from __future__ import annotations

import argparse
import hashlib
import json
import sys
from pathlib import Path


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def main() -> int:
    parser = argparse.ArgumentParser(description="Assemble completed Task 12 page packages.")
    parser.add_argument("--runtime-scripts", type=Path, required=True)
    parser.add_argument("--project", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--pages", type=int, nargs="+", required=True)
    args = parser.parse_args()

    sys.path.insert(0, str(args.runtime_scripts.resolve()))
    from pptx import Presentation
    from workflow_v6_reconstruction import _copy_page_slide

    completed: list[tuple[int, Path]] = []
    missing: list[int] = []
    for page_number in args.pages:
        path = args.project / "06_v6" / "pages" / f"page_{page_number:03d}" / "page.pptx"
        if path.is_file():
            completed.append((page_number, path))
        else:
            missing.append(page_number)
    if not completed:
        raise SystemExit("no completed Task 12 pages are available")

    deck = Presentation(completed[0][1])
    layout = deck.slides[0].slide_layout
    for page_number, source in completed[1:]:
        _copy_page_slide(source, deck, layout, page_number)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    deck.save(args.output)
    reopened = Presentation(args.output)
    if len(reopened.slides) != len(completed):
        raise SystemExit("assembled sample deck slide count is incorrect")
    if any(not any(shape.has_text_frame or shape.has_table for shape in slide.shapes) for slide in reopened.slides):
        raise SystemExit("assembled sample deck contains a slide without editable text or table objects")

    report = {
        "artifact_version": "task12-selected-page-deck-v1",
        "requested_pages": args.pages,
        "assembled_pages": [number for number, _path in completed],
        "missing_pages": missing,
        "slide_count": len(reopened.slides),
        "output": str(args.output.resolve()),
        "sha256": sha256(args.output),
        "editable_text_or_table_on_every_slide": True,
    }
    report_path = args.output.with_suffix(".json")
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
