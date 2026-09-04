#!/usr/bin/env python3
"""Read final editable authority evidence from one preserved Task 12 page."""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation


def _object_id(shape: object) -> str | None:
    properties = shape._element.xpath(".//p:cNvPr")
    if not properties:
        return None
    description = properties[0].get("descr", "")
    return description.removeprefix("object_id:") if description.startswith("object_id:") else None


def _inside(point: tuple[int, int], shape: object) -> bool:
    x, y = point
    return (
        shape.left - 1 <= x <= shape.left + shape.width + 1
        and shape.top - 1 <= y <= shape.top + shape.height + 1
    )


def collect(page_dir: Path, runtime_scripts: Path) -> dict[str, object]:
    runtime_scripts = runtime_scripts.resolve()
    runtime = runtime_scripts.parents[1] / "reconstruct-editable-slide/cli/editppt/runtime"
    sys.path[:0] = [str(runtime_scripts), str(runtime)]
    from validate_pptx import (  # noqa: PLC0415
        _connector_endpoints,
        _shape_arrowheads,
        _shape_kind,
        quantitative_chart_readback_violations,
    )

    request_path = page_dir / "accepted_reconstruction_request.json"
    manifest_path = page_dir / "manifest.json"
    pptx_path = page_dir / "page.pptx"
    missing = [path.name for path in (request_path, manifest_path, pptx_path) if not path.is_file()]
    if missing:
        return {"available": False, "applicable": False, "compliant": None, "missing": missing}

    request = json.loads(request_path.read_text(encoding="utf-8"))
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    relationship = request.get("page_plan", {}).get("primary_relationship", {})
    nodes = relationship.get("nodes", [])
    edges = relationship.get("edges", [])
    numeric = request.get("numeric_authority")
    applicable = bool(nodes or edges or numeric is not None)

    manifest_by_id: dict[str, list[dict[str, object]]] = {}
    for section in ("text_boxes", "tables", "images", "shapes", "charts"):
        for item in manifest.get(section, []):
            if isinstance(item, dict) and isinstance(item.get("object_id"), str):
                manifest_by_id.setdefault(item["object_id"], []).append({"section": section, **item})

    deck = Presentation(pptx_path)
    shapes_by_id: dict[str, list[object]] = {}
    for shape in deck.slides[0].shapes:
        object_id = _object_id(shape)
        if object_id:
            shapes_by_id.setdefault(object_id, []).append(shape)

    node_results = []
    for node in nodes:
        node_id = node.get("node_id")
        manifest_count = len(manifest_by_id.get(node_id, []))
        pptx_count = len(shapes_by_id.get(node_id, []))
        node_results.append({
            "node_id": node_id,
            "manifest_count": manifest_count,
            "pptx_count": pptx_count,
            "compliant": manifest_count == 1 and pptx_count == 1,
        })

    edge_results = []
    for edge in edges:
        source_id, target_id = edge.get("from_node"), edge.get("to_node")
        edge_id = f"edge:{source_id}->{target_id}"
        manifest_matches = manifest_by_id.get(edge_id, [])
        pptx_matches = shapes_by_id.get(edge_id, [])
        result: dict[str, object] = {
            "edge_id": edge_id,
            "from_node": source_id,
            "to_node": target_id,
            "manifest_count": len(manifest_matches),
            "manifest_object_type": manifest_matches[0].get("type") if len(manifest_matches) == 1 else None,
            "pptx_count": len(pptx_matches),
            "pptx_object_type": None,
            "endpoints_emu": None,
            "source_endpoint_inside": False,
            "target_endpoint_inside": False,
            "direction_match": False,
            "arrowheads": {},
            "arrowhead_match": False,
            "compliant": False,
        }
        if len(pptx_matches) == 1:
            connector = pptx_matches[0]
            result["pptx_object_type"] = _shape_kind(connector)
            endpoints = _connector_endpoints(connector)
            result["endpoints_emu"] = list(endpoints)
            if len(shapes_by_id.get(source_id, [])) == len(shapes_by_id.get(target_id, [])) == 1:
                source_inside = _inside(endpoints[:2], shapes_by_id[source_id][0])
                target_inside = _inside(endpoints[2:], shapes_by_id[target_id][0])
                result["source_endpoint_inside"] = source_inside
                result["target_endpoint_inside"] = target_inside
                result["direction_match"] = source_inside and target_inside
            arrowheads = _shape_arrowheads(connector)
            result["arrowheads"] = arrowheads
            result["arrowhead_match"] = arrowheads == {"tailEnd": "triangle"}
        result["compliant"] = (
            len(manifest_matches) == 1
            and len(pptx_matches) == 1
            and result["pptx_object_type"] in {"line", "connector"}
            and result["direction_match"]
            and result["arrowhead_match"]
        )
        edge_results.append(result)

    numeric_result: dict[str, object] = {"present": numeric is not None, "compliant": None}
    if numeric is not None:
        matching = [
            chart for chart in manifest.get("charts", [])
            if isinstance(chart, dict) and all(chart.get(key) == value for key, value in numeric.items())
        ]
        violations = quantitative_chart_readback_violations(pptx_path, [manifest])
        numeric_result.update({
            "sealed_authority": numeric,
            "matching_manifest_charts": len(matching),
            "readback_violations": violations,
            "compliant": len(matching) == 1 and not violations,
        })

    compliant = all(item["compliant"] for item in node_results + edge_results) if applicable else None
    if applicable and numeric_result["compliant"] is not None:
        compliant = bool(compliant) and bool(numeric_result["compliant"])
    final_object_ids = sorted(shapes_by_id)
    return {
        "available": True,
        "applicable": applicable,
        "compliant": compliant,
        "contract": "sealed page_plan.primary_relationship and numeric_authority",
        "sealed_node_ids": [item.get("node_id") for item in nodes],
        "final_editable_node_ids": sorted(set(final_object_ids).intersection(item.get("node_id") for item in nodes)),
        "final_editable_object_ids": final_object_ids,
        "nodes": node_results,
        "edges": edge_results,
        "numeric_authority": numeric_result,
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--page-dir", type=Path, required=True)
    parser.add_argument("--runtime-scripts", type=Path, required=True)
    args = parser.parse_args()
    print(json.dumps(collect(args.page_dir, args.runtime_scripts), ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
