from __future__ import annotations

import hashlib
import json
import subprocess
import sys
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
PLUGIN = ROOT.parents[1]
EDITPPT_RUNTIME = PLUGIN / "skills" / "reconstruct-editable-slide" / "cli" / "editppt" / "runtime"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from test_workflow_v6_reconstruction import _body, _project, _write_signed_receipt  # noqa: E402
import workflow_v6_pipeline  # noqa: E402
import workflow_v6_reconstruction_worker as worker_module  # noqa: E402
from workflow_v6_pipeline import (  # noqa: E402
    PipelineConfiguration,
    PipelineDependencies,
    production_pipeline_dependencies,
    run_pages,
)
from workflow_v6_reconstruction_worker import (  # noqa: E402
    PageWorkerResult,
    reconstruct_accepted_page,
)
from workflow_v6_reconstruction import assemble_v6_deck  # noqa: E402
from workflow_v6_media import normalized_raster_pixel_seal  # noqa: E402
from workflow_v6_state import load, save  # noqa: E402
from test_quantitative_chart_v123_e2e import (  # noqa: E402
    _production_worker,
    _qualitative_manifest,
    _relationship_manifest,
    _with_relationship,
)


def _accepted_outcome(project: Path, page_number: int = 1):
    receipt = json.loads(
        (project / "04_v6" / "images" / f"page_{page_number:03d}.json").read_text(
            encoding="utf-8"
        )
    )
    selected = receipt.get("candidate", receipt.get("selected"))
    candidate = SimpleNamespace(
        path=project / selected["path"],
        attempt=selected["attempt"],
    )
    return SimpleNamespace(
        status="accepted",
        accepted=SimpleNamespace(candidate=candidate),
    )


def _workspace(project: Path, page_number: int = 1):
    return SimpleNamespace(project_copy=project, page_number=page_number)


def _successful_worker(calls: list, text: str = "Editable worker output"):
    def manifest(page_request):
        value = _with_relationship(
            _qualitative_manifest("flow", "timeline_roadmap"), page_request,
        )
        value["text_boxes"].append({
            "object_id": "worker-output-text",
            "name": "worker-output-text",
            "box_px": [200, 300, 800, 100],
            "text": text,
            "font_size": 20,
        })
        return value

    sealed_worker = _production_worker(manifest, [])

    def invoke(request):
        calls.append(request)
        return sealed_worker(request)

    return invoke


def _numeric_authority() -> dict[str, object]:
    return {
        "title": "Revenue",
        "relationship": "change_over_time",
        "rendering_primitive": "line_point",
        "chart_variant": "line",
        "unit": "USD m",
        "basis": "reported revenue",
        "period": "FY2024-FY2025",
        "series": [{"name": "Revenue", "categories": ["2024", "2025"], "values": [12, 18]}],
    }


def test_direct_codex_worker_success_uses_zero_paddle_and_recovers_with_zero_calls(
    tmp_path: Path,
):
    project = _project(tmp_path, 1)
    worker_calls: list = []
    paddle_calls: list = []

    first = reconstruct_accepted_page(
        _workspace(project),
        _accepted_outcome(project),
        page_worker=_successful_worker(worker_calls),
        paddle_runner=lambda request: paddle_calls.append(request),
    )
    recovered = reconstruct_accepted_page(
        _workspace(project),
        _accepted_outcome(project),
        page_worker=lambda request: pytest.fail("recovery called the page worker"),
        paddle_runner=lambda request: pytest.fail("recovery called Paddle"),
    )
    assembly = assemble_v6_deck(project)

    assert first["reconstruction_mode"] == "codex_direct_reconstruction"
    assert recovered["recovered"] is True
    assert len(worker_calls) == 1
    assert paddle_calls == []
    prompt = worker_calls[0].prompt_file.read_text(encoding="utf-8")
    assert str(EDITPPT_RUNTIME / "main.py") in prompt
    assert "do not rely on a separately installed CLI" in prompt
    assert load(project)["pages"][0]["state"] == "page_complete"
    page_request = json.loads((worker_calls[0].page_dir / "page_request.json").read_text(encoding="utf-8"))
    accepted_request = json.loads(
        (worker_calls[0].page_dir / "accepted_reconstruction_request.json").read_text(encoding="utf-8")
    )
    final = json.loads((project / "06_v6/pages/page_001/page.json").read_text(encoding="utf-8"))
    reconstruction = json.loads(
        (project / "05_v6/reconstruction_runs/page_001/reconstruction.json").read_text(encoding="utf-8")
    )
    assert page_request["accepted_source_body"] == accepted_request["source_body"]
    assert page_request["worker_source_body"]["normalized_pixel_sha256"] == (
        accepted_request["source_body"]["normalized_pixel_sha256"]
    )
    assert final["accepted_source_body"] == accepted_request["source_body"]
    assert reconstruction["accepted_source_body"] == accepted_request["source_body"]
    assert reconstruction["worker_source_body"] == page_request["worker_source_body"]
    assert [
        {key: item[key] for key in ("page_number", "status", "authority_mode")}
        for item in assembly["page_authority"]
    ] == [{
        "page_number": 1, "status": "verified", "authority_mode": "sealed_reconstruction",
    }]
    assert assembly["page_authority"][0]["visual_qa"]["status"] in {
        "passed", "unavailable",
    }
    if assembly["release_ready"]:
        assert assembly["assembled_visual_qa"]["status"] == "passed"
        assert assembly["sha256"] == hashlib.sha256(
            (project / assembly["output"]).read_bytes()
        ).hexdigest()
    else:
        assert assembly["status"] == "validation_incomplete"
        assert assembly["release_status"] == "not_release_ready"
        assert assembly["final_output"] is None
        assert "output" not in assembly
        assert (project / assembly["candidate_output"]["relative_path"]).is_file()


@pytest.mark.parametrize(
    "target,mutate",
    [
        (
            "accepted_reconstruction_request.json",
            lambda value: value["page_plan"].update({"page_purpose": "tampered"}),
        ),
        (
            "page_request.json",
            lambda value: value["page_plan"].update({"page_purpose": "tampered"}),
        ),
        (
            "page_jobs.json",
            lambda value: value["pages"][0]["dispatch"].update(
                {"page_request_sha256": "0" * 64}
            ),
        ),
        (
            "manifest.json",
            lambda value: value.update({"shapes": []}),
        ),
        (
            "reconstruction.json",
            lambda value: value.update({"final_page_sha256": "0" * 64}),
        ),
        (
            "acceptance_receipt",
            lambda value: value["page_plan"].update({"page_purpose": "tampered"}),
        ),
        (
            "final_page_receipt",
            lambda value: value.update({"artifact_version": "tampered"}),
        ),
        (
            "final_page_receipt",
            lambda value: value.update({"page_pptx": "06_v6/pages/page_999/page.pptx"}),
        ),
        (
            "final_page_receipt",
            lambda value: value["fixed_frame"].update({"passed": False}),
        ),
    ],
)
def test_assembly_revalidates_the_complete_sealed_page_authority_chain(
    tmp_path: Path, target: str, mutate,
):
    project = _project(tmp_path, 1)
    calls: list = []
    reconstruct_accepted_page(
        _workspace(project), _accepted_outcome(project),
        page_worker=_successful_worker(calls),
    )
    run_dir = project / "05_v6/reconstruction_runs/page_001"
    path = {
        "acceptance_receipt": project / "04_v6/images/page_001.json",
        "final_page_receipt": project / "06_v6/pages/page_001/page.json",
        "reconstruction.json": run_dir / "reconstruction.json",
        "page_jobs.json": run_dir / "page_jobs.json",
    }.get(target, run_dir / "pages/page_001" / target)
    value = json.loads(path.read_text(encoding="utf-8"))
    mutate(value)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    with pytest.raises(
        (RuntimeError, ValueError),
        match="authority|receipt|request|relationship|edge|signature|artifact",
    ):
        assemble_v6_deck(project)

    assert not (project / "08_final/deck.pptx").exists()


def test_assembly_rejects_changed_final_page_bytes(tmp_path: Path):
    project = _project(tmp_path, 1)
    reconstruct_accepted_page(
        _workspace(project), _accepted_outcome(project),
        page_worker=_successful_worker([]),
    )
    page = project / "06_v6/pages/page_001/page.pptx"
    page.write_bytes(page.read_bytes() + b"tampered")

    with pytest.raises(RuntimeError, match="completed reconstructed page changed"):
        assemble_v6_deck(project)

    assert not (project / "08_final/deck.pptx").exists()


@pytest.mark.parametrize("page_role", ["content", "appendix"])
def test_assembly_rejects_non_special_page_disguised_as_special_page(
    tmp_path: Path, page_role: str,
):
    project = _project(tmp_path, 1)
    composition_path = project / "02_v6/page_composition.json"
    composition = json.loads(composition_path.read_text(encoding="utf-8"))
    composition["pages"][0]["page_role"] = page_role
    composition_path.write_text(
        json.dumps(composition, ensure_ascii=False, indent=2) + "\n", encoding="utf-8",
    )
    reconstruct_accepted_page(
        _workspace(project), _accepted_outcome(project),
        page_worker=_successful_worker([]),
    )
    receipt_path = project / "06_v6/pages/page_001/page.json"
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    receipt["artifact_version"] = "special-page-v6"
    receipt_path.write_text(
        json.dumps(receipt, ensure_ascii=False, indent=2) + "\n", encoding="utf-8",
    )

    with pytest.raises(RuntimeError, match="artifact|role|authority"):
        assemble_v6_deck(project)


def test_assembly_rejects_unsigned_source_chain_swapped_away_from_signed_candidate(
    tmp_path: Path,
):
    project = _project(tmp_path, 1)
    reconstruct_accepted_page(
        _workspace(project), _accepted_outcome(project),
        page_worker=_successful_worker([]),
    )
    swapped_path = project / "04_v6/images/swapped.png"
    Image.new("RGB", (1904, 896), "black").save(swapped_path)
    swapped_bytes = swapped_path.read_bytes()
    accepted_source = {
        "path": swapped_path.relative_to(project).as_posix(),
        "sha256": hashlib.sha256(swapped_bytes).hexdigest(),
        **normalized_raster_pixel_seal(swapped_bytes),
    }
    run_dir = project / "05_v6/reconstruction_runs/page_001"
    page_dir = run_dir / "pages/page_001"
    worker_path = page_dir / "source.png"
    worker_path.write_bytes(swapped_bytes)
    worker_source = {
        "path": worker_path.relative_to(project).as_posix(),
        "sha256": hashlib.sha256(swapped_bytes).hexdigest(),
        **normalized_raster_pixel_seal(swapped_bytes),
    }
    request_paths = [
        project / "05_v6/reconstruction_requests/page_001.json",
        page_dir / "accepted_reconstruction_request.json",
    ]
    for path in request_paths:
        value = json.loads(path.read_text(encoding="utf-8"))
        value["source_body"] = accepted_source
        path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    page_request_path = page_dir / "page_request.json"
    page_request = json.loads(page_request_path.read_text(encoding="utf-8"))
    page_request["accepted_source_body"] = accepted_source
    page_request["worker_source_body"] = worker_source
    page_request_path.write_text(
        json.dumps(page_request, ensure_ascii=False, indent=2) + "\n", encoding="utf-8",
    )
    jobs_path = run_dir / "page_jobs.json"
    jobs = json.loads(jobs_path.read_text(encoding="utf-8"))
    jobs["pages"][0]["dispatch"]["page_request_sha256"] = hashlib.sha256(
        page_request_path.read_bytes()
    ).hexdigest()
    jobs_path.write_text(json.dumps(jobs, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    reconstruction_path = run_dir / "reconstruction.json"
    reconstruction = json.loads(reconstruction_path.read_text(encoding="utf-8"))
    reconstruction.update({
        "accepted_image_sha256": accepted_source["sha256"],
        "accepted_image_pixel_sha256": accepted_source["normalized_pixel_sha256"],
        "accepted_source_body": accepted_source,
        "worker_source_body": worker_source,
    })
    reconstruction_path.write_text(
        json.dumps(reconstruction, ensure_ascii=False, indent=2) + "\n", encoding="utf-8",
    )
    final_path = project / "06_v6/pages/page_001/page.json"
    final = json.loads(final_path.read_text(encoding="utf-8"))
    final["accepted_source_body"] = accepted_source
    final["worker_source_body"] = worker_source
    final_path.write_text(json.dumps(final, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    with pytest.raises((RuntimeError, ValueError), match="accepted|candidate|source|authority"):
        assemble_v6_deck(project)


def test_assembly_rejects_manifest_bytes_changed_after_worker_record(tmp_path: Path):
    project = _project(tmp_path, 1)
    reconstruct_accepted_page(
        _workspace(project), _accepted_outcome(project),
        page_worker=_successful_worker([]),
    )
    manifest_path = project / "05_v6/reconstruction_runs/page_001/pages/page_001/manifest.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    manifest["harmless_note"] = "changed after record"
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8",
    )

    with pytest.raises((RuntimeError, ValueError), match="manifest|record|authority|digest"):
        assemble_v6_deck(project)


def test_assembly_rejects_worker_pptx_with_undeclared_shape_after_record(tmp_path: Path):
    project = _project(tmp_path, 1)
    reconstruct_accepted_page(
        _workspace(project), _accepted_outcome(project),
        page_worker=_successful_worker([]),
    )
    worker_path = project / "05_v6/reconstruction_runs/page_001/pages/page_001/page.pptx"
    deck = Presentation(worker_path)
    deck.slides[0].shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text = (
        "undeclared"
    )
    deck.save(worker_path)

    with pytest.raises((RuntimeError, ValueError), match="worker|record|authority|digest|PPTX"):
        assemble_v6_deck(project)


def test_page_worker_prompt_enforces_sealed_text_repairs(tmp_path: Path):
    project = _project(tmp_path, 1)
    receipt_path = project / "04_v6/images/page_001.json"
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    receipt["reconstruction_repairs"] = [
        {
            "category": "severe_usability",
            "detail": "将错字“清出”修正为“退出”，其余构图保持不变。",
        }
    ]
    _write_signed_receipt(project, 1, receipt)
    calls: list = []

    reconstruct_accepted_page(
        _workspace(project),
        _accepted_outcome(project),
        page_worker=_successful_worker(calls, text="退出"),
    )

    prompt = calls[0].prompt_file.read_text(encoding="utf-8")
    assert "SEALED TEXT REPAIRS" in prompt
    assert "将错字“清出”修正为“退出”" in prompt


def test_page_worker_request_copies_numeric_authority_before_whole_request_hash(
    tmp_path: Path,
):
    project = _project(tmp_path, 1)
    authority = _numeric_authority()
    materials = project / "02_v6/page_materials/page_001.json"
    materials.parent.mkdir(parents=True, exist_ok=True)
    materials.write_text(
        json.dumps({"chart_facts": [authority]}, ensure_ascii=False), encoding="utf-8"
    )
    calls: list = []

    def worker(request):
        calls.append(request)
        return _production_worker(_relationship_manifest, [])(request)

    reconstruct_accepted_page(
        _workspace(project), _accepted_outcome(project), page_worker=worker,
    )

    page_request_path = calls[0].page_dir / "page_request.json"
    page_request = json.loads(page_request_path.read_text(encoding="utf-8"))
    receipt = json.loads((project / "04_v6/images/page_001.json").read_text(encoding="utf-8"))
    jobs = json.loads((calls[0].run_dir / "page_jobs.json").read_text(encoding="utf-8"))
    assert page_request["numeric_authority"] == authority
    assert page_request["page_plan"] == receipt["page_plan"]
    assert jobs["pages"][0]["dispatch"]["page_request_sha256"] == hashlib.sha256(
        page_request_path.read_bytes()
    ).hexdigest()
    assert not list(calls[0].page_dir.glob("numeric_authority*.json"))


def test_qualitative_relationship_never_creates_numeric_authority_in_page_request(
    tmp_path: Path,
):
    project = _project(tmp_path, 1)
    materials = project / "02_v6/page_materials/page_001.json"
    materials.parent.mkdir(parents=True, exist_ok=True)
    materials.write_text(
        json.dumps(
            {
                "chart_facts": [{
                    "title": "Change over time",
                    "relationship": "change_over_time",
                    "source_wording": "The source states a sequence but no complete values.",
                    "disabled_primitive": "line_point",
                    "fallback": "timeline_roadmap",
                    "series": [],
                }]
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    calls: list = []

    reconstruct_accepted_page(
        _workspace(project),
        _accepted_outcome(project),
        page_worker=_successful_worker(calls),
    )

    page_request = json.loads(
        (calls[0].page_dir / "page_request.json").read_text(encoding="utf-8")
    )
    assert "numeric_authority" not in page_request


def test_interrupted_prepare_validates_accepted_request_before_resyncing_authority(
    tmp_path: Path,
):
    project = _project(tmp_path, 1)
    accepted_request = worker_module.build_reconstruction_request(project, page_number=1)
    run_dir, page_dir, _prompt_file = worker_module._prepare_run(
        project, accepted_request, 1
    )
    page_request_path = page_dir / "page_request.json"
    stale_authority = _numeric_authority()
    page_request = json.loads(page_request_path.read_text(encoding="utf-8"))
    page_request["numeric_authority"] = stale_authority
    page_request_path.write_text(
        json.dumps(page_request, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )

    changed_request = {
        **accepted_request,
        "numeric_authority": {**stale_authority, "unit": "EUR m"},
    }
    with pytest.raises(RuntimeError, match="accepted reconstruction request changed"):
        worker_module._prepare_run(project, changed_request, 1)
    assert json.loads(page_request_path.read_text(encoding="utf-8"))[
        "numeric_authority"
    ] == stale_authority

    recovered_run, recovered_page, _prompt_file = worker_module._prepare_run(
        project, accepted_request, 1
    )
    assert recovered_run == run_dir
    assert recovered_page == page_dir
    assert "numeric_authority" not in json.loads(
        page_request_path.read_text(encoding="utf-8")
    )


def test_interrupted_prepare_rejects_stale_page_plan_before_resync(tmp_path: Path):
    project = _project(tmp_path, 1)
    accepted_request = worker_module.build_reconstruction_request(project, page_number=1)
    _run_dir, page_dir, _prompt_file = worker_module._prepare_run(project, accepted_request, 1)
    page_request_path = page_dir / "page_request.json"
    original_page_request = page_request_path.read_bytes()
    changed_request = json.loads(json.dumps(accepted_request))
    changed_request["page_plan"]["primary_relationship"]["edges"][0]["to_node"] = "source"

    with pytest.raises(RuntimeError, match="accepted reconstruction request changed"):
        worker_module._prepare_run(project, changed_request, 1)

    assert page_request_path.read_bytes() == original_page_request


def test_unreadable_text_uses_paddle_once_then_same_page_worker(tmp_path: Path):
    project = _project(tmp_path, 1)
    worker_calls: list = []
    paddle_calls: list = []

    def worker(request):
        worker_calls.append(request)
        if len(worker_calls) == 1:
            return PageWorkerResult(
                status="needs_paddle",
                reason="accepted image text is too small to transcribe reliably",
            )
        return _successful_worker([], "Paddle assisted editable output")(request)

    def paddle(request):
        paddle_calls.append(request)
        hints = request.page_dir / "text_hints.json"
        hints.write_text('{"backend":"paddleocr-vl","lines":[]}', encoding="utf-8")
        return hints

    result = reconstruct_accepted_page(
        _workspace(project),
        _accepted_outcome(project),
        page_worker=worker,
        paddle_runner=paddle,
        paddle_token="configured-token",
        paddle_authorized=True,
    )

    assert result["reconstruction_mode"] == "paddle_assisted_reconstruction"
    assert len(worker_calls) == 2
    assert len(paddle_calls) == 1
    assert worker_calls[0].page_number == worker_calls[1].page_number == 1
    assert worker_calls[0].page_dir == worker_calls[1].page_dir
    assert worker_calls[0].text_hints is None
    assert worker_calls[1].text_hints.name == "text_hints.json"


@pytest.mark.parametrize(
    ("token", "authorized", "paddle_fails"),
    [
        ("", True, False),
        ("configured-token", False, False),
        ("configured-token", True, True),
    ],
)
def test_unreadable_text_without_permitted_working_paddle_stops_page(
    tmp_path: Path, token: str | None, authorized: bool, paddle_fails: bool,
):
    project = _project(tmp_path, 1)
    paddle_calls: list = []

    def worker(request):
        return PageWorkerResult(status="needs_paddle", reason="text unreadable")

    def paddle(request):
        paddle_calls.append(request)
        if paddle_fails:
            raise RuntimeError("Paddle unavailable")
        return request.page_dir / "text_hints.json"

    with pytest.raises(RuntimeError, match="Paddle|text unreadable"):
        reconstruct_accepted_page(
            _workspace(project),
            _accepted_outcome(project),
            page_worker=worker,
            paddle_runner=paddle,
            paddle_token=token,
            paddle_authorized=authorized,
        )

    assert not (project / "06_v6" / "pages" / "page_001" / "page.pptx").exists()
    assert len(paddle_calls) == (1 if token and authorized and paddle_fails else 0)


def test_failed_paddle_assisted_worker_does_not_publish_page(tmp_path: Path):
    project = _project(tmp_path, 1)
    calls = 0

    def worker(request):
        nonlocal calls
        calls += 1
        if calls == 1:
            return PageWorkerResult(status="needs_paddle", reason="dense unreadable text")
        return PageWorkerResult(status="failed", reason="could not reconstruct after Paddle")

    def paddle(request):
        hints = request.page_dir / "text_hints.json"
        hints.write_text('{"backend":"paddleocr-vl","lines":[]}', encoding="utf-8")
        return hints

    with pytest.raises(RuntimeError, match="could not reconstruct"):
        reconstruct_accepted_page(
            _workspace(project),
            _accepted_outcome(project),
            page_worker=worker,
            paddle_runner=paddle,
            paddle_token="configured-token",
            paddle_authorized=True,
        )

    assert calls == 2
    assert not (project / "06_v6" / "pages" / "page_001" / "page.pptx").exists()


def test_failed_dispatched_worker_requires_explicit_reset_before_resubmission(tmp_path: Path):
    project = _project(tmp_path, 1)

    def failed_worker(request):
        jobs_path = request.run_dir / "page_jobs.json"
        jobs = json.loads(jobs_path.read_text(encoding="utf-8"))
        jobs["pages"][0]["status"] = "dispatched"
        jobs["pages"][0]["dispatch"] = {"agent_id": "codex-page-worker"}
        jobs_path.write_text(json.dumps(jobs), encoding="utf-8")
        (request.page_dir / "worker-last-message.txt").write_text(
            "worker could not execute commands", encoding="utf-8",
        )
        return PageWorkerResult(status="failed", reason="worker unavailable")

    with pytest.raises(RuntimeError, match="worker unavailable"):
        reconstruct_accepted_page(
            _workspace(project), _accepted_outcome(project), page_worker=failed_worker,
        )
    with pytest.raises(RuntimeError, match="reset required"):
        reconstruct_accepted_page(
            _workspace(project), _accepted_outcome(project),
            page_worker=lambda request: pytest.fail("failed worker must not be resubmitted"),
        )


def test_nonaccepted_and_legacy_fallback_states_cannot_reconstruct(tmp_path: Path):
    project = _project(tmp_path, 1)
    outcome = _accepted_outcome(project)
    outcome.status = "failed"
    outcome.accepted = None
    with pytest.raises(ValueError, match="independently accepted"):
        reconstruct_accepted_page(
            _workspace(project), outcome,
            page_worker=lambda request: pytest.fail("worker must not run"),
        )

    state = load(project)
    state["pages"][0]["state"] = "accepted_fallback_first"
    with pytest.raises(ValueError, match="state is invalid"):
        save(project, state)


def test_production_pipeline_defaults_auto_reconstruct_and_assemble():
    dependencies = production_pipeline_dependencies()
    assert dependencies.reconstruct_page is reconstruct_accepted_page
    assert callable(dependencies.assemble_project)


def test_run_pages_without_dependency_override_automatically_reconstructs_acceptance(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
):
    project = _project(tmp_path, 1)
    reconstructed: list[tuple[object, object]] = []
    assembled: list[dict[int, object]] = []
    accepted = SimpleNamespace(status="accepted", accepted=SimpleNamespace(candidate=object()))
    dependencies = PipelineDependencies(
        open_workspace=lambda root, page: SimpleNamespace(project_copy=root, page_number=page),
        evidence_recorder=lambda workspace: object(),
        candidate_loop=lambda workspace, **kwargs: accepted,
        reconstruct_page=lambda workspace, outcome: reconstructed.append((workspace, outcome)),
        assemble_project=lambda root, outcomes: assembled.append(outcomes),
    )
    monkeypatch.setattr(workflow_v6_pipeline, "production_pipeline_dependencies", lambda: dependencies)

    report = run_pages(
        project,
        [1],
        configuration=PipelineConfiguration(
            page_workers=1,
            initial_page_concurrency=1,
            maximum_page_concurrency=1,
        ),
    )

    assert report.completed_pages == (1,)
    assert len(reconstructed) == 1
    assert reconstructed[0][1] is accepted
    assert assembled == [{1: accepted}]


def test_single_page_prepare_dispatches_worker_without_ocr_or_local_mode(tmp_path: Path):
    source = tmp_path / "accepted.png"
    Image.new("RGB", (1904, 896), "white").save(source)
    run_dir = tmp_path / "run"
    prepare = subprocess.run(
        [
            sys.executable,
            str(EDITPPT_RUNTIME / "main.py"),
            "prepare",
            str(source),
            "--job-dir", str(run_dir),
            "--max-concurrent-pages", "1",
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    assert prepare.returncode == 0, prepare.stderr
    assert not (run_dir / "pages" / "page_001" / "text_hints.json").exists()
    assert not (run_dir / "pages" / "page_001" / "text_hints.png").exists()

    next_step = subprocess.run(
        [sys.executable, str(EDITPPT_RUNTIME / "main.py"), "run", "next", str(run_dir), "--json"],
        capture_output=True,
        text=True,
        check=False,
    )
    assert next_step.returncode == 0, next_step.stderr
    payload = json.loads(next_step.stdout)
    assert payload["stage"] == "dispatch_pages"
    assert payload["suggested_pages"] == ["page_001"]
    assert "--local" not in next_step.stdout


def test_shipped_runtime_has_no_retired_low_quality_or_api_fallback_surface():
    forbidden = (
        "builtin-ink",
        "rebuild_page_locally",
        "--local",
        "accepted_fallback_first",
        "codex_fallback",
        "--no-text-hints",
        "OPENAI_API_KEY",
        "OPENAI_BASE_URL",
        "IMAGE_TO_EDITABLE_PPT_IMAGE_MODEL",
        "openai-compatible-api",
    )
    findings: list[str] = []
    for path in PLUGIN.rglob("*"):
        if not path.is_file() or "__pycache__" in path.parts or "tests" in path.parts:
            continue
        if path.suffix.lower() not in {".py", ".md", ".json", ".yaml", ".toml"}:
            continue
        text = path.read_text(encoding="utf-8", errors="replace")
        for token in forbidden:
            if token in text:
                findings.append(f"{path.relative_to(PLUGIN)}: {token}")
    assert findings == []

    help_result = subprocess.run(
        [sys.executable, str(EDITPPT_RUNTIME / "main.py"), "--help"],
        capture_output=True,
        text=True,
        check=False,
    )
    assert help_result.returncode == 0
    assert all(token not in help_result.stdout for token in forbidden)


def test_page_worker_must_preserve_source_line_fitting_before_passing():
    prompt = (
        PLUGIN
        / "skills"
        / "reconstruct-editable-slide"
        / "prompts"
        / "page-worker.md"
    ).read_text(encoding="utf-8")

    assert "no clipped, truncated, or unintended wrapped text" in prompt
    assert "line breaks visible in source.png" in prompt
    assert "preview.png at the source image dimensions" in prompt


def test_codex_process_failure_reports_transport_error_before_missing_validation(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
):
    run_dir = tmp_path / "run"
    page_dir = run_dir / "pages" / "page_001"
    page_dir.mkdir(parents=True)
    (run_dir / "page_jobs.json").write_text(
        json.dumps({"pages": [{"status": "dispatched"}]}), encoding="utf-8",
    )
    prompt = page_dir / "worker-prompt.md"
    source = page_dir / "source.png"
    prompt.write_text("reconstruct", encoding="utf-8")
    Image.new("RGB", (1904, 896), "white").save(source)
    (page_dir / "validation.json").write_text('{"passed":true}', encoding="utf-8")
    (page_dir / "manifest.json").write_text('{"text_boxes":[]}', encoding="utf-8")
    _body(page_dir / "page.pptx", "stale passed body")
    monkeypatch.setattr(worker_module, "_codex_executable", lambda: "codex")
    monkeypatch.setattr(
        worker_module.subprocess,
        "run",
        lambda *args, **kwargs: subprocess.CompletedProcess(
            args[0], 1, stdout="", stderr="remote worker transport failed",
        ),
    )
    request = worker_module.PageWorkerRequest(
        project=tmp_path,
        page_number=1,
        run_dir=run_dir,
        page_dir=page_dir,
        source_image=source,
        prompt_file=prompt,
        text_hints=None,
        timeout=30,
    )

    result = worker_module._default_page_worker(request)

    assert result.status == "failed"
    assert result.reason == "remote worker transport failed"
    assert not (page_dir / "validation.json").exists()
    assert not (page_dir / "manifest.json").exists()
    assert not (page_dir / "page.pptx").exists()


def test_formal_reconstruction_commits_state_after_all_receipts(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    project = _project(tmp_path, 1)
    calls: list = []
    original_write = worker_module._atomic_json

    def fail_reconstruction_receipt(root, path, value):
        if Path(path).name == "reconstruction.json":
            raise RuntimeError("simulated reconstruction receipt failure")
        return original_write(root, path, value)

    monkeypatch.setattr(worker_module, "_atomic_json", fail_reconstruction_receipt)
    with pytest.raises(RuntimeError, match="receipt failure"):
        reconstruct_accepted_page(
            _workspace(project), _accepted_outcome(project), page_worker=_successful_worker(calls),
        )

    assert load(project)["pages"][0]["state"] != "page_complete"
    assert not (project / "06_v6/pages/page_001/page.pptx").exists()
    assert not (project / "06_v6/pages/page_001/page.json").exists()
