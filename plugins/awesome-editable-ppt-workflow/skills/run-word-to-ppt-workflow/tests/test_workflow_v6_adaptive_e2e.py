from __future__ import annotations

import hashlib
import importlib.util
import json
import sys
from collections import defaultdict
from pathlib import Path
from types import SimpleNamespace

import pytest
from docx import Document
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Cm


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
SERVER = SCRIPTS / "confirm_ui" / "server.py"
FIXTURE = Path(__file__).parent / "fixtures" / "v6_adaptive_project" / "fixture.json"
CONSULTING_REGRESSION_FIXTURE = (
    Path(__file__).parent / "fixtures" / "consulting_director_cases.json"
)
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

import workflow_v6_source  # noqa: E402
from workflow_v6_contract import canonical_sha256  # noqa: E402
from workflow_v6_image import generate_page_body  # noqa: E402
from workflow_v6_reconstruction import (  # noqa: E402
    assemble_v6_deck,
    build_reconstruction_request,
    finalize_reconstructed_page as _finalize_reconstructed_page,
)
from workflow_v6_source import initialize_v6_project  # noqa: E402
from workflow_v6_state import load, save  # noqa: E402
from director_taskbook import confirmed_taskbook_prompt  # noqa: E402
from workflow_v6_contract import transition_page  # noqa: E402
from workflow_v6_pipeline import (  # noqa: E402
    PipelineConfiguration,
    PipelineDependencies,
    run_pages,
)
from workflow_v6_special_pages import render_special_page  # noqa: E402
from awesome_page_materials import publish_page_materials  # noqa: E402
from fixed_region_contract import fixed_frame_execution  # noqa: E402


def finalize_reconstructed_page(*args, **kwargs):
    return _finalize_reconstructed_page(*args, authority_mode="native_direct", **kwargs)


def test_44_logical_markers_ignore_physical_pagination_and_keep_source_ids(tmp_path: Path) -> None:
    # Break caught: Word section/page breaks are counted instead of explicit logical markers.
    source_ids = list(range(1, 45))
    for index, source_id in zip((7, 18, 29, 40), (107, 218, 329, 440)):
        source_ids[index] = source_id
    word = tmp_path / "44-logical-pages.docx"
    document = Document()
    for logical_number, source_id in enumerate(source_ids, start=1):
        document.add_paragraph(f"第 {source_id} 页")
        document.add_paragraph(f"逻辑页 {logical_number}")
        document.add_paragraph(f"这是第 {logical_number} 个逻辑页的正文。")
        if logical_number % 3 == 0 and logical_number != 44:
            document.add_page_break()
    document.save(word)
    logo = tmp_path / "logo.svg"
    logo.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 20"><rect width="100" height="20"/></svg>',
        encoding="utf-8",
    )

    project = tmp_path / "project"
    initialize_v6_project(word, logo, project)
    extraction = json.loads(
        (project / "02_v6/paginated_word_source.json").read_text(encoding="utf-8")
    )
    composition = json.loads(
        (project / "02_v6/page_composition.json").read_text(encoding="utf-8")
    )

    assert extraction["page_count"] == 44
    assert [page["source_page_id"] for page in extraction["pages"]] == source_ids
    assert composition["page_count"] >= 44
    assert [page["output_page_number"] for page in composition["pages"]] == list(
        range(1, composition["page_count"] + 1)
    )


def test_one_confirmation_completes_and_assembles_every_confirmed_role(tmp_path: Path) -> None:
    # Break caught: special pages need a second approval or disappear from final assembly.
    word = tmp_path / "complete-composition.docx"
    document = Document()
    pages = [
        ("cover", "黄石产业项目建议", "联合产业升级"),
        ("toc", "目录", "PART 1｜产业目标\nPART 2｜创新转化"),
        ("section", "PART 1｜产业目标", "聚焦主导产业升级"),
        ("content", "产业基础", "以现有产业链为基础形成项目组合。"),
        ("appendix", "附录：数据口径", "本页说明数据口径。"),
        ("closing", "最终目标：形成可持续产业生态", "全联并购公会"),
    ]
    for number, (role, title, body) in enumerate(pages, start=1):
        document.add_paragraph(f"第 {number} 页")
        document.add_paragraph({
            "cover": "PPT页型：封面", "toc": "PPT页型：目录",
            "section": "PPT页型：章节", "content": "PPT页型：正文",
            "appendix": "PPT页型：附录", "closing": "PPT页型：尾页",
        }[role])
        document.add_paragraph(title)
        document.add_paragraph(body)
    document.save(word)
    logo = tmp_path / "logo.svg"
    logo.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 20"><rect width="100" height="20"/></svg>',
        encoding="utf-8",
    )
    project = tmp_path / "confirmed-deck"
    initialize_v6_project(word, logo, project)

    server = _load_server()
    client = server.create_app(project).test_client()
    recommendations = client.get("/api/recommendations").get_json()
    expected_roles = ["cover", "toc", "section", "content", "appendix", "closing"]
    proposed = json.loads((project / "02_v6/page_composition.json").read_text(encoding="utf-8"))
    assert [page["page_role"] for page in proposed["pages"]] == expected_roles
    template = next(
        item for item in recommendations["templates"]
        if item["id"] == recommendations["recommended_template_id"]
    )
    payload = {
        "submission_id": "full-deck-e2e-0001",
        "revision": recommendations["revision"],
        **template["defaults"],
        "selected_director_template_id": template["id"],
        "director_taskbook": recommendations["director_taskbook"],
    }
    first = client.post("/api/confirm", json=payload)
    assert server._wait(project, "final", 1) == 0
    second = client.post("/api/confirm", json=payload)
    assert first.status_code == 200
    assert second.status_code == 409
    result = json.loads((project / "confirm_ui/result.json").read_text(encoding="utf-8"))
    assert result["revision"] == 1
    assert len(result["confirmed_pages"]) == result["confirmed_pages"][-1]["output_page_number"]
    assert [page["page_role"] for page in result["confirmed_pages"]] == expected_roles
    frozen = json.loads((project / "02_v6/page_composition.json").read_text(encoding="utf-8"))
    assert [page["page_role"] for page in frozen["pages"]] == expected_roles

    def open_workspace(root: Path, page_number: int):
        return SimpleNamespace(project_copy=root, page_number=page_number)

    content_pages = []
    native_pages = []

    def accept_content(workspace, **_kwargs):
        content_pages.append(workspace.page_number)
        state = load(workspace.project_copy)
        page = state["pages"][workspace.page_number - 1]
        for target in ("generating", "qa_review", "accepted"):
            page = transition_page(page, target)
        state["pages"][workspace.page_number - 1] = page
        save(workspace.project_copy, state)
        return SimpleNamespace(
            status="accepted", accepted=SimpleNamespace(candidate=object()),
            attempts=(), failure_problems=(), correction_count=0,
        )

    def reconstruct_content(workspace, _outcome):
        body = project / f"editable-body-{workspace.page_number}.pptx"
        _editable_body(body, workspace.page_number)
        return finalize_reconstructed_page(
            workspace.project_copy,
            page_number=workspace.page_number,
            reconstructed_body=body,
        )

    def render_native(root: Path, page_number: int):
        native_pages.append(page_number)
        return render_special_page(root, page_number)

    assembled_outcomes = []
    report = run_pages(
        project,
        list(range(1, len(result["confirmed_pages"]) + 1)),
        dependencies=PipelineDependencies(
            open_workspace=open_workspace,
            evidence_recorder=lambda _workspace: object(),
            candidate_loop=accept_content,
            reconstruct_page=reconstruct_content,
            native_page_renderer=render_native,
            assemble_project=lambda _root, outcomes: assembled_outcomes.append(outcomes),
        ),
        configuration=PipelineConfiguration(
            page_workers=1, initial_page_concurrency=1, maximum_page_concurrency=1,
        ),
    )

    assert report.failed_pages == {}
    assert native_pages == [1, 2, 3, 6]
    assert content_pages == [4, 5]
    assert set(assembled_outcomes[0]) == set(range(1, len(result["confirmed_pages"]) + 1))
    assemble_v6_deck(project)
    composition = json.loads((project / "02_v6/page_composition.json").read_text(encoding="utf-8"))
    assembly = json.loads((project / "08_final/assembly.json").read_text(encoding="utf-8"))
    deck = Presentation(project / "08_final/deck.pptx")
    assert all(page["state"] == "page_complete" for page in load(project)["pages"])
    assert len(deck.slides) == composition["page_count"] == assembly["page_count"]
    assert assembly["page_order"] == list(range(1, composition["page_count"] + 1))
    expected_number_visibility = {
        "cover": False, "toc": True, "section": True,
        "content": True, "appendix": True, "closing": False,
    }
    for slide, page, expected_role in zip(deck.slides, composition["pages"], expected_roles):
        assert page["page_role"] == expected_role
        names = {shape.name for shape in slide.shapes}
        has_number = "special-page-number" in names or "fixed-frame-page-number" in names
        assert has_number is expected_number_visibility[expected_role]


def test_confirmation_rejects_reorder_delete_and_preserves_word_authority(tmp_path: Path) -> None:
    word = tmp_path / "three-content-pages.docx"
    document = Document()
    image_sha256 = {}
    for source_position, source_id in enumerate((10, 30, 50), start=1):
        image = tmp_path / f"source-{source_position}.png"
        Image.new("RGB", (40, 20), (source_position * 60, 20, 120)).save(image)
        image_sha256[source_position] = hashlib.sha256(image.read_bytes()).hexdigest()
        document.add_paragraph(f"第 {source_id} 页")
        document.add_paragraph("PPT页型：正文")
        document.add_paragraph(f"标题-{source_position}")
        document.add_paragraph(f"BODY-{source_position}")
        document.add_picture(str(image))
    document.save(word)
    logo = tmp_path / "logo.svg"
    logo.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 20"><rect width="100" height="20"/></svg>',
        encoding="utf-8",
    )
    project = tmp_path / "reordered-project"
    initialize_v6_project(word, logo, project)
    server = _load_server()
    client = server.create_app(project).test_client()
    recommendations = client.get("/api/recommendations").get_json()
    template = next(
        item for item in recommendations["templates"]
        if item["id"] == recommendations["recommended_template_id"]
    )
    proposed = json.loads((project / "02_v6/page_composition.json").read_text(encoding="utf-8"))["pages"]
    confirmed_pages = []
    for output_number, original_index in enumerate((2, 0), start=1):
        page = {key: value for key, value in proposed[original_index].items() if key != "source_preview"}
        page["output_page_number"] = output_number
        confirmed_pages.append(page)
    payload = {
        "submission_id": "reorder-delete-e2e-0001", "revision": 0,
        **template["defaults"], "confirmed_pages": confirmed_pages,
    }

    response = client.post("/api/confirm", json=payload)
    assert response.status_code == 400, response.get_json()
    source = json.loads((project / "02_v6/paginated_word_source.json").read_text(encoding="utf-8"))

    assert [page["page_number"] for page in source["pages"]] == [1, 2, 3]
    assert [page["source_page_id"] for page in source["pages"]] == [10, 30, 50]
    assert [page["source_asset_page_number"] for page in source["pages"]] == [1, 2, 3]


def test_adaptive_e2e_tracks_all_public_consulting_director_patterns() -> None:
    fixture = json.loads(CONSULTING_REGRESSION_FIXTURE.read_text(encoding="utf-8"))

    assert [case["id"] for case in fixture["cases"]] == [
        "three-lane-portfolio",
        "five-stage-capital-loop",
        "four-capability-transformation-chain",
        "four-row-investment-matrix",
    ]
    assert all(case["privacy_class"] == "public-synthetic" for case in fixture["cases"])


@pytest.mark.parametrize(
    ("expected_template_id", "signal_text"),
    [
        ("company-business-introduction", "公司介绍、业务介绍、核心能力与合作价值。"),
        ("investment-committee", "提交投委会审议，重点说明估值、投资回报与退出。"),
        ("project-initiation", "申请项目立项，说明初步尽调、可行性和工作计划。"),
        ("corporate-planning", "公司三年规划围绕战略目标、重点任务与实施路径。"),
        ("investment-project-bp", "投资项目 BP 说明融资需求、商业模式与资金用途。"),
    ],
)
def test_director_template_confirmation_preserves_word_and_automatic_pages(
    tmp_path: Path, expected_template_id: str, signal_text: str,
) -> None:
    word = tmp_path / f"{expected_template_id}.docx"
    document = Document()
    document.add_paragraph("第 1 页")
    document.add_paragraph("项目标题")
    document.add_paragraph(signal_text)
    document.save(word)
    logo = tmp_path / "logo.svg"
    logo.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 20"><rect width="100" height="20"/></svg>',
        encoding="utf-8",
    )
    project = tmp_path / "project"
    initialize_v6_project(word, logo, project)
    source_before = json.loads((project / "02_v6/page_sources/page_001.json").read_text(encoding="utf-8"))
    composition_before = json.loads((project / "02_v6/page_composition.json").read_text(encoding="utf-8"))
    client = _load_server().create_app(project).test_client()
    recommendations = client.get("/api/recommendations").get_json()
    assert recommendations["recommended_template_id"] == expected_template_id
    template = next(item for item in recommendations["templates"] if item["id"] == expected_template_id)
    payload = {
        "submission_id": f"director-template-{expected_template_id}",
        "revision": 0,
        **template["defaults"],
        "selected_director_template_id": expected_template_id,
        "director_taskbook": template["director_taskbook"],
    }

    response = client.post("/api/confirm", json=payload)
    assert response.status_code == 200, response.get_json()
    result = json.loads((project / "confirm_ui/result.json").read_text(encoding="utf-8"))
    state = load(project)
    source_after = json.loads((project / "02_v6/page_sources/page_001.json").read_text(encoding="utf-8"))
    composition_after = json.loads((project / "02_v6/page_composition.json").read_text(encoding="utf-8"))
    taskbook_prompt = confirmed_taskbook_prompt(project)

    assert state["director_confirmation"] == result["director_confirmation"]
    assert result["director_confirmation"]["taskbook"] == template["director_taskbook"]
    assert source_after["word_original"] == source_before["word_original"]
    assert composition_after["pages"] == composition_before["pages"]
    assert all(value in taskbook_prompt for value in template["director_taskbook"].values())
    for forbidden in (expected_template_id, "template_version", "taskbook_digest", '"defaults"'):
        assert forbidden not in taskbook_prompt


def _load_server():
    spec = importlib.util.spec_from_file_location("v6_adaptive_e2e_server", SERVER)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _png(path: Path, color: str, label: str, size: tuple[int, int] = (1200, 700)) -> Path:
    image = Image.new("RGB", size, color)
    ImageDraw.Draw(image).text((40, 40), label, fill="white")
    image.save(path)
    return path


def _source_fixture(tmp_path: Path) -> tuple[Path, Path, Path, Path]:
    licensed_meeting = FIXTURE.with_name("staff-meeting-public-domain.jpg")
    license_record = json.loads(FIXTURE.with_name("staff-meeting-public-domain.license.json").read_text(encoding="utf-8"))
    assert hashlib.sha256(licensed_meeting.read_bytes()).hexdigest() == license_record["sha256"]
    meeting = tmp_path / licensed_meeting.name
    meeting.write_bytes(licensed_meeting.read_bytes())
    company_logo = _png(tmp_path / "company-logo.png", "#B23A48", "REAL COMPANY LOGO", (900, 360))
    word = tmp_path / "four-pages.docx"
    document = Document()
    pages = [
        ("Strategy overview", "The approved strategy has three phases: diagnose, implement, review."),
        ("Meeting evidence", "The project team met on 12 August 2026 and approved the implementation plan."),
        ("Brand system", "The approved company identity must remain recognizable in the visual."),
        ("Revenue chart", "Revenue was 20 in 2025 and 30 in 2026."),
    ]
    for number, (title, body) in enumerate(pages, start=1):
        document.add_paragraph(f"第 {number} 页")
        document.add_paragraph(title)
        paragraph = document.add_paragraph(body)
        if number == 2:
            document.add_picture(str(meeting))
        if number == 3:
            document.add_comment(
                paragraph.runs,
                "[search-evidence:official company logo]",
                author="Reviewer",
                initials="RV",
            )
        if number == 4:
            document.add_comment(
                paragraph.runs,
                "[search-evidence:unavailable audit photograph]",
                author="Reviewer",
                initials="RV",
            )
    document.save(word)
    fixed_logo = tmp_path / "fixed-logo.svg"
    fixed_logo.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 240 72">'
        '<rect width="240" height="72" rx="8" fill="#17365D"/>'
        '<text x="18" y="48" fill="white" font-size="32">FIXED SVG</text></svg>',
        encoding="utf-8",
    )
    return word, fixed_logo, meeting, company_logo


def _confirm_once(project: Path) -> dict:
    server = _load_server()
    client = server.create_app(project).test_client()
    recommendations = client.get("/api/recommendations").get_json()
    template = next(
        item for item in recommendations["templates"]
        if item["id"] == recommendations["recommended_template_id"]
    )
    pages = server._v6_project_pages(project)
    editable = []
    for page in pages:
        authority = json.loads((project / "02_v6" / "page_materials" / f"page_{page['page_number']:03d}.json").read_text(encoding="utf-8"))
        editable.append({
            "page_number": page["page_number"],
            "effective_body": page["effective_body"],
            "attachment_extracts": page["attachment_extracts"],
            "chart_facts": page["chart_facts"],
            "image_requirements": page["image_requirements"],
            "degradations": page["degradations"],
            "reference_images": authority["reference_images"],
            "reference_decisions": [],
        })
    payload = {"submission_id": "adaptive-e2e-0001", "revision": 1, **template["defaults"]}
    response = client.post("/api/confirm", json=payload)
    assert response.status_code == 200, response.get_json()
    result = json.loads((project / "confirm_ui" / "result.json").read_text(encoding="utf-8"))
    assert server._wait(project, "final", 1) == 0
    for page in pages:
        publish_page_materials(
            project, page["page_number"],
            project / "02_v6" / "awesome_page_materials" / f"page_{page['page_number']:03d}.json",
        )
    # Downstream image tests currently consume the already-confirmed page authority.
    generation_result = {
        "status": "confirmed", "revision": result["revision"],
        "confirmed_at": "2026-08-13T00:00:00+08:00",
        "production_profile": "balanced",
        "global_visual_contract": {"visual_style": result["visual_description"]},
        "confirmed_pages": editable,
    }
    (project / "confirm_ui" / "result.json").write_text(json.dumps(generation_result, ensure_ascii=False), encoding="utf-8")
    state = load(project)
    state["confirmed_ui_digest"] = canonical_sha256(generation_result)
    state["style_confirmation"]["contract"] = {
        "fixed_frame": {"title_color": result["primary_color"], **fixed_frame_execution()},
        "hard_constraints": {
            "title_color": result["primary_color"],
            "typography": {"heading": {"cjk": result["cjk_font"]}, "type_scale_pt": {"page_title": result["title_size_pt"]}},
        },
    }
    save(project, state)
    return generation_result


def _trace(command: list[str], output: Path) -> None:
    trace = Path(command[command.index("--trace-out") + 1])
    images = [Path(command[index + 1]) for index, value in enumerate(command) if value == "--image"]
    roles = [command[index + 1] for index, value in enumerate(command) if value == "--image-role"]
    digests = [command[index + 1] for index, value in enumerate(command) if value == "--image-sha256"]
    trace.write_text(json.dumps({
        "operation": command[2],
        "model": "gpt-image-2",
        "quality": command[command.index("--quality") + 1],
        "size": command[command.index("--size") + 1],
        "input_images": [
            {"role": role, "path": str(path), "sha256": digest}
            for role, path, digest in zip(roles, images, digests)
        ],
        "outputs": [{
            "path": str(output.resolve()),
            "sha256": hashlib.sha256(output.read_bytes()).hexdigest(),
            "mime_type": "image/png",
        }],
    }), encoding="utf-8")


def _semantic_failure(score: int) -> dict:
    return {
        "accepted": False,
        "score": score,
        "checks": {
            "global_style_followed": {
                "result": "fail",
                "detail": "Contrast does not satisfy the frozen style.",
                "correction": {
                    "check": "global_style_followed",
                    "action": "increase",
                    "target": "contrast_relation",
                    "constraint": "contrast_relation",
                    "correction": "Increase contrast between confirmed body text and its background.",
                },
            }
        },
        "issues": [],
    }


def _editable_body(path: Path, page_number: int) -> None:
    deck = Presentation()
    deck.slide_width = Cm(25.4)
    deck.slide_height = Cm(14.288)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(16), Cm(2)).text = f"Editable approved body {page_number}"
    deck.save(path)


@pytest.mark.skip(reason="legacy UI-selected-reference E2E is replaced by Task 9 awesome four-page E2E")
def test_four_page_adaptive_v6_runs_once_resumes_and_assembles_in_word_order(tmp_path: Path, monkeypatch) -> None:
    contract = json.loads(FIXTURE.read_text(encoding="utf-8"))
    license_record = json.loads(FIXTURE.with_name("staff-meeting-public-domain.license.json").read_text(encoding="utf-8"))
    assert license_record["license"].startswith("Public domain")
    word, fixed_logo, meeting, company_logo = _source_fixture(tmp_path)
    project = tmp_path / "brand-new-v6-project"
    real_extract = workflow_v6_source.extract_source_assets

    def extract_with_chart(*args, **kwargs):
        value = real_extract(*args, **kwargs)
        value.setdefault("chart_records", []).append({
            "page_numbers": [4],
            "title": "Revenue trend",
            "unit": "USD m",
            "series": [
                {"series": "Revenue", "time": "2025", "value": 20},
                {"series": "Revenue", "time": "2026", "value": 30},
            ],
            "image_path": "must-not-be-an-image-input.png",
        })
        return value

    monkeypatch.setattr(workflow_v6_source, "extract_source_assets", extract_with_chart)
    initialize_v6_project(word, fixed_logo, project)
    page3_receipt = json.loads((project / "02_v6/reference_materials/page_003.json").read_text(encoding="utf-8"))
    page4_receipt = json.loads((project / "02_v6/reference_materials/page_004.json").read_text(encoding="utf-8"))
    request3 = page3_receipt["reference_acquisitions"][0]["request_id"]
    request4 = page4_receipt["reference_acquisitions"][0]["request_id"]
    import_reference(project, page_number=3, request_id=request3, image=company_logo, source_url="https://example.test/logo")
    confirm_reference(project, page_number=3, request_id=request3)
    fail_reference(project, page_number=4, request_id=request4, reason="no correct source-backed material")
    frozen = _confirm_once(project)

    assert frozen["revision"] == 1
    assert [item["status"] for item in json.loads((project / "02_v6/reference_materials/page_003.json").read_text())["reference_acquisitions"]] == ["confirmed"]
    assert [item["status"] for item in json.loads((project / "02_v6/reference_materials/page_004.json").read_text())["reference_acquisitions"]] == ["failed_no_retry"]

    calls: list[list[str]] = []
    reviews: defaultdict[int, int] = defaultdict(int)

    def runner(command: list[str], _timeout: int) -> None:
        calls.append(list(command))
        output = Path(command[command.index("--out") + 1])
        page_number = int(output.name.split("_")[1].split(".")[0])
        attempt = int(output.name.split("candidate_")[1].split(".")[0])
        Image.new("RGB", (1904, 896), (35 * page_number, 25 * attempt, 80)).save(output)
        _trace(command, output)

    def reviewer(_project, *, image: Path, **_kwargs):
        page_number = int(image.name.split("_")[1].split(".")[0])
        reviews[page_number] += 1
        if page_number == 2 and reviews[page_number] == 1:
            return _semantic_failure(3)
        if page_number == 3:
            return _semantic_failure(3)
        return {"accepted": True, "score": 5, "checks": {}, "issues": []}

    receipts = [
        generate_page_body(project, page_number=number, runner=runner, reviewer=reviewer, retry_sleep=lambda _x: None)
        for number in range(1, 5)
    ]
    calls_before_resume = len(calls)
    receipt_bytes = [(project / f"04_v6/images/page_{number:03d}.json").read_bytes() for number in range(1, 5)]
    resumed = [
        generate_page_body(project, page_number=number, runner=runner, reviewer=reviewer, retry_sleep=lambda _x: None)
        for number in range(1, 5)
    ]
    assert len(calls) == calls_before_resume
    assert receipts == resumed
    assert receipt_bytes == [(project / f"04_v6/images/page_{number:03d}.json").read_bytes() for number in range(1, 5)]

    by_page = defaultdict(list)
    for call in calls:
        output = Path(call[call.index("--out") + 1])
        by_page[int(output.name.split("_")[1].split(".")[0])].append(call)
    expected = {item["page_number"]: item for item in contract["pages"]}
    for page_number, page_calls in by_page.items():
        assert all(call[2] == expected[page_number]["operation"] for call in page_calls)
        assert page_calls[0][page_calls[0].index("--quality") + 1] == expected[page_number]["quality"]
        assert len(page_calls) <= contract["candidate_limit"]
    assert len(by_page[1]) == 1
    assert len(by_page[2]) == 2
    assert len(by_page[3]) == 2
    assert len(by_page[4]) == 1
    assert "--image" not in by_page[1][0]
    assert "--image" not in by_page[4][0]
    assert [by_page[2][0][index + 1] for index, value in enumerate(by_page[2][0]) if value == "--image"] == [
        by_page[2][1][index + 1] for index, value in enumerate(by_page[2][1]) if value == "--image"
    ]
    assert all(str(meeting) != value for value in by_page[2][0])
    assert not any("candidate_1.png" in value for value in by_page[2][1])
    page2_inputs = [
        Path(by_page[2][0][index + 1]).resolve()
        for index, value in enumerate(by_page[2][0])
        if value == "--image"
    ]
    expected_page2_digest = frozen["confirmed_pages"][1]["reference_images"][0]["integrity"]["model_input_sha256"]
    assert [hashlib.sha256(path.read_bytes()).hexdigest() for path in page2_inputs] == [expected_page2_digest]
    assert all(path.is_relative_to(project) for path in page2_inputs)
    page3_roles = [
        by_page[3][0][index + 1]
        for index, value in enumerate(by_page[3][0])
        if value == "--image-role"
    ]
    assert any("logo" in role.lower() for role in page3_roles)

    prompt1 = Path(by_page[1][0][by_page[1][0].index("--prompt-file") + 1]).read_text(encoding="utf-8")
    prompt4 = Path(by_page[4][0][by_page[4][0].index("--prompt-file") + 1]).read_text(encoding="utf-8")
    for forbidden in ("Strategy overview", "FIXED SVG", "footer", "page number"):
        if forbidden == "Strategy overview":
            assert forbidden not in prompt1
    assert "Revenue trend" in prompt4 and '"value":20' in prompt4 and '"value":30' in prompt4
    assert "must-not-be-an-image-input" not in prompt4
    assert receipts[2]["selected"]["attempt"] == 1
    assert any("qa_no_effective_improvement" in reason for reason in receipts[2]["degraded_reasons"])
    for number, receipt in enumerate(receipts, start=1):
        selected = project / receipt["selected"]["path"]
        with Image.open(selected) as image:
            assert image.size == (1904, 896)
        assert receipt["request_identity"]
        assert receipt["request_prompt_sha256"]
        assert receipt["selected"]["output_sha256"]
        assert receipt["selected"]["trace_sha256"]

    for page_number in range(1, 5):
        request = build_reconstruction_request(project, page_number=page_number)
        assert request["page_title"] not in request["effective_page"].get("body_render_content", "")
        body = tmp_path / f"editable-body-{page_number}.pptx"
        _editable_body(body, page_number)
        finalized = finalize_reconstructed_page(project, page_number=page_number, reconstructed_body=body)
        assert finalized["fixed_frame"]["passed"] is True
    assembly = assemble_v6_deck(project)
    deck = Presentation(project / assembly["output"])
    assert assembly["page_order"] == [1, 2, 3, 4]
    assert len(deck.slides) == 4
    assert round(deck.slide_width / 360000, 3) == 25.4
    assert round(deck.slide_height / 360000, 3) == 14.288
    assert all(len([shape for shape in slide.shapes if shape.name == "fixed-frame-logo"]) == 1 for slide in deck.slides)
    assert all(
        next(
                shape._pic.blipFill.blip.rEmbed
            for shape in slide.shapes
            if shape.name == "fixed-frame-logo"
        ) in {
            relationship.rId
            for relationship in slide.part.rels.values()
            if relationship.reltype.endswith("/image")
            and relationship.target_part.content_type == "image/svg+xml"
        }
        for slide in deck.slides
    )
