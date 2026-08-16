from __future__ import annotations

import hashlib
import json
import os
import subprocess
import sys
import time
import zipfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

import pytest
from PIL import Image
from docx import Document
from pptx import Presentation
from pypdf import PdfWriter
import xlsxwriter


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

WORKFLOW_RUNTIME = Path.home() / ".codex/plugin-runtimes/editable-ppt-workflow-fixed-canvas-cm-v2/workflow/Scripts/python.exe"


def _sha(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _fixtures(root: Path) -> list[Path]:
    source = root / "01_source_assets"
    source.mkdir(parents=True)

    image = source / "photo.png"
    Image.new("RGB", (640, 360), "navy").save(image)

    pdf = source / "brief.pdf"
    drawing = PdfWriter()
    drawing.add_blank_page(width=640, height=360)
    drawing.add_blank_page(width=640, height=360)
    with pdf.open("wb") as stream:
        drawing.write(stream)

    docx = source / "brief.docx"
    document = Document()
    document.add_heading("DOCX PAGE", level=1)
    document.add_paragraph("attachment visual evidence")
    document.save(docx)

    pptx = source / "brief.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(pptx)

    xlsx = source / "brief.xlsx"
    workbook = xlsxwriter.Workbook(str(xlsx))
    workbook.add_worksheet().write("A1", "XLSX PAGE")
    workbook.close()
    return [image, pdf, docx, pptx, xlsx]


@pytest.mark.skipif(sys.platform != "win32", reason="Awesome attachment rendering is Windows-only")
def test_render_supported_attachments_is_lossless_ordered_and_reusable(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    from awesome_attachment_render import render_attachment

    if WORKFLOW_RUNTIME.is_file():
        monkeypatch.setattr(sys, "executable", str(WORKFLOW_RUNTIME))

    project = tmp_path / "project"
    project.mkdir()
    for attachment in _fixtures(project):
        before = attachment.read_bytes()
        receipt = render_attachment(project, attachment)
        assert attachment.read_bytes() == before
        assert receipt.original_sha256 == hashlib.sha256(before).hexdigest()
        assert receipt.renderer_identity
        assert receipt.pages
        assert [page.page_number for page in receipt.pages] == list(range(1, len(receipt.pages) + 1))
        assert all(page.width > 0 and page.height > 0 and len(page.sha256) == 64 for page in receipt.pages)
        assert len(receipt.contact_sheet.sha256) == 64
        payload = receipt.to_dict()
        forbidden = {"text", "summary", "chart_facts", "semantic_extract", "hyperlinks"}
        assert not forbidden.intersection(payload)
        assert render_attachment(project, attachment).to_dict() == payload


@pytest.mark.skipif(sys.platform != "win32", reason="Awesome attachment rendering is Windows-only")
def test_render_receipt_is_stably_reused_by_another_process(tmp_path: Path):
    project = tmp_path / "project"
    project.mkdir()
    attachment = _fixtures(project)[1]
    script = (
        "import json,sys; from pathlib import Path; "
        f"sys.path.insert(0,{str(SCRIPTS)!r}); "
        "from awesome_attachment_render import render_attachment; "
        "print(json.dumps(render_attachment(Path(sys.argv[1]),Path(sys.argv[2])).to_dict(),sort_keys=True))"
    )
    first = subprocess.run([sys.executable, "-c", script, str(project), str(attachment)], check=True, capture_output=True, text=True)
    second = subprocess.run([sys.executable, "-c", script, str(project), str(attachment)], check=True, capture_output=True, text=True)
    assert json.loads(first.stdout) == json.loads(second.stdout)


@pytest.mark.skipif(sys.platform != "win32", reason="Awesome attachment rendering is Windows-only")
def test_corrupt_and_unsupported_attachment_only_fail_owning_page(tmp_path: Path):
    from awesome_attachment_render import AttachmentRenderError, render_page_attachments
    from workflow_v6_contract import new_page, new_project
    from workflow_v6_state import load, save

    project = tmp_path / "project"
    source = project / "01_source_assets"
    source.mkdir(parents=True)
    bad = source / "bad.pdf"
    bad.write_bytes(b"not a pdf")
    unsupported = source / "payload.exe"
    unsupported.write_bytes(b"MZ")
    state = new_project(
        word_source={"path": "word.docx", "sha256": "1" * 64},
        logo_source={"path": "logo.svg", "sha256": "2" * 64},
        pages=[new_page(1, title="one"), new_page(2, title="two")],
    )
    save(project, state)

    for page_number, attachment in ((1, bad), (2, unsupported)):
        with pytest.raises(AttachmentRenderError, match=attachment.name):
            render_page_attachments(project, page_number, [attachment])
        updated = load(project)
        assert updated["pages"][page_number - 1]["state"] == "technical_failed"
        other = 2 if page_number == 1 else 1
        if other > page_number:
            assert updated["pages"][other - 1]["state"] == "prepared"
        assert attachment.name in updated["pages"][page_number - 1]["technical_failure"]["detail"]


def test_attachment_renderer_rejects_non_windows_before_writes(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    import awesome_attachment_render as renderer

    project = tmp_path / "project"
    source = project / "01_source_assets"
    source.mkdir(parents=True)
    attachment = source / "brief.pdf"
    attachment.write_bytes(b"pdf")
    monkeypatch.setattr(renderer.os, "name", "posix")
    with pytest.raises(renderer.AttachmentRenderError, match="Windows-only"):
        renderer.render_attachment(project, attachment)
    assert not (project / "02_v6" / "attachment_renders").exists()


def test_source_has_no_network_or_semantic_extraction_paths():
    source = (SCRIPTS / "awesome_attachment_render.py").read_text(encoding="utf-8")
    forbidden = ("requests", "urlopen", "http://", "https://", "extract_text", "chart_facts", "summary")
    assert all(token not in source for token in forbidden)


@pytest.mark.skipif(sys.platform != "win32", reason="Windows cache custody")
def test_same_identity_concurrent_render_converges_and_source_path_is_not_aliased(tmp_path: Path):
    from awesome_attachment_render import render_attachment

    project = tmp_path / "project"
    project.mkdir()
    first = _fixtures(project)[0]
    second = first.with_name("same.png")
    second.write_bytes(first.read_bytes())
    with ThreadPoolExecutor(max_workers=2) as pool:
        receipts = list(pool.map(lambda _: render_attachment(project, first), range(2)))
    assert receipts[0].to_dict() == receipts[1].to_dict()
    other = render_attachment(project, second)
    assert other.original_path != receipts[0].original_path
    assert other.contact_sheet.path != receipts[0].contact_sheet.path


@pytest.mark.skipif(sys.platform != "win32", reason="Windows cache custody")
def test_cache_poison_is_rejected_without_overwrite(tmp_path: Path):
    from awesome_attachment_render import AttachmentRenderError, render_attachment

    project = tmp_path / "project"
    project.mkdir()
    attachment = _fixtures(project)[0]
    receipt = render_attachment(project, attachment)
    page = project / receipt.pages[0].path
    page.write_bytes(b"poison")
    with pytest.raises(AttachmentRenderError, match="incomplete or changed"):
        render_attachment(project, attachment)
    assert page.read_bytes() == b"poison"


def test_receipt_validator_is_closed_and_rejects_nonconsecutive_pages(tmp_path: Path):
    from awesome_attachment_render import validate_attachment_receipt

    base = {
        "schema_version": "awesome-attachment-render-v1",
        "original_path": "01_source_assets/a.pdf",
        "original_sha256": "a" * 64,
        "original_byte_size": 1,
        "renderer_identity": "renderer/build;params=x",
        "pages": [
            {"page_number": 2, "path": "02_v6/attachment_renders/x/page_0002.png", "width": 10, "height": 10, "byte_size": 5, "sha256": "b" * 64}
        ],
        "contact_sheet": {"page_number": 0, "path": "02_v6/attachment_renders/x/contact_sheet.png", "width": 10, "height": 10, "byte_size": 5, "sha256": "c" * 64},
    }
    with pytest.raises(ValueError, match="consecutive"):
        validate_attachment_receipt(base)
    base["pages"][0]["page_number"] = 1
    base["summary"] = "forbidden"
    with pytest.raises(ValueError, match="fields"):
        validate_attachment_receipt(base)


@pytest.mark.skipif(sys.platform != "win32", reason="Windows renderer boundaries")
def test_pdf_preflight_rejects_oversized_page_before_rasterization(tmp_path: Path):
    from pypdf import PdfWriter
    from awesome_attachment_render import AttachmentRenderError, render_attachment

    project = tmp_path / "project"
    source = project / "01_source_assets/huge.pdf"
    source.parent.mkdir(parents=True)
    writer = PdfWriter()
    writer.add_blank_page(width=20_000, height=20_000)
    with source.open("wb") as stream:
        writer.write(stream)
    with pytest.raises(AttachmentRenderError, match="dimensions exceed"):
        render_attachment(project, source)


def test_office_timeout_kills_only_started_pid_tree(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    import awesome_attachment_render as renderer

    class HungChild:
        pid = 424242
        returncode = None
        calls = 0

        def communicate(self, timeout=None):
            self.calls += 1
            if self.calls == 1:
                raise subprocess.TimeoutExpired("office-child", timeout)
            return "", ""

        def terminate(self):
            self.returncode = -1

        def kill(self):
            self.returncode = -9

    killed: list[list[str]] = []
    monkeypatch.setattr(renderer.subprocess, "Popen", lambda *args, **kwargs: HungChild())
    monkeypatch.setattr(renderer, "_assign_kill_on_close_job", lambda child: None)
    monkeypatch.setattr(renderer, "_resume_suspended_process", lambda child: None)
    monkeypatch.setattr(
        renderer.subprocess, "run",
        lambda command, **kwargs: killed.append(command) or subprocess.CompletedProcess(command, 0, "", ""),
    )
    with pytest.raises(renderer.AttachmentRenderError, match="timed out"):
        renderer._office_to_pdf(tmp_path / "snapshot.docx", ".docx", tmp_path / "out.pdf")
    assert killed == []


@pytest.mark.skipif(sys.platform != "win32", reason="Windows junction custody")
def test_render_cache_junction_is_rejected_without_outside_publication(tmp_path: Path):
    from awesome_attachment_render import AttachmentRenderError, render_attachment

    project = tmp_path / "project"
    project.mkdir()
    source = _fixtures(project)[0]
    outside = tmp_path / "outside"
    outside.mkdir()
    cache = project / "02_v6/attachment_renders"
    cache.parent.mkdir(parents=True)
    completed = subprocess.run(
        ["powershell", "-NoProfile", "-Command", f"New-Item -ItemType Junction -Path '{cache}' -Target '{outside}' | Out-Null"],
        capture_output=True, text=True, check=False,
    )
    assert completed.returncode == 0, completed.stdout + completed.stderr
    try:
        with pytest.raises(AttachmentRenderError, match="reparse"):
            render_attachment(project, source)
        assert list(outside.iterdir()) == []
    finally:
        if cache.exists():
            cache.rmdir()


@pytest.mark.parametrize("suffix", [".docx", ".pptx", ".xlsx"])
def test_ooxml_external_relationships_are_rejected_before_office(tmp_path: Path, suffix: str):
    from awesome_attachment_render import AttachmentRenderError, render_attachment

    project = tmp_path / "project"
    source = project / f"01_source_assets/canary{suffix}"
    source.parent.mkdir(parents=True)
    with zipfile.ZipFile(source, "w") as package:
        package.writestr(
            "_rels/.rels",
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="x" Target="https://example.invalid/canary" TargetMode="External"/>'
            '</Relationships>',
        )
    with pytest.raises(AttachmentRenderError, match="external relationship"):
        render_attachment(project, source)


def test_receipt_validator_binds_identity_paths_and_lowercase_digests():
    from awesome_attachment_render import validate_attachment_receipt

    identity = "d" * 64
    base = {
        "schema_version": "awesome-attachment-render-v1",
        "original_path": "01_source_assets/a.pdf",
        "original_sha256": "a" * 64,
        "original_byte_size": 1,
        "renderer_identity": "renderer/build;params=x",
        "pages": [{"page_number": 1, "path": f"02_v6/attachment_renders/{identity}/page_0001.png", "width": 10, "height": 10, "byte_size": 5, "sha256": "b" * 64}],
        "contact_sheet": {"page_number": 0, "path": f"02_v6/attachment_renders/{identity}/contact_sheet.png", "width": 10, "height": 10, "byte_size": 5, "sha256": "c" * 64},
    }
    validate_attachment_receipt(base, expected_identity=identity, expected_source_path="01_source_assets/a.pdf", expected_source_size=1)
    bad = json.loads(json.dumps(base)); bad["original_sha256"] = "A" * 64
    with pytest.raises(ValueError, match="digest"):
        validate_attachment_receipt(bad, expected_identity=identity, expected_source_path="01_source_assets/a.pdf", expected_source_size=1)
    bad = json.loads(json.dumps(base)); bad["pages"][0]["path"] = f"02_v6/attachment_renders/{identity}/page_0002.png"
    with pytest.raises(ValueError, match="canonical"):
        validate_attachment_receipt(bad, expected_identity=identity, expected_source_path="01_source_assets/a.pdf", expected_source_size=1)


@pytest.mark.skipif(sys.platform != "win32", reason="Windows stable receipt reuse")
def test_valid_png_cache_poison_is_rejected_by_recorded_digest(tmp_path: Path):
    from PIL import Image
    from awesome_attachment_render import AttachmentRenderError, render_attachment

    project = tmp_path / "project"
    project.mkdir()
    source = _fixtures(project)[0]
    receipt = render_attachment(project, source)
    page = project / receipt.pages[0].path
    Image.new("RGB", (64, 48), "blue").save(page, "PNG")
    with pytest.raises(AttachmentRenderError, match="incomplete or changed"):
        render_attachment(project, source)


@pytest.mark.skipif(sys.platform != "win32", reason="Windows staging custody")
def test_staging_root_junction_is_rejected_and_outside_remains_empty(tmp_path: Path):
    from awesome_attachment_render import AttachmentRenderError, render_attachment

    project = tmp_path / "project"
    project.mkdir()
    source = _fixtures(project)[0]
    outside = tmp_path / "outside"
    outside.mkdir()
    staging = project / "02_v6"
    completed = subprocess.run(
        ["powershell", "-NoProfile", "-Command", f"New-Item -ItemType Junction -Path '{staging}' -Target '{outside}' | Out-Null"],
        capture_output=True, text=True, check=False,
    )
    assert completed.returncode == 0, completed.stdout + completed.stderr
    try:
        with pytest.raises(AttachmentRenderError, match="reparse"):
            render_attachment(project, source)
        assert list(outside.iterdir()) == []
    finally:
        if staging.exists():
            staging.rmdir()


def test_office_renderer_identity_binds_reported_application_version(monkeypatch: pytest.MonkeyPatch):
    import awesome_attachment_render as renderer

    monkeypatch.setattr(renderer, "_office_application_build", lambda suffix: ("16.0", "17928"))
    identity = renderer._renderer_identity(".docx")
    assert "word/version=16.0/build=17928" in identity
    assert "automation-security=3" in identity


@pytest.mark.skipif(sys.platform != "win32", reason="Windows page state transaction")
def test_attachment_entry_rejects_non_prepared_page_before_render(tmp_path: Path):
    from awesome_attachment_render import AttachmentRenderError, render_page_attachments
    from workflow_v6_contract import new_page, new_project
    from workflow_v6_state import save

    project = tmp_path / "project"
    source = project / "01_source_assets"
    source.mkdir(parents=True)
    image = source / "a.png"
    Image.new("RGB", (10, 10), "white").save(image)
    page = new_page(1, title="one")
    page["state"] = "generating"
    state = new_project(
        word_source={"path": "word.docx", "sha256": "1" * 64},
        logo_source={"path": "logo.svg", "sha256": "2" * 64}, pages=[page],
    )
    save(project, state)
    with pytest.raises(AttachmentRenderError, match="requires prepared"):
        render_page_attachments(project, 1, [image])
    assert not (project / "02_v6/attachment_renders").exists()


@pytest.mark.skipif(sys.platform != "win32", reason="Windows atomic page failure")
def test_multi_attachment_failure_publishes_no_partial_page_authority_and_retry_is_explicit(tmp_path: Path):
    from awesome_attachment_render import AttachmentRenderError, render_page_attachments
    from workflow_v6_contract import new_page, new_project
    from workflow_v6_state import load, save

    project = tmp_path / "project"
    source = project / "01_source_assets"
    source.mkdir(parents=True)
    good = source / "good.png"
    bad = source / "bad.pdf"
    Image.new("RGB", (10, 10), "white").save(good)
    bad.write_bytes(b"broken")
    save(project, new_project(
        word_source={"path": "word.docx", "sha256": "1" * 64},
        logo_source={"path": "logo.svg", "sha256": "2" * 64}, pages=[new_page(1, title="one")],
    ))
    with pytest.raises(AttachmentRenderError, match="bad.pdf"):
        render_page_attachments(project, 1, [good, bad])
    page = load(project)["pages"][0]
    assert page["state"] == "technical_failed"
    assert page["material_receipt"] is None
    with pytest.raises(AttachmentRenderError, match="requires prepared"):
        render_page_attachments(project, 1, [good])


def test_source_never_releases_staging_or_uses_path_replace_for_publication():
    source = (SCRIPTS / "awesome_attachment_render.py").read_text(encoding="utf-8")
    assert "temporary.replace(final_dir)" not in source
    assert "os.close(staging_descriptor)\n        staging_descriptor = -1\n        with mutation_lock" not in source


def test_zip_relationship_preflight_rejects_duplicate_and_ratio_bomb(tmp_path: Path):
    from awesome_attachment_render import AttachmentRenderError, _reject_external_ooxml_relationships

    duplicate = tmp_path / "duplicate.docx"
    with zipfile.ZipFile(duplicate, "w") as package:
        package.writestr("_rels/.rels", "<Relationships/>")
        package.writestr("_RELS/.RELS", "<Relationships/>")
    with pytest.raises(AttachmentRenderError, match="duplicate"):
        _reject_external_ooxml_relationships(duplicate)

    bomb = tmp_path / "bomb.docx"
    with zipfile.ZipFile(bomb, "w", compression=zipfile.ZIP_DEFLATED) as package:
        package.writestr("_rels/.rels", b"A" * (2 * 1024 * 1024))
    with pytest.raises(AttachmentRenderError, match="ratio|large"):
        _reject_external_ooxml_relationships(bomb)


def test_owned_runner_starts_suspended_and_has_no_taskkill_primary():
    source = (SCRIPTS / "awesome_attachment_render.py").read_text(encoding="utf-8")
    assert "CREATE_SUSPENDED" in source
    assert "NtResumeProcess" in source
    assert '["taskkill"' not in source


def test_receipt_rejects_more_than_renderer_page_limit():
    from awesome_attachment_render import validate_attachment_receipt

    identity = "d" * 64
    pages = [
        {"page_number": i, "path": f"02_v6/attachment_renders/{identity}/page_{i:04d}.png", "width": 1, "height": 1, "byte_size": 1, "sha256": "b" * 64}
        for i in range(1, 202)
    ]
    receipt = {
        "schema_version": "awesome-attachment-render-v1", "original_path": "01_source_assets/a.pdf",
        "original_sha256": "a" * 64, "original_byte_size": 1, "renderer_identity": "r",
        "pages": pages, "contact_sheet": {"page_number": 0, "path": f"02_v6/attachment_renders/{identity}/contact_sheet.png", "width": 1, "height": 1, "byte_size": 1, "sha256": "c" * 64},
    }
    with pytest.raises(ValueError, match="page count"):
        validate_attachment_receipt(receipt, expected_identity=identity)


@pytest.mark.skipif(sys.platform != "win32", reason="Windows publication custody")
def test_publication_verifies_from_held_handles_and_blocks_receipt_replacement(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    import awesome_attachment_render as renderer

    project = tmp_path / "project"
    project.mkdir()
    attachment = _fixtures(project)[0]
    attempts: list[str] = []

    def attack(pending: Path, _identity: str) -> None:
        receipt = pending / "receipt.json"
        replacement = tmp_path / "replacement.json"
        replacement.write_text("{}", encoding="utf-8")
        with pytest.raises(PermissionError):
            os.replace(replacement, receipt)
        attempts.append("blocked")

    monkeypatch.setattr(renderer, "_after_pending_receipt_close", attack)
    receipt = renderer.render_attachment(project, attachment)
    assert attempts == ["blocked"]
    assert json.loads((project / Path(receipt.contact_sheet.path).parent / "receipt.json").read_text(encoding="utf-8")) == receipt.to_dict()


@pytest.mark.skipif(sys.platform != "win32", reason="Windows pending directory custody")
def test_pending_directory_cannot_move_outside_after_first_write(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    import awesome_attachment_render as renderer

    project = tmp_path / "project"
    project.mkdir()
    attachment = _fixtures(project)[0]
    outside = tmp_path / "outside"
    outside.mkdir()
    attempted = False

    def attack(pending: Path, filename: str) -> None:
        nonlocal attempted
        if attempted:
            return
        attempted = True
        completed = subprocess.run(
            ["powershell", "-NoProfile", "-Command", f"Move-Item -LiteralPath '{pending}' -Destination '{outside / pending.name}' -ErrorAction Stop"],
            capture_output=True, text=True, check=False,
        )
        assert completed.returncode != 0
        assert not any(outside.iterdir())

    monkeypatch.setattr(renderer, "_after_pending_write", attack)
    renderer.render_attachment(project, attachment)
    assert attempted


@pytest.mark.skipif(sys.platform != "win32", reason="Windows crash-safe cache")
def test_uncommitted_pending_cache_is_recovered_but_committed_corruption_fails_closed(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    import awesome_attachment_render as renderer

    project = tmp_path / "project"
    project.mkdir()
    attachment = _fixtures(project)[0]
    monkeypatch.setattr(renderer, "_after_pending_receipt_close", lambda *_: (_ for _ in ()).throw(SystemExit("crash")))
    with pytest.raises(SystemExit, match="crash"):
        renderer.render_attachment(project, attachment)
    assert len(list((project / renderer.ROOT).glob(".pending-*"))) == 1

    monkeypatch.setattr(renderer, "_after_pending_receipt_close", lambda *_: None)
    receipt = renderer.render_attachment(project, attachment)
    assert not list((project / renderer.ROOT).glob(".pending-*"))
    page = project / receipt.pages[0].path
    page.write_bytes(b"corrupt")
    with pytest.raises(renderer.AttachmentRenderError, match="incomplete or changed"):
        renderer.render_attachment(project, attachment)
    assert page.read_bytes() == b"corrupt"


@pytest.mark.skipif(sys.platform != "win32", reason="Windows crash-safe cache")
@pytest.mark.parametrize("crash_after", [1, 2, 3])
def test_every_pending_write_crash_is_recovered_on_next_run(tmp_path: Path, monkeypatch: pytest.MonkeyPatch, crash_after: int):
    import awesome_attachment_render as renderer

    project = tmp_path / "project"
    project.mkdir()
    attachment = _fixtures(project)[0]
    calls = 0

    def crash(_pending: Path, _filename: str) -> None:
        nonlocal calls
        calls += 1
        if calls == crash_after:
            raise SystemExit("power loss")

    monkeypatch.setattr(renderer, "_after_pending_write", crash)
    with pytest.raises(SystemExit, match="power loss"):
        renderer.render_attachment(project, attachment)
    monkeypatch.setattr(renderer, "_after_pending_write", lambda *_: None)
    receipt = renderer.render_attachment(project, attachment)
    assert (project / receipt.contact_sheet.path).is_file()
    assert not list((project / renderer.ROOT).glob(".pending-*"))
    assert len(list((project / renderer.ROOT).glob(".orphan-*"))) == 1


@pytest.mark.skipif(sys.platform != "win32", reason="Windows exact quarantine cleanup")
def test_quarantine_replacement_is_blocked_while_exact_cleanup_handle_is_held(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    import awesome_attachment_render as renderer

    project = tmp_path / "project"
    project.mkdir()
    attachment = _fixtures(project)[0]
    monkeypatch.setattr(renderer, "_after_pending_write", lambda *_: (_ for _ in ()).throw(SystemExit("crash")))
    with pytest.raises(SystemExit):
        renderer.render_attachment(project, attachment)
    replacement = tmp_path / "replacement"
    replacement.mkdir()
    canary = replacement / "canary.txt"
    canary.write_text("must survive")
    attempts: list[bool] = []

    def attack(quarantine: Path) -> None:
        moved = tmp_path / "moved-quarantine"
        first = subprocess.run(
            ["powershell", "-NoProfile", "-Command", f"Move-Item -LiteralPath '{quarantine}' -Destination '{moved}' -ErrorAction Stop"],
            capture_output=True, text=True, check=False,
        )
        attempts.append(first.returncode != 0)

    monkeypatch.setattr(renderer, "_after_quarantine_rename", attack)
    monkeypatch.setattr(renderer, "_after_pending_write", lambda *_: None)
    renderer.render_attachment(project, attachment)
    assert attempts == [True]
    assert canary.read_text() == "must survive"


def test_office_identity_probe_uses_owned_suspended_runner(monkeypatch: pytest.MonkeyPatch):
    import awesome_attachment_render as renderer

    calls: list[tuple[list[str], int, str]] = []
    monkeypatch.setattr(
        renderer, "_run_owned_process",
        lambda command, timeout, label: calls.append((command, timeout, label)) or (0, '["16.0","18000"]\n', ""),
    )
    renderer._office_application_build.cache_clear()
    assert renderer._office_application_build(".docx") == ("16.0", "18000")
    assert len(calls) == 1
    assert calls[0][1:] == (30, "Office renderer identity probe")


@pytest.mark.skipif(sys.platform != "win32", reason="Windows Job Object runtime")
def test_owned_runner_timeout_kills_real_child_descendant_and_keeps_unrelated_sentinel(tmp_path: Path):
    import awesome_attachment_render as renderer

    def alive(pid: int) -> bool:
        completed = subprocess.run(
            ["powershell", "-NoProfile", "-Command", f"if (Get-Process -Id {pid} -ErrorAction SilentlyContinue) {{ exit 0 }} else {{ exit 1 }}"],
            capture_output=True, check=False,
        )
        return completed.returncode == 0

    sentinel = subprocess.Popen([sys.executable, "-c", "import time; time.sleep(30)"])
    child_pid = tmp_path / "child.pid"
    command = [
        sys.executable, "-c",
        "import subprocess,sys,time; p=subprocess.Popen([sys.executable,'-c','import time; time.sleep(30)']); "
        "open(sys.argv[1],'w').write(str(p.pid)); time.sleep(30)", str(child_pid),
    ]
    try:
        with pytest.raises(renderer.AttachmentRenderError, match="timed out"):
            renderer._run_owned_process(command, 1, "owned process test")
        assert sentinel.poll() is None
        descendant = int(child_pid.read_text())
        limit = time.time() + 5
        while alive(descendant) and time.time() < limit:
            time.sleep(.05)
        assert not alive(descendant)
    finally:
        sentinel.terminate()
        sentinel.wait(5)


@pytest.mark.skipif(sys.platform != "win32", reason="Windows page lease")
def test_same_page_processes_have_exclusive_render_ownership_and_loser_cannot_commit(tmp_path: Path):
    from workflow_v6_contract import new_page, new_project
    from workflow_v6_state import load, save

    project = tmp_path / "project"
    source = project / "01_source_assets"
    source.mkdir(parents=True)
    attachment = source / "a.png"
    Image.new("RGB", (10, 10), "white").save(attachment)
    save(project, new_project(
        word_source={"path": "word.docx", "sha256": "1" * 64},
        logo_source={"path": "logo.svg", "sha256": "2" * 64}, pages=[new_page(1, title="one")],
    ))
    gate = tmp_path / "gate"
    gate.mkdir()
    worker = "\n".join([
        "import sys,time",
        "from pathlib import Path",
        "sys.path.insert(0,sys.argv[1])",
        "import awesome_attachment_render as r",
        "project,attachment,gate,mode=Path(sys.argv[2]),Path(sys.argv[3]),Path(sys.argv[4]),sys.argv[5]",
        "def fake(*_):",
        " (gate/(mode+'.inside')).write_text('1')",
        " if mode=='first':",
        "  limit=time.time()+10",
        "  while not (gate/'release').exists() and time.time()<limit: time.sleep(.01)",
        "  raise r.AttachmentRenderError('first failed')",
        " raise AssertionError('second renderer must never run')",
        "r.render_attachment=fake",
        "(gate/(mode+'.started')).write_text('1')",
        "try: r.render_page_attachments(project,1,[attachment])",
        "except Exception as exc: (gate/(mode+'.result')).write_text(type(exc).__name__+':'+str(exc))",
    ])
    first = subprocess.Popen([sys.executable, "-c", worker, str(SCRIPTS), str(project), str(attachment), str(gate), "first"])
    limit = time.time() + 10
    while not (gate / "first.inside").exists() and time.time() < limit:
        time.sleep(.01)
    second = subprocess.Popen([sys.executable, "-c", worker, str(SCRIPTS), str(project), str(attachment), str(gate), "second"])
    while not (gate / "second.started").exists() and time.time() < limit:
        time.sleep(.01)
    time.sleep(.2)
    assert not (gate / "second.inside").exists()
    (gate / "release").write_text("1")
    assert first.wait(15) == 0
    assert second.wait(15) == 0
    page = load(project)["pages"][0]
    assert page["state"] == "technical_failed"
    assert page["material_state"] != "available"
    assert "requires prepared" in (gate / "second.result").read_text()


@pytest.mark.skipif(sys.platform != "win32", reason="Windows opaque page lease")
def test_page_render_worker_rejects_forged_mismatched_and_stale_leases(tmp_path: Path):
    import awesome_attachment_render as renderer
    from workflow_v6_contract import new_page, new_project
    from workflow_v6_state import save

    project = tmp_path / "project"
    source = project / "01_source_assets"
    source.mkdir(parents=True)
    image = source / "a.png"
    Image.new("RGB", (10, 10), "white").save(image)
    save(project, new_project(
        word_source={"path": "word.docx", "sha256": "1" * 64},
        logo_source={"path": "logo.svg", "sha256": "2" * 64}, pages=[new_page(1, title="one"), new_page(2, title="two")],
    ))
    with pytest.raises(renderer.AttachmentRenderError, match="lease"):
        renderer._render_page_attachments_owned(project, 1, [image], "forged")
    with renderer._page_render_lease(project, 1) as lease:
        with pytest.raises(renderer.AttachmentRenderError, match="lease"):
            renderer._render_page_attachments_owned(project, 2, [image], lease)
    with pytest.raises(renderer.AttachmentRenderError, match="lease"):
        renderer._render_page_attachments_owned(project, 1, [image], lease)


@pytest.mark.skipif(sys.platform != "win32", reason="Windows lease root custody")
def test_page_lease_root_junction_swap_cannot_write_outside(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    import awesome_attachment_render as renderer

    project = tmp_path / "project"
    (project / "02_v6").mkdir(parents=True)
    outside = tmp_path / "outside"
    outside.mkdir()

    def attack(_descriptor: int) -> None:
        lease_root = project / "02_v6/page_render_leases"
        moved = tmp_path / "moved-leases"
        completed = subprocess.run(
            ["powershell", "-NoProfile", "-Command", f"Move-Item -LiteralPath '{lease_root}' -Destination '{moved}' -ErrorAction Stop"],
            capture_output=True, text=True, check=False,
        )
        assert completed.returncode != 0

    monkeypatch.setattr(renderer, "_after_lease_root_open", attack)
    with renderer._page_render_lease(project, 1):
        pass
    assert not any(outside.iterdir())
