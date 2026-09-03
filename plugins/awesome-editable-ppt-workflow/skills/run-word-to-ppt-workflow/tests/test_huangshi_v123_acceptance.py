from __future__ import annotations

import base64
import hashlib
import json
import os
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


TESTS = Path(__file__).resolve().parent
PLUGIN = TESTS.parents[2]
REPO = TESTS.parents[4]
SCRIPTS = PLUGIN / "skills/run-word-to-ppt-workflow/scripts"
RUNTIME = PLUGIN / "skills/reconstruct-editable-slide/cli/editppt/runtime"
sys.path[:0] = [str(SCRIPTS), str(RUNTIME)]

from awesome_page_materials import publish_page_materials  # noqa: E402
from codex_subscription_runtime import CodexStructuredResult  # noqa: E402
from complex_page_experiment.director import direct_page  # noqa: E402
from complex_page_experiment.materials import build_complete_page_material_view  # noqa: E402
from complex_page_experiment.workspace import open_live_page_workspace  # noqa: E402
from director_taskbook import taskbook_digest  # noqa: E402
from extract_docx_pages import DEFAULT_MARKER, extract  # noqa: E402
from fixed_region_runtime import CONTENT_BOX, SLIDE  # noqa: E402
from workflow_v6_contract import new_page, new_project  # noqa: E402
from workflow_v6_reconstruction import (  # noqa: E402
    assemble_v6_deck,
    build_reconstruction_request,
    finalize_reconstructed_page,
)
from workflow_v6_reconstruction_worker import PageWorkerResult, reconstruct_accepted_page  # noqa: E402
from workflow_v6_state import create, load  # noqa: E402
from test_quantitative_chart_v123_e2e import _connector_endpoints  # noqa: E402
from test_workflow_v6_reconstruction import _write_signed_receipt  # noqa: E402


DESKTOP = Path.home() / "Desktop"
WORD = DESKTOP / "黄石市产业创新与母基金专业化管理合作建议_PPT生成专用Word副本_V3.docx"
LOGO = DESKTOP / "尚融logo.png"
WORD_SHA256 = "519FC2C5DAA0B4A2E65954E6FA20DF461E04587749C69AFB5952C6535A4A4A11"
LOGO_SHA256 = "9681840BACFBA51E87E47D687C1CA1F9C542F9C235577280447E96070726BCF0"
SELECTED = (5, 10, 14, 20, 21, 31, 40, 41)
OUTPUT = REPO / "tmp/v1.2.3-acceptance/huangshi"
PREVIEW_FONT = str(Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts/msyh.ttc")

PAGE_CONTRACTS = {
    5: ("geography", ("武汉", "上海", "深圳", "北京", "黄石")),
    10: ("quantitative_chart", ("420亿元", "100亿元", "450家", "100个以上", "50家", "3—5家")),
    14: ("flow", ("项目发现", "技术评价", "商业验证", "中试放大", "落地承接", "成长赋能", "并购退出")),
    20: ("hierarchy", ("1+4+N", "产业发展基金", "天使类", "科创类", "产业类", "专项类", "N只细分子基金")),
    21: ("causality", ("总规模百亿元", "总投资68.2亿元", "没有完整披露")),
    31: ("composition_architecture", ("产业研究与创新导入", "母基金管理", "项目与并购", "能源/AI及颠覆性技术", "专业服务")),
    40: ("flow", ("0—30天", "31—60天", "61—90天", "12个月")),
    41: ("analytical_table", ("资本形成", "产业投资", "科技转化", "企业成长", "绿色转型", "退出循环")),
}

STRUCTURAL_RELATIONSHIPS = {
    5: (
        (("external-centres", "武汉、上海、深圳、北京", ("word-block-000053",)), ("local-growth", "本地产业增量", ("word-block-000052",))),
        (("external-centres", "local-growth", "资源转化", ("word-block-000052", "word-block-000053")),),
    ),
    14: (
        tuple((f"stage-{index}", label, ("word-block-000125",)) for index, label in enumerate(PAGE_CONTRACTS[14][1], start=1)),
        tuple((f"stage-{index}", f"stage-{index + 1}", None, ("word-block-000125",)) for index in range(1, 7)),
    ),
    20: (
        (("mother-fund", "政府投资母基金", ("word-block-000166",)), ("angel", "天使类", ("word-block-000166",)), ("innovation", "科创类", ("word-block-000166",)), ("industry", "产业类", ("word-block-000166",)), ("special", "专项类", ("word-block-000166",)), ("subfunds", "N只细分子基金", ("word-block-000166",))),
        tuple(("mother-fund", child, None, ("word-block-000166",)) for child in ("angel", "innovation", "industry", "special", "subfunds")),
    ),
    21: (
        (("operating-evidence", "已经进入基金和项目落地阶段", ("word-block-000172", "word-block-000173", "word-block-000174", "word-block-000175", "word-block-000176")), ("disclosure-gap", "没有完整披露", ("word-block-000179",)), ("cautious-judgment", "专业化、市场化运营能力仍有进一步提升空间", ("word-block-000179",))),
        (("operating-evidence", "cautious-judgment", "支持判断", ("word-block-000173", "word-block-000174", "word-block-000175", "word-block-000176", "word-block-000179")), ("disclosure-gap", "cautious-judgment", "约束判断", ("word-block-000179",))),
    ),
    40: (
        (("days-0-30", "0—30天", ("word-block-000338",)), ("days-31-60", "31—60天", ("word-block-000338",)), ("days-61-90", "61—90天", ("word-block-000338",)), ("month-12", "12个月", ("word-block-000340", "word-block-000341", "word-block-000342", "word-block-000343", "word-block-000344", "word-block-000345", "word-block-000346"))),
        (("days-0-30", "days-31-60", None, ("word-block-000338",)), ("days-31-60", "days-61-90", None, ("word-block-000338",)), ("days-61-90", "month-12", None, ("word-block-000338", "word-block-000340"))),
    ),
}


def _box_px(left: float, top: float, width: float, height: float) -> list[int]:
    return [
        round((left - CONTENT_BOX["left"]) / CONTENT_BOX["width"] * 1904),
        round((top - CONTENT_BOX["top"]) / CONTENT_BOX["height"] * 896),
        round(width / CONTENT_BOX["width"] * 1904),
        round(height / CONTENT_BOX["height"] * 896),
    ]


def _point_px(left: float, top: float) -> list[int]:
    return _box_px(left, top, 0, 0)[:2]


def _accepted_outcome(project: Path, page_number: int) -> SimpleNamespace:
    receipt = json.loads((project / "04_v6/images" / f"page_{page_number:03d}.json").read_text(encoding="utf-8"))
    selected = receipt.get("candidate", receipt.get("selected"))
    candidate = SimpleNamespace(path=project / selected["path"], attempt=selected["attempt"])
    return SimpleNamespace(status="accepted", accepted=SimpleNamespace(candidate=candidate))


def _production_worker(manifest: dict, calls: list[dict], director_prompt: str | None = None):
    def worker(request):
        page_request = json.loads((request.page_dir / "page_request.json").read_text(encoding="utf-8"))
        accepted_request = json.loads((request.page_dir / "accepted_reconstruction_request.json").read_text(encoding="utf-8"))
        assert page_request == json.loads((request.page_dir / "page_request.json").read_text(encoding="utf-8"))
        assert page_request.get("numeric_authority") == accepted_request.get("numeric_authority")
        manifest_path = request.page_dir / "manifest.json"
        manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
        dispatch = subprocess.run(
            [sys.executable, str(RUNTIME / "record_page_dispatch.py"), str(request.run_dir), "--page", "page_001", "--agent-id", "deterministic-worker", "--prompt-file", str(request.prompt_file)],
            capture_output=True, text=True, check=False,
        )
        assert dispatch.returncode == 0, dispatch.stderr
        for command in (
            [sys.executable, str(RUNTIME / "main.py"), "page", "build", str(request.page_dir)],
            [sys.executable, str(RUNTIME / "main.py"), "page", "validate", str(request.page_dir), "--report", "validation.json"],
        ):
            completed = subprocess.run(command, capture_output=True, text=True, check=False)
            assert completed.returncode == 0, completed.stdout or completed.stderr
        validation = json.loads((request.page_dir / "validation.json").read_text(encoding="utf-8"))
        assert validation["passed"] is True
        calls.append({"page_request": page_request, "accepted_request": accepted_request, "manifest": json.loads(manifest_path.read_text(encoding="utf-8")), "prompt": request.prompt_file.read_text(encoding="utf-8")})
        return PageWorkerResult(status="completed", reconstructed_body=request.page_dir / "page.pptx")

    return worker


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest().upper()


def _page_text(page: dict) -> str:
    return "\n".join(
        "\n".join(cell for row in block["rows"] for cell in row)
        if block["type"] == "table" else block.get("text", "")
        for block in page["blocks"]
    )


def _page_title(page: dict) -> tuple[str, str]:
    for block in page["blocks"][1:]:
        text = block.get("text")
        if isinstance(text, str) and text.strip():
            return text.strip(), block["source_block_id"]
    return f"Source page {page['page_number']}", page["blocks"][0]["source_block_id"]


def _style() -> dict:
    return {
        "primary_color": "#17365D", "secondary_color": "#C7352B",
        "background_color": "#FFFFFF", "cjk_font": "Microsoft YaHei",
        "latin_font": "Arial", "title_size_pt": 28, "body_size_pt": 12,
        "caption_size_pt": 9, "regional_characteristics": "黄石产业投资",
        "visual_description": "Formal investment committee presentation.",
    }


def _director_value(
    page_number: int,
    material_view,
    grammar: str,
) -> dict:
    fact_ids = [
        str(block["source_block_id"])
        for block in material_view.value["complete_word_content"]
    ]
    relationship = STRUCTURAL_RELATIONSHIPS.get(page_number, ((), ()))
    nodes = [
        {"node_id": node_id, "label": label, "fact_ids": list(node_fact_ids)}
        for node_id, label, node_fact_ids in relationship[0]
    ]
    edges = [
        {
            "from_node": from_node,
            "to_node": to_node,
            **({"label": label} if label else {}),
            "fact_ids": list(edge_fact_ids),
        }
        for from_node, to_node, label, edge_fact_ids in relationship[1]
    ]
    return {
        "schema_version": "awesome-consulting-page-director-v3",
        "page_number": page_number,
        "quality": "high",
        "page_plan": {
            "page_purpose": "Present the complete source page through one source-bound main relationship.",
            "primary_relationship": {
                "grammar": grammar,
                "description": f"Use one {grammar} relationship grounded only in the mapped Word facts.",
                "fact_ids": fact_ids,
                "visual_instruction": f"Express the {grammar} relationship with source facts mapped to its visible parts.",
                "nodes": nodes,
                "edges": edges,
            },
            "core_exhibit": {
                "grammar": grammar,
                "description": f"One source-bound {grammar} exhibit.",
                "fact_ids": fact_ids,
            },
            "support_groups": [],
            "reading_path": "Read the complete source facts through the single main relationship.",
            "local_visuals": [],
        },
        "selected_references": [],
    }


def _director_result(value: dict) -> CodexStructuredResult:
    return CodexStructuredResult(
        value=value, thread_id="huangshi-test", turn_id="deterministic-director",
        model="deterministic-test-stub", model_provider="local-test", auth_mode="test",
        plan_type="test", usage={"input_tokens": 0, "output_tokens": 0},
        safe_trace={"boundary": "external-model-stub"}, effort="high",
        duration_seconds=0.01, startup_reused=True,
    )


def _chart_fact(page_number: int, source_page: dict, source_text: str, grammar: str) -> dict:
    if page_number == 10:
        return {
            "title": "2030年产业增加值目标",
            "relationship": "target_by_industry",
            "rendering_primitive": "column_bar",
            "chart_variant": "bar",
            "unit": "亿元",
            "basis": "2030年产业增加值目标",
            "period": "2030年",
            "source_wording": source_page["blocks"][5]["text"],
            "series": [{
                "name": "产业增加值目标",
                "categories": ["数字经济核心产业", "通用人工智能产业"],
                "values": [420, 100],
            }],
        }
    return {
        "title": source_page["blocks"][1]["text"],
        "relationship": grammar,
        "source_wording": source_text,
        "disabled_primitive": "quantitative_encoding",
        "fallback": grammar,
        "series": [],
    }


def _manifest(source_page: int, labels: tuple[str, ...], page_plan: dict | None = None) -> dict:
    if source_page == 10:
        boxes = [(0.65 + column * 3.1, 1.35 + row * 1.65, 2.75, 1.1) for row in range(2) for column in range(3)]
    elif source_page == 20:
        boxes = [(4.175, 1.0, 1.65, 0.7), (4.175, 1.9, 1.65, 0.7), *[(0.55 + index * 2.25, 2.95, 1.65, 0.7) for index in range(4)], (4.175, 4.05, 1.65, 0.7)]
    elif source_page == 21:
        boxes = [(0.8, 1.2, 4.0, 1.0), (5.2, 1.2, 4.0, 1.0), (0.8, 3.2, 8.4, 1.0)]
    elif source_page == 40:
        boxes = [*( (0.8 + index * 2.9, 1.4, 2.55, 1.0) for index in range(3)), (3.25, 3.5, 3.5, 1.0)]
    else:
        width = 8.4 / len(labels)
        boxes = [(0.8 + index * width, 1.75, width - 0.12, 1.5) for index in range(len(labels))]
    relationship = page_plan["primary_relationship"] if page_plan else {"nodes": [], "edges": []}
    if source_page == 20 and relationship["nodes"]:
        relationship_boxes = [(4.175, 1.0, 1.65, 0.7), *[(0.55 + index * 1.8, 3.0, 1.45, 0.7) for index in range(5)]]
    elif len(relationship["nodes"]) == len(boxes):
        relationship_boxes = boxes
    else:
        width = 8.4 / max(1, len(relationship["nodes"]))
        relationship_boxes = [(0.8 + index * width, 3.6, width - 0.12, 0.9) for index in range(len(relationship["nodes"]))]
    node_boxes = {
        node["node_id"]: box
        for node, box in zip(relationship["nodes"], relationship_boxes, strict=True)
    }
    connectors = []
    for edge in relationship["edges"]:
        sx, sy, sw, sh = node_boxes[edge["from_node"]]
        ex, ey, ew, eh = node_boxes[edge["to_node"]]
        edge_id = f"edge:{edge['from_node']}->{edge['to_node']}"
        connectors.append({
            "object_id": edge_id,
            "name": edge_id,
            "type": "line",
            "points_px": [*_point_px(sx + sw / 2, sy + sh / 2), *_point_px(ex + ew / 2, ey + eh / 2)],
            "stroke": "#6B7A90",
        })
    return {
        "workflow_contract_version": "fixed-canvas-cm-v2", "reconstruction_contract_version": "editable-image-v3",
        "slide": dict(SLIDE), "content_box": dict(CONTENT_BOX), "source": {"width_px": 1904, "height_px": 896},
        "text_inventory": [], "visual_inventory": [], "background_strategy": "native white body background",
        "quality_checks": {"font_size_calibrated": True, "visual_inventory_matched": True, "background_strategy_checked": True, "shape_corner_geometry_checked": True},
        "text_boxes": [
            {"object_id": f"page-{source_page}-source-label-{index}", "name": f"page-{source_page}-source-label-{index}", "box_px": _box_px(x + 0.08, y + 0.15, w - 0.16, min(0.7, h - 0.2)), "text": label, "font_size": 12, "preview_font": PREVIEW_FONT, "align": "center"}
            for index, (label, (x, y, w, h)) in enumerate(zip(labels, boxes, strict=True))
        ],
        "tables": [],
        "shapes": ([
            {"object_id": node["node_id"], "name": node["node_id"], "type": "rect", "box_px": _box_px(*node_boxes[node["node_id"]]), "fill": "#EAF2F8", "stroke": "#6B7A90"}
            for node in relationship["nodes"]
        ] if relationship["nodes"] else [
            {"object_id": f"page-{source_page}-source-node-{index}", "name": f"page-{source_page}-source-node-{index}", "type": "rect", "box_px": _box_px(x, y, w, h), "fill": "#EAF2F8", "stroke": "#6B7A90"}
            for index, (x, y, w, h) in enumerate(boxes)
        ]) + connectors,
        "images": [], "charts": [], "asset_provenance": [],
    }


def _build_project(root: Path, source: dict, logo_svg: Path) -> Path:
    project = root / "project"
    (project / "00_source").mkdir(parents=True)
    shutil.copy2(WORD, project / "00_source/source.docx")
    shutil.copy2(logo_svg, project / "00_source/logo.svg")
    pages = []
    for source_page in source["pages"]:
        page_number = source_page["page_number"]
        title, _title_block = _page_title(source_page)
        page = new_page(page_number, title=title)
        page["state"] = "accepted"
        image = project / "04_v6/images" / f"page_{page_number:03d}.png"
        image.parent.mkdir(parents=True, exist_ok=True)
        Image.new("RGB", (1904, 896), "white").save(image)  # external Image2 boundary stub
        digest = hashlib.sha256(image.read_bytes()).hexdigest()
        candidate = {"path": image.relative_to(project).as_posix(), "attempt": 1, "operation": "generate"}
        page["first_candidate"] = dict(candidate)
        page["selected_candidate"] = dict(candidate)
        pages.append(page)

    state = new_project(word_source={"path": "00_source/source.docx"}, logo_source={"path": "00_source/logo.svg"}, pages=pages)
    state["style_confirmation"] = {"status": "confirmed", "contract": _style()}
    state["confirmed_ui_revision"] = 1
    state["confirmed_ui_digest"] = "a" * 64
    state["page_materials_status"] = "pending"
    taskbook = {
        "use_scenario": "黄石市产业创新与母基金专业化管理建议汇报", "presenter": "尚融财富项目团队",
        "primary_audience": "政府与基金决策者", "audience_prior_knowledge": "熟悉黄石产业与基金基础",
        "desired_outcome": "评估合作建议与实施路径", "emphasis": "母基金、实施路径、前90天", "deemphasis": "背景性重复介绍",
    }
    state["director_confirmation"] = {
        "template_id": "investment-committee", "template_version": "1.0",
        "taskbook": taskbook, "taskbook_digest": taskbook_digest(taskbook),
    }
    create(project, state)
    for page_number in range(1, source["page_count"] + 1):
        _write_signed_receipt(project, page_number)

    paginated = {"pages": []}
    composition = {"artifact_version": "page-composition-v1", "page_count": source["page_count"], "warnings": [], "pages": []}
    for source_page in source["pages"]:
        page_number = source_page["page_number"]
        blocks = [dict(block) for block in source_page["blocks"]]
        title, title_block = _page_title(source_page)
        paginated["pages"].append({**source_page, "fixed_page_title": title, "fixed_page_title_source_block_id": title_block})
        composition["pages"].append({"output_page_number": page_number, "source_page_id": page_number, "page_role": "content", "role_source": "explicit", "chapter_title": "", "fixed_page_title": title, "source_page_number": page_number, "material_source_block_ids": [block["source_block_id"] for block in blocks], "visible_page_number": True})
        effective = project / "02_v6/effective_pages" / f"page_{page_number:03d}.json"
        effective.parent.mkdir(parents=True, exist_ok=True)
        effective.write_text(json.dumps({"page_number": page_number, "word_original": _page_text(source_page)}, ensure_ascii=False), encoding="utf-8")
    (project / "02_v6/paginated_word_source.json").write_text(json.dumps(paginated, ensure_ascii=False), encoding="utf-8")
    (project / "02_v6/page_composition.json").write_text(json.dumps(composition, ensure_ascii=False), encoding="utf-8")
    (project / "02_v6/source_assets.json").write_text('{"assets":[]}', encoding="utf-8")
    for page_number in range(1, source["page_count"] + 1):
        publish_page_materials(project, page_number, project / "02_v6/awesome_page_materials" / f"page_{page_number:03d}.json")
    return project


def _preview_has_ink(image: Image.Image, shape, slide_width: int, slide_height: int) -> bool:
    left = max(0, round(shape.left / slide_width * image.width))
    top = max(0, round(shape.top / slide_height * image.height))
    right = min(image.width, round((shape.left + shape.width) / slide_width * image.width))
    bottom = min(image.height, round((shape.top + shape.height) / slide_height * image.height))
    crop = image.crop((left, top, right, bottom)).convert("RGB")
    return any(max(pixel) - min(pixel) > 12 or sum(pixel) < 690 for pixel in crop.get_flattened_data())


def _assert_relationship_shapes(slide, page_plan: dict) -> None:
    relationship = page_plan["primary_relationship"]
    named = {shape.name: shape for shape in slide.shapes}
    for node in relationship["nodes"]:
        assert sum(shape.name == node["node_id"] for shape in slide.shapes) == 1

    def inside(point: tuple[int, int], node) -> bool:
        x, y = point
        return node.left - 1 <= x <= node.left + node.width + 1 and node.top - 1 <= y <= node.top + node.height + 1

    for edge in relationship["edges"]:
        edge_name = f"edge:{edge['from_node']}->{edge['to_node']}"
        connectors = [shape for shape in slide.shapes if shape.name == edge_name]
        assert len(connectors) == 1
        connector = connectors[0]
        assert connector.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        assert connector._element.spPr.prstGeom.get("prst") == "line"
        x1, y1, x2, y2 = _connector_endpoints(connector)
        assert inside((x1, y1), named[edge["from_node"]])
        assert inside((x2, y2), named[edge["to_node"]])


def test_huangshi_controlled_acceptance_runs_real_production_path_without_ui(tmp_path: Path) -> None:
    if not WORD.is_file() or not LOGO.is_file():
        pytest.skip("user-supplied Huangshi acceptance files are not present")
    assert _sha256(WORD) == WORD_SHA256
    assert _sha256(LOGO) == LOGO_SHA256
    source = extract(WORD, DEFAULT_MARKER)
    assert source["page_count"] == 42

    OUTPUT.mkdir(parents=True, exist_ok=True)
    svg = OUTPUT / "shangrong-logo-test-wrapper.svg"
    encoded_logo = base64.b64encode(LOGO.read_bytes()).decode("ascii")
    svg.write_text(f'<svg xmlns="http://www.w3.org/2000/svg" width="220" height="46"><image width="220" height="46" href="data:image/png;base64,{encoded_logo}"/></svg>', encoding="utf-8")
    project = _build_project(tmp_path, source, svg)
    source_pages = {page["page_number"]: page for page in source["pages"]}
    composition = json.loads((project / "02_v6/page_composition.json").read_text(encoding="utf-8"))
    page_mapping = {page["output_page_number"]: page["source_page_number"] for page in composition["pages"]}
    assert {page_number: page_mapping[page_number] for page_number in SELECTED} == {
        page_number: page_number for page_number in SELECTED
    }
    director_prompts = []
    director_artifacts = {}
    worker_calls: list[dict] = []

    for page_number, source_page in source_pages.items():
        director_prompt = None
        if page_number in SELECTED:
            grammar, labels = PAGE_CONTRACTS[page_number]
            source_text = _page_text(source_page)
            assert all(label in source_text for label in labels)
            workspace = open_live_page_workspace(project, page_number)
            material_view = build_complete_page_material_view(workspace)
            _title, title_source_id = _page_title(source_page)
            source_ids = [
                block["source_block_id"] for block in source_page["blocks"]
                if block["source_block_id"] != title_source_id
            ]
            assert material_view.value["page_number"] == page_number
            assert [block["source_block_id"] for block in material_view.value["complete_word_content"]] == source_ids
            value = _director_value(page_number, material_view, grammar)
            artifact = direct_page(workspace, material_view, timeout=30, invoke=lambda *_args, v=value, **_kwargs: _director_result(v))
            plan = artifact.page_plan
            assert plan["primary_relationship"]["grammar"] == grammar
            assert plan["core_exhibit"]["grammar"] == grammar
            assert plan["primary_relationship"]["fact_ids"] == source_ids
            assert plan["core_exhibit"]["fact_ids"] == source_ids
            assert not plan["support_groups"]
            assert artifact.actual_prompt.count("## ") == 6
            assert grammar in artifact.actual_prompt
            assert all(source_id in artifact.actual_prompt for source_id in source_ids)
            receipt_path = project / "04_v6/images" / f"page_{page_number:03d}.json"
            receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
            receipt["page_plan"] = plan
            _write_signed_receipt(project, page_number, receipt)
            director_prompts.append(artifact.actual_prompt)
            director_artifacts[page_number] = artifact
            director_prompt = artifact.actual_prompt
            materials = project / "02_v6/page_materials" / f"page_{page_number:03d}.json"
            materials.parent.mkdir(parents=True, exist_ok=True)
            chart_fact = _chart_fact(page_number, source_page, source_text, grammar)
            materials.write_text(json.dumps({"chart_facts": [chart_fact]}, ensure_ascii=False), encoding="utf-8")
            reconstruction_request = build_reconstruction_request(project, page_number=page_number)
            if page_number == 10:
                assert reconstruction_request["numeric_authority"] == chart_fact
            else:
                assert "numeric_authority" not in reconstruction_request
            assert reconstruction_request["page_plan"] == plan
            manifest = _manifest(page_number, labels, plan)
        else:
            manifest = _manifest(page_number, (source_page["blocks"][2].get("text", "Source page"),))
        sealed_request = build_reconstruction_request(project, page_number=page_number)
        before = len(worker_calls)
        reconstruct_accepted_page(
            SimpleNamespace(project_copy=project, page_number=page_number),
            _accepted_outcome(project, page_number),
            page_worker=_production_worker(manifest, worker_calls, director_prompt),
        )
        assert len(worker_calls) == before + 1
        assert worker_calls[-1]["manifest"] == manifest
        assert worker_calls[-1]["accepted_request"] == sealed_request

    assembly = assemble_v6_deck(project)
    assert assembly["page_order"] == list(range(1, 43))
    deck_path = OUTPUT / "huangshi-full-42-pages-v1.2.3.pptx"
    shutil.copy2(project / assembly["output"], deck_path)
    deck = Presentation(deck_path)
    assert len(deck.slides) == 42
    assert set(director_artifacts) == set(SELECTED)
    assert {artifact.page_plan["primary_relationship"]["grammar"] for artifact in director_artifacts.values()} == {
        "analytical_table", "flow", "hierarchy", "geography", "causality", "quantitative_chart", "composition_architecture",
    }
    assert all(not any(shape.has_chart for shape in slide.shapes) for slide in deck.slides)
    for source_number in SELECTED:
        slide = deck.slides[source_number - 1]
        _grammar, labels = PAGE_CONTRACTS[source_number]
        by_name = {shape.name: shape for shape in slide.shapes}
        names = set(by_name)
        assert "fixed-frame-logo" in names
        assert not any(term in name.casefold() for name in names for term in ("axis", "gantt", "mekko", "target line", "difference arrow"))
        slide_text = "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
        assert all(label in slide_text for label in labels)
        assert by_name["fixed-frame-title"].text == source_pages[source_number]["blocks"][1]["text"]
        assert by_name["fixed-frame-page-number"].text == str(source_number)
        relationship = director_artifacts[source_number].page_plan["primary_relationship"]
        _assert_relationship_shapes(slide, director_artifacts[source_number].page_plan)
        label_shapes = [shape for name, shape in by_name.items() if name.startswith(f"page-{source_number}-source-label-")]
        nodes = (
            [by_name[node["node_id"]] for node in relationship["nodes"]]
            if relationship["nodes"] else
            [shape for name, shape in by_name.items() if name.startswith(f"page-{source_number}-source-node-")]
        )
        assert len(label_shapes) == len(labels)
        assert all(shape.has_text_frame for shape in label_shapes)
        assert len(nodes) == len(relationship["nodes"] or label_shapes)
        if source_number == 14:
            assert len(nodes) == 7
            assert len({shape.width for shape in nodes}) == len({shape.top for shape in nodes}) == 1
        if source_number == 20:
            sealed_edges = {(edge["from_node"], edge["to_node"]) for edge in relationship["edges"]}
            assert ("mother-fund", "subfunds") in sealed_edges
            assert ("innovation", "subfunds") not in sealed_edges
            assert len({shape.top for shape in nodes}) == 2
        if source_number == 21:
            assert nodes[0].width == nodes[1].width
            assert "没有完整披露" in slide_text
            assert all(token not in slide_text for token in ("差额", "占比", "=", "+"))
        if source_number == 40:
            assert len({shape.width for shape in nodes[:3]}) == len({shape.height for shape in nodes[:3]}) == len({shape.top for shape in nodes[:3]}) == 1
            assert nodes[3].top > nodes[0].top

    with zipfile.ZipFile(deck_path) as package:
        embedded_svg = [package.read(name) for name in package.namelist() if name.startswith("ppt/media/") and name.endswith(".svg")]
        slide_xml = b"\n".join(package.read(name) for name in package.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
        assert not any(name.startswith("ppt/charts/") for name in package.namelist())
        assert b"<c:chart" not in slide_xml and b"<c:plotArea" not in slide_xml
        assert all(term not in slide_xml.lower() for term in (b"axis", b"area", b"bubble", b"target line", b"difference arrow"))
    assert svg.read_bytes() in embedded_svg
    assert encoded_logo.encode("ascii") in svg.read_bytes()
    assert all("source facts mapped to its visible parts" in prompt and "numeric axes" in prompt for prompt in director_prompts)

    findings = {
        "covered_grammars": {str(page): PAGE_CONTRACTS[page][0] for page in SELECTED},
        "unsupported_quantitative_primitives": ["line", "scatter", "bubble", "waterfall", "true_mekko"],
        "reason": "the selected source facts support one bar comparison on page 10, but not the other quantitative primitives",
        "production_path": ["extract_docx_pages.extract", "build_complete_page_material_view", "direct_page", "build_reconstruction_request", "write_pptx", "finalize_reconstructed_page", "assemble_v6_deck"],
        "optional_preview": "set EDITPPT_OPTIONAL_OFFICECLI_VALIDATION=1 and run the optional assembled-preview test",
        "remaining_limitations": ["external Image2 and director model calls are deterministic boundary stubs; this test proves contracts and production plumbing, not generated aesthetics"],
    }
    (OUTPUT / "acceptance-findings.json").write_text(json.dumps(findings, ensure_ascii=False, indent=2), encoding="utf-8")
    assert all(page["state"] == "page_complete" for page in load(project)["pages"])


def test_huangshi_optional_assembled_powerpoint_preview(tmp_path: Path) -> None:
    if os.getenv("EDITPPT_OPTIONAL_OFFICECLI_VALIDATION") != "1":
        pytest.skip("set EDITPPT_OPTIONAL_OFFICECLI_VALIDATION=1 to enable assembled PowerPoint previews")
    if not WORD.is_file() or not LOGO.is_file():
        pytest.skip("user-supplied Huangshi acceptance files are not present")
    from awesome_attachment_render import _office_to_pdf, _render_pdf
    from doctor import powerpoint_status

    status = powerpoint_status()
    if not status["available"]:
        pytest.skip(f"assembled PowerPoint preview unavailable: {status['detail']}")
    deck_path = OUTPUT / "huangshi-full-42-pages-v1.2.3.pptx"
    if not deck_path.is_file():
        pytest.skip("run the core Huangshi assembled acceptance first")
    rendered_pdf = tmp_path / "huangshi-assembled.pdf"
    rendered_pages = tmp_path / "assembled-previews"
    rendered_pages.mkdir()
    _office_to_pdf(deck_path, ".pptx", rendered_pdf)
    preview_paths = _render_pdf(rendered_pdf, rendered_pages)
    deck = Presentation(deck_path)
    assert len(preview_paths) == len(deck.slides) == 42
    for source_number in SELECTED:
        slide = deck.slides[source_number - 1]
        preview = Image.open(preview_paths[source_number - 1]).convert("RGB")
        preview_path = OUTPUT / "previews" / f"page-{source_number:02d}.png"
        preview_path.parent.mkdir(parents=True, exist_ok=True)
        preview.save(preview_path)
        by_name = {shape.name: shape for shape in slide.shapes}
        assert _preview_has_ink(preview, by_name["fixed-frame-title"], deck.slide_width, deck.slide_height)
        assert _preview_has_ink(preview, by_name["fixed-frame-logo"], deck.slide_width, deck.slide_height)
        assert _preview_has_ink(preview, by_name["fixed-frame-page-number"], deck.slide_width, deck.slide_height)
        label_shapes = [shape for name, shape in by_name.items() if name.startswith(f"page-{source_number}-source-label-")]
        assert all(_preview_has_ink(preview, shape, deck.slide_width, deck.slide_height) for shape in label_shapes)
