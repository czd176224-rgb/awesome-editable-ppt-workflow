from __future__ import annotations

import json
import os
import subprocess
import sys
from pathlib import Path

import pytest
from pptx import Presentation


SCRIPTS = Path(__file__).resolve().parents[1] / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))


from workflow_v6_contract import new_page, new_project
from workflow_v6_state import create, load, save


ROLE_TEXT = {
    "cover": ["黄石项目建议", "联合产业升级", "全联并购公会"],
    "toc": ["目录", "PART 1｜产业目标", "PART 2｜创新转化"],
    "section": ["PART 1｜产业目标", "产业目标", "聚焦主导产业升级"],
    "closing": ["最终目标：形成可持续产业生态", "全联并购公会"],
}


def _extend_materials(project: Path, total: int) -> list[str]:
    source_path = project / "02_v6/paginated_word_source.json"
    source = json.loads(source_path.read_text(encoding="utf-8"))
    blocks = source["pages"][0]["blocks"]
    for index in range(len(blocks) + 1, total + 1):
        blocks.append({
            "type": "paragraph",
            "text": f"补充材料{index}",
            "source_block_id": f"block-{index}",
        })
    source_path.write_text(json.dumps(source, ensure_ascii=False), encoding="utf-8")
    ids = [f"block-{index}" for index in range(1, total + 1)]
    composition_path = project / "02_v6/page_composition.json"
    composition = json.loads(composition_path.read_text(encoding="utf-8"))
    composition["pages"][0]["material_source_block_ids"] = ids
    composition_path.write_text(json.dumps(composition, ensure_ascii=False), encoding="utf-8")
    return ids


def prepared_confirmed_project(tmp_path: Path, *, role: str, visible_page_number: bool) -> Path:
    root = tmp_path / role
    (root / "00_source").mkdir(parents=True)
    (root / "00_source" / "logo.svg").write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 20"><rect width="100" height="20"/></svg>',
        encoding="utf-8",
    )
    page = new_page(1, title=ROLE_TEXT[role][0])
    state = new_project(
        word_source={"path": "00_source/source.docx"},
        logo_source={"path": "00_source/logo.svg"},
        pages=[page],
    )
    state["style_confirmation"] = {
        "status": "confirmed",
        "contract": {
            "primary_color": "#17365D",
            "secondary_color": "#C7352B",
            "background_color": "#F1F2F3",
            "cjk_font": "Microsoft YaHei",
            "latin_font": "Arial",
            "title_size_pt": 28,
            "body_size_pt": 16,
            "caption_size_pt": 10,
        },
    }
    state["confirmed_ui_revision"] = 1
    state["confirmed_ui_digest"] = "a" * 64
    state["page_materials_status"] = "pending"
    create(root, state)

    blocks = [
        {"type": "paragraph", "text": text, "source_block_id": f"block-{index}"}
        for index, text in enumerate(ROLE_TEXT[role], start=1)
    ]
    source_dir = root / "02_v6"
    source_dir.mkdir(parents=True)
    (source_dir / "paginated_word_source.json").write_text(
        json.dumps({"pages": [{"page_number": 1, "blocks": blocks}]}, ensure_ascii=False),
        encoding="utf-8",
    )
    composition_page = {
        "output_page_number": 1,
        "source_page_id": 1,
        "page_role": role,
        "role_source": "explicit",
        "chapter_title": "产业目标" if role == "section" else "",
        "fixed_page_title": ROLE_TEXT[role][0],
        "source_page_number": 1,
        "material_source_block_ids": [block["source_block_id"] for block in blocks],
        "visible_page_number": visible_page_number,
    }
    (source_dir / "page_composition.json").write_text(
        json.dumps({
            "artifact_version": "page-composition-v1",
            "page_count": 1,
            "pages": [composition_page],
            "warnings": [],
        }, ensure_ascii=False),
        encoding="utf-8",
    )
    return root


def _reordered_toc_project(tmp_path: Path) -> Path:
    root = tmp_path / "reordered"
    (root / "00_source").mkdir(parents=True)
    (root / "00_source" / "logo.svg").write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 20"><rect width="100" height="20"/></svg>',
        encoding="utf-8",
    )
    state = new_project(
        word_source={"path": "00_source/source.docx"},
        logo_source={"path": "00_source/logo.svg"},
        pages=[new_page(1, title="Body"), new_page(2, title="确认目录")],
    )
    state["style_confirmation"] = {"status": "confirmed", "contract": {
        "primary_color": "#17365D", "secondary_color": "#C7352B",
        "background_color": "#FFFFFF", "cjk_font": "Microsoft YaHei",
        "title_size_pt": 28, "body_size_pt": 16,
    }}
    state["confirmed_ui_revision"] = 1
    state["confirmed_ui_digest"] = "a" * 64
    state["page_materials_status"] = "pending"
    create(root, state)
    source = root / "02_v6"
    source.mkdir()
    source.joinpath("paginated_word_source.json").write_text(json.dumps({"pages": [{
        "page_number": 9, "blocks": [
            {"type": "paragraph", "text": "旧目录", "source_block_id": "toc-title"},
            {"type": "paragraph", "text": "PART 1｜产业目标", "source_block_id": "toc-1"},
        ],
    }]}, ensure_ascii=False), encoding="utf-8")
    source.joinpath("page_composition.json").write_text(json.dumps({
        "artifact_version": "page-composition-v1", "page_count": 2, "warnings": [], "pages": [
            {"output_page_number": 1, "source_page_id": 1, "page_role": "content", "role_source": "explicit", "chapter_title": "", "fixed_page_title": "Body", "source_page_number": 1, "material_source_block_ids": ["body"], "visible_page_number": True},
            {"output_page_number": 2, "source_page_id": 9, "page_role": "toc", "role_source": "explicit", "chapter_title": "", "fixed_page_title": "确认目录", "source_page_number": 9, "material_source_block_ids": ["toc-title", "toc-1"], "visible_page_number": True},
        ],
    }, ensure_ascii=False), encoding="utf-8")
    return root


@pytest.mark.parametrize("role,visible", [
    ("cover", False), ("toc", True), ("section", True), ("closing", False),
])
def test_special_page_is_editable_and_obeys_page_number_policy(
    tmp_path: Path, role: str, visible: bool,
) -> None:
    # Break caught: a special role falls back to an image or uses the wrong page-number policy.
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role=role, visible_page_number=visible)

    receipt = render_special_page(project, 1)
    deck = Presentation(project / receipt["page_pptx"])
    names = [shape.name for shape in deck.slides[0].shapes]

    assert len(deck.slides) == 1
    assert any(shape.has_text_frame for shape in deck.slides[0].shapes)
    assert ("special-page-number" in names) is visible
    assert receipt["page_role"] == role
    assert load(project)["pages"][0]["state"] == "page_complete"
    assert str(deck.slides[0].background.fill.fore_color.rgb) == "F1F2F3"


@pytest.mark.parametrize("role,visible", [
    ("cover", False), ("toc", True), ("section", True), ("closing", False),
])
def test_special_page_displayed_text_is_source_traced(
    tmp_path: Path, role: str, visible: bool,
) -> None:
    # Break caught: the renderer invents a contact detail, institution claim, or other display copy.
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role=role, visible_page_number=visible)
    receipt = render_special_page(project, 1)
    slide = Presentation(project / receipt["page_pptx"]).slides[0]
    allowed = "\n".join(ROLE_TEXT[role] + (["1"] if visible else []))

    displayed = [shape.text.strip() for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
    assert displayed
    assert all(text in allowed for text in displayed)
    assert set(displayed) == set(ROLE_TEXT[role] + (["1"] if visible else []))


@pytest.mark.parametrize("role", ["cover", "section", "closing"])
def test_special_page_displays_every_frozen_source_block_once(
    tmp_path: Path, role: str,
) -> None:
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(
        tmp_path, role=role, visible_page_number=role == "section",
    )
    source_path = project / "02_v6/paginated_word_source.json"
    source = json.loads(source_path.read_text(encoding="utf-8"))
    first_extra = len(source["pages"][0]["blocks"]) + 1
    extra = [
        {"type": "paragraph", "text": f"补充材料{index}", "source_block_id": f"block-{index}"}
        for index in range(first_extra, 7)
    ]
    source["pages"][0]["blocks"].extend(extra)
    source_path.write_text(json.dumps(source, ensure_ascii=False), encoding="utf-8")
    composition_path = project / "02_v6/page_composition.json"
    composition = json.loads(composition_path.read_text(encoding="utf-8"))
    composition["pages"][0]["material_source_block_ids"] = [f"block-{index}" for index in range(1, 7)]
    composition_path.write_text(json.dumps(composition, ensure_ascii=False), encoding="utf-8")

    receipt = render_special_page(project, 1)
    slide = Presentation(project / receipt["page_pptx"]).slides[0]
    source_names = [shape.name for shape in slide.shapes if shape.name.startswith("special-source-")]

    assert receipt["displayed_source_block_ids"] == [f"block-{index}" for index in range(1, 7)]
    assert source_names == [f"special-source-block-{index}" for index in range(1, 7)]


@pytest.mark.parametrize("role", ["cover", "section", "closing"])
def test_special_page_overflow_fails_with_exact_source_block_id(
    tmp_path: Path, role: str,
) -> None:
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(
        tmp_path, role=role, visible_page_number=role == "section",
    )
    _extend_materials(project, 7)

    with pytest.raises(ValueError, match=rf"{role}.*block-7"):
        render_special_page(project, 1)
    assert not (project / "06_v6/pages/page_001/page.pptx").exists()


@pytest.mark.parametrize("role", ["cover", "section", "closing"])
def test_special_page_text_density_overflow_names_the_exact_block(
    tmp_path: Path, role: str,
) -> None:
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(
        tmp_path, role=role, visible_page_number=role == "section",
    )
    path = project / "02_v6/paginated_word_source.json"
    source = json.loads(path.read_text(encoding="utf-8"))
    overflow_id = source["pages"][0]["blocks"][-1]["source_block_id"]
    source["pages"][0]["blocks"][-1]["text"] = "超长材料" * 160
    path.write_text(json.dumps(source, ensure_ascii=False), encoding="utf-8")

    with pytest.raises(ValueError, match=rf"{role}.*{overflow_id}"):
        render_special_page(project, 1)


def test_toc_text_density_overflow_names_the_exact_entry(tmp_path: Path) -> None:
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role="toc", visible_page_number=True)
    path = project / "02_v6/paginated_word_source.json"
    source = json.loads(path.read_text(encoding="utf-8"))
    source["pages"][0]["blocks"][1]["text"] = "超长目录项" * 300
    path.write_text(json.dumps(source, ensure_ascii=False), encoding="utf-8")

    with pytest.raises(ValueError, match=r"toc.*block-2"):
        render_special_page(project, 1)


@pytest.mark.parametrize("role,visible", [
    ("cover", False), ("toc", True), ("section", True), ("closing", False),
])
def test_special_page_uses_readable_fonts_and_safe_source_bounds(
    tmp_path: Path, role: str, visible: bool,
) -> None:
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role=role, visible_page_number=visible)
    receipt = render_special_page(project, 1)
    deck = Presentation(project / receipt["page_pptx"])
    source_shapes = [
        shape for shape in deck.slides[0].shapes if shape.name.startswith("special-source-")
    ]
    font_sizes = [
        run.font.size.pt
        for shape in source_shapes
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    ]

    assert min(font_sizes) >= 12
    assert max(font_sizes) >= 24
    assert all(
        shape.left >= 0 and shape.top >= 0
        and shape.left + shape.width <= deck.slide_width
        and shape.top + shape.height <= deck.slide_height
        for shape in source_shapes
    )


@pytest.mark.parametrize("entry_count,expected_variant,expected_columns", [
    (4, "toc-one-column", 1),
    (8, "toc-two-column", 2),
])
def test_toc_layout_switches_columns_by_entry_count(
    tmp_path: Path, entry_count: int, expected_variant: str, expected_columns: int,
) -> None:
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role="toc", visible_page_number=True)
    _extend_materials(project, entry_count + 1)

    receipt = render_special_page(project, 1)
    deck = Presentation(project / receipt["page_pptx"])
    entries = [
        shape for shape in deck.slides[0].shapes
        if shape.name.startswith("special-source-") and shape.name != "special-source-block-1"
    ]

    assert receipt["layout_variant"] == expected_variant
    assert len({shape.left for shape in entries}) == expected_columns


def test_toc_continuation_displays_title_without_duplicating_title_source_id(tmp_path: Path) -> None:
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role="toc", visible_page_number=True)
    ids = _extend_materials(project, 9)[1:]
    composition_path = project / "02_v6/page_composition.json"
    composition = json.loads(composition_path.read_text(encoding="utf-8"))
    composition["pages"][0].update({
        "source_page_id": None,
        "source_page_number": None,
        "role_source": "synthesized",
        "composition_page_id": "toc-continuation:block-1:2",
        "material_source_block_ids": ids,
    })
    composition_path.write_text(json.dumps(composition, ensure_ascii=False), encoding="utf-8")

    receipt = render_special_page(project, 1)
    slide = Presentation(project / receipt["page_pptx"]).slides[0]
    names = [shape.name for shape in slide.shapes]

    assert receipt["displayed_source_block_ids"] == ids
    assert receipt["layout_variant"] == "toc-continuation-two-column"
    assert "special-toc-title" in names
    assert "special-source-block-1" not in names


@pytest.mark.parametrize("role,total", [
    ("cover", 6), ("toc", 13), ("section", 6), ("closing", 6),
])
def test_special_page_source_shapes_do_not_overlap(
    tmp_path: Path, role: str, total: int,
) -> None:
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(
        tmp_path, role=role, visible_page_number=role in {"toc", "section"},
    )
    _extend_materials(project, total)
    receipt = render_special_page(project, 1)
    slide = Presentation(project / receipt["page_pptx"]).slides[0]
    shapes = [shape for shape in slide.shapes if shape.name.startswith("special-source-")]

    for index, left in enumerate(shapes):
        for right in shapes[index + 1:]:
            horizontal = min(left.left + left.width, right.left + right.width) - max(left.left, right.left)
            vertical = min(left.top + left.height, right.top + right.height) - max(left.top, right.top)
            assert horizontal <= 0 or vertical <= 0, (left.name, right.name)


def test_special_page_clamps_confirmed_sizes_to_readable_minimums(tmp_path: Path) -> None:
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role="toc", visible_page_number=True)
    state = load(project)
    state["style_confirmation"]["contract"]["title_size_pt"] = 10
    state["style_confirmation"]["contract"]["body_size_pt"] = 8
    save(project, state)

    receipt = render_special_page(project, 1)
    slide = Presentation(project / receipt["page_pptx"]).slides[0]
    sizes = {
        shape.name: min(
            run.font.size.pt
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
        )
        for shape in slide.shapes
        if shape.name.startswith("special-source-")
    }

    assert sizes["special-source-block-1"] >= 24
    assert min(value for name, value in sizes.items() if name != "special-source-block-1") >= 12


def test_synthesized_closing_resolves_material_by_block_trace_not_output_number(tmp_path: Path) -> None:
    # Break caught: a synthesized closing has no original source page matching its output position.
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role="closing", visible_page_number=False)
    source_path = project / "02_v6/paginated_word_source.json"
    source = json.loads(source_path.read_text(encoding="utf-8"))
    source["pages"][0]["page_number"] = 43
    source_path.write_text(json.dumps(source, ensure_ascii=False), encoding="utf-8")
    composition_path = project / "02_v6/page_composition.json"
    composition = json.loads(composition_path.read_text(encoding="utf-8"))
    composition["pages"][0].update({"source_page_id": None, "source_page_number": None, "role_source": "synthesized"})
    composition_path.write_text(json.dumps(composition, ensure_ascii=False), encoding="utf-8")

    receipt = render_special_page(project, 1)
    text = "\n".join(shape.text for shape in Presentation(project / receipt["page_pptx"]).slides[0].shapes if shape.has_text_frame)

    assert "最终目标：形成可持续产业生态" in text


def test_reordered_special_page_resolves_original_source_trace(tmp_path: Path) -> None:
    # Break caught: an inserted page shifts a TOC output number away from its original source page number.
    from workflow_v6_special_pages import render_special_page

    project = _reordered_toc_project(tmp_path)
    receipt = render_special_page(project, 2)
    text = "\n".join(shape.text for shape in Presentation(project / receipt["page_pptx"]).slides[0].shapes if shape.has_text_frame)

    assert "确认目录" in text
    assert "PART 1｜产业目标" in text


@pytest.mark.parametrize("duplicate_text", ["冲突的另一份材料", "黄石项目建议"])
def test_special_page_rejects_duplicate_traced_block_ids(
    tmp_path: Path, duplicate_text: str,
) -> None:
    # Break caught: an ambiguous global source block ID silently chooses the first occurrence.
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role="cover", visible_page_number=False)
    path = project / "02_v6/paginated_word_source.json"
    source = json.loads(path.read_text(encoding="utf-8"))
    source["pages"].append({"page_number": 99, "blocks": [
        {"type": "paragraph", "text": duplicate_text, "source_block_id": "block-1"},
    ]})
    path.write_text(json.dumps(source, ensure_ascii=False), encoding="utf-8")

    with pytest.raises(ValueError, match="duplicate traced source block ID"):
        render_special_page(project, 1)


def test_special_page_rejects_duplicate_required_source_block_ids(tmp_path: Path) -> None:
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role="cover", visible_page_number=False)
    path = project / "02_v6/page_composition.json"
    composition = json.loads(path.read_text(encoding="utf-8"))
    composition["pages"][0]["material_source_block_ids"].append("block-2")
    path.write_text(json.dumps(composition, ensure_ascii=False), encoding="utf-8")

    with pytest.raises(ValueError, match="duplicate source block IDs"):
        render_special_page(project, 1)


@pytest.mark.parametrize("role,confirmed_title,confirmed_chapter", [
    ("cover", "确认后的封面标题", ""),
    ("section", "确认后的 PART 1", "确认后的章节标题"),
])
def test_special_page_uses_frozen_confirmed_titles(
    tmp_path: Path, role: str, confirmed_title: str, confirmed_chapter: str,
) -> None:
    # Break caught: one-time UI title edits are silently replaced by old Word wording.
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role=role, visible_page_number=role == "section")
    path = project / "02_v6/page_composition.json"
    composition = json.loads(path.read_text(encoding="utf-8"))
    composition["pages"][0]["fixed_page_title"] = confirmed_title
    composition["pages"][0]["chapter_title"] = confirmed_chapter
    path.write_text(json.dumps(composition, ensure_ascii=False), encoding="utf-8")

    receipt = render_special_page(project, 1)
    displayed = [shape.text for shape in Presentation(project / receipt["page_pptx"]).slides[0].shapes if shape.has_text_frame]

    assert confirmed_title in displayed
    if confirmed_chapter:
        assert confirmed_chapter in displayed
    assert ROLE_TEXT[role][0] not in displayed


def test_special_page_rejects_logo_path_outside_project(tmp_path: Path) -> None:
    # Break caught: a corrupted logo source path changes the trusted secure-I/O root.
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role="cover", visible_page_number=False)
    outside = tmp_path / "outside.svg"
    outside.write_text('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 10 2"/>', encoding="utf-8")
    state = load(project)
    replacement = new_project(
        word_source=state["word_source"], logo_source={"path": str(outside.resolve())}, pages=state["pages"],
    )
    for field in ("style_confirmation", "confirmed_ui_revision", "confirmed_ui_digest", "page_materials_status"):
        replacement[field] = state[field]
    (project / "workflow_v6.json").write_text(json.dumps(replacement, ensure_ascii=False), encoding="utf-8")

    with pytest.raises(ValueError, match="logo"):
        render_special_page(project, 1)


@pytest.mark.parametrize("view_box", ["0 0 100 20 1", "0 0 NaN 20", "0 0 Inf 20"])
def test_special_page_rejects_malformed_or_nonfinite_svg_viewbox(
    tmp_path: Path, view_box: str,
) -> None:
    # Break caught: malformed geometry reaches coordinate conversion instead of the SVG boundary.
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role="cover", visible_page_number=False)
    (project / "00_source/logo.svg").write_text(
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="{view_box}"/>', encoding="utf-8",
    )

    with pytest.raises(ValueError, match="exactly four finite numbers with positive width and height"):
        render_special_page(project, 1)


@pytest.mark.skipif(os.name != "nt", reason="real Windows junction test")
def test_special_page_rejects_locked_logo_directory_junction(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    # Break caught: a linked 00_source directory supplies an outside logo after state validation.
    from workflow_v6_special_pages import render_special_page

    project = prepared_confirmed_project(tmp_path, role="cover", visible_page_number=False)
    outside = tmp_path / "outside-logo"
    outside.mkdir()
    (outside / "logo.svg").write_text('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 10 2"/>', encoding="utf-8")
    (project / "00_source/logo.svg").unlink()
    os.rmdir(project / "00_source")
    result = subprocess.run(["cmd", "/c", "mklink", "/J", str(project / "00_source"), str(outside)], capture_output=True, text=True, check=False)
    assert result.returncode == 0, result.stdout + result.stderr
    import workflow_v6_secure_io as secure_io
    real_read = secure_io.read_bytes
    safe_svg = b'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 20"/>'
    monkeypatch.setattr(
        secure_io, "read_bytes",
        lambda root, relative, **kwargs: safe_svg
        if Path(relative) == Path("00_source/logo.svg")
        else real_read(root, relative, **kwargs),
    )

    with pytest.raises((OSError, ValueError), match="reparse"):
        render_special_page(project, 1)
