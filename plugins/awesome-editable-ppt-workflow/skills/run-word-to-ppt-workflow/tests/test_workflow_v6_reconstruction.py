from __future__ import annotations

import json
import hashlib
import hmac
import os
import sys
from pathlib import Path

import pytest
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Cm


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from workflow_v6_contract import canonical_sha256, new_page, new_project  # noqa: E402
from workflow_v6_reconstruction import (  # noqa: E402
    _compare_body_images,
    _render_reconstructed_body,
    assemble_v6_deck,
    build_reconstruction_request,
    finalize_reconstructed_page as _finalize_reconstructed_page,
)
import workflow_v6_reconstruction as reconstruction_module  # noqa: E402
from editppt.runtime.fixed_region_runtime import CONTENT_BOX, SLIDE  # noqa: E402
from workflow_v6_state import create, load, save  # noqa: E402
from awesome_page_materials import publish_page_materials  # noqa: E402
from director_taskbook import project_emphasis_pages, taskbook_digest  # noqa: E402
from complex_page_experiment.loop import signing_key  # noqa: E402


def finalize_reconstructed_page(*args, **kwargs):
    project = Path(args[0] if args else kwargs["project"])
    page_number = kwargs["page_number"]
    state = load(project)
    state["word_source"]["authority_mode"] = "legacy_non_word"
    state["source_identity"] = canonical_sha256({
        "word_source": state["word_source"], "logo_source": state["logo_source"],
    })
    state["pages"][page_number - 1]["selected_candidate"] = None
    save(project, state)
    (project / "04_v6" / "images" / f"page_{page_number:03d}.json").unlink(missing_ok=True)
    return _finalize_reconstructed_page(*args, authority_mode="native_direct", **kwargs)


@pytest.fixture(autouse=True)
def _deterministic_assembled_deck_renderer(monkeypatch):
    def render(_pptx: Path, expected_pages: int, render_dir: Path) -> dict:
        render_dir.mkdir(parents=True, exist_ok=True)
        for page_number in range(1, expected_pages + 1):
            Image.new("RGB", (1904, 1071), "white").save(
                render_dir / f"page-{page_number:03d}.png"
            )
        return {
            "available": True,
            "status": "passed",
            "detail": f"test renderer produced {expected_pages} slides",
            "slides": expected_pages,
            "rendered_slides": expected_pages,
        }

    monkeypatch.setattr(reconstruction_module, "_render_powerpoint_deck", render)


def _page_plan(page_number: int) -> dict[str, object]:
    return {
        "page_purpose": f"Explain page {page_number}.",
        "primary_relationship": {
            "grammar": "flow",
            "description": "Source feeds destination.",
            "fact_ids": [f"body-{page_number}"],
            "visual_instruction": "Keep the accepted flow and its connectors.",
            "nodes": [
                {"node_id": "source", "label": "Source", "fact_ids": [f"body-{page_number}"]},
                {"node_id": "destination", "label": "Destination", "fact_ids": [f"body-{page_number}"]},
            ],
            "edges": [{
                "from_node": "source",
                "to_node": "destination",
                "label": "feeds",
                "fact_ids": [f"body-{page_number}"],
            }],
        },
        "core_exhibit": {
            "grammar": "flow",
            "description": "Accepted flow exhibit.",
            "fact_ids": [f"body-{page_number}"],
        },
        "support_groups": [],
        "reading_path": "Read source, then destination.",
        "local_visuals": [],
    }


def _write_signed_receipt(project: Path, page_number: int, value: dict | None = None) -> dict:
    state = load(project)
    image = project / "04_v6" / "images" / f"page_{page_number:03d}.png"
    digest = hashlib.sha256(image.read_bytes()).hexdigest()
    receipt = value or {
        "schema_version": "awesome-complex-page-acceptance-v1",
        "status": "accepted",
        "experiment_id": f"live-page-{page_number:03d}",
        "page_number": page_number,
        "source_snapshot_sha256": state["source_identity"],
        "source_identity": state["source_identity"],
        "ui_revision": 1,
        "ui_digest": "a" * 64,
        "material_view_sha256": "b" * 64,
        "page_plan": _page_plan(page_number),
        "candidate": {
            "attempt": 1,
            "path": image.relative_to(project).as_posix(),
            "sha256": digest,
            "byte_size": image.stat().st_size,
            "width": 1904,
            "height": 896,
            "trace_path": "04_v6/trace.json",
            "prompt_path": "04_v6/prompt.md",
            "operation": "generate",
            "quality": "high",
            "selected_reference_ids": [],
            "input_sha256s": [],
            "prompt_sha256": "c" * 64,
            "request_identity": "d" * 64,
            "duration_seconds": 1,
            "duration_unavailable_reason": None,
        },
        "candidate_history_sha256": "e" * 64,
        "prompt_history_sha256": "f" * 64,
        "selected_real_reference_ids": [],
        "accepted_review": {
            "decision": "accept",
            "problems": [],
            "model": "test-model",
            "effort": None,
            "duration_seconds": 1,
            "authority_path": "04_v6/review.json",
            "authority_sha256": "1" * 64,
        },
        "provider_authority": {
            "trace_path": "04_v6/trace.json",
            "trace_sha256": "2" * 64,
            "capability_path": "04_v6/capability.json",
            "capability_sha256": "3" * 64,
            "capability_nonce": "test-nonce",
            "journal_path": "04_v6/journal.jsonl",
            "journal_sha256": "4" * 64,
            "request_identity": "d" * 64,
        },
        "fixed_frame": {},
        "evidence_checkpoint": {
            "schema_version": "awesome-complex-page-candidate-acceptance-checkpoint-v1",
            "experiment_id": f"live-page-{page_number:03d}",
            "workspace_identity_sha256": "5" * 64,
            "source_snapshot_sha256": state["source_identity"],
            "page_number": page_number,
            "selected_attempt": 1,
            "candidate_sha256": digest,
            "request_identity": "d" * 64,
            "review_authority_sha256": "1" * 64,
            "event_count": 3,
            "terminal_event_index": 2,
            "evidence_prefix_sha256": "6" * 64,
            "causal_events": [
                {"event_index": index, "value": {}} for index in range(3)
            ],
            "checkpoint_sha256": "7" * 64,
        },
    }
    authority = receipt["provider_authority"]
    for name, payload in (
        ("trace", b'{"trace":"test"}\n'),
        ("capability", b'{"capability":"test"}\n'),
        ("journal", b'{"event":"test"}\n'),
    ):
        authority_path = project / authority[f"{name}_path"]
        authority_path.parent.mkdir(parents=True, exist_ok=True)
        authority_path.write_bytes(payload)
        authority[f"{name}_sha256"] = hashlib.sha256(payload).hexdigest()
    unsigned = {key: item for key, item in receipt.items() if key not in {"key_id", "hmac_sha256"}}
    key_id, key = signing_key()
    unsigned["key_id"] = key_id
    payload = json.dumps(unsigned, ensure_ascii=False, sort_keys=True, separators=(",", ":")).encode()
    signed = {**unsigned, "hmac_sha256": hmac.new(key, payload, hashlib.sha256).hexdigest()}
    path = project / "04_v6" / "images" / f"page_{page_number:03d}.json"
    path.write_text(
        json.dumps(signed, ensure_ascii=False, sort_keys=True, separators=(",", ":")) + "\n",
        encoding="utf-8",
    )
    return signed


def _style():
    return {
        "primary_color": "#17365D",
        "secondary_color": "#C7352B",
        "background_color": "#E7F1FA",
        "fixed_frame": {
            "geometry_version": "fixed-canvas-cm-v2",
            "body_bounds_cm": {"x": 0.81, "y": 2.3, "w": 23.78, "h": 11.18},
            "body_bounds": {"x": 0.81 / 25.4, "y": 2.3 / 14.288, "w": 23.78 / 25.4, "h": 11.18 / 14.288},
            "title_bounds_cm": {"x": 0.9, "y": 0.5, "w": 20.066, "h": 1.4288},
            "title_bounds": {"x": 0.9 / 25.4, "y": 0.5 / 14.288, "w": 20.066 / 25.4, "h": 1.4288 / 14.288},
            "logo_bounds_cm": {"x": 21.844, "y": 0.57152, "w": 2.667, "h": 1.0716},
            "logo_bounds": {"x": 21.844 / 25.4, "y": 0.57152 / 14.288, "w": 2.667 / 25.4, "h": 1.0716 / 14.288},
            "footer_line": {"x": 0.9, "y": 13.64504, "w": 23.6, "h": 0.028576, "color": "#B8C0CC"},
            "page_number_bounds_cm": {"x": 23.368, "y": 13.687904, "w": 1.143, "h": 0.3572},
            "page_number_bounds": {"x": 23.368 / 25.4, "y": 13.687904 / 14.288, "w": 1.143 / 25.4, "h": 0.3572 / 14.288},
            "page_number_style": {"font": "Microsoft YaHei", "size_pt": 9, "color": "#6B7280"},
            "title_color": "#0B1727",
        },
        "hard_constraints": {
            "title_color": "#0B1727",
            "typography": {
                "heading": {"cjk": "Microsoft YaHei"},
                "type_scale_pt": {"page_title": 28},
            },
        },
    }


def _body(path: Path, text: str, *, color: str | None = None, bold: bool = False):
    deck = Presentation()
    deck.slide_width = Cm(25.4)
    deck.slide_height = Cm(14.288)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(Cm(2), Cm(3), Cm(10), Cm(2))
    box.text = text
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
    deck.save(path)


def _empty_manifest() -> dict:
    return {
        "workflow_contract_version": "fixed-canvas-cm-v2",
        "reconstruction_contract_version": "editable-image-v3",
        "slide": dict(SLIDE),
        "content_box": dict(CONTENT_BOX),
        "source": {"width_px": 1904, "height_px": 896},
        "text_inventory": [],
        "visual_inventory": [],
        "background_strategy": "native white body background",
        "quality_checks": {
            "font_size_calibrated": True,
            "visual_inventory_matched": True,
            "background_strategy_checked": True,
            "shape_corner_geometry_checked": True,
        },
        "text_boxes": [], "tables": [], "shapes": [], "images": [], "charts": [],
    }


def _project(tmp_path: Path, page_count: int = 2) -> Path:
    root = tmp_path / "project"
    (root / "00_source").mkdir(parents=True)
    (root / "00_source" / "logo.svg").write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 20"><rect width="100" height="20"/></svg>',
        encoding="utf-8",
    )
    pages = [new_page(number, title=f"标题{number}") for number in range(1, page_count + 1)]
    for page in pages:
        page["state"] = "accepted"
        image = root / "04_v6" / "images" / f"page_{page['page_number']:03d}.png"
        image.parent.mkdir(parents=True, exist_ok=True)
        Image.new("RGB", (1904, 896), "white").save(image)
        page["first_candidate"] = {"path": image.relative_to(root).as_posix(), "attempt": 1, "operation": "generate"}
        page["selected_candidate"] = dict(page["first_candidate"])
        effective = root / "02_v6" / "effective_pages" / f"page_{page['page_number']:03d}.json"
        effective.parent.mkdir(parents=True, exist_ok=True)
        effective.write_text(json.dumps({"page_number": page["page_number"], "word_original": "正文"}), encoding="utf-8")
    state = new_project(
        word_source={"path": "00_source/source.docx"},
        logo_source={"path": "00_source/logo.svg"},
        pages=pages,
    )
    material_style = {
        "primary_color": "#17365D", "secondary_color": "#C7352B",
        "background_color": "#FFFFFF", "cjk_font": "Microsoft YaHei",
        "latin_font": "Arial", "title_size_pt": 28, "body_size_pt": 12,
        "caption_size_pt": 9, "regional_characteristics": "",
        "visual_description": "Formal editorial presentation.",
    }
    state["style_confirmation"] = {"status": "confirmed", "contract": material_style}
    state["confirmed_ui_revision"] = 1
    state["confirmed_ui_digest"] = "a" * 64
    state["page_materials_status"] = "pending"
    create(root, state)
    for number in range(1, page_count + 1):
        _write_signed_receipt(root, number)
    source_dir = root / "02_v6"
    (source_dir / "paginated_word_source.json").write_text(json.dumps({"pages": [
        {"page_number": number, "fixed_page_title": f"标题{number}", "fixed_page_title_source_block_id": f"title-{number}", "page_comments": [], "blocks": [
            {"type": "paragraph", "text": f"标题{number}", "source_block_id": f"title-{number}", "source_block_index": (number - 1) * 2, "source_order": 1, "relationship_ids": [], "comment_ids": []},
            {"type": "paragraph", "text": "正文", "source_block_id": f"body-{number}", "source_block_index": (number - 1) * 2 + 1, "source_order": 2, "relationship_ids": [], "comment_ids": []},
        ]} for number in range(1, page_count + 1)
    ]}, ensure_ascii=False), encoding="utf-8")
    (source_dir / "page_composition.json").write_text(json.dumps({
        "artifact_version": "page-composition-v1", "page_count": page_count, "warnings": [],
        "pages": [
            {"output_page_number": number, "source_page_id": number, "page_role": "content", "role_source": "explicit", "chapter_title": "", "fixed_page_title": f"标题{number}", "source_page_number": number, "material_source_block_ids": [f"title-{number}", f"body-{number}"], "visible_page_number": True}
            for number in range(1, page_count + 1)
        ],
    }, ensure_ascii=False), encoding="utf-8")
    (source_dir / "source_assets.json").write_text('{"assets":[]}', encoding="utf-8")
    for number in range(1, page_count + 1):
        publish_page_materials(root, number, root / "02_v6" / "awesome_page_materials" / f"page_{number:03d}.json")
    state = load(root)
    state["style_confirmation"]["contract"] = _style()
    save(root, state)
    return root


def test_v6_reconstruction_request_has_no_exact_material_and_requires_post_visual_qa(tmp_path: Path):
    project = _project(tmp_path, 1)
    receipt = json.loads((project / "04_v6/images/page_001.json").read_text(encoding="utf-8"))
    request = build_reconstruction_request(project, page_number=1)
    assert request["workflow_contract_version"] == "awesome-word-ppt-workflow-v1"
    assert request["requirements"]["exact_reference_material_custody"] is False
    assert request["requirements"]["post_reconstruction_visual_qa"] is True
    assert request["sealed_image_edits"] == []
    assert request["sealed_text_repairs"] == []
    assert request["page_plan"] == receipt["page_plan"]


def test_reconstruction_request_hash_changes_with_resigned_page_plan(tmp_path: Path):
    project = _project(tmp_path, 1)
    build_reconstruction_request(project, page_number=1)
    request_path = project / "05_v6/reconstruction_requests/page_001.json"
    first_hash = hashlib.sha256(request_path.read_bytes()).hexdigest()
    receipt_path = project / "04_v6/images/page_001.json"
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    receipt["page_plan"]["primary_relationship"]["edges"][0]["to_node"] = "source"
    _write_signed_receipt(project, 1, receipt)

    second = build_reconstruction_request(project, page_number=1)
    second_hash = hashlib.sha256(request_path.read_bytes()).hexdigest()

    assert second["page_plan"] == receipt["page_plan"]
    assert second_hash != first_hash


def test_reconstruction_request_rejects_tampered_page_plan_before_copy(tmp_path: Path):
    project = _project(tmp_path, 1)
    receipt_path = project / "04_v6/images/page_001.json"
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    receipt["page_plan"]["primary_relationship"]["nodes"][0]["node_id"] = "tampered"
    receipt_path.write_text(json.dumps(receipt, ensure_ascii=False), encoding="utf-8")

    with pytest.raises(ValueError, match="signature"):
        build_reconstruction_request(project, page_number=1)

    assert not (project / "05_v6/reconstruction_requests/page_001.json").exists()


def test_reconstruction_request_preserves_numeric_authority(tmp_path: Path):
    project = _project(tmp_path, 1)
    authority = {
        "title": "项目计划",
        "rendering_primitive": "time_interval",
        "period": "2026-09",
        "series": [{
            "name": "计划",
            "categories": ["尽调", "投决"],
            "start_dates": ["2026-09-01", "2026-09-11"],
            "end_dates": ["2026-09-10", "2026-09-15"],
        }],
    }
    path = project / "02_v6/page_materials/page_001.json"
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps({"chart_facts": [authority]}, ensure_ascii=False), encoding="utf-8")

    request = build_reconstruction_request(project, page_number=1)

    assert request["numeric_authority"] == authority


def test_reconstruction_request_refuses_incomplete_or_ambiguous_authority(tmp_path: Path):
    project = _project(tmp_path, 1)
    complete = {
        "title": "Revenue",
        "rendering_primitive": "column_bar",
        "chart_variant": "column",
        "unit": "USD m",
        "basis": "FY2025 revenue",
        "series": [{"name": "Revenue", "categories": ["A"], "values": [10]}],
    }
    incomplete = {
        "title": "Priorities",
        "disabled_primitive": "column_bar",
        "fallback": "native_table",
        "source_wording": "A is high priority and B is medium priority.",
        "series": [{"name": "Priority", "categories": ["A", "B"], "values": ["high", "medium"]}],
    }
    path = project / "02_v6/page_materials/page_001.json"
    path.parent.mkdir(parents=True, exist_ok=True)

    path.write_text(json.dumps({"chart_facts": [incomplete]}, ensure_ascii=False), encoding="utf-8")
    assert "numeric_authority" not in build_reconstruction_request(project, page_number=1)
    assert json.loads(path.read_text(encoding="utf-8"))["chart_facts"][0]["source_wording"]

    path.write_text(json.dumps({"chart_facts": [complete, complete]}, ensure_ascii=False), encoding="utf-8")
    assert "numeric_authority" not in build_reconstruction_request(project, page_number=1)


def test_reconstruction_request_refuses_legacy_times_numeric_authority(tmp_path: Path):
    project = _project(tmp_path, 1)
    legacy = {
        "title": "Revenue trend", "rendering_primitive": "line_point", "chart_variant": "line",
        "unit": "USD m", "basis": "FY2025 revenue",
        "series": [{"name": "Revenue", "times": ["2024", "2025"], "values": [12, 18]}],
    }
    path = project / "02_v6/page_materials/page_001.json"
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps({"chart_facts": [legacy]}, ensure_ascii=False), encoding="utf-8")

    request = build_reconstruction_request(project, page_number=1)

    assert "numeric_authority" not in request


def test_reconstruction_request_carries_signed_text_repairs(tmp_path: Path):
    project = _project(tmp_path, 1)
    receipt_path = project / "04_v6/images/page_001.json"
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    receipt["reconstruction_repairs"] = [
        {
            "category": "severe_usability",
            "detail": "将错字“清出”修正为“退出”，其余构图保持不变。",
        }
    ]
    _write_signed_receipt(project, 1, receipt)

    request = build_reconstruction_request(project, page_number=1)

    assert request["sealed_text_repairs"] == receipt["reconstruction_repairs"]


def test_finalize_requires_exact_native_text_repair(tmp_path: Path):
    project = _project(tmp_path, 1)
    receipt_path = project / "04_v6/images/page_001.json"
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    receipt["reconstruction_repairs"] = [
        {
            "category": "severe_usability",
            "detail": "将错字“清出”修正为“退出”，其余构图保持不变。",
            "find": "清出",
            "replace": "退出",
        }
    ]
    _write_signed_receipt(project, 1, receipt)
    wrong = tmp_path / "wrong.pptx"
    _body(wrong, "清出表现挂钩")

    with pytest.raises(ValueError, match="did not apply a sealed repair"):
        _finalize_reconstructed_page(
            project,
            page_number=1,
            reconstructed_body=wrong,
        )

    correct = tmp_path / "correct.pptx"
    _body(correct, "退出表现挂钩")
    report = finalize_reconstructed_page(
        project, page_number=1, reconstructed_body=correct,
    )
    assert report["fixed_frame"]["passed"] is True


def test_finalize_and_assemble_add_fixed_layers_without_office_or_visual_qa(tmp_path: Path):
    project = _project(tmp_path, 2)
    for page in (1, 2):
        body = tmp_path / f"body-{page}.pptx"
        _body(body, f"可编辑正文{page}")
        report = finalize_reconstructed_page(project, page_number=page, reconstructed_body=body)
        assert report["post_reconstruction_visual_qa"]["status"] == "skipped"
        assert report["fixed_frame"]["passed"] is True

    report = assemble_v6_deck(project)
    output = project / report["output"]
    deck = Presentation(output)
    assert len(deck.slides) == 2
    assert report["post_reconstruction_visual_qa"]["status"] == "skipped"
    assert report["release_ready"] is True
    assert report["release_status"] == "release_ready"
    assert report["openxml_validation"]["status"] == "passed"
    assert report["final_output"] == {
        "path": str(output),
        "relative_path": report["output"],
        "sha256": report["sha256"],
    }
    assert {report["enhanced_validation"][name]["status"] for name in ("officecli", "powerpoint")} <= {
        "passed", "failed", "skipped",
    }
    assert all(page["state"] == "page_complete" for page in load(project)["pages"])


def test_post_reconstruction_visual_comparison_rejects_severe_body_loss(tmp_path: Path):
    source = Image.new("RGB", (1904, 896), "white")
    for left in (80, 520, 960, 1400):
        for top in (80, 330, 580):
            for x in range(left, left + 320):
                for y in range(top, top + 170):
                    source.putpixel((x, y), (23, 54, 93))
    source_path = tmp_path / "source.png"
    preview_path = tmp_path / "preview.png"
    source.save(source_path)
    Image.new("RGB", (1200, 675), "white").save(preview_path)

    report = _compare_body_images(source_path, preview_path)

    assert report["passed"] is False
    assert report["reason"] == "severe_body_content_loss"
    assert report["metrics"]["foreground_retention"] < report["thresholds"]["minimum_foreground_retention"]


@pytest.mark.skipif(os.name != "nt", reason="actual PowerPoint render requires Windows")
def test_actual_pptx_render_detects_removed_body_content(tmp_path: Path):
    populated = tmp_path / "populated.pptx"
    partial = tmp_path / "partial.pptx"
    for path, boxes in ((populated, 12), (partial, 1)):
        deck = Presentation()
        deck.slide_width = Cm(25.4)
        deck.slide_height = Cm(14.288)
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        for index in range(boxes):
            row, column = divmod(index, 4)
            shape = slide.shapes.add_shape(
                1, Cm(1.2 + column * 5.8), Cm(2.8 + row * 3.2), Cm(4.8), Cm(2.2),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string("17365D")
        deck.save(path)
    source = tmp_path / "source.png"
    missing = tmp_path / "missing.png"

    populated_render = _render_reconstructed_body(populated, source)
    if populated_render["status"] == "unavailable":
        pytest.skip(populated_render["detail"])
    assert populated_render["status"] == "passed"
    missing_render = _render_reconstructed_body(partial, missing)
    assert missing_render["status"] == "passed"

    comparison = _compare_body_images(source, missing)
    assert comparison["passed"] is False
    assert comparison["reason"] == "severe_body_content_loss"


def test_sealed_page_visual_failure_blocks_publication(tmp_path: Path, monkeypatch):
    project = _project(tmp_path, 1)
    body = tmp_path / "sealed" / "page.pptx"
    body.parent.mkdir()
    _body(body, "sealed body")
    (body.parent / "manifest.json").write_text("{}", encoding="utf-8")
    source = project / "04_v6/images/page_001.png"
    monkeypatch.setattr(
        reconstruction_module,
        "_require_final_authority",
        lambda *_args, **_kwargs: {
            "accepted_receipt": {},
            "accepted_source_body": {"path": source.relative_to(project).as_posix(), "sha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "worker_source_body": {},
        },
    )
    monkeypatch.setattr(
        reconstruction_module,
        "_run_post_reconstruction_visual_qa",
        lambda *_args, **_kwargs: {"status": "failed", "passed": False, "reason": "severe_body_content_loss"},
    )

    with pytest.raises(ValueError, match="post-reconstruction visual QA failed"):
        _finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)

    assert not (project / "06_v6/pages/page_001/page.pptx").exists()


def test_sealed_page_without_actual_renderer_remains_assembly_ready(tmp_path: Path, monkeypatch):
    project = _project(tmp_path, 1)
    body = tmp_path / "sealed" / "page.pptx"
    body.parent.mkdir()
    _body(body, "sealed body")
    source = project / "04_v6/images/page_001.png"
    monkeypatch.setattr(
        reconstruction_module,
        "_require_final_authority",
        lambda *_args, **_kwargs: {
            "accepted_receipt": {},
            "accepted_source_body": {"path": source.relative_to(project).as_posix(), "sha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "worker_source_body": {},
        },
    )
    monkeypatch.setattr(
        reconstruction_module,
        "_run_post_reconstruction_visual_qa",
        lambda *_args, **_kwargs: {
            "status": "unavailable", "passed": False,
            "reason": "actual_pptx_render_unavailable",
        },
    )

    report = _finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)

    assert report["artifact_version"] == "final-page-v6"
    assert report["post_reconstruction_visual_qa"]["status"] == "unavailable"
    assert (project / report["page_pptx"]).is_file()
    assert load(project)["pages"][0]["state"] == "page_complete"
    manifest = project / "05_v6/reconstruction_runs/page_001/pages/page_001/manifest.json"
    manifest.parent.mkdir(parents=True)
    manifest.write_text(json.dumps(_empty_manifest()), encoding="utf-8")
    monkeypatch.setattr(
        reconstruction_module,
        "verify_completed_page_authority",
        lambda *_args, **_kwargs: {
            "status": "verified",
            "authority_mode": "sealed_reconstruction",
            "visual_qa": report["post_reconstruction_visual_qa"],
        },
    )
    monkeypatch.setattr(
        reconstruction_module,
        "_render_powerpoint_deck",
        lambda *_args, **_kwargs: {
            "available": False, "status": "skipped", "detail": "not installed",
        },
    )

    assembly = assemble_v6_deck(project)

    assert assembly["status"] == "validation_incomplete"
    assert assembly["release_ready"] is False
    assert (project / assembly["candidate_output"]["relative_path"]).is_file()


def test_final_openxml_failure_preserves_previous_deck(tmp_path: Path, monkeypatch):
    project = _project(tmp_path, 1)
    body = tmp_path / "body.pptx"
    _body(body, "editable body")
    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)
    first = assemble_v6_deck(project)
    output = project / first["output"]
    old_bytes = output.read_bytes()
    assembly_report = project / first["assembly_report"]
    old_report = assembly_report.read_bytes()
    monkeypatch.setattr(
        reconstruction_module,
        "_validate_final_openxml",
        lambda *_args, **_kwargs: {"status": "failed", "passed": False, "reason": "broken package"},
    )

    with pytest.raises(ValueError, match="OpenXML validation failed"):
        assemble_v6_deck(project)

    assert output.read_bytes() == old_bytes
    assert assembly_report.read_bytes() == old_report
    assert not list(project.glob(".08_final.*.tmp"))


def test_available_enhanced_validation_failure_preserves_previous_deck(tmp_path: Path, monkeypatch):
    project = _project(tmp_path, 1)
    body = tmp_path / "body.pptx"
    _body(body, "editable body")
    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)
    monkeypatch.setattr(
        reconstruction_module, "_officecli_validation",
        lambda *_args, **_kwargs: {"available": False, "status": "skipped", "detail": "not installed"},
    )
    calls = 0

    def powerpoint_result(_pptx, expected_pages, render_dir):
        nonlocal calls
        calls += 1
        render_dir.mkdir(parents=True, exist_ok=True)
        for page_number in range(1, expected_pages + 1):
            Image.new("RGB", (1904, 1071), "white").save(
                render_dir / f"page-{page_number:03d}.png"
            )
        return {
            "available": True,
            "status": "passed" if calls == 1 else "failed",
            "detail": "simulated render validation",
        }

    monkeypatch.setattr(reconstruction_module, "_render_powerpoint_deck", powerpoint_result)
    first = assemble_v6_deck(project)
    output = project / first["output"]
    old_bytes = output.read_bytes()
    assembly_report = project / first["assembly_report"]
    old_report = assembly_report.read_bytes()

    with pytest.raises(ValueError, match="enhanced final validation failed"):
        assemble_v6_deck(project)

    assert output.read_bytes() == old_bytes
    assert assembly_report.read_bytes() == old_report
    assert not list(project.glob(".08_final.*.tmp"))


def test_unavailable_powerpoint_creates_candidate_without_replacing_final(tmp_path: Path, monkeypatch):
    project = _project(tmp_path, 1)
    body = tmp_path / "body.pptx"
    _body(body, "editable body")
    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)
    monkeypatch.setattr(
        reconstruction_module, "_officecli_validation",
        lambda *_args, **_kwargs: {"available": False, "status": "skipped", "detail": "not installed"},
    )
    calls = 0

    def powerpoint_result(_pptx, expected_pages, render_dir):
        nonlocal calls
        calls += 1
        if calls == 1:
            render_dir.mkdir(parents=True, exist_ok=True)
            for page_number in range(1, expected_pages + 1):
                Image.new("RGB", (1904, 1071), "white").save(
                    render_dir / f"page-{page_number:03d}.png"
                )
        return (
            {"available": True, "status": "passed", "detail": "rendered"}
            if calls == 1 else
            {"available": False, "status": "skipped", "detail": "not installed"}
        )

    monkeypatch.setattr(reconstruction_module, "_render_powerpoint_deck", powerpoint_result)
    first = assemble_v6_deck(project)
    final = project / first["output"]
    old_final = final.read_bytes()
    assembly_report = project / first["assembly_report"]
    old_report = assembly_report.read_bytes()

    candidate = assemble_v6_deck(project)

    assert candidate["status"] == "validation_incomplete"
    assert candidate["release_ready"] is False
    assert candidate["final_output"] is None
    assert (project / candidate["candidate_output"]["relative_path"]).is_file()
    assert final.read_bytes() == old_final
    assert assembly_report.read_bytes() == old_report


def test_assembled_visual_qa_rejects_one_localized_body_block(
    tmp_path: Path, monkeypatch,
):
    project = _project(tmp_path, 1)
    body = tmp_path / "body.pptx"
    _body(body, "editable body")
    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)
    source = project / "04_v6/images/page_001.png"
    source_image = Image.new("RGB", (1904, 896), "white")
    source_draw = ImageDraw.Draw(source_image)
    for row in range(4):
        for column in range(8):
            left = column * 230 + 20
            top = row * 210 + 20
            source_draw.rectangle((left, top, left + 150, top + 95), fill="#17365D")
    source_image.save(source)
    page_report_path = project / "06_v6/pages/page_001/page.json"
    page_report = json.loads(page_report_path.read_text(encoding="utf-8"))
    page_report["accepted_source_body"] = {"path": source.relative_to(project).as_posix()}
    page_report_path.write_text(json.dumps(page_report), encoding="utf-8")
    manifest = project / "05_v6/reconstruction_runs/page_001/pages/page_001/manifest.json"
    manifest.parent.mkdir(parents=True)
    manifest.write_text(json.dumps(_empty_manifest()), encoding="utf-8")
    monkeypatch.setattr(
        reconstruction_module,
        "verify_completed_page_authority",
        lambda *_args, **_kwargs: {
            "status": "verified",
            "authority_mode": "sealed_reconstruction",
            "visual_qa": {"status": "unavailable", "passed": False},
        },
    )

    def localized_render(_pptx: Path, expected_pages: int, render_dir: Path) -> dict:
        assert expected_pages == 1
        render_dir.mkdir(parents=True)
        image = Image.new("RGB", (1904, 1071), "white")
        ImageDraw.Draw(image).rectangle((180, 240, 380, 390), fill="#17365D")
        image.save(render_dir / "page-001.png")
        return {"available": True, "status": "passed", "detail": "rendered"}

    monkeypatch.setattr(reconstruction_module, "_render_powerpoint_deck", localized_render)

    with pytest.raises(ValueError, match="assembled deck visual QA failed"):
        assemble_v6_deck(project)

    assert not (project / "08_final/deck.pptx").exists()
    assert not list(project.glob(".08_final.*.tmp"))


def test_assembly_report_write_failure_preserves_previous_final_package(
    tmp_path: Path, monkeypatch,
):
    project = _project(tmp_path, 1)
    body = tmp_path / "body.pptx"
    _body(body, "editable body")
    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)
    first = assemble_v6_deck(project)
    output = project / first["output"]
    assembly_report = project / first["assembly_report"]
    old_deck = output.read_bytes()
    old_report = assembly_report.read_bytes()

    monkeypatch.setattr(
        reconstruction_module,
        "_write_staged_json",
        lambda *_args, **_kwargs: (_ for _ in ()).throw(OSError("report write failed")),
    )

    with pytest.raises(OSError, match="report write failed"):
        assemble_v6_deck(project)

    assert output.read_bytes() == old_deck
    assert assembly_report.read_bytes() == old_report
    assert not list(project.glob(".08_final.*.tmp"))


def test_assembly_allows_manifestless_genuine_native_direct_page(tmp_path: Path):
    project = _project(tmp_path, 1)
    body = tmp_path / "native-direct-body.pptx"
    _body(body, "native-direct editable body")
    final_report = finalize_reconstructed_page(
        project, page_number=1, reconstructed_body=body,
    )

    assert "accepted_receipt" not in final_report
    assert load(project)["pages"][0]["selected_candidate"] is None
    assert not (project / "04_v6/images/page_001.json").exists()
    assert not (
        project / "05_v6/reconstruction_runs/page_001/pages/page_001/manifest.json"
    ).exists()

    report = assemble_v6_deck(project)

    assert report["status"] == "complete"
    assert (project / report["output"]).is_file()


def test_formal_word_page_cannot_downgrade_to_native_direct_by_deleting_seals(tmp_path: Path):
    project = _project(tmp_path, 1)
    body = tmp_path / "native-direct-body.pptx"
    _body(body, "native-direct editable body")
    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)
    state = load(project)
    state["word_source"].pop("authority_mode", None)
    state["source_identity"] = canonical_sha256({
        "word_source": state["word_source"], "logo_source": state["logo_source"],
    })
    save(project, state)

    with pytest.raises(RuntimeError, match="native-direct.*legacy"):
        assemble_v6_deck(project)

    assert not (project / "08_final/deck.pptx").exists()


def test_assembly_rejects_manifestless_page_2_native_chart_before_output(tmp_path: Path):
    project = _project(tmp_path, 2)
    for page_number in (1, 2):
        body = tmp_path / f"native-direct-{page_number}.pptx"
        _body(body, f"native-direct editable body {page_number}")
        finalize_reconstructed_page(project, page_number=page_number, reconstructed_body=body)

    page_path = project / "06_v6/pages/page_002/page.pptx"
    page_deck = Presentation(page_path)
    chart_data = ChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("目标", (420, 100))
    page_deck.slides[0].shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Cm(13), Cm(3), Cm(8), Cm(5),
        chart_data,
    )
    page_deck.save(page_path)
    final_report_path = project / "06_v6/pages/page_002/page.json"
    final_report = json.loads(final_report_path.read_text(encoding="utf-8"))
    final_report["sha256"] = hashlib.sha256(page_path.read_bytes()).hexdigest()
    final_report_path.write_text(json.dumps(final_report), encoding="utf-8")

    with pytest.raises(ValueError, match="manifestless native-direct page contains an undeclared native chart"):
        assemble_v6_deck(project)

    assert not (project / "08_final/deck.pptx").exists()
    assert sum(shape.has_chart for shape in Presentation(page_path).slides[0].shapes) == 1


def test_assembly_rejects_sealed_page_when_manifest_is_missing(tmp_path: Path):
    project = _project(tmp_path, 2)
    for page_number in (1, 2):
        body = tmp_path / f"body-{page_number}.pptx"
        _body(body, f"editable body {page_number}")
        finalize_reconstructed_page(
            project, page_number=page_number, reconstructed_body=body,
        )

    page_path = project / "06_v6/pages/page_002/page.pptx"
    page_deck = Presentation(page_path)
    chart_data = ChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("目标", (420, 100))
    page_deck.slides[0].shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Cm(13), Cm(3), Cm(8), Cm(5),
        chart_data,
    )
    page_deck.save(page_path)

    state = load(project)
    state["pages"][1]["selected_candidate"] = dict(state["pages"][1]["first_candidate"])
    save(project, state)
    receipt = _write_signed_receipt(project, 2)
    final_report_path = project / "06_v6/pages/page_002/page.json"
    final_report = json.loads(final_report_path.read_text(encoding="utf-8"))
    receipt_path = project / "04_v6/images/page_002.json"
    final_report["accepted_receipt"] = {
        "path": receipt_path.relative_to(project).as_posix(),
        "sha256": hashlib.sha256(receipt_path.read_bytes()).hexdigest(),
    }
    final_report["sha256"] = hashlib.sha256(page_path.read_bytes()).hexdigest()
    final_report_path.write_text(json.dumps(final_report), encoding="utf-8")
    manifest_path = (
        project / "05_v6/reconstruction_runs/page_002/pages/page_001/manifest.json"
    )
    assert not manifest_path.exists()
    assert receipt["status"] == "accepted"

    with pytest.raises(RuntimeError, match="authority is incomplete"):
        assemble_v6_deck(project)

    assert not (project / "08_final/deck.pptx").exists()
    source_page = Presentation(page_path)
    assert sum(shape.has_chart for shape in source_page.slides[0].shapes) == 1


def test_finalize_sets_whole_slide_background_and_recolors_non_emphasis_text_only(tmp_path: Path):
    project = _project(tmp_path, 1)
    body = tmp_path / "accent-body.pptx"
    _body(body, "非重点页结论", color="#C7352B", bold=True)

    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)

    deck = Presentation(project / "06_v6/pages/page_001/page.pptx")
    slide = deck.slides[0]
    body_shape = next(shape for shape in slide.shapes if shape.name == "TextBox 1")
    run = body_shape.text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == "17365D"
    assert run.font.bold is True
    assert str(slide.background.fill.fore_color.rgb) == "E7F1FA"


def test_finalize_preserves_accent_text_on_automatically_matched_emphasis_page(tmp_path: Path):
    project = _project(tmp_path, 1)
    state = load(project)
    taskbook = {
        "use_scenario": "汇报",
        "presenter": "项目团队",
        "primary_audience": "管理层",
        "audience_prior_knowledge": "已阅读材料",
        "desired_outcome": "形成判断",
        "emphasis": "正文",
        "deemphasis": "重复背景",
    }
    state["director_confirmation"] = {
        "template_id": "investment-committee",
        "template_version": "1.0",
        "taskbook": taskbook,
        "taskbook_digest": taskbook_digest(taskbook),
    }
    save(project, state)
    body = tmp_path / "emphasis-body.pptx"
    _body(body, "重点页结论", color="#C7352B")

    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)

    deck = Presentation(project / "06_v6/pages/page_001/page.pptx")
    body_shape = next(shape for shape in deck.slides[0].shapes if shape.name == "TextBox 1")
    assert str(body_shape.text_frame.paragraphs[0].runs[0].font.color.rgb) == "C7352B"


def test_emphasis_pages_are_recomputed_after_paginated_word_changes(tmp_path: Path):
    project = _project(tmp_path, 2)
    state = load(project)
    taskbook = {
        "use_scenario": "汇报", "presenter": "项目团队", "primary_audience": "管理层",
        "audience_prior_knowledge": "已阅读材料", "desired_outcome": "形成判断",
        "emphasis": "新增成果", "deemphasis": "重复背景",
    }
    state["director_confirmation"] = {
        "template_id": "investment-committee", "template_version": "1.0",
        "taskbook": taskbook, "taskbook_digest": taskbook_digest(taskbook),
    }
    save(project, state)
    assert project_emphasis_pages(project) == set()
    source_path = project / "02_v6/paginated_word_source.json"
    source = json.loads(source_path.read_text(encoding="utf-8"))
    source["pages"][1]["blocks"][1]["text"] = "新增成果"
    source_path.write_text(json.dumps(source, ensure_ascii=False), encoding="utf-8")
    assert project_emphasis_pages(project) == {2}


def test_non_emphasis_theme_color_is_replaced_only_when_it_resolves_to_secondary_family(
    tmp_path: Path,
):
    project = _project(tmp_path, 2)
    for page, theme in ((1, MSO_THEME_COLOR.ACCENT_1), (2, MSO_THEME_COLOR.ACCENT_2)):
        body = tmp_path / f"theme-{page}.pptx"
        _body(body, f"主题文字{page}")
        deck = Presentation(body)
        run = deck.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
        run.font.color.theme_color = theme
        deck.save(body)
        finalize_reconstructed_page(project, page_number=page, reconstructed_body=body)

    primary_theme = Presentation(project / "06_v6/pages/page_001/page.pptx")
    secondary_theme = Presentation(project / "06_v6/pages/page_002/page.pptx")
    primary_color = primary_theme.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.color
    secondary_color = secondary_theme.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.color
    assert primary_color.theme_color == MSO_THEME_COLOR.ACCENT_1
    assert str(secondary_color.rgb) == "17365D"


def test_non_emphasis_replacement_uses_theme_fill_for_contrast(tmp_path: Path):
    project = _project(tmp_path, 1)
    body = tmp_path / "theme-fill.pptx"
    _body(body, "深色底文字", color="#C7352B")
    deck = Presentation(body)
    shape = deck.slides[0].shapes[0]
    shape.fill.solid()
    shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    deck.save(body)

    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)

    final = Presentation(project / "06_v6/pages/page_001/page.pptx")
    run = final.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == "000000"


def test_assembly_rejects_page_changed_after_finalization(tmp_path: Path):
    project = _project(tmp_path, 1)
    body = tmp_path / "body.pptx"
    _body(body, "最终校验")
    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)
    page_path = project / "06_v6/pages/page_001/page.pptx"
    page_deck = Presentation(page_path)
    slide = page_deck.slides[0]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    slide.shapes[0].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string("C7352B")
    page_deck.save(page_path)

    with pytest.raises(RuntimeError, match="completed reconstructed page changed"):
        assemble_v6_deck(project)

    assert not (project / "08_final/deck.pptx").exists()


def test_current_confirmed_project_missing_composition_fails_closed_at_assembly(tmp_path: Path):
    project = _project(tmp_path, 1)
    body = tmp_path / "body.pptx"
    _body(body, "current body")
    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)
    (project / "02_v6/page_composition.json").unlink()

    with pytest.raises(ValueError, match="composition"):
        assemble_v6_deck(project)


def test_exact_legacy_content_only_confirmation_can_assemble_without_composition(tmp_path: Path):
    project = _project(tmp_path, 1)
    body = tmp_path / "legacy-body.pptx"
    _body(body, "legacy body")
    finalize_reconstructed_page(project, page_number=1, reconstructed_body=body)
    (project / "02_v6/page_composition.json").unlink()
    (project / "confirm_ui").mkdir(exist_ok=True)
    (project / "confirm_ui/result.json").write_text(json.dumps({
        "status": "confirmed", "revision": 1,
        "confirmed_at": "2026-08-23T00:00:00+08:00",
        "production_profile": "balanced", "global_visual_contract": {},
        "confirmed_pages": [{"page_number": 1, "effective_body": "legacy body"}],
    }), encoding="utf-8")

    report = assemble_v6_deck(project)

    assert report["page_count"] == 1


def test_assembly_accepts_role_specific_special_shapes_with_content_fixed_frame(tmp_path: Path):
    # Break caught: assembly requires the four content fixed-frame objects on a native cover.
    from workflow_v6_special_pages import render_special_page

    project = _project(tmp_path, 2)
    composition = {
        "artifact_version": "page-composition-v1", "page_count": 2, "warnings": [],
        "pages": [
            {"output_page_number": 1, "source_page_id": 1, "page_role": "cover", "role_source": "explicit", "chapter_title": "", "fixed_page_title": "标题1", "source_page_number": 1, "material_source_block_ids": ["title-1", "body-1"], "visible_page_number": False},
            {"output_page_number": 2, "source_page_id": 2, "page_role": "content", "role_source": "explicit", "chapter_title": "", "fixed_page_title": "标题2", "source_page_number": 2, "material_source_block_ids": ["title-2", "body-2"], "visible_page_number": True},
        ],
    }
    (project / "02_v6" / "page_composition.json").write_text(
        json.dumps(composition, ensure_ascii=False), encoding="utf-8",
    )
    state = load(project)
    state["pages"][0]["selected_candidate"] = None
    save(project, state)
    (project / "04_v6/images/page_001.json").unlink()
    render_special_page(project, 1)
    body = tmp_path / "body-2.pptx"
    _body(body, "可编辑正文2")
    finalize_reconstructed_page(project, page_number=2, reconstructed_body=body)

    report = assemble_v6_deck(project)
    deck = Presentation(project / report["output"])

    assert len(deck.slides) == 2
    assert any(shape.name.startswith("special-") for shape in deck.slides[0].shapes)
    assert len([shape for shape in deck.slides[1].shapes if shape.name.startswith("fixed-frame-")]) == 4


def test_finalize_allows_verified_same_accepted_page_repair(tmp_path: Path):
    project = _project(tmp_path, 1)
    first = tmp_path / "body-first.pptx"
    repaired = tmp_path / "body-repaired.pptx"
    _body(first, "first reconstruction")
    _body(repaired, "repaired reconstruction")

    initial = finalize_reconstructed_page(project, page_number=1, reconstructed_body=first)
    result = finalize_reconstructed_page(project, page_number=1, reconstructed_body=repaired)

    assert initial["sha256"] != result["sha256"]
    assert result["fixed_frame"]["passed"] is True
    assert load(project)["pages"][0]["state"] == "page_complete"
    output = project / result["page_pptx"]
    assert "repaired reconstruction" in "".join(
        shape.text for shape in Presentation(output).slides[0].shapes if hasattr(shape, "text")
    )

    receipt = output.with_name("page.json")
    value = json.loads(receipt.read_text(encoding="utf-8"))
    value["sha256"] = "0" * 64
    receipt.write_text(json.dumps(value), encoding="utf-8")
    with pytest.raises(ValueError, match="existing finalized page authority"):
        finalize_reconstructed_page(project, page_number=1, reconstructed_body=first)


def test_completed_page_repair_receipt_failure_preserves_old_package(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
):
    project = _project(tmp_path, 1)
    first = tmp_path / "body-first.pptx"
    repaired = tmp_path / "body-repaired.pptx"
    _body(first, "first reconstruction")
    _body(repaired, "repaired reconstruction")
    finalize_reconstructed_page(project, page_number=1, reconstructed_body=first)
    page_path = project / "06_v6/pages/page_001/page.pptx"
    receipt_path = project / "06_v6/pages/page_001/page.json"
    old_page = page_path.read_bytes()
    old_receipt = receipt_path.read_bytes()
    old_state = (project / "workflow_v6.json").read_bytes()

    original_write = reconstruction_module._write_json

    def fail_page_receipt(path, value):
        if Path(path).name == "page.json":
            raise RuntimeError("simulated repair receipt failure")
        return original_write(path, value)

    monkeypatch.setattr(reconstruction_module, "_write_json", fail_page_receipt)

    with pytest.raises(RuntimeError, match="repair receipt failure"):
        finalize_reconstructed_page(project, page_number=1, reconstructed_body=repaired)

    assert page_path.read_bytes() == old_page
    assert receipt_path.read_bytes() == old_receipt
    assert (project / "workflow_v6.json").read_bytes() == old_state
