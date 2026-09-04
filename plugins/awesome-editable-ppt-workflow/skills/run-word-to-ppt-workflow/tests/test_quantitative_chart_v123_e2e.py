from __future__ import annotations

import json
import re
import subprocess
import sys
import zipfile
from copy import deepcopy
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches


TESTS = Path(__file__).resolve().parent
PLUGIN = TESTS.parents[2]
REPO = TESTS.parents[4]
RUNTIME = PLUGIN / "skills/reconstruct-editable-slide/cli/editppt/runtime"
SCRIPTS = PLUGIN / "skills/run-word-to-ppt-workflow/scripts"
sys.path[:0] = [str(RUNTIME), str(SCRIPTS)]

from build_pptx_from_manifest import render_preview, write_deck  # noqa: E402
from fixed_region_runtime import CONTENT_BOX, SLIDE  # noqa: E402
from workflow_v6_materials import select_numeric_authority  # noqa: E402
from workflow_v6_reconstruction import build_reconstruction_request  # noqa: E402
from workflow_v6_reconstruction_worker import (  # noqa: E402
    PageWorkerResult,
    reconstruct_accepted_page,
)
from test_workflow_v6_reconstruction import _project, _write_signed_receipt  # noqa: E402


OUTPUT = REPO / "tmp/v1.2.3-acceptance/synthetic"


def _box_px(left: float, top: float, width: float, height: float) -> list[int]:
    return [
        round((left - CONTENT_BOX["left"]) / CONTENT_BOX["width"] * 1904),
        round((top - CONTENT_BOX["top"]) / CONTENT_BOX["height"] * 896),
        round(width / CONTENT_BOX["width"] * 1904),
        round(height / CONTENT_BOX["height"] * 896),
    ]


def _accepted_outcome(project: Path) -> SimpleNamespace:
    receipt = json.loads((project / "04_v6/images/page_001.json").read_text(encoding="utf-8"))
    selected = receipt.get("candidate", receipt.get("selected"))
    candidate = SimpleNamespace(path=project / selected["path"], attempt=selected["attempt"])
    return SimpleNamespace(status="accepted", accepted=SimpleNamespace(candidate=candidate))


def _production_worker(manifest_factory, calls: list[dict], post_validate=None):
    def worker(request):
        page_request = json.loads((request.page_dir / "page_request.json").read_text(encoding="utf-8"))
        accepted_request = json.loads((request.page_dir / "accepted_reconstruction_request.json").read_text(encoding="utf-8"))
        assert page_request.get("numeric_authority") == accepted_request.get("numeric_authority")
        manifest = manifest_factory(page_request)
        manifest_path = request.page_dir / "manifest.json"
        manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
        dispatch = subprocess.run(
            [sys.executable, str(RUNTIME / "record_page_dispatch.py"), str(request.run_dir), "--page", "page_001", "--agent-id", "deterministic-worker", "--prompt-file", str(request.prompt_file)],
            capture_output=True, text=True, check=False,
        )
        assert dispatch.returncode == 0, dispatch.stderr
        commands = (
            [sys.executable, str(RUNTIME / "main.py"), "page", "build", str(request.page_dir)],
            [sys.executable, str(RUNTIME / "main.py"), "page", "validate", str(request.page_dir), "--report", "validation.json"],
        )
        for command in commands:
            completed = subprocess.run(command, capture_output=True, text=True, check=False)
            assert completed.returncode == 0, completed.stderr
        validation = json.loads((request.page_dir / "validation.json").read_text(encoding="utf-8"))
        assert validation["passed"] is True
        if post_validate is not None:
            post_validate(request.page_dir / "page.pptx")
        calls.append({"request": page_request, "prompt": request.prompt_file.read_text(encoding="utf-8"), "manifest": manifest, "page_dir": request.page_dir})
        return PageWorkerResult(status="completed", reconstructed_body=request.page_dir / "page.pptx")

    return worker


def _worker_chart(page_request: dict) -> dict:
    authority = page_request["numeric_authority"]
    return _manifest({"object_id": "sealed-chart", "name": "sealed-chart", "box_px": [190, 90, 1142, 620], **authority})


def _relationship_manifest(page_request: dict) -> dict:
    return _with_relationship(_worker_chart(page_request), page_request)


def _with_relationship(manifest: dict, page_request: dict) -> dict:
    relationship = page_request["page_plan"]["primary_relationship"]
    boxes = {
        relationship["nodes"][0]["node_id"]: [80, 90, 260, 120],
        relationship["nodes"][1]["node_id"]: [1560, 90, 260, 120],
    }
    manifest["shapes"] = list(manifest.get("shapes", [])) + [
        {
            "object_id": node_id,
            "name": node_id,
            "type": "rect",
            "box_px": box,
            "fill": "#EAF2F8",
            "stroke": "#6B7A90",
        }
        for node_id, box in boxes.items()
    ]
    for index, edge in enumerate(relationship["edges"]):
        start = boxes[edge["from_node"]]
        end = boxes[edge["to_node"]]
        manifest["shapes"].append({
            "object_id": f"edge:{edge['from_node']}->{edge['to_node']}",
            "name": f"edge:{edge['from_node']}->{edge['to_node']}",
            "type": "line",
            "points_px": [
                start[0] + start[2] / 2,
                start[1] + start[3] / 2 + index * 10,
                end[0] + end[2] / 2,
                end[1] + end[3] / 2 + index * 10,
            ],
            "stroke": "#6B7A90",
        })
    return manifest


def _connector_endpoints(shape) -> tuple[int, int, int, int]:
    xfrm = shape._element.xpath(".//a:xfrm")[0]
    return (
        shape.left + shape.width if xfrm.get("flipH") in {"1", "true"} else shape.left,
        shape.top + shape.height if xfrm.get("flipV") in {"1", "true"} else shape.top,
        shape.left if xfrm.get("flipH") in {"1", "true"} else shape.left + shape.width,
        shape.top if xfrm.get("flipV") in {"1", "true"} else shape.top + shape.height,
    )


def _assert_relationship_geometry(slide, page_plan: dict) -> None:
    named = {shape.name: shape for shape in slide.shapes}
    relationship = page_plan["primary_relationship"]
    assert {node["node_id"] for node in relationship["nodes"]} <= named.keys()

    def inside(point: tuple[int, int], node) -> bool:
        x, y = point
        return (
            node.left - 1 <= x <= node.left + node.width + 1
            and node.top - 1 <= y <= node.top + node.height + 1
        )

    for edge in relationship["edges"]:
        edge_name = f"edge:{edge['from_node']}->{edge['to_node']}"
        matches = [shape for shape in slide.shapes if shape.name == edge_name]
        assert len(matches) == 1
        connector = matches[0]
        preset = connector._element.xpath(".//a:prstGeom")
        assert (
            connector._element.tag.rsplit("}", 1)[-1] == "cxnSp"
            or (preset and preset[0].get("prst") == "line")
        )
        x1, y1, x2, y2 = _connector_endpoints(connector)
        assert inside((x1, y1), named[edge["from_node"]])
        assert inside((x2, y2), named[edge["to_node"]])
        assert {
            item.tag.rsplit("}", 1)[-1]: item.get("type")
            for item in connector._element.spPr.get_or_add_ln()
            if item.tag.rsplit("}", 1)[-1] in {"headEnd", "tailEnd"}
        } == {"tailEnd": "triangle"}


def test_relationship_geometry_rejects_named_rectangle_with_arrow_xml(tmp_path: Path) -> None:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    source = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    source.name = "source"
    destination = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(1), Inches(2), Inches(1))
    destination.name = "destination"
    impostor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.49), Inches(5), Inches(0.02))
    impostor.name = "edge:source->destination"
    arrow = OxmlElement("a:tailEnd")
    arrow.set("type", "triangle")
    impostor._element.spPr.get_or_add_ln().append(arrow)
    path = tmp_path / "named-rectangle-edge.pptx"
    deck.save(path)

    page_plan = {
        "primary_relationship": {
            "nodes": [{"node_id": "source"}, {"node_id": "destination"}],
            "edges": [{"from_node": "source", "to_node": "destination"}],
        },
    }
    with pytest.raises(AssertionError):
        _assert_relationship_geometry(Presentation(path).slides[0], page_plan)


def _manifest(chart: dict) -> dict:
    return {
        "workflow_contract_version": "fixed-canvas-cm-v2",
        "reconstruction_contract_version": "editable-image-v3",
        "slide": dict(SLIDE),
        "content_box": dict(CONTENT_BOX),
        "source": {"width_px": 1904, "height_px": 896},
        "text_inventory": [],
        "visual_inventory": [],
        "background_strategy": "native white body background",
        "quality_checks": {
            "font_size_calibrated": True,
            "visual_inventory_matched": True,
            "background_strategy_checked": True,
            "shape_corner_geometry_checked": True,
        },
        "text_boxes": [],
        "tables": [],
        "shapes": [],
        "images": [],
        "charts": [chart],
        "asset_provenance": [],
    }


def _qualitative_manifest(relationship: str, fallback: str) -> dict:
    manifest = _manifest({})
    manifest["charts"] = []
    layouts = {
        "driver_bridge": [(0.8 + index * 2.2, 2.4, 1.8, 0.9) for index in range(4)],
        "timeline_roadmap": [(0.8 + index * 2.2, 2.4, 1.8, 0.9) for index in range(4)],
        "qualitative_quadrant_or_comparison_table": [(1.2 + column * 4.2, 1.5 + row * 1.8, 3.2, 1.2) for row in range(2) for column in range(2)],
        "uniform_nodes": [(0.8 + index * 2.2, 2.4, 1.8, 0.9) for index in range(4)],
        "equal_width_hierarchy": [(4.1, 1.2, 1.8, 0.9), *[(1.3 + index * 2.8, 3.0, 1.8, 0.9) for index in range(3)]],
        "roadmap_milestones": [(0.8 + index * 2.2, 2.4, 1.8, 0.9) for index in range(4)],
        "comparison_table": [(1.0, 1.6, 3.7, 2.4), (5.3, 1.6, 3.7, 2.4)],
        "goal_current_gap": [(0.8 + index * 3.0, 2.1, 2.4, 1.4) for index in range(3)],
    }
    boxes = layouts[fallback]
    manifest["shapes"] = [
        {
            "object_id": f"{fallback}-node-{index}",
            "name": f"{fallback} node {index}",
            "type": "rect",
            "box_px": _box_px(left, top, width, height),
            "fill": "#EAF2F8",
            "stroke": "#6B7A90",
        }
        for index, (left, top, width, height) in enumerate(boxes)
    ]
    manifest["text_boxes"] = [
        {
            "object_id": f"{fallback}-title",
            "name": f"{fallback} title",
            "box_px": _box_px(0.7, 1.0, 8.6, 0.5),
            "text": f"{relationship}: {fallback} (qualitative, non-scaled)",
            "font_size": 22,
        },
        *[
            {
                "object_id": f"{fallback}-label-{index}",
                "name": f"{fallback} label {index}",
                "box_px": _box_px(left + 0.15, top + 0.25, width - 0.3, min(0.7, height - 0.3)),
                "text": f"Source-backed item {chr(65 + index)}",
                "font_size": 14,
                "align": "center",
            }
            for index, (left, top, width, height) in enumerate(boxes)
        ],
    ]
    return manifest


def _one_dimensional(name: str, variant: str, *, target: bool = False) -> dict:
    chart = {
        "object_id": name,
        "name": name,
        "box_px": [190, 90, 1142, 620],
        "rendering_primitive": "column_bar" if variant in {"column", "bar"} else "line_point",
        "chart_variant": variant,
        "title": name,
        "unit": "USD m",
        "period": "FY2025",
        "basis": "same portfolio companies",
        "series": [{"name": "Revenue", "categories": ["A", "B"], "values": [12, 18]}],
    }
    if target:
        chart.update({"target_value": 20, "actual_value": 18})
    return chart


def _xy(name: str, variant: str) -> dict:
    chart = {
        "object_id": name,
        "name": name,
        "box_px": [190, 90, 1142, 620],
        "rendering_primitive": "xy",
        "chart_variant": variant,
        "title": name,
        "period": "FY2025",
        "x_label": "Growth",
        "x_unit": "%",
        "x_basis": "FY2025 revenue growth",
        "y_label": "Margin",
        "y_unit": "% pts",
        "y_basis": "FY2025 EBITDA margin",
        "series": [{"name": "Companies", "x_values": [1.2, 2.4], "y_values": [8, 11]}],
    }
    if variant == "bubble":
        chart.update({"size_label": "Revenue", "size_unit": "USD m", "size_basis": "FY2025 revenue"})
        chart["series"][0]["size_values"] = [30, 50]
    return chart


def _waterfall() -> dict:
    return {
        "object_id": "waterfall", "name": "waterfall", "box_px": [190, 90, 1142, 620],
        "rendering_primitive": "cumulative_bridge", "title": "waterfall", "unit": "USD m",
        "basis": "FY2025 EBITDA", "period": "FY2025",
        "series": [{"name": "Bridge", "categories": ["Pricing", "Volume"], "start": 100, "changes": [20, -5], "end": 115}],
    }


def _gantt() -> dict:
    return {
        "object_id": "gantt", "name": "gantt", "box_px": [190, 90, 1142, 620],
        "rendering_primitive": "time_interval", "title": "gantt", "period": "September 2026",
        "series": [{"name": "Plan", "categories": ["Diligence", "IC"],
                    "start_dates": ["2026-09-01", "2026-09-11"], "end_dates": ["2026-09-10", "2026-09-15"]}],
    }


def _mekko() -> dict:
    return {
        "object_id": "variable-rectangle", "name": "variable-rectangle", "box_px": [190, 90, 1142, 620],
        "rendering_primitive": "variable_rectangle", "title": "variable rectangle", "period": "FY2025",
        "series": [{"name": "Markets", "categories": ["A", "B"], "width_values": [40, 60],
                    "width_label": "Market size", "width_unit": "USD m", "width_basis": "2025 market",
                    "share_values": [[25, 75], [40, 60]], "share_label": "Share", "share_unit": "%",
                    "share_basis": "2025 composition", "share_denominator": 100}],
    }


QUANTITATIVE_CASES = [
    _one_dimensional("column-bar", "column"),
    _one_dimensional("line", "line"),
    _xy("scatter", "scatter"),
    _xy("bubble", "bubble"),
    _one_dimensional("dot", "dot"),
    _waterfall(),
    _gantt(),
    _mekko(),
    _one_dimensional("target-actual-variance", "dot", target=True),
]


MATRIX = [
    ("drivers", _waterfall(), "driver_bridge"),
    ("time_change", _one_dimensional("time-change", "line"), "timeline_roadmap"),
    ("two_variables", _xy("two-variables", "scatter"), "qualitative_quadrant_or_comparison_table"),
    ("third_variable", _xy("third-variable", "bubble"), "uniform_nodes"),
    ("market_size_share", _mekko(), "equal_width_hierarchy"),
    ("project_stage_time", _gantt(), "roadmap_milestones"),
    ("option_comparison", _one_dimensional("options", "dot"), "comparison_table"),
    ("target_actual_variance", _one_dimensional("target", "dot", target=True), "goal_current_gap"),
]


def test_final_editable_relationship_geometry_and_numeric_authority_match_sealed_request(
    tmp_path: Path,
) -> None:
    project = _project(tmp_path / "relationship-and-numeric", 1)
    receipt_path = project / "04_v6/images/page_001.json"
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    relationship = receipt["page_plan"]["primary_relationship"]
    relationship.update({
        "grammar": "hierarchy",
        "description": "Source and destination exchange a governed resource.",
        "nodes": [
            {"node_id": "source", "label": "Source", "fact_ids": ["body-1"]},
            {"node_id": "destination", "label": "Destination", "fact_ids": ["body-1"]},
        ],
        "edges": [
            {"from_node": "source", "to_node": "destination", "label": "feeds", "fact_ids": ["body-1"]},
            {"from_node": "destination", "to_node": "source", "label": "returns", "fact_ids": ["body-1"]},
        ],
    })
    _write_signed_receipt(project, 1, receipt)

    chart = {**_one_dimensional("relationship-chart", "dot"), "relationship": "time_change"}
    materials = project / "02_v6/page_materials/page_001.json"
    materials.parent.mkdir(parents=True, exist_ok=True)
    materials.write_text(json.dumps({"chart_facts": [chart]}), encoding="utf-8")
    expected_authority = select_numeric_authority([chart])
    assert expected_authority is not None

    calls: list[dict] = []
    result = reconstruct_accepted_page(
        SimpleNamespace(project_copy=project, page_number=1),
        _accepted_outcome(project),
        page_worker=_production_worker(
            _relationship_manifest,
            calls,
        ),
    )

    request = calls[0]["request"]
    assert request["numeric_authority"] == expected_authority
    slide = Presentation(project / result["final_page"]).slides[0]
    _assert_relationship_geometry(slide, request["page_plan"])
    named = {shape.name: shape for shape in slide.shapes}
    assert named["sealed-chart Unit"].text == expected_authority["unit"]
    assert named["sealed-chart Period"].text == expected_authority["period"]
    assert [named[f"sealed-chart Category {index}"].text for index in (1, 2)] == ["A", "B"]
    assert [named[f"sealed-chart Value {index}"].text for index in (1, 2)] == ["12", "18"]


def _relationship_project(tmp_path: Path) -> tuple[Path, dict]:
    project = _project(tmp_path, 1)
    receipt_path = project / "04_v6/images/page_001.json"
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    receipt["page_plan"]["primary_relationship"]["edges"].append({
        "from_node": "destination", "to_node": "source", "label": "returns", "fact_ids": ["body-1"],
    })
    _write_signed_receipt(project, 1, receipt)
    chart = {**_one_dimensional("relationship-chart", "dot"), "relationship": "time_change"}
    materials = project / "02_v6/page_materials/page_001.json"
    materials.parent.mkdir(parents=True, exist_ok=True)
    materials.write_text(json.dumps({"chart_facts": [chart]}), encoding="utf-8")
    return project, receipt


@pytest.mark.parametrize("defect,error", [
    ("missing_node", "V6 sealed relationship node is missing or duplicated: destination"),
    ("wrong_direction", "V6 sealed relationship edge direction is invalid: edge:source->destination"),
    ("non_line", "V6 sealed relationship edge is not a real line: edge:source->destination"),
])
def test_host_finalization_rejects_tampered_sealed_relationship(
    tmp_path: Path, defect: str, error: str,
) -> None:
    project, _receipt = _relationship_project(tmp_path / defect)

    def tamper_built_pptx(path: Path) -> None:
        deck = Presentation(path)
        slide = deck.slides[0]
        if defect == "missing_node":
            node = next(item for item in slide.shapes if item.name == "destination")
            node._element.getparent().remove(node._element)
        elif defect == "wrong_direction":
            edge = next(item for item in slide.shapes if item.name == "edge:source->destination")
            edge._element.xpath(".//a:xfrm")[0].set("flipH", "1")
        else:
            edge = next(item for item in slide.shapes if item.name == "edge:source->destination")
            edge._element.xpath(".//a:prstGeom")[0].set("prst", "rect")
        deck.save(path)

    with pytest.raises(ValueError) as exc_info:
        reconstruct_accepted_page(
            SimpleNamespace(project_copy=project, page_number=1),
            _accepted_outcome(project),
            page_worker=_production_worker(
                _relationship_manifest, [], post_validate=tamper_built_pptx,
            ),
        )
    assert str(exc_info.value) == error


@pytest.mark.parametrize("shape_name", [
    "sealed-chart Unit", "sealed-chart Period", "sealed-chart Category 1",
])
def test_finalization_rejects_missing_numeric_label_unit_or_period(
    tmp_path: Path, shape_name: str,
) -> None:
    project, _receipt = _relationship_project(tmp_path / shape_name.replace(" ", "-"))

    def remove_authority_text(path: Path) -> None:
        deck = Presentation(path)
        shape = next(item for item in deck.slides[0].shapes if item.name == shape_name)
        shape.text = ""
        deck.save(path)

    with pytest.raises(ValueError, match="sealed numeric authority"):
        reconstruct_accepted_page(
            SimpleNamespace(project_copy=project, page_number=1),
            _accepted_outcome(project),
            page_worker=_production_worker(
                _relationship_manifest,
                [],
                post_validate=remove_authority_text,
            ),
        )


@pytest.mark.parametrize("missing", [
    "accepted_reconstruction_request.json", "manifest.json",
])
def test_finalization_fails_closed_when_sealed_worker_artifact_is_missing(
    tmp_path: Path, missing: str,
) -> None:
    project, _receipt = _relationship_project(tmp_path / missing.removesuffix(".json"))

    def remove_authority(path: Path) -> None:
        deck = Presentation(path)
        node = next(item for item in deck.slides[0].shapes if item.name == "destination")
        node._element.getparent().remove(node._element)
        deck.save(path)
        (path.parent / missing).unlink()

    with pytest.raises(ValueError, match="sealed reconstruction"):
        reconstruct_accepted_page(
            SimpleNamespace(project_copy=project, page_number=1),
            _accepted_outcome(project),
            page_worker=_production_worker(
                _relationship_manifest,
                [],
                post_validate=remove_authority,
            ),
        )


def test_nine_quantitative_cases_cover_ten_visual_marks_in_editable_pptx_and_previews() -> None:
    OUTPUT.mkdir(parents=True, exist_ok=True)
    entries = []
    for index, chart in enumerate(QUANTITATIVE_CASES, start=1):
        manifest = _manifest(chart)
        manifest_path = OUTPUT / f"quantitative-{index:02d}.json"
        manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
        preview = OUTPUT / f"quantitative-{index:02d}.png"
        render_preview(manifest, manifest_path, preview)
        assert Image.open(preview).size == (1200, 675)
        entries.append({"manifest": manifest, "manifest_path": str(manifest_path)})

    deck_path = OUTPUT / "quantitative-nine-cases.pptx"
    write_deck({"workflow_contract_version": "fixed-canvas-cm-v2", "slide": dict(SLIDE)}, entries, deck_path, [])
    deck = Presentation(deck_path)

    assert len(deck.slides) == 9
    assert all(any(shape.has_chart for shape in deck.slides[index].shapes) for index in range(4))
    assert not any(shape.has_chart for shape in deck.slides[4].shapes)  # dot uses editable shapes
    assert not any(shape.has_chart for index in range(5, 8) for shape in deck.slides[index].shapes)
    for index in (4, 5, 6, 7, 8):
        assert any(shape.shape_type != 13 for shape in deck.slides[index].shapes)
    comparison_names = {shape.name for shape in deck.slides[8].shapes}
    assert "target-actual-variance Target Line" in comparison_names
    assert "target-actual-variance Difference Arrow" in comparison_names
    assert "target-actual-variance Difference" in comparison_names


def test_eight_relationships_use_real_selector_reconstruction_and_renderer_contracts(
    tmp_path: Path,
) -> None:
    OUTPUT.mkdir(parents=True, exist_ok=True)
    entries = []
    for index, (relationship, quantitative, fallback) in enumerate(MATRIX, start=1):
        quantitative = deepcopy(quantitative)
        quantitative["relationship"] = relationship
        authority = select_numeric_authority([quantitative])
        assert authority is not None
        assert authority["rendering_primitive"] == quantitative["rendering_primitive"]
        quantitative_project = _project(tmp_path / f"{relationship}-quantitative", 1)
        quantitative_materials = quantitative_project / "02_v6/page_materials/page_001.json"
        quantitative_materials.parent.mkdir(parents=True, exist_ok=True)
        quantitative_materials.write_text(json.dumps({"chart_facts": [quantitative]}, ensure_ascii=False), encoding="utf-8")
        quantitative_calls: list[dict] = []
        reconstruct_accepted_page(
            SimpleNamespace(project_copy=quantitative_project, page_number=1),
            _accepted_outcome(quantitative_project),
            page_worker=_production_worker(
                _relationship_manifest, quantitative_calls,
            ),
        )
        assert quantitative_calls[0]["request"]["numeric_authority"] == authority
        assert "sealed numeric authority owns quantitative mark size" in quantitative_calls[0]["prompt"]
        assert (quantitative_calls[0]["page_dir"] / "page.pptx").is_file()

        qualitative = {
            "title": relationship,
            "relationship": relationship,
            "source_wording": "Source describes the relationship but supplies no complete numeric dimensions.",
            "disabled_primitive": quantitative["rendering_primitive"],
            "fallback": fallback,
            "series": [],
        }
        assert select_numeric_authority([qualitative]) is None
        project = _project(tmp_path / f"{relationship}-qualitative", 1)
        material_path = project / "02_v6/page_materials/page_001.json"
        material_path.parent.mkdir(parents=True, exist_ok=True)
        material_path.write_text(
            json.dumps({"chart_facts": [qualitative]}, ensure_ascii=False), encoding="utf-8",
        )
        request = build_reconstruction_request(project, page_number=1)
        assert "numeric_authority" not in request
        qualitative_calls: list[dict] = []
        reconstruct_accepted_page(
            SimpleNamespace(project_copy=project, page_number=1),
            _accepted_outcome(project),
            page_worker=_production_worker(
                lambda page_request, r=relationship, f=fallback: _with_relationship(
                    _qualitative_manifest(r, f), page_request,
                ),
                qualitative_calls,
            ),
        )
        assert "numeric_authority" not in qualitative_calls[0]["request"]
        assert "No sealed numeric authority is present" in qualitative_calls[0]["prompt"]
        manifest = qualitative_calls[0]["manifest"]
        manifest_path = OUTPUT / f"qualitative-{index:02d}.json"
        manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
        preview = OUTPUT / f"qualitative-{index:02d}.png"
        render_preview(manifest, manifest_path, preview)
        assert Image.open(preview).size == (1200, 675)
        entries.append({"manifest": manifest, "manifest_path": str(manifest_path)})

    path = OUTPUT / "qualitative-eight-modes.pptx"
    write_deck({"workflow_contract_version": "fixed-canvas-cm-v2", "slide": dict(SLIDE)}, entries, path, [])

    reopened = Presentation(path)
    assert len(reopened.slides) == 8
    assert all(not any(shape.has_chart for shape in slide.shapes) for slide in reopened.slides)
    with zipfile.ZipFile(path) as package:
        assert not any(name.startswith("ppt/charts/") for name in package.namelist())
        assert not any(b"<c:chart" in package.read(name) for name in package.namelist() if name.endswith(".xml"))
    expected_counts = (4, 4, 4, 4, 4, 4, 2, 3)
    for slide, (relationship, _quantitative, fallback), expected_count in zip(reopened.slides, MATRIX, expected_counts, strict=True):
        names = {shape.name for shape in slide.shapes}
        assert any(fallback in name for name in names)
        nodes = [shape for shape in slide.shapes if shape.name.startswith(f"{fallback} node")]
        assert len(nodes) == expected_count
        assert len({shape.width for shape in nodes}) == 1
        assert len({shape.height for shape in nodes}) == 1
        assert len({shape.width * shape.height for shape in nodes}) == 1
        assert all(shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.auto_shape_type == MSO_SHAPE.RECTANGLE for shape in nodes)
        assert {"source", "destination"} <= {shape.name for shape in slide.shapes}
        assert all(shape.shape_type != MSO_SHAPE_TYPE.CHART for shape in slide.shapes)
        slide_text = "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
        assert not re.search(r"\d", slide_text)
        if relationship in {"drivers", "time_change", "third_variable", "project_stage_time", "goal_current_gap"}:
            assert len({shape.top for shape in nodes}) == 1
        if relationship == "two_variables":
            assert len({shape.left for shape in nodes}) == 2 and len({shape.top for shape in nodes}) == 2
        if relationship == "market_size_share":
            assert len({shape.top for shape in nodes}) == 2 and len([shape for shape in nodes if shape.top == min(node.top for node in nodes)]) == 1
