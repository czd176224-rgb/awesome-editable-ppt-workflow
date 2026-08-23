"""Render confirmed cover, TOC, section, and closing pages as native PPT objects."""

from __future__ import annotations

import hashlib
import json
import math
import re
import uuid
from pathlib import Path
from typing import Any, Mapping
from xml.etree import ElementTree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Cm, Pt

import workflow_v6_secure_io as secure_io
from workflow_v6_composition import validate_composition
from workflow_v6_contract import validate_project
from workflow_v6_state import mutation_lock, save


SPECIAL_ROLES = frozenset({"cover", "toc", "section", "closing"})


def _color(value: Any, fallback: str) -> RGBColor:
    text = value if isinstance(value, str) and re.fullmatch(r"#[0-9A-Fa-f]{6}", value) else fallback
    return RGBColor.from_string(text[1:].upper())


def _add_text(slide, *, name: str, text: str, box: tuple[float, float, float, float],
              font: str, size: float, color: str, bold: bool = False,
              align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(*(Cm(value) for value in box))
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Cm(0.05)
    frame.margin_top = frame.margin_bottom = Cm(0.03)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _color(color, "#17365D")
    return shape


def _add_rule(slide, *, name: str, box: tuple[float, float, float, float], color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *(Cm(value) for value in box))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = _color(color, "#C7352B")
    shape.line.fill.background()


def _svg_aspect_ratio(payload: bytes) -> float:
    error = "locked logo SVG viewBox must contain exactly four finite numbers with positive width and height"
    try:
        root = ElementTree.fromstring(payload)
        values = re.split(r"[\s,]+", root.get("viewBox", "").strip())
        numbers = [float(value) for value in values]
    except (ElementTree.ParseError, IndexError, ValueError) as exc:
        raise ValueError(error) from exc
    if len(numbers) != 4 or not all(math.isfinite(value) for value in numbers):
        raise ValueError(error)
    width, height = numbers[2], numbers[3]
    if width <= 0 or height <= 0:
        raise ValueError(error)
    return width / height


def _add_logo(slide, payload: bytes, *, box: tuple[float, float, float, float]) -> None:
    digest = hashlib.sha256(payload).hexdigest()[:16]
    partname = PackURI(f"/ppt/media/special-logo-{digest}.svg")
    package = slide.part.package
    part = next((item for item in package.iter_parts() if item.partname == partname), None)
    if part is None:
        part = Part(partname, "image/svg+xml", package, payload)
    relationship = slide.part.relate_to(part, RT.IMAGE)
    x, y, maximum_width, maximum_height = box
    ratio = _svg_aspect_ratio(payload)
    width, height = maximum_width, maximum_width / ratio
    if height > maximum_height:
        height, width = maximum_height, maximum_height * ratio
    x += maximum_width - width
    y += (maximum_height - height) / 2
    picture = parse_xml(
        f'<p:pic {nsdecls("a", "p", "r")}>'
        f'<p:nvPicPr><p:cNvPr id="{slide.shapes._next_shape_id}" name="special-logo"/>'
        '<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
        f'<p:blipFill><a:blip r:embed="{relationship}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        f'<p:spPr><a:xfrm><a:off x="{int(Cm(x))}" y="{int(Cm(y))}"/>'
        f'<a:ext cx="{int(Cm(width))}" cy="{int(Cm(height))}"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
    )
    slide.shapes._spTree.insert_element_before(picture, "p:extLst")


def _add_page_number(slide, page_number: int, *, font: str, color: str) -> None:
    _add_text(
        slide, name="special-page-number", text=str(page_number),
        box=(22.8, 13.35, 1.1, 0.45), font=font, size=9, color=color,
        align=PP_ALIGN.RIGHT,
    )


def _read_json(root: Path, relative: Path) -> dict[str, Any]:
    value = json.loads(secure_io.read_bytes(root, relative).decode("utf-8"))
    if not isinstance(value, dict):
        raise ValueError(f"expected JSON object: {relative.as_posix()}")
    return value


def _block_text(block: Mapping[str, Any]) -> str:
    if block.get("type") in {"paragraph", "list"}:
        value = block.get("text")
        return value.strip() if isinstance(value, str) else ""
    if block.get("type") == "table" and isinstance(block.get("rows"), list):
        return "\n".join(" | ".join(str(cell) for cell in row) for row in block["rows"] if isinstance(row, list)).strip()
    value = block.get("markdown")
    return value.strip() if isinstance(value, str) else ""


def _style(contract: Mapping[str, Any]) -> tuple[str, str, str, str, str, float, float]:
    typography = contract.get("typography") if isinstance(contract.get("typography"), Mapping) else {}
    heading = typography.get("heading") if isinstance(typography.get("heading"), Mapping) else {}
    body = typography.get("body") if isinstance(typography.get("body"), Mapping) else {}
    scale = typography.get("type_scale_pt") if isinstance(typography.get("type_scale_pt"), Mapping) else {}
    color = contract.get("color") if isinstance(contract.get("color"), Mapping) else {}
    palette = color.get("palette") if isinstance(color.get("palette"), Mapping) else {}
    return (
        str(heading.get("cjk") or contract.get("cjk_font") or "Microsoft YaHei"),
        str(body.get("cjk") or contract.get("cjk_font") or "Microsoft YaHei"),
        str(palette.get("primary") or contract.get("primary_color") or "#17365D"),
        str(palette.get("accent") or contract.get("secondary_color") or "#C7352B"),
        str(palette.get("background") or contract.get("background_color") or "#FFFFFF"),
        float(scale.get("page_title") or contract.get("title_size_pt") or 28),
        float(scale.get("body") or contract.get("body_size_pt") or 16),
    )


def _row_heights(
    texts: list[str], *, width: float, font_size: float, available: float,
) -> tuple[list[float], int | None]:
    gap = 0.18
    characters_per_line = max(8, int(width / (font_size * 0.026)))
    heights: list[float] = []
    used = 0.0
    for index, text in enumerate(texts):
        line_count = sum(
            max(1, math.ceil(len(line) / characters_per_line))
            for line in text.splitlines() or [""]
        )
        height = max(0.7, line_count * font_size * 0.043 + 0.12)
        required = height + (gap if heights else 0.0)
        if used + required > available:
            return heights, index
        heights.append(height)
        used += required
    return heights, None


def render_special_page(project: Path, page_number: int) -> dict[str, Any]:
    """Render one frozen special page and mark its canonical page package complete."""
    secure_io.reject_reparse_chain(Path(project))
    root = Path(project).resolve(strict=True)
    state = _read_json(root, Path("workflow_v6.json"))
    validate_project(state)
    if type(page_number) is not int or not 1 <= page_number <= len(state["pages"]):
        raise ValueError("V6 special page number is out of range")
    composition = _read_json(root, Path("02_v6/page_composition.json"))
    validate_composition(composition, confirmed=True)
    page = composition["pages"][page_number - 1]
    role = page["page_role"]
    if role not in SPECIAL_ROLES:
        raise ValueError("V6 page is not a native special role")
    if page["output_page_number"] != page_number:
        raise ValueError("V6 special page identity is invalid")
    required_ids = list(page["material_source_block_ids"])
    if len(required_ids) != len(set(required_ids)):
        raise ValueError("V6 special page contains duplicate source block IDs")

    source = _read_json(root, Path("02_v6/paginated_word_source.json"))
    allowed_ids = set(required_ids)
    traced: dict[str, str] = {}
    occurrences: set[str] = set()
    for source_page in source.get("pages", []):
        if not isinstance(source_page, Mapping):
            continue
        for block in source_page.get("blocks", []):
            if not isinstance(block, Mapping):
                continue
            block_id = block.get("source_block_id")
            text = _block_text(block)
            if block_id in allowed_ids:
                if block_id in occurrences:
                    raise ValueError(f"duplicate traced source block ID: {block_id}")
                occurrences.add(block_id)
                if text:
                    traced[block_id] = text
    if set(traced) != allowed_ids:
        raise ValueError("V6 special page has no traced display material")
    lines = [traced[block_id] for block_id in required_ids]
    continuation = str(page.get("composition_page_id") or "").startswith("toc-continuation:")
    capacity = (12 if continuation else 13) if role == "toc" else 6
    if len(required_ids) > capacity:
        overflow = ", ".join(required_ids[capacity:])
        raise ValueError(
            f"V6 {role} page {page_number} overflow source block IDs: {overflow}"
        )
    contract = state.get("style_confirmation", {}).get("contract")
    if not isinstance(contract, Mapping):
        raise ValueError("V6 confirmed visual contract is missing")
    heading_font, body_font, primary, accent, background, title_size, body_size = _style(contract)
    title_size = max(title_size, 24)
    body_size = max(body_size, 12)
    chapter_title = page["chapter_title"].strip()
    section_start = 2 if (
        role == "section"
        and len(required_ids) > 1
        and chapter_title
        and chapter_title != page["fixed_page_title"].strip()
    ) else 1
    row_heights: list[float] = []
    toc_start = 0 if continuation else 1
    toc_two_columns = role == "toc" and len(lines[toc_start:]) > 6
    toc_row_heights: list[float] = []
    if role in {"cover", "closing", "section"}:
        start = section_start if role == "section" else 1
        width = 18.5 if role == "section" else 16.0
        available = {"cover": 4.85, "section": 4.55, "closing": 4.95}[role]
        row_heights, overflow_index = _row_heights(
            lines[start:], width=width, font_size=body_size, available=available,
        )
        if overflow_index is not None:
            overflow = ", ".join(required_ids[start + overflow_index:])
            raise ValueError(
                f"V6 {role} page {page_number} overflow source block IDs: {overflow}"
            )
    elif role == "toc":
        entry_width = 10.0 if toc_two_columns else 20.2
        entries = lines[toc_start:]
        groups = [entries[:6], entries[6:]] if toc_two_columns else [entries]
        for column, group in enumerate(groups):
            heights, overflow_index = _row_heights(
                group, width=entry_width, font_size=body_size, available=8.4,
            )
            if overflow_index is not None:
                item_index = column * 6 + overflow_index
                raise ValueError(
                    f"V6 toc page {page_number} overflow source block IDs: "
                    f"{required_ids[toc_start + item_index]}"
                )
            toc_row_heights.extend(heights)
    logo_relative = Path(state["logo_source"].get("path", ""))
    if logo_relative != Path("00_source/logo.svg"):
        raise ValueError("V6 locked logo path must be 00_source/logo.svg")
    logo_path = root / logo_relative
    secure_io.reject_reparse_chain(logo_path)
    logo_payload = secure_io.read_bytes(root, logo_relative)

    deck = Presentation()
    deck.slide_width, deck.slide_height = Cm(25.4), Cm(14.288)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _color(background, "#FFFFFF")
    title = page["fixed_page_title"].strip()
    displayed_ids: list[str] = []
    layout_variant = role

    def add_source(block_id: str, *, text: str,
                   box: tuple[float, float, float, float], font: str,
                   size: float, color: str, bold: bool = False) -> None:
        _add_text(
            slide, name=f"special-source-{block_id}", text=text, box=box,
            font=font, size=size, color=color, bold=bold,
        )
        displayed_ids.append(block_id)

    if role == "cover":
        _add_rule(slide, name="special-accent-rule", box=(2.0, 2.65, 3.4, 0.12), color=accent)
        add_source(required_ids[0], text=title, box=(2.0, 3.0, 18.5, 2.4), font=heading_font, size=max(title_size, 30), color=primary, bold=True)
        y = 5.65
        for block_id, text, height in zip(required_ids[1:], lines[1:], row_heights):
            add_source(block_id, text=text, box=(2.0, y, 16.0, height), font=body_font, size=body_size, color=primary)
            y += height + 0.18
        _add_logo(slide, logo_payload, box=(19.0, 11.5, 4.3, 1.2))
    elif role == "toc":
        entry_ids = required_ids if continuation else required_ids[1:]
        entry_lines = lines if continuation else lines[1:]
        two_columns = toc_two_columns
        layout_variant = (
            f"toc-continuation-{'two' if two_columns else 'one'}-column"
            if continuation else f"toc-{'two' if two_columns else 'one'}-column"
        )
        _add_rule(slide, name="special-accent-rule", box=(1.8, 2.45, 21.8, 0.08), color=accent)
        if continuation:
            _add_text(slide, name="special-toc-title", text=title, box=(1.8, 1.1, 21.8, 1.15), font=heading_font, size=max(title_size, 28), color=primary, bold=True)
        else:
            add_source(required_ids[0], text=title, box=(1.8, 1.1, 21.8, 1.15), font=heading_font, size=max(title_size, 28), color=primary, bold=True)
        y_by_column = [3.0, 3.0]
        for index, (block_id, text, height) in enumerate(zip(entry_ids, entry_lines, toc_row_heights)):
            if two_columns:
                column = index // 6
                box = (2.0 + column * 11.0, y_by_column[column], 10.0, height)
            else:
                column = 0
                box = (2.2, y_by_column[0], 20.2, height)
            add_source(block_id, text=text, box=box, font=body_font, size=body_size, color=primary)
            y_by_column[column] += height + 0.18
    elif role == "section":
        start = section_start
        if start == 2:
            add_source(required_ids[0], text=title, box=(2.0, 2.8, 18.0, 0.75), font=body_font, size=body_size, color=accent, bold=True)
            add_source(required_ids[1], text=chapter_title, box=(2.0, 3.7, 20.0, 2.2), font=heading_font, size=max(title_size, 30), color=primary, bold=True)
        else:
            add_source(required_ids[0], text=title, box=(2.0, 3.4, 20.0, 2.2), font=heading_font, size=max(title_size, 30), color=primary, bold=True)
        _add_rule(slide, name="special-accent-rule", box=(2.0, 6.05, 4.2, 0.1), color=accent)
        y = 6.45
        for block_id, text, height in zip(required_ids[start:], lines[start:], row_heights):
            add_source(block_id, text=text, box=(2.0, y, 18.5, height), font=body_font, size=body_size, color=accent)
            y += height + 0.18
    else:
        _add_rule(slide, name="special-accent-rule", box=(2.0, 3.35, 3.4, 0.12), color=accent)
        add_source(required_ids[0], text=title, box=(2.0, 3.75, 19.5, 2.5), font=heading_font, size=max(title_size, 30), color=primary, bold=True)
        y = 6.5
        for block_id, text, height in zip(required_ids[1:], lines[1:], row_heights):
            add_source(block_id, text=text, box=(2.0, y, 16.0, height), font=body_font, size=body_size, color=primary)
            y += height + 0.18
        _add_logo(slide, logo_payload, box=(19.0, 11.5, 4.3, 1.2))
    if displayed_ids != required_ids:
        raise ValueError(f"V6 special page source/display mismatch: required={required_ids}, displayed={displayed_ids}")
    if page["visible_page_number"]:
        _add_page_number(slide, page_number, font=body_font, color=primary)

    output = root / "06_v6" / "pages" / f"page_{page_number:03d}" / "page.pptx"
    temporary = output.with_name(f".special-{uuid.uuid4().hex[:8]}.pptx")
    with secure_io.hold_parent(root, temporary.relative_to(root), create=True):
        deck.save(temporary)
        payload = temporary.read_bytes()
    secure_io.atomic_write_bytes(root, output.relative_to(root), payload, replace=output.exists())
    temporary.unlink(missing_ok=True)
    receipt = {
        "artifact_version": "special-page-v6",
        "status": "page_complete",
        "page_number": page_number,
        "page_role": role,
        "page_pptx": output.relative_to(root).as_posix(),
        "sha256": hashlib.sha256(payload).hexdigest(),
        "visible_page_number": page["visible_page_number"],
        "displayed_source_block_ids": displayed_ids,
        "layout_variant": layout_variant,
    }
    secure_io.atomic_write_bytes(
        root, output.with_name("page.json").relative_to(root),
        (json.dumps(receipt, ensure_ascii=False, indent=2) + "\n").encode("utf-8"),
        replace=output.with_name("page.json").exists(),
    )
    with mutation_lock(root):
        current = _read_json(root, Path("workflow_v6.json"))
        current["pages"][page_number - 1]["state"] = "page_complete"
        save(root, current)
    return receipt


__all__ = ["SPECIAL_ROLES", "render_special_page"]
