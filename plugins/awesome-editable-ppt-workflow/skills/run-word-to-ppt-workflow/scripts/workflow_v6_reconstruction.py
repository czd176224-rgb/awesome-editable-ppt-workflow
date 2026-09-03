"""V6 reconstruction requests, fixed-layer finalization, and deck assembly."""

from __future__ import annotations

import hashlib
import json
import os
import shutil
import sys
import uuid
import zipfile
import copy
import colorsys
from io import BytesIO
from pathlib import Path
from typing import Any, Mapping
from xml.etree import ElementTree

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI

from fixed_frame import apply_fixed_frame, inspect_fixed_frame
from fixed_region_contract import fixed_frame_execution
from director_taskbook import project_emphasis_pages
from workflow_v6_contract import geometry_contract, transition_page, validate_project
from workflow_v6_materials import select_numeric_authority
from workflow_v6_composition import load_composition_authority
from workflow_v6_state import mutation_lock, save
import workflow_v6_secure_io as secure_io
from complex_page_experiment import (
    open_live_page_workspace,
    verify_signed_acceptance_receipt,
)


_EDITPPT_CLI = (
    Path(__file__).resolve().parents[2]
    / "reconstruct-editable-slide"
    / "cli"
)
if _EDITPPT_CLI.is_dir() and str(_EDITPPT_CLI) not in sys.path:
    sys.path.insert(0, str(_EDITPPT_CLI))
from editppt.runtime.validate_pptx import (  # noqa: E402
    _connector_endpoints,
    _shape_arrowheads,
    _shape_kind,
    quantitative_chart_readback_violations,
)


_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_RELATIONSHIP_ATTRIBUTES = {f"{{{_R}}}embed", f"{{{_R}}}id", f"{{{_R}}}link"}


def _require_current_fixed_frame(slide) -> None:
    names = [shape.name for shape in slide.shapes if shape.name.startswith("fixed-frame-")]
    if sorted(names) != sorted((
        "fixed-frame-title", "fixed-frame-logo", "fixed-frame-footer", "fixed-frame-page-number",
    )):
        raise ValueError("assembled V6 fixed-layer inventory is incorrect")


def _require_special_page_shapes(slide, role: str, visible_page_number: bool) -> None:
    names = [shape.name for shape in slide.shapes]
    if any(name.startswith("fixed-frame-") for name in names):
        raise ValueError("assembled V6 special page contains a content fixed frame")
    if f"special-{role}-title" not in names and not any(
        name.startswith("special-source-") for name in names
    ):
        raise ValueError("assembled V6 special-page title is missing")
    if ("special-page-number" in names) != visible_page_number:
        raise ValueError("assembled V6 special-page number policy is incorrect")
    if role in {"cover", "closing"} and names.count("special-logo") != 1:
        raise ValueError("assembled V6 special-page logo is missing")


def _copy_image_relationship(relationship, destination_slide) -> str:
    source_part = relationship.target_part
    if source_part.content_type != "image/svg+xml":
        _part, new_id = destination_slide.part.get_or_add_image_part(BytesIO(source_part.blob))
        return new_id
    digest = hashlib.sha256(source_part.blob).hexdigest()[:16]
    partname = PackURI(f"/ppt/media/fixed-logo-{digest}.svg")
    package = destination_slide.part.package
    target = next((part for part in package.iter_parts() if part.partname == partname), None)
    if target is None:
        target = Part(partname, "image/svg+xml", package, source_part.blob)
    return destination_slide.part.relate_to(target, RT.IMAGE)


def _copy_page_slide(source_path: Path, destination: Presentation, destination_layout, page_number: int) -> None:
    source = Presentation(source_path)
    if len(source.slides) != 1:
        raise ValueError("V6 reconstructed page package must contain one slide")
    if source.slide_width != destination.slide_width or source.slide_height != destination.slide_height:
        raise ValueError("all V6 reconstructed pages must share one slide size")
    source_slide = source.slides[0]
    destination_slide = destination.slides.add_slide(destination_layout)
    mapping: dict[str, str] = {}
    for relationship in source_slide.part.rels.values():
        if relationship.reltype == RT.SLIDE_LAYOUT:
            continue
        if relationship.reltype == RT.IMAGE:
            mapping[relationship.rId] = _copy_image_relationship(relationship, destination_slide)
        elif relationship.is_external:
            mapping[relationship.rId] = destination_slide.part.relate_to(
                relationship.target_ref, relationship.reltype, is_external=True
            )
        else:
            raise ValueError(f"unsupported V6 slide relationship: {relationship.reltype}")
    copied_content = copy.deepcopy(source_slide.element.cSld)
    copied_content.set("name", f"editable-ppt-v6-page:{page_number}")
    for node in copied_content.iter():
        for attribute, old_id in tuple(node.attrib.items()):
            if attribute in _RELATIONSHIP_ATTRIBUTES:
                if old_id not in mapping:
                    raise ValueError(f"unresolved V6 slide relationship: {old_id}")
                node.set(attribute, mapping[old_id])
    destination_slide.element.replace(destination_slide.element.cSld, copied_content)
    source_color_map = source_slide.element.clrMapOvr
    if source_color_map is not None:
        destination_color_map = destination_slide.element.clrMapOvr
        copied_color_map = copy.deepcopy(source_color_map)
        if destination_color_map is None:
            destination_slide.element.append(copied_color_map)
        else:
            destination_slide.element.replace(destination_color_map, copied_color_map)


def _sha256(path: Path) -> str:
    try:
        root = _project_for_artifact(path)
    except ValueError:
        data = path.read_bytes()
    else:
        data = secure_io.read_bytes(root, path.relative_to(root))
    return hashlib.sha256(data).hexdigest()


def _read_json(path: Path) -> dict[str, Any]:
    root = _project_for_artifact(path)
    value = json.loads(secure_io.read_bytes(root, path.relative_to(root)).decode("utf-8"))
    if not isinstance(value, dict):
        raise ValueError(f"expected JSON object: {path}")
    return value


def _write_json(path: Path, value: Mapping[str, Any]) -> None:
    root = _project_for_artifact(path)
    secure_io.atomic_write_bytes(
        root, path.relative_to(root),
        (json.dumps(value, ensure_ascii=False, indent=2) + "\n").encode("utf-8"),
        replace=path.exists(),
    )


def _project_for_artifact(path: Path) -> Path:
    for marker in ("02_v6", "04_v6", "05_v6", "06_v6", "08_final"):
        if marker in path.parts:
            return Path(*path.parts[:path.parts.index(marker)])
    raise ValueError("reconstruction artifact is outside canonical project storage")


def _load_reconstruction_state(root: Path) -> dict[str, Any]:
    value = json.loads(
        secure_io.read_bytes(root, Path("workflow_v6.json")).decode("utf-8")
    )
    if not isinstance(value, dict):
        raise ValueError("V6 state root must be an object")
    validate_project(value)
    return value


def _update_reconstruction_page(
    root: Path, page_number: int, page: Mapping[str, Any],
) -> Path:
    with mutation_lock(root):
        state = _load_reconstruction_state(root)
        if page_number < 1 or page_number > len(state["pages"]):
            raise ValueError("V6 page number is out of range")
        if page.get("page_number") != page_number:
            raise ValueError("V6 page update identity is invalid")
        state["pages"][page_number - 1] = dict(page)
        return save(root, state)


def _fixed_frame_style(style: Mapping[str, Any]) -> Mapping[str, Any]:
    """Map the current confirmed UI contract onto the existing fixed-frame input."""
    if isinstance(style.get("fixed_frame"), Mapping) and isinstance(
        style.get("hard_constraints"), Mapping
    ):
        return style
    cjk_font = style.get("cjk_font")
    title_size = style.get("title_size_pt")
    title_color = style.get("primary_color")
    if (
        not isinstance(cjk_font, str)
        or not cjk_font.strip()
        or not isinstance(title_size, (int, float))
        or float(title_size) <= 0
        or not isinstance(title_color, str)
    ):
        raise ValueError("V6 confirmed UI contract cannot drive the fixed frame")
    return {
        **dict(style),
        "fixed_frame": {"title_color": title_color, **fixed_frame_execution()},
        "hard_constraints": {
            "title_color": title_color,
            "typography": {
                "heading": {"cjk": cjk_font.strip()},
                "type_scale_pt": {"page_title": float(title_size)},
            },
        },
    }


def build_reconstruction_request(project: Path, *, page_number: int) -> dict[str, Any]:
    secure_io.reject_reparse_chain(Path(project))
    root = Path(project).resolve()
    state = _load_reconstruction_state(root)
    if page_number < 1 or page_number > len(state["pages"]):
        raise ValueError("V6 page number is out of range")
    page = state["pages"][page_number - 1]
    if page["state"] != "accepted":
        raise ValueError("V6 page must have a selected Image2 body before reconstruction")
    page_selected = page.get("selected_candidate")
    if not isinstance(page_selected, Mapping):
        raise ValueError("V6 selected body is missing")
    accepted_receipt = root / "04_v6" / "images" / f"page_{page_number:03d}.json"
    if not accepted_receipt.is_file():
        raise ValueError("V6 accepted Image2 receipt is missing")
    receipt_bytes = secure_io.read_bytes(root, accepted_receipt.relative_to(root))
    workspace = open_live_page_workspace(root, page_number)
    receipt = verify_signed_acceptance_receipt(workspace, receipt_bytes)
    if not isinstance(receipt, Mapping) or receipt.get("page_number") != page_number:
        raise ValueError("V6 accepted Image2 receipt identity is invalid")
    selected = receipt.get("candidate", receipt.get("selected"))
    if not isinstance(selected, Mapping):
        raise ValueError("V6 accepted Image2 receipt has no selected body")
    relative_value = selected.get("path")
    receipt_image_digest = selected.get("output_sha256", selected.get("sha256"))
    if not isinstance(relative_value, str) or not isinstance(receipt_image_digest, str):
        raise ValueError("V6 accepted Image2 receipt image authority is incomplete")
    relative = Path(relative_value)
    if (
        relative.is_absolute()
        or ".." in relative.parts
        or relative.suffix.casefold() != ".png"
        or tuple(relative.parts[:2]) != ("04_v6", "images")
    ):
        raise ValueError("V6 accepted body path is invalid")
    if any(
        page_selected.get(field) != selected.get(field)
        for field in ("attempt", "path")
    ) or (
        "operation" in page_selected
        and page_selected.get("operation") != selected.get("operation")
    ) or (
        "sha256" in page_selected
        and page_selected.get("sha256") != receipt_image_digest
    ):
        raise ValueError("V6 accepted receipt and selected body do not match")
    image_bytes = secure_io.read_bytes(root, relative)
    image_digest = hashlib.sha256(image_bytes).hexdigest()
    if image_digest != receipt_image_digest:
        raise ValueError("V6 accepted body digest does not match its receipt")
    try:
        with Image.open(BytesIO(image_bytes)) as decoded:
            image_format = decoded.format
            image_size = decoded.size
    except OSError as exc:
        raise ValueError("V6 accepted body is not a readable PNG") from exc
    if image_format != "PNG" or image_size != (1904, 896):
        raise ValueError("V6 accepted body must be a 1904x896 PNG")
    repairs = receipt.get("reconstruction_repairs", [])
    if not isinstance(repairs, list) or any(
        not isinstance(item, Mapping)
        or not isinstance(item.get("category"), str)
        or not item["category"].strip()
        or not isinstance(item.get("detail"), str)
        or not item["detail"].strip()
        for item in repairs
    ):
        raise ValueError("V6 accepted text repair authority is invalid")
    request = {
        "artifact_version": "reconstruction-request-v6",
        "workflow_contract_version": "awesome-word-ppt-workflow-v1",
        "operation": "reconstruct_editable_slide",
        "page_number": page_number,
        "page_title": page["title"],
        "accepted_receipt": {
            "path": accepted_receipt.relative_to(root).as_posix(),
            "sha256": hashlib.sha256(receipt_bytes).hexdigest(),
        },
        "source_body": {
            "path": relative.as_posix(),
            "sha256": image_digest,
            "pixels": {"width": 1904, "height": 896},
        },
        "sealed_image_edits": [],
        "sealed_text_repairs": [dict(item) for item in repairs],
        "page_plan": receipt["page_plan"],
        "geometry": geometry_contract(),
        "requirements": {
            "object_level_editable": True,
            "body_only": True,
            "fixed_layers_added_after_reconstruction": True,
            "post_reconstruction_visual_qa": False,
            "exact_reference_material_custody": False,
        },
    }
    materials_path = root / "02_v6" / "page_materials" / f"page_{page_number:03d}.json"
    if materials_path.is_file():
        materials = json.loads(
            secure_io.read_bytes(root, materials_path.relative_to(root)).decode("utf-8")
        )
        authority = select_numeric_authority(materials.get("chart_facts", []))
        if authority:
            request["numeric_authority"] = authority
    path = root / "05_v6" / "reconstruction_requests" / f"page_{page_number:03d}.json"
    _write_json(path, request)
    return request


def _editable_slide_text(deck: Presentation) -> str:
    values: list[str] = []
    for shape in deck.slides[0].shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            values.append(shape.text)
        if getattr(shape, "has_table", False):
            values.extend(
                cell.text
                for row in shape.table.rows
                for cell in row.cells
                if cell.text.strip()
            )
    return "\n".join(values)


def _rgb(value: str) -> tuple[int, int, int]:
    text = value.strip().lstrip("#")
    if len(text) != 6:
        raise ValueError("confirmed color must be six-digit hex")
    return tuple(int(text[index : index + 2], 16) for index in (0, 2, 4))


def _mix(color: tuple[int, int, int], target: int, amount: float) -> tuple[int, int, int]:
    return tuple(round(channel * (1 - amount) + target * amount) for channel in color)


def _accent_family(candidate: tuple[int, int, int], accent: tuple[int, int, int]) -> bool:
    known = {accent, _mix(accent, 0, 0.20), *(_mix(accent, 255, amount) for amount in (0.40, 0.70, 0.88))}
    if candidate in known:
        return True
    accent_hue, accent_saturation, _ = colorsys.rgb_to_hsv(*(value / 255 for value in accent))
    hue, saturation, _ = colorsys.rgb_to_hsv(*(value / 255 for value in candidate))
    if accent_saturation < 0.15:
        return False
    return min(abs(hue - accent_hue), 1 - abs(hue - accent_hue)) <= 0.04 and saturation >= max(
        0.12, accent_saturation * 0.20
    )


def _luminance(color: tuple[int, int, int]) -> float:
    channels = [
        value / 12.92 if value <= 0.04045 else ((value + 0.055) / 1.055) ** 2.4
        for value in (channel / 255 for channel in color)
    ]
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def _contrast(first: tuple[int, int, int], second: tuple[int, int, int]) -> float:
    lighter, darker = sorted((_luminance(first), _luminance(second)), reverse=True)
    return (lighter + 0.05) / (darker + 0.05)


def _direct_rgb(color: object) -> tuple[int, int, int] | None:
    try:
        value = getattr(color, "rgb")
    except (AttributeError, ValueError):
        return None
    if value is None:
        return None
    return tuple(int(str(value)[index : index + 2], 16) for index in (0, 2, 4))


def _theme_colors(slide: object) -> dict[str, tuple[int, int, int]]:
    parts = (
        part for part in slide.part.package.iter_parts()
        if "theme" in str(part.partname)
    )
    try:
        root = ElementTree.fromstring(next(parts).blob)
    except (StopIteration, ElementTree.ParseError):
        return {}
    scheme = root.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme")
    if scheme is None:
        return {}
    result: dict[str, tuple[int, int, int]] = {}
    for slot in scheme:
        color = next(iter(slot), None)
        if color is None:
            continue
        value = color.get("val") or color.get("lastClr")
        if value and len(value) == 6:
            try:
                result[slot.tag.rsplit("}", 1)[-1]] = _rgb(value)
            except ValueError:
                continue
    return result


def _font_rgb(
    color: object, theme_colors: Mapping[str, tuple[int, int, int]],
) -> tuple[int, int, int] | None:
    direct = _direct_rgb(color)
    if direct is not None:
        return direct
    try:
        slot = str(color.theme_color).split(" ", 1)[0].lower().replace("_", "")
    except (AttributeError, TypeError, ValueError):
        return None
    # ponytail: theme tint/shade transforms are ignored; resolve them if authored decks depend on them.
    return theme_colors.get(slot)


def _fill_rgb(
    owner: object,
    fallback: tuple[int, int, int],
    theme_colors: Mapping[str, tuple[int, int, int]],
) -> tuple[int, int, int]:
    try:
        value = _font_rgb(getattr(owner, "fill").fore_color, theme_colors)
    except (AttributeError, TypeError, ValueError):
        return fallback
    return value or fallback


def _replace_text_frame_colors(
    text_frame: object,
    *,
    background: tuple[int, int, int],
    accent: tuple[int, int, int],
    default: tuple[int, int, int],
    theme_colors: Mapping[str, tuple[int, int, int]],
) -> None:
    replacement = default
    if _accent_family(default, accent) or _contrast(default, background) < 4.5:
        replacement = max(((0, 0, 0), (255, 255, 255)), key=lambda color: _contrast(color, background))
    for paragraph in getattr(text_frame, "paragraphs", []):
        for run in paragraph.runs:
            _replace_font_color(
                run.font,
                background=background,
                accent=accent,
                default=default,
                replacement=replacement,
                theme_colors=theme_colors,
            )


def _replace_font_color(
    font: object,
    *,
    background: tuple[int, int, int],
    accent: tuple[int, int, int],
    default: tuple[int, int, int],
    replacement: tuple[int, int, int] | None = None,
    theme_colors: Mapping[str, tuple[int, int, int]],
) -> None:
    if replacement is None:
        replacement = default
        if _accent_family(default, accent) or _contrast(default, background) < 4.5:
            replacement = max(
                ((0, 0, 0), (255, 255, 255)),
                key=lambda color: _contrast(color, background),
            )
    color = getattr(font, "color", None)
    actual = _font_rgb(color, theme_colors) if color is not None else None
    if actual is not None and _accent_family(actual, accent):
        font.color.rgb = RGBColor(*replacement)


def _replace_non_emphasis_slide_text_colors(slide: object, style: Mapping[str, Any]) -> None:
    accent = _rgb(str(style.get("secondary_color", "#C7352B")))
    default = _rgb(str(style.get("primary_color", "#1F2937")))
    slide_background = _rgb(str(style.get("background_color", "#FFFFFF")))
    theme_colors = _theme_colors(slide)

    def visit(shape: object, inherited_background: tuple[int, int, int]) -> None:
        background = _fill_rgb(shape, inherited_background, theme_colors)
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is not None:
            _replace_text_frame_colors(
                text_frame, background=background, accent=accent, default=default,
                theme_colors=theme_colors,
            )
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    _replace_text_frame_colors(
                        cell.text_frame,
                        background=_fill_rgb(cell, background, theme_colors),
                        accent=accent,
                        default=default,
                        theme_colors=theme_colors,
                    )
        if getattr(shape, "has_chart", False):
            chart = shape.chart
            if getattr(chart, "has_title", False):
                _replace_text_frame_colors(
                    chart.chart_title.text_frame,
                    background=background,
                    accent=accent,
                    default=default,
                    theme_colors=theme_colors,
                )
            if getattr(chart, "has_legend", False):
                _replace_font_color(
                    chart.legend.font,
                    background=background,
                    accent=accent,
                    default=default,
                    theme_colors=theme_colors,
                )
            for axis_name in ("category_axis", "value_axis", "series_axis"):
                try:
                    axis = getattr(chart, axis_name)
                except (AttributeError, ValueError):
                    continue
                _replace_font_color(
                    axis.tick_labels.font,
                    background=background,
                    accent=accent,
                    default=default,
                    theme_colors=theme_colors,
                )
                if getattr(axis, "has_title", False):
                    _replace_text_frame_colors(
                        axis.axis_title.text_frame,
                        background=background,
                        accent=accent,
                        default=default,
                        theme_colors=theme_colors,
                    )
            for plot in chart.plots:
                if getattr(plot, "has_data_labels", False):
                    _replace_font_color(
                        plot.data_labels.font,
                        background=background,
                        accent=accent,
                        default=default,
                        theme_colors=theme_colors,
                    )
        for child in getattr(shape, "shapes", []):
            visit(child, background)

    for shape in slide.shapes:
        visit(shape, slide_background)


def _replace_non_emphasis_text_colors(deck: Presentation, style: Mapping[str, Any]) -> None:
    for slide in deck.slides:
        _replace_non_emphasis_slide_text_colors(slide, style)


def _validate_reconstructed_text_repairs(
    root: Path, page_number: int, deck: Presentation,
) -> None:
    receipt_path = root / "04_v6" / "images" / f"page_{page_number:03d}.json"
    # Legacy/native direct-finalization paths predate accepted-image receipts.
    # They carry no sealed repair authority, so there is nothing to enforce.
    if not receipt_path.exists():
        return
    receipt = _read_json(receipt_path)
    repairs = receipt.get("reconstruction_repairs", [])
    if not isinstance(repairs, list):
        raise ValueError("V6 accepted text repair authority is invalid")
    text = _editable_slide_text(deck)
    for repair in repairs:
        if not isinstance(repair, Mapping):
            raise ValueError("V6 accepted text repair authority is invalid")
        find = repair.get("find")
        replace = repair.get("replace")
        if find is None and replace is None:
            continue
        if (
            not isinstance(find, str)
            or not isinstance(replace, str)
            or find in text
            or replace not in text
        ):
            raise ValueError("V6 reconstructed native text did not apply a sealed repair")


def _shape_object_id(shape: object) -> str | None:
    properties = shape._element.xpath(".//p:cNvPr")
    if not properties:
        return None
    description = properties[0].get("descr", "")
    return description.removeprefix("object_id:") if description.startswith("object_id:") else None


def _inside(point: tuple[int, int], shape: object) -> bool:
    x, y = point
    return (
        shape.left - 1 <= x <= shape.left + shape.width + 1
        and shape.top - 1 <= y <= shape.top + shape.height + 1
    )


def _require_final_authority(
    root: Path,
    page_number: int,
    reconstructed_body: Path,
    deck: Presentation,
    authority_mode: str,
) -> Mapping[str, Any] | None:
    """Verify sealed worker authority before the host publishes the editable page."""
    if authority_mode == "native_direct":
        page = _load_reconstruction_state(root)["pages"][page_number - 1]
        if page.get("selected_candidate") is not None:
            raise ValueError("V6 native-direct finalization cannot use a selected candidate")
        receipt_path = root / "04_v6" / "images" / f"page_{page_number:03d}.json"
        if receipt_path.exists():
            raise ValueError("V6 native-direct finalization cannot use an acceptance receipt")
        if page.get("state") not in {"accepted", "reconstructing", "page_complete"}:
            raise ValueError("V6 native-direct acceptance authority is missing")
        return None
    if authority_mode != "sealed_reconstruction":
        raise ValueError("V6 finalization authority mode is invalid")
    request_path = reconstructed_body.parent / "accepted_reconstruction_request.json"
    manifest_path = reconstructed_body.parent / "manifest.json"
    if not request_path.is_file():
        raise ValueError("V6 sealed reconstruction request is missing")
    if not manifest_path.is_file():
        raise ValueError("V6 sealed reconstruction manifest is missing")
    request = _read_json(request_path)
    accepted = request.get("accepted_receipt")
    if not isinstance(accepted, Mapping) or accepted.get("path") != (
        Path("04_v6") / "images" / f"page_{page_number:03d}.json"
    ).as_posix():
        raise ValueError("V6 sealed acceptance receipt relationship is invalid")
    receipt_path = root / str(accepted["path"])
    if not receipt_path.is_file():
        raise ValueError("V6 sealed acceptance receipt is missing")
    receipt_bytes = secure_io.read_bytes(root, receipt_path.relative_to(root))
    receipt = verify_signed_acceptance_receipt(
        open_live_page_workspace(root, page_number), receipt_bytes,
    )
    if (
        accepted.get("sha256") != hashlib.sha256(receipt_bytes).hexdigest()
        or receipt.get("page_number") != page_number
        or request.get("page_plan") != receipt.get("page_plan")
    ):
        raise ValueError("V6 sealed acceptance receipt relationship is invalid")
    canonical_request = _read_json(
        root / "05_v6" / "reconstruction_requests" / f"page_{page_number:03d}.json"
    )
    page_request_path = reconstructed_body.parent / "page_request.json"
    jobs_path = reconstructed_body.parents[2] / "page_jobs.json"
    page_request = _read_json(page_request_path)
    jobs = _read_json(jobs_path)
    job = next(
        (
            item for item in jobs.get("pages", [])
            if isinstance(item, Mapping) and item.get("page_id") == "page_001"
        ),
        None,
    )
    dispatch = job.get("dispatch") if isinstance(job, Mapping) else None
    if (
        request != canonical_request
        or not isinstance(job, Mapping)
        or not isinstance(dispatch, Mapping)
        or dispatch.get("page_request_sha256") != _sha256(page_request_path)
        or any(
            page_request.get(field) != request.get(field)
            for field in ("page_plan", "numeric_authority")
        )
    ):
        raise ValueError("V6 sealed reconstruction request relationship is invalid")
    manifest = _read_json(manifest_path)

    relationship = request.get("page_plan", {}).get("primary_relationship", {})
    nodes = relationship.get("nodes", [])
    edges = relationship.get("edges", [])
    shapes_by_id: dict[str, list[object]] = {}
    for shape in deck.slides[0].shapes:
        object_id = _shape_object_id(shape)
        if object_id:
            shapes_by_id.setdefault(object_id, []).append(shape)
    manifest_ids = [
        item.get("object_id")
        for section in ("text_boxes", "tables", "images", "shapes", "charts")
        for item in manifest.get(section, [])
        if isinstance(item, Mapping)
    ]
    for node in nodes:
        node_id = node.get("node_id") if isinstance(node, Mapping) else None
        if not isinstance(node_id, str) or manifest_ids.count(node_id) != 1 or len(shapes_by_id.get(node_id, [])) != 1:
            raise ValueError(f"V6 sealed relationship node is missing or duplicated: {node_id}")
    for edge in edges:
        if not isinstance(edge, Mapping):
            raise ValueError("V6 sealed relationship edge is invalid")
        source_id, target_id = edge.get("from_node"), edge.get("to_node")
        edge_id = f"edge:{source_id}->{target_id}"
        matches = shapes_by_id.get(edge_id, [])
        if (
            source_id not in shapes_by_id
            or target_id not in shapes_by_id
            or manifest_ids.count(edge_id) != 1
            or len(matches) != 1
        ):
            raise ValueError(f"V6 sealed relationship edge is missing or duplicated: {edge_id}")
        connector = matches[0]
        if _shape_kind(connector) not in {"line", "connector"}:
            raise ValueError(f"V6 sealed relationship edge is not a real line: {edge_id}")
        start_x, start_y, end_x, end_y = _connector_endpoints(connector)
        if (
            not _inside((start_x, start_y), shapes_by_id[source_id][0])
            or not _inside((end_x, end_y), shapes_by_id[target_id][0])
            or _shape_arrowheads(connector) != {"tailEnd": "triangle"}
        ):
            raise ValueError(f"V6 sealed relationship edge direction is invalid: {edge_id}")

    authority = request.get("numeric_authority")
    if authority is not None:
        charts = manifest.get("charts", [])
        matching = [
            chart for chart in charts
            if isinstance(chart, Mapping)
            and all(chart.get(key) == value for key, value in authority.items())
        ]
        if len(matching) != 1:
            raise ValueError("V6 sealed numeric authority is missing or changed")
        violations = quantitative_chart_readback_violations(reconstructed_body, [manifest])
        if violations:
            raise ValueError("V6 sealed numeric authority failed readback: " + json.dumps(violations, ensure_ascii=False))
    return accepted


def finalize_reconstructed_page(
    project: Path,
    *,
    page_number: int,
    reconstructed_body: Path,
    authority_mode: str = "sealed_reconstruction",
) -> dict[str, Any]:
    secure_io.reject_reparse_chain(Path(project))
    root = Path(project).resolve()
    state = _load_reconstruction_state(root)
    reconstructed_body = Path(reconstructed_body).resolve()
    if not reconstructed_body.is_file() or reconstructed_body.suffix.lower() != ".pptx":
        raise ValueError("V6 reconstructed body must be an existing PPTX")
    opened = Presentation(reconstructed_body)
    if len(opened.slides) != 1:
        raise ValueError("V6 reconstructed body must contain exactly one slide")
    _validate_reconstructed_text_repairs(root, page_number, opened)
    accepted_receipt = _require_final_authority(
        root, page_number, reconstructed_body, opened, authority_mode,
    )
    page_index = page_number - 1
    page = state["pages"][page_index]
    if page["state"] not in {"accepted", "reconstructing", "page_complete"}:
        raise ValueError("V6 page is not ready for reconstruction finalization")
    repairing_complete_page = page["state"] == "page_complete"
    if not repairing_complete_page and page["state"] != "reconstructing":
        page = transition_page(page, "reconstructing")
    output_dir = root / "06_v6" / "pages" / f"page_{page_number:03d}"
    output = output_dir / "page.pptx"
    style = state["style_confirmation"]["contract"]
    if not isinstance(style, Mapping):
        raise ValueError("V6 confirmed style contract is missing")
    style_execution = _fixed_frame_style(style)
    logo = root / state["logo_source"]["path"]
    if repairing_complete_page:
        receipt_path = output_dir / "page.json"
        if not output.is_file() or not receipt_path.is_file():
            raise ValueError("V6 existing finalized page authority is missing")
        receipt = _read_json(receipt_path)
        existing_fixed = inspect_fixed_frame(
            output,
            expected_title=page["title"],
            expected_page_number=page_number,
            style_execution=style_execution,
            logo_svg=logo,
        )
        if (
            receipt.get("artifact_version") != "final-page-v6"
            or receipt.get("page_number") != page_number
            or receipt.get("page_pptx") != output.relative_to(root).as_posix()
            or receipt.get("sha256") != _sha256(output)
            or not isinstance(receipt.get("fixed_frame"), Mapping)
            or receipt["fixed_frame"].get("passed") is not True
            or existing_fixed.get("passed") is not True
        ):
            raise ValueError("V6 existing finalized page authority is invalid")
    if page_number not in project_emphasis_pages(root):
        _replace_non_emphasis_text_colors(opened, style)
        buffer = BytesIO()
        opened.save(buffer)
        reconstructed_bytes = buffer.getvalue()
    else:
        reconstructed_bytes = reconstructed_body.read_bytes()
    repair_target: Path | None = None
    if repairing_complete_page:
        repair_target = output_dir / f".page-repair-{uuid.uuid4().hex[:8]}.pptx"
        secure_io.atomic_write_bytes(root, repair_target.relative_to(root), reconstructed_bytes)
        finalization_output = repair_target
    elif output.is_file():
        existing_bytes = secure_io.read_bytes(root, output.relative_to(root))
        if existing_bytes != reconstructed_bytes:
            raise ValueError("V6 reconstructed page output already contains different bytes")
        finalization_output = output
    else:
        secure_io.atomic_write_bytes(root, output.relative_to(root), reconstructed_bytes)
        finalization_output = output
    try:
        apply_fixed_frame(
            finalization_output,
            page_title=page["title"],
            page_number=page_number,
            style_execution=style_execution,
            logo_svg=logo,
        )
        fixed = inspect_fixed_frame(
            finalization_output,
            expected_title=page["title"],
            expected_page_number=page_number,
            style_execution=style_execution,
            logo_svg=logo,
        )
        if fixed.get("passed") is not True:
            raise ValueError("V6 fixed-layer validation failed: " + "; ".join(fixed.get("issues", [])))
        if repairing_complete_page:
            secure_io.atomic_write_bytes(
                root,
                output.relative_to(root),
                secure_io.read_bytes(root, finalization_output.relative_to(root)),
                replace=True,
            )
    finally:
        if repair_target is not None and repair_target.is_file():
            repair_target.unlink()
    if not repairing_complete_page:
        page = transition_page(page, "page_complete")
        _update_reconstruction_page(root, page_number, page)
    report = {
        "artifact_version": "final-page-v6",
        "page_number": page_number,
        "page_pptx": output.relative_to(root).as_posix(),
        "sha256": _sha256(output),
        "fixed_frame": fixed,
        "post_reconstruction_visual_qa": False,
    }
    if accepted_receipt is not None:
        report["accepted_receipt"] = dict(accepted_receipt)
    _write_json(output_dir / "page.json", report)
    return report


def assemble_v6_deck(project: Path) -> dict[str, Any]:
    secure_io.reject_reparse_chain(Path(project))
    root = Path(project).resolve()
    state = _load_reconstruction_state(root)
    composition = load_composition_authority(root)
    if composition is None:
        page_contracts = [
            {"page_role": "content", "visible_page_number": True} for _page in state["pages"]
        ]
    else:
        page_contracts = composition["pages"]
    if len(page_contracts) != len(state["pages"]):
        raise ValueError("assembled V6 composition page count is incorrect")
    if any(page["state"] != "page_complete" for page in state["pages"]):
        raise ValueError("every V6 page must be complete before assembly")
    pages = [
        root / "06_v6" / "pages" / f"page_{page['page_number']:03d}" / "page.pptx"
        for page in state["pages"]
    ]
    if any(not path.is_file() for path in pages):
        raise ValueError("a V6 finalized page package is missing")
    output_dir = root / "08_final"
    output = output_dir / "deck.pptx"
    temporary = output_dir / f".deck-v6-{uuid.uuid4().hex[:8]}.tmp"
    deck = Presentation(pages[0])
    layout = deck.slides[0].slide_layout
    for page_number, path in enumerate(pages[1:], start=2):
        _copy_page_slide(path, deck, layout, page_number)
    with secure_io.hold_parent(root, temporary.relative_to(root), create=True):
        deck.save(temporary)
        reopened = Presentation(temporary)
        if len(reopened.slides) != len(pages):
            raise ValueError("assembled V6 slide count is incorrect")
        style = state["style_confirmation"]["contract"]
        if not isinstance(style, Mapping):
            raise ValueError("V6 confirmed style contract is missing")
        emphasis_pages = project_emphasis_pages(root)
        background = _rgb(str(style.get("background_color", "#FFFFFF")))
        for page_number, (slide, contract) in enumerate(
            zip(reopened.slides, page_contracts), start=1
        ):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*background)
            if page_number not in emphasis_pages:
                _replace_non_emphasis_slide_text_colors(slide, style)
            role = contract["page_role"]
            if role in {"content", "appendix"}:
                _require_current_fixed_frame(slide)
            else:
                _require_special_page_shapes(slide, role, contract["visible_page_number"])
        if any(
            not any(shape.has_text_frame or shape.has_table for shape in slide.shapes)
            for slide in reopened.slides
        ):
            raise ValueError("assembled V6 slide has no editable text or table object")
        reopened.save(temporary)
        if not zipfile.is_zipfile(temporary):
            raise ValueError("assembled V6 output is not an OpenXML package")
        temporary_bytes = temporary.read_bytes()
    secure_io.atomic_write_bytes(root, output.relative_to(root), temporary_bytes, replace=output.exists())
    temporary.unlink(missing_ok=True)
    report = {
        "artifact_version": "final-assembly-v6",
        "workflow_contract_version": "awesome-word-ppt-workflow-v1",
        "status": "complete",
        "page_count": len(pages),
        "page_order": [page["page_number"] for page in state["pages"]],
        "output": output.relative_to(root).as_posix(),
        "sha256": _sha256(output),
        "mechanical_validation": {
            "openxml_package": True,
            "slide_count": True,
            "fixed_layers": True,
            "editable_objects": True,
        },
        "office_render_required": False,
        "post_reconstruction_visual_qa": False,
    }
    _write_json(output_dir / "assembly.json", report)
    return report
