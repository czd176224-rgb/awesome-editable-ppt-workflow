from __future__ import annotations

import importlib.util
import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from editppt.runtime.editable_page_cache import (
    PackageValidationError,
    create_page_package,
)


PROMPT_BUILDER = Path(__file__).resolve().parents[2] / "scripts" / "build-page-worker-prompt.py"
MANIFEST_SCHEMA = Path(__file__).resolve().parents[2] / "references" / "manifest-schema.md"


def _load_prompt_builder():
    spec = importlib.util.spec_from_file_location("page_worker_prompt_builder", PROMPT_BUILDER)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _worker_prompt(
    tmp_path: Path,
    numeric_authority: dict[str, object] | None,
    page_plan: dict[str, object] | None = None,
) -> str:
    module = _load_prompt_builder()
    run_dir = tmp_path / "run"
    page_dir = run_dir / "pages" / "page_001"
    page_dir.mkdir(parents=True)
    request: dict[str, object] = {"source_image": str(page_dir / "source.png")}
    if numeric_authority is not None:
        request["numeric_authority"] = numeric_authority
    if page_plan is not None:
        request["page_plan"] = page_plan
    (page_dir / "page_request.json").write_text(
        json.dumps(request, ensure_ascii=False), encoding="utf-8"
    )
    return module.build_prompt(run_dir, {"page_id": "page_001"}, page_dir)


def _write_editable_pptx(path, *, slide_count: int) -> None:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    for index in range(1, slide_count):
        presentation.slides.add_slide(presentation.slide_layouts[6])
    for index, slide in enumerate(presentation.slides, start=1):
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        text_box.text = f"Editable page {index}"
    presentation.save(path)


def test_create_page_package_records_one_real_editable_slide(tmp_path):
    project = tmp_path / "project"
    project.mkdir()
    pptx = project / "page.pptx"
    descriptor = project / "packages" / "page-001.json"
    _write_editable_pptx(pptx, slide_count=1)

    result = create_page_package(
        project,
        page_number=1,
        cache_key="a" * 64,
        pptx=pptx,
        output=descriptor,
    )

    payload = json.loads(result.read_text(encoding="utf-8"))
    assert payload == {
        "schema_version": "editable-page-package-v1",
        "cache_key": "a" * 64,
        "pptx": "page.pptx",
        "pptx_sha256": payload["pptx_sha256"],
        "editable_object_count": 1,
        "slide_fingerprint": payload["slide_fingerprint"],
    }
    assert len(payload["pptx_sha256"]) == 64
    assert len(payload["slide_fingerprint"]) == 64


def test_create_page_package_rejects_a_multi_slide_reconstruction(tmp_path):
    project = tmp_path / "project"
    project.mkdir()
    pptx = project / "two-slides.pptx"
    _write_editable_pptx(pptx, slide_count=2)

    with pytest.raises(PackageValidationError, match="exactly one slide"):
        create_page_package(
            project,
            page_number=1,
            cache_key="b" * 64,
            pptx=pptx,
            output=project / "packages" / "page-001.json",
        )


def test_installed_runtime_exposes_only_manifest_record_and_finalize_routes():
    """The shipped CLI must not retain V4 record/finalize dispatches."""
    runtime = Path(__file__).resolve().parents[1] / "editppt" / "runtime"

    for name in ("editable_page_cache.py",):
        source = (runtime / name).read_text(encoding="utf-8")
        assert "parents[4]" not in source
        assert '"run-word-to-ppt-workflow" / "scripts"' not in source
    main = (runtime / "main.py").read_text(encoding="utf-8")
    assert '"record_page_result.py"' not in main
    assert '"finalize_deck_run.py"' not in main
    assert '"record_manifest_page_result.py"' in main
    assert '"finalize_manifest_deck_run.py"' in main
    for name in ("record_page_result.py", "finalize_deck_run.py", "final_assembler.py"):
        assert not (runtime / name).exists()


def test_page_worker_prompt_seals_composition_and_numeric_ownership(tmp_path):
    authority = {
        "title": "Revenue",
        "rendering_primitive": "column_bar",
        "chart_variant": "column",
        "unit": "USD m",
        "basis": "reported revenue",
        "period": "FY2025",
        "series": [{"name": "Revenue", "categories": ["A"], "values": [10]}],
    }

    prompt = _worker_prompt(tmp_path, authority)

    assert "accepted source image owns chart container placement, composition, and style" in prompt
    assert "sealed numeric authority owns quantitative mark size, position, and labels" in prompt
    assert "Keep chart_variant unchanged for native charts" in prompt
    assert "special rendering primitives omit chart_variant" in prompt
    assert "do not calculate new metrics" in prompt
    assert json.dumps(authority, ensure_ascii=False, sort_keys=True) in prompt


def test_page_worker_prompt_requires_special_authority_to_omit_chart_variant(tmp_path):
    authority = {
        "title": "Bridge", "rendering_primitive": "cumulative_bridge",
        "unit": "USD m", "basis": "reported EBITDA", "period": "FY2025",
        "series": [{"name": "Bridge", "categories": ["Price"], "start": 10, "changes": [2], "end": 12}],
    }

    prompt = _worker_prompt(tmp_path, authority)

    assert "special rendering primitives omit chart_variant" in prompt
    assert "any added variant must be rejected" in prompt
    assert '"chart_variant"' not in json.dumps(authority)

    authority["chart_variant"] = "bubble"
    with pytest.raises(SystemExit, match="must omit chart_variant"):
        _worker_prompt(tmp_path / "invalid", authority)


def test_manifest_schema_matches_worker_chart_variant_and_special_field_contracts(tmp_path):
    schema = MANIFEST_SCHEMA.read_text(encoding="utf-8")
    native_prompt = _worker_prompt(tmp_path / "native", {
        "title": "Revenue", "rendering_primitive": "column_bar", "chart_variant": "column",
        "unit": "USD m", "basis": "reported revenue", "period": "FY2025",
        "series": [{"name": "Revenue", "categories": ["A"], "values": [10]}],
    })
    special_prompt = _worker_prompt(tmp_path / "special", {
        "title": "Bridge", "rendering_primitive": "cumulative_bridge",
        "unit": "USD m", "basis": "reported EBITDA", "period": "FY2025",
        "series": [{"name": "Bridge", "categories": ["Price"], "start": 10,
        "start_label": "FY2024", "changes": [2], "end": 12, "end_label": "FY2025"}],
    })

    assert "Standard chart primitives must include the matching explicit `chart_variant`" in schema
    assert "special shape-based primitives must omit `chart_variant`" in schema
    assert "Keep chart_variant unchanged for native charts" in native_prompt
    assert "special rendering primitives omit chart_variant" in special_prompt
    for field in (
        "cumulative_bridge", "start", "changes", "end", "start_label", "end_label",
        "time_interval", "start_dates", "end_dates", "variable_rectangle", "width_values",
        "width_label", "width_unit", "width_basis", "share_values", "share_denominator",
        "share_label", "share_unit", "share_basis",
    ):
        assert f"`{field}`" in schema


def test_page_worker_prompt_without_numeric_authority_forbids_quantitative_geometry(tmp_path):
    prompt = _worker_prompt(tmp_path, None)

    assert "No sealed numeric authority is present" in prompt
    for forbidden_geometry in (
        "numeric axes",
        "proportional geometry",
        "bubble-size ranking",
        "target-line magnitude",
        "difference magnitude",
    ):
        assert forbidden_geometry in prompt


def test_page_worker_prompt_contains_sealed_relationship_nodes_edges_and_direction(tmp_path):
    page_plan = {
        "page_purpose": "Explain the accepted relationship.",
        "primary_relationship": {
            "nodes": [
                {"node_id": "regional-resources", "label": "Regional resources"},
                {"node_id": "fund-platform", "label": "Fund platform"},
            ],
            "edges": [{
                "from_node": "regional-resources",
                "to_node": "fund-platform",
                "label": "enter",
            }],
        },
    }

    prompt = _worker_prompt(tmp_path, None, page_plan)

    assert "SEALED PAGE PLAN RELATIONSHIP AUTHORITY" in prompt
    assert '"node_id": "regional-resources"' in prompt
    assert '"from_node": "regional-resources"' in prompt
    assert '"to_node": "fund-platform"' in prompt
    assert "regional-resources -> fund-platform" in prompt
    assert (
        'Required node object_id values (verbatim, one object each): '
        '["regional-resources", "fund-platform"]'
    ) in prompt
    assert (
        'Required connector object_id values (verbatim, one connector each): '
        '["edge:regional-resources->fund-platform"]'
    ) in prompt
    assert "Do not alias, rename, prefix, suffix, or split any required ID" in prompt
    assert (
        "Before reporting success, verify each required ID appears exactly once in "
        "manifest.json as object_id and exactly once in page.pptx cNvPr descr as "
        "object_id:<ID>."
    ) in prompt
