from __future__ import annotations

import sys
from copy import deepcopy
from pathlib import Path

import pytest
from PIL import ImageDraw
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE


RUNTIME = Path(__file__).resolve().parents[1] / "editppt" / "runtime"
sys.path.insert(0, str(RUNTIME))
WORKFLOW_SCRIPTS = Path(__file__).resolve().parents[3] / "run-word-to-ppt-workflow" / "scripts"
sys.path.insert(0, str(WORKFLOW_SCRIPTS))

from build_pptx_from_manifest import normalize_manifest, px_to_inches, render_preview, write_deck, write_pptx  # noqa: E402
from fixed_region_runtime import CONTENT_BOX, SLIDE  # noqa: E402
from source_assets import _chart_record  # noqa: E402
from workflow_v6_materials import select_numeric_authority  # noqa: E402
import validate_pptx  # noqa: E402


CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "bubble": XL_CHART_TYPE.BUBBLE,
}


def _manifest(chart: dict) -> dict:
    return {
        "workflow_contract_version": "fixed-canvas-cm-v2",
        "reconstruction_contract_version": "editable-image-v3",
        "slide": dict(SLIDE),
        "content_box": dict(CONTENT_BOX),
        "source": {"width_px": 1904, "height_px": 896},
        "text_inventory": [],
        "visual_inventory": [],
        "background_strategy": "native white body background",
        "quality_checks": {
            "font_size_calibrated": True,
            "visual_inventory_matched": True,
            "background_strategy_checked": True,
            "shape_corner_geometry_checked": True,
        },
        "text_boxes": [],
        "tables": [],
        "shapes": [],
        "images": [],
        "charts": [chart],
        "asset_provenance": [],
    }


def _one_dimensional(variant: str = "column") -> dict:
    return {
        "object_id": "chart-1",
        "name": "Revenue chart",
        "box_px": [190, 90, 1142, 620],
        "rendering_primitive": "column_bar" if variant in {"column", "bar"} else "line_point",
        "chart_variant": variant,
        "title": "Revenue",
        "unit": "USD m",
        "period": "FY2025",
        "basis": "same portfolio companies",
        "series": [{"name": "Revenue", "categories": ["A", "B"], "values": [12, 18]}],
    }


def _xy(variant: str) -> dict:
    chart = {
        "object_id": "chart-1",
        "name": "Portfolio chart",
        "box_px": [190, 90, 1142, 620],
        "rendering_primitive": "xy",
        "chart_variant": variant,
        "title": "Portfolio",
        "period": "FY2025",
        "x_label": "Growth",
        "x_unit": "%",
        "x_basis": "FY2025 revenue growth",
        "y_label": "Margin",
        "y_unit": "% pts",
        "y_basis": "FY2025 EBITDA margin",
        "series": [{"name": "Companies", "x_values": [1.2, 2.4], "y_values": [8, 11]}],
    }
    if variant == "bubble":
        chart.update({"size_label": "Revenue", "size_unit": "USD m", "size_basis": "FY2025 revenue"})
        chart["series"][0]["size_values"] = [30, 50]
    return chart


def _cumulative_bridge() -> dict:
    return {
        "object_id": "chart-1",
        "name": "Value bridge",
        "box_px": [190, 90, 1142, 620],
        "rendering_primitive": "cumulative_bridge",
        "title": "Value bridge",
        "unit": "USD m",
        "basis": "FY2025 EBITDA",
        "period": "FY2025",
        "series": [{
            "name": "Bridge",
            "categories": ["Pricing", "Volume"],
            "start": 100,
            "changes": [20, -5],
            "end": 115,
        }],
    }


def _time_interval() -> dict:
    return {
        "object_id": "chart-1",
        "name": "Project schedule",
        "box_px": [190, 90, 1142, 620],
        "rendering_primitive": "time_interval",
        "title": "Project schedule",
        "period": "September 2026",
        "series": [{
            "name": "Plan",
            "categories": ["Diligence", "IC"],
            "start_dates": ["2026-09-01", "2026-09-11"],
            "end_dates": ["2026-09-10", "2026-09-15"],
        }],
    }


def _variable_rectangle() -> dict:
    return {
        "object_id": "chart-1",
        "name": "Market composition",
        "box_px": [190, 90, 1142, 620],
        "rendering_primitive": "variable_rectangle",
        "title": "Market composition",
        "period": "FY2025",
        "series": [{
            "name": "Markets",
            "categories": ["A", "B"],
            "width_values": [40, 60],
            "width_label": "Market size",
            "width_unit": "USD m",
            "width_basis": "2025 addressable market",
            "share_values": [[25, 75], [40, 60]],
            "share_label": "Portfolio share",
            "share_unit": "%",
            "share_basis": "2025 composition",
            "share_denominator": 100,
        }],
    }


def _description(shape) -> str | None:
    return shape._element.xpath(".//p:cNvPr")[0].get("descr")


def _preset(shape) -> str | None:
    geometry = shape._element.xpath(".//a:prstGeom")
    return geometry[0].get("prst") if geometry else None


@pytest.mark.parametrize(
    ("chart", "expected_type"),
    [
        (_one_dimensional("column"), XL_CHART_TYPE.COLUMN_CLUSTERED),
        (_one_dimensional("bar"), XL_CHART_TYPE.BAR_CLUSTERED),
        (_one_dimensional("line"), XL_CHART_TYPE.LINE),
        (_xy("scatter"), XL_CHART_TYPE.XY_SCATTER),
        (_xy("bubble"), XL_CHART_TYPE.BUBBLE),
    ],
)
def test_native_variants_preserve_exact_type_data_labels_and_fixed_canvas_box(
    tmp_path: Path, chart: dict, expected_type: XL_CHART_TYPE
) -> None:
    manifest = _manifest(chart)
    out = tmp_path / f"{chart['chart_variant']}.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")

    presentation = Presentation(out)
    chart_shape = next(shape for shape in presentation.slides[0].shapes if shape.has_chart)
    expected_box = px_to_inches(manifest, *chart["box_px"])
    metadata = {shape.name: shape.text for shape in presentation.slides[0].shapes if shape.has_text_frame}

    assert chart_shape.chart.chart_type == expected_type
    assert chart_shape.name == chart["name"]
    assert chart_shape.left == pytest.approx(expected_box["left"] * 914400, abs=1)
    assert chart_shape.top == pytest.approx(expected_box["top"] * 914400, abs=1)
    assert chart_shape.width == pytest.approx(expected_box["width"] * 914400, abs=1)
    assert chart_shape.height == pytest.approx(expected_box["height"] * 914400, abs=1)
    assert chart_shape.chart.chart_title.text_frame.text == chart["title"]
    if chart["rendering_primitive"] == "xy":
        assert chart_shape.chart.category_axis.axis_title.text_frame.text == chart["x_label"]
        assert chart_shape.chart.value_axis.axis_title.text_frame.text == chart["y_label"]
        assert metadata[f"{chart['name']} X Basis"] == chart["x_basis"]
        assert metadata[f"{chart['name']} Y Basis"] == chart["y_basis"]
        if chart["chart_variant"] == "bubble":
            assert metadata[f"{chart['name']} Size Basis"] == chart["size_basis"]
    else:
        assert metadata[f"{chart['name']} Basis"] == chart["basis"]
    assert metadata[f"{chart['name']} Unit"] == chart.get("unit", "x: % | y: % pts | size: USD m" if chart["chart_variant"] == "bubble" else "x: % | y: % pts")
    assert metadata[f"{chart['name']} Period"] == "FY2025"
    assert validate_pptx.quantitative_chart_readback_violations(out, [manifest]) == []


def test_cumulative_bridge_uses_exact_editable_waterfall_geometry_and_labels(tmp_path: Path) -> None:
    chart = _cumulative_bridge()
    manifest = _manifest(chart)
    out = tmp_path / "bridge.pptx"

    write_pptx(manifest, out, tmp_path / "manifest.json")

    slide = Presentation(out).slides[0]
    named = {shape.name: shape for shape in slide.shapes}
    assert not any(shape.has_chart for shape in slide.shapes)
    assert _description(named[chart["name"]]) == "object_id:chart-1"
    assert _preset(named[chart["name"]]) == "rect"
    assert (named[chart["name"]].left, named[chart["name"]].top, named[chart["name"]].width, named[chart["name"]].height) == (
        1145882, 1232277, 5134682, 2785018,
    )
    assert (named["Value bridge Bar 1"].left, named["Value bridge Bar 1"].top, named["Value bridge Bar 1"].width, named["Value bridge Bar 1"].height) == (
        2108634, 2021365, 577652, 1438926,
    )
    assert (named["Value bridge Bar 3"].left, named["Value bridge Bar 3"].top, named["Value bridge Bar 3"].width, named["Value bridge Bar 3"].height) == (
        4034140, 1733580, 577652, 71946,
    )
    assert _preset(named["Value bridge Connector 2"]) == "line"
    assert (named["Value bridge Connector 2"].left, named["Value bridge Connector 2"].top, named["Value bridge Connector 2"].width, named["Value bridge Connector 2"].height) == (
        3649039, 1733580, 385101, 0,
    )
    texts = {shape.text for shape in slide.shapes if shape.has_text_frame}
    assert {"Value bridge", "Bridge", "Pricing", "Volume", "100", "20", "-5", "115", "USD m", "FY2025", "FY2025 EBITDA"}.issubset(texts)
    assert {"Start", "End"}.isdisjoint(texts)
    assert validate_pptx.quantitative_chart_readback_violations(out, [manifest]) == []


def test_time_interval_uses_exact_editable_gantt_geometry_dates_and_labels(tmp_path: Path) -> None:
    chart = _time_interval()
    manifest = _manifest(chart)
    out = tmp_path / "gantt.pptx"

    write_pptx(manifest, out, tmp_path / "manifest.json")

    slide = Presentation(out).slides[0]
    named = {shape.name: shape for shape in slide.shapes}
    assert not any(shape.has_chart for shape in slide.shapes)
    assert (named["Project schedule Bar 1"].left, named["Project schedule Bar 1"].top, named["Project schedule Bar 1"].width, named["Project schedule Bar 1"].height) == (
        2429552, 1998157, 2396185, 417753,
    )
    assert (named["Project schedule Bar 2"].left, named["Project schedule Bar 2"].top, named["Project schedule Bar 2"].width, named["Project schedule Bar 2"].height) == (
        4825737, 2833662, 1198092, 417753,
    )
    assert (named["Project schedule Axis"].left, named["Project schedule Axis"].top, named["Project schedule Axis"].width, named["Project schedule Axis"].height) == (
        2429552, 3460291, 3594277, 0,
    )
    assert _preset(named["Project schedule Axis"]) == "line"
    texts = {shape.text for shape in slide.shapes if shape.has_text_frame}
    assert {"Project schedule", "Plan", "Diligence", "IC", "2026-09-01 – 2026-09-10", "2026-09-11 – 2026-09-15", "September 2026"}.issubset(texts)
    assert validate_pptx.quantitative_chart_readback_violations(out, [manifest]) == []


def test_variable_rectangle_uses_exact_normalized_widths_shares_and_labels(tmp_path: Path) -> None:
    chart = _variable_rectangle()
    manifest = _manifest(chart)
    out = tmp_path / "variable-rectangle.pptx"

    write_pptx(manifest, out, tmp_path / "manifest.json")

    slide = Presentation(out).slides[0]
    named = {shape.name: shape for shape in slide.shapes}
    expected_segments = [
        (1762043, 1789280, 1643098, 431678),
        (1762043, 2220958, 1643098, 1295033),
        (3405141, 1789280, 2464647, 690684),
        (3405141, 2479965, 2464647, 1036027),
    ]
    assert not any(shape.has_chart for shape in slide.shapes)
    assert [
        (named[f"Market composition Segment {index}"].left, named[f"Market composition Segment {index}"].top,
         named[f"Market composition Segment {index}"].width, named[f"Market composition Segment {index}"].height)
        for index in range(1, 5)
    ] == expected_segments
    texts = {shape.text for shape in slide.shapes if shape.has_text_frame}
    assert {"Market composition", "Markets", "A", "B", "40", "60", "25", "75", "100", "Market size", "USD m", "2025 addressable market", "Portfolio share", "%", "2025 composition", "FY2025"}.issubset(texts)
    assert validate_pptx.quantitative_chart_readback_violations(out, [manifest]) == []


@pytest.mark.parametrize(
    ("chart", "expected"),
    [
        ({**_cumulative_bridge(), "series": [{
            **_cumulative_bridge()["series"][0], "end": 115.00000001,
        }]}, "end"),
        ({**_variable_rectangle(), "series": [{
            **_variable_rectangle()["series"][0], "share_values": [[25, 75.00000001], [40, 60]],
        }]}, "share_denominator"),
    ],
)
def test_special_source_totals_require_exact_decimal_equality(
    tmp_path: Path, chart: dict, expected: str,
) -> None:
    with pytest.raises(ValueError, match=expected):
        write_pptx(_manifest(chart), tmp_path / "decimal-mismatch.pptx", tmp_path / "manifest.json")


@pytest.mark.parametrize(
    "chart",
    [
        {**_cumulative_bridge(), "series": [{**_cumulative_bridge()["series"][0], "end": 115.00000001}]},
        {**_variable_rectangle(), "series": [{**_variable_rectangle()["series"][0], "share_values": [[25, 75.00000001], [40, 60]]}]},
    ],
)
def test_selector_to_renderer_rejects_the_same_inexact_special_totals(tmp_path: Path, chart: dict) -> None:
    assert select_numeric_authority([chart]) is None
    with pytest.raises(ValueError):
        write_pptx(_manifest(chart), tmp_path / "rejected.pptx", tmp_path / "manifest.json")


def test_waterfall_uses_source_endpoint_labels_only(tmp_path: Path) -> None:
    chart = _cumulative_bridge()
    chart["series"][0].update({"start_label": "FY2024", "end_label": "FY2025"})
    out = tmp_path / "source-endpoints.pptx"

    write_pptx(_manifest(chart), out, tmp_path / "manifest.json")

    texts = {shape.text for shape in Presentation(out).slides[0].shapes if shape.has_text_frame}
    assert {"FY2024", "FY2025"}.issubset(texts)
    assert {"Start", "End"}.isdisjoint(texts)


def test_zero_waterfall_values_use_boundary_markers_and_positive_external_labels(tmp_path: Path) -> None:
    chart = _cumulative_bridge()
    chart["series"][0].update({
        "categories": ["No change", "Increase", "Decrease"],
        "start": 0,
        "changes": [0, 5, -5],
        "end": 0,
    })
    manifest = _manifest(chart)
    normalized = normalize_manifest(manifest)
    named_shapes = {item["name"]: item for item in normalized["shapes"]}
    named_text = {item["name"]: item for item in normalized["text_boxes"]}

    assert {"Value bridge Boundary 1", "Value bridge Boundary 2", "Value bridge Boundary 5"}.issubset(named_shapes)
    assert all(named_shapes[name]["type"] == "line" and named_shapes[name]["width"] > 0 for name in (
        "Value bridge Boundary 1", "Value bridge Boundary 2", "Value bridge Boundary 5",
    ))
    assert all(item["width"] > 0 and item["height"] > 0 for item in normalized["shapes"] if item["type"] == "rect")
    assert all(item["width"] > 0 and item["height"] > 0 for item in normalized["text_boxes"])
    assert named_text["Value bridge Value 1"]["text"] == "0"
    assert named_text["Value bridge Value 2"]["text"] == "0"
    assert named_text["Value bridge Value 5"]["text"] == "0"
    assert named_text["Value bridge Value 1"]["top"] + named_text["Value bridge Value 1"]["height"] <= named_shapes["Value bridge Boundary 1"]["top"]

    out = tmp_path / "zero-bridge.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")
    assert validate_pptx.quantitative_chart_readback_violations(out, [manifest]) == []


def test_zero_share_uses_boundary_marker_and_positive_external_label(tmp_path: Path) -> None:
    chart = _variable_rectangle()
    chart["series"][0]["share_values"][0] = [0, 100]
    manifest = _manifest(chart)
    normalized = normalize_manifest(manifest)
    named_shapes = {item["name"]: item for item in normalized["shapes"]}
    named_text = {item["name"]: item for item in normalized["text_boxes"]}

    assert "Market composition Segment 1" not in named_shapes
    assert named_shapes["Market composition Boundary 1"]["type"] == "line"
    assert named_shapes["Market composition Boundary 1"]["width"] > 0
    assert named_text["Market composition Share 1"]["text"] == "0"
    assert named_text["Market composition Share 1"]["width"] > 0
    assert named_text["Market composition Share 1"]["height"] > 0
    assert named_text["Market composition Share 1"]["left"] >= named_shapes["Market composition Boundary 1"]["left"] + named_shapes["Market composition Boundary 1"]["width"]
    assert all(item["width"] > 0 and item["height"] > 0 for item in normalized["shapes"] if item["type"] == "rect")
    assert all(item["width"] > 0 and item["height"] > 0 for item in normalized["text_boxes"])

    out = tmp_path / "zero-share.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")
    assert validate_pptx.quantitative_chart_readback_violations(out, [manifest]) == []


def test_tiny_positive_share_and_short_gantt_keep_data_geometry_but_externalize_readable_labels() -> None:
    market = _variable_rectangle()
    market["series"][0]["share_values"][0] = [0.1, 99.9]
    market_normalized = normalize_manifest(_manifest(market))
    market_chart = market_normalized["charts"][0]
    market_shapes = {item["name"]: item for item in market_normalized["shapes"]}
    market_text = {item["name"]: item for item in market_normalized["text_boxes"]}

    assert market_shapes["Market composition Segment 1"]["height"] > 0
    assert market_text["Market composition Share 1"]["left"] >= market_shapes["Market composition Segment 1"]["left"] + market_shapes["Market composition Segment 1"]["width"]
    assert market_text["Market composition Share 1"]["height"] >= market_chart["height"] * 0.04

    gantt = _time_interval()
    gantt["series"][0].update({
        "start_dates": ["2026-09-01", "2026-12-31"],
        "end_dates": ["2026-09-01", "2026-12-31"],
    })
    gantt_normalized = normalize_manifest(_manifest(gantt))
    gantt_chart = gantt_normalized["charts"][0]
    gantt_shapes = {item["name"]: item for item in gantt_normalized["shapes"]}
    gantt_text = {item["name"]: item for item in gantt_normalized["text_boxes"]}

    assert gantt_shapes["Project schedule Bar 1"]["width"] > 0
    assert gantt_text["Project schedule Date 1"]["width"] >= gantt_chart["width"] * 0.12
    assert gantt_text["Project schedule Date 1"]["width"] > gantt_shapes["Project schedule Bar 1"]["width"]


@pytest.mark.parametrize(
    ("chart", "mutate", "expected"),
    [
        (_cumulative_bridge(), lambda chart: chart["series"][0].update({"end": 999}), "end"),
        (_cumulative_bridge(), lambda chart: chart.pop("basis"), "basis"),
        (_time_interval(), lambda chart: chart["series"][0]["start_dates"].__setitem__(0, "09/01/2026"), "ISO"),
        (_time_interval(), lambda chart: chart["series"][0]["end_dates"].__setitem__(0, "2026-08-31"), "before"),
        (_variable_rectangle(), lambda chart: chart["series"][0]["width_values"].__setitem__(0, 0), "positive"),
        (_variable_rectangle(), lambda chart: chart["series"][0]["share_values"].__setitem__(0, [25, 70]), "share_denominator"),
    ],
)
def test_special_quantitative_forms_reject_unreadable_source_dimensions(
    tmp_path: Path, chart: dict, mutate, expected: str,
) -> None:
    chart = deepcopy(chart)
    mutate(chart)

    with pytest.raises(ValueError, match=expected):
        write_pptx(_manifest(chart), tmp_path / "invalid-special.pptx", tmp_path / "manifest.json")


@pytest.mark.parametrize("variant", ["bubble", "column"])
def test_special_manifest_rejects_any_extra_chart_variant(tmp_path: Path, variant: str) -> None:
    chart = _cumulative_bridge()
    chart["chart_variant"] = variant

    with pytest.raises(ValueError, match="must be omitted"):
        write_pptx(_manifest(chart), tmp_path / "variant-special.pptx", tmp_path / "manifest.json")


@pytest.mark.parametrize(
    ("chart", "shape_name", "damage", "field"),
    [
        (_cumulative_bridge(), "Value bridge Bar 1", "type", "charts[0].bars[0].type"),
        (_time_interval(), "Project schedule Bar 1", "geometry", "charts[0].bars[0].geometry"),
        (_variable_rectangle(), "Market composition Segment 1", "identity", "charts[0].segments[0].object_id"),
        (_cumulative_bridge(), "Value bridge Value 2", "label", "charts[0].values[1]"),
    ],
)
def test_special_chart_readback_rejects_wrong_type_geometry_identity_or_label(
    tmp_path: Path, chart: dict, shape_name: str, damage: str, field: str,
) -> None:
    manifest = _manifest(chart)
    out = tmp_path / f"special-{damage}.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")
    presentation = Presentation(out)
    shape = next(item for item in presentation.slides[0].shapes if item.name == shape_name)
    if damage == "type":
        shape._element.xpath(".//a:prstGeom")[0].set("prst", "ellipse")
    elif damage == "geometry":
        off = shape._element.xpath(".//a:xfrm/a:off")[0]
        off.set("x", str(int(off.get("x")) + 1000))
    elif damage == "identity":
        shape._element.xpath(".//p:cNvPr")[0].set("descr", "object_id:wrong")
    else:
        shape.text = "999"
    presentation.save(out)

    violations = validate_pptx.quantitative_chart_readback_violations(out, [manifest])

    assert any(item["field"] == field for item in violations)


def test_final_deck_rebuild_preserves_all_special_chart_readback(tmp_path: Path) -> None:
    manifests = [_manifest(chart) for chart in (_cumulative_bridge(), _time_interval(), _variable_rectangle())]
    out = tmp_path / "special-final.pptx"

    write_deck(
        {"workflow_contract_version": "fixed-canvas-cm-v2", "slide": dict(SLIDE)},
        [
            {"manifest": deepcopy(manifest), "manifest_path": tmp_path / f"manifest-{index}.json"}
            for index, manifest in enumerate(manifests, start=1)
        ],
        out,
        [],
    )

    assert validate_pptx.quantitative_chart_readback_violations(out, manifests) == []


def test_dot_uses_editable_points_connectors_and_exact_source_labels(tmp_path: Path) -> None:
    chart = _one_dimensional("dot")
    manifest = _manifest(chart)
    out = tmp_path / "dot.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")

    slide = Presentation(out).slides[0]
    names = [shape.name for shape in slide.shapes]
    texts = [shape.text for shape in slide.shapes if shape.has_text_frame]

    assert not any(shape.has_chart for shape in slide.shapes)
    assert chart["name"] in names
    assert next(shape.text for shape in slide.shapes if shape.name == "Revenue chart Title") == "Revenue"
    assert sum(name.startswith("Revenue chart Point ") for name in names) == 2
    assert sum(name.startswith("Revenue chart Connector ") for name in names) == 2
    assert next(shape.text for shape in slide.shapes if shape.name == "Revenue chart Series 1") == "Revenue"
    assert {"A", "B", "12", "18", "USD m", "FY2025"}.issubset(texts)
    assert validate_pptx.quantitative_chart_readback_violations(out, [manifest]) == []


def test_explicit_target_actual_add_target_line_and_direct_difference_arrow(tmp_path: Path) -> None:
    chart = {**_one_dimensional("dot"), "target_value": 0.3, "actual_value": 0.2}
    manifest = _manifest(chart)
    out = tmp_path / "target.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")

    slide = Presentation(out).slides[0]
    named = {shape.name: shape for shape in slide.shapes}

    assert "Revenue chart Target Line" in named
    assert "Revenue chart Difference Arrow" in named
    assert named["Revenue chart Target"].text == "Target: 0.3"
    assert named["Revenue chart Actual"].text == "Actual: 0.2"
    assert named["Revenue chart Difference"].text == "Difference: -0.1"
    assert validate_pptx.quantitative_chart_readback_violations(out, [manifest]) == []


def test_one_dimensional_chart_accepts_explicit_shared_series_unit_and_basis(tmp_path: Path) -> None:
    chart = _one_dimensional()
    chart["series"][0].update({"unit": chart.pop("unit"), "basis": chart.pop("basis")})
    manifest = _manifest(chart)
    out = tmp_path / "series-metadata.pptx"

    write_pptx(manifest, out, tmp_path / "manifest.json")

    slide = Presentation(out).slides[0]
    assert next(shape.text for shape in slide.shapes if shape.name == "Revenue chart Unit") == "USD m"
    assert next(shape.text for shape in slide.shapes if shape.name == "Revenue chart Basis") == "same portfolio companies"
    assert validate_pptx.quantitative_chart_readback_violations(out, [manifest]) == []


@pytest.mark.parametrize(
    ("chart", "role", "changed", "field"),
    [
        (_one_dimensional(), "Basis", "wrong basis", "charts[0].basis"),
        (_one_dimensional(), "Basis", None, "charts[0].basis"),
        (_xy("scatter"), "X Basis", "wrong x basis", "charts[0].x_basis"),
        (_xy("scatter"), "Y Basis", "wrong y basis", "charts[0].y_basis"),
        (_xy("bubble"), "Size Basis", "wrong size basis", "charts[0].size_basis"),
    ],
)
def test_readback_rejects_changed_or_missing_dimension_basis(
    tmp_path: Path, chart: dict, role: str, changed: str | None, field: str
) -> None:
    manifest = _manifest(chart)
    out = tmp_path / "basis.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")
    presentation = Presentation(out)
    named = {shape.name: shape for shape in presentation.slides[0].shapes}
    basis_shape = named[f"{chart['name']} {role}"]
    if changed is None:
        basis_shape._element.getparent().remove(basis_shape._element)
    else:
        basis_shape.text = changed
    presentation.save(out)

    violations = validate_pptx.quantitative_chart_readback_violations(out, [manifest])

    assert any(item["field"] == field for item in violations)


def test_dot_readback_rejects_changed_title(tmp_path: Path) -> None:
    chart = _one_dimensional("dot")
    manifest = _manifest(chart)
    out = tmp_path / "dot-title.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")
    presentation = Presentation(out)
    named = {shape.name: shape for shape in presentation.slides[0].shapes}
    named["Revenue chart Title"].text = "Wrong title"
    presentation.save(out)

    violations = validate_pptx.quantitative_chart_readback_violations(out, [manifest])

    assert any(item["field"] == "charts[0].title" for item in violations)


@pytest.mark.parametrize(
    ("damage", "field"),
    [
        ("point_identity", "charts[0].points[0].object_id"),
        ("point_type", "charts[0].points[0].type"),
        ("connector_geometry", "charts[0].connectors[0].geometry"),
    ],
)
def test_dot_mark_readback_rejects_wrong_identity_type_or_geometry(
    tmp_path: Path, damage: str, field: str
) -> None:
    chart = _one_dimensional("dot")
    manifest = _manifest(chart)
    out = tmp_path / f"dot-{damage}.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")
    presentation = Presentation(out)
    named = {shape.name: shape for shape in presentation.slides[0].shapes}
    if damage == "point_identity":
        named["Revenue chart Point 1"]._element.xpath(".//p:cNvPr")[0].set("descr", "object_id:wrong")
    elif damage == "point_type":
        named["Revenue chart Point 1"]._element.xpath(".//a:prstGeom")[0].set("prst", "rect")
    else:
        off = named["Revenue chart Connector 1"]._element.xpath(".//a:xfrm/a:off")[0]
        off.set("x", str(float(off.get("x")) + 1000))
    presentation.save(out)

    violations = validate_pptx.quantitative_chart_readback_violations(out, [manifest])

    assert any(item["field"] == field for item in violations)


@pytest.mark.parametrize(
    ("damage", "field"),
    [
        ("target_identity", "charts[0].target_line.object_id"),
        ("target_geometry", "charts[0].target_line.geometry"),
        ("difference_type", "charts[0].difference_arrow.type"),
        ("difference_arrowheads", "charts[0].difference_arrow.arrowheads"),
        ("difference_label", "charts[0].difference"),
    ],
)
def test_target_mark_readback_rejects_wrong_identity_geometry_type_or_arrowheads(
    tmp_path: Path, damage: str, field: str
) -> None:
    chart = {**_one_dimensional("dot"), "target_value": 20, "actual_value": 18}
    manifest = _manifest(chart)
    out = tmp_path / f"target-{damage}.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")
    presentation = Presentation(out)
    slide = presentation.slides[0]
    named = {shape.name: shape for shape in slide.shapes}
    if damage == "target_identity":
        named["Revenue chart Target Line"]._element.xpath(".//p:cNvPr")[0].set("descr", "object_id:wrong")
    elif damage == "target_geometry":
        off = named["Revenue chart Target Line"]._element.xpath(".//a:xfrm/a:off")[0]
        off.set("y", str(float(off.get("y")) + 1000))
    elif damage == "difference_type":
        old = named["Revenue chart Difference Arrow"]
        xfrm = old._element.xpath(".//a:xfrm")[0]
        off, ext = xfrm.xpath("./a:off")[0], xfrm.xpath("./a:ext")[0]
        replacement = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(round(float(off.get("x")))),
            int(round(float(off.get("y")))),
            int(round(float(ext.get("cx")))),
            int(round(float(ext.get("cy")))),
        )
        replacement.name = old.name
        replacement._element.xpath(".//p:cNvPr")[0].set(
            "descr", old._element.xpath(".//p:cNvPr")[0].get("descr")
        )
        old._element.getparent().remove(old._element)
    elif damage == "difference_arrowheads":
        line = named["Revenue chart Difference Arrow"]._element.xpath(".//a:ln")[0]
        for arrow in line.xpath("./a:headEnd | ./a:tailEnd"):
            line.remove(arrow)
    else:
        named["Revenue chart Difference"].text = "Difference: 999"
    presentation.save(out)

    violations = validate_pptx.quantitative_chart_readback_violations(out, [manifest])

    assert any(item["field"] == field for item in violations)


@pytest.mark.parametrize(
    "mutate, expected",
    [
        (lambda manifest: manifest["charts"][0].pop("chart_variant"), "chart_variant"),
        (lambda manifest: manifest["charts"][0].update({"box_px": [-1, 90, 1142, 620]}), "box_px"),
        (lambda manifest: manifest["charts"][0].update({"anchor": "2cm,2cm,18cm,10cm"}), "anchor"),
        (lambda manifest: manifest["charts"][0].update({"target_value": 20}), "target_value"),
        (lambda manifest: manifest["text_boxes"].append({"object_id": "chart-1", "name": "duplicate", "text": "x", "box_px": [1, 1, 10, 10]}), "unique"),
    ],
)
def test_manifest_rejects_ambiguous_or_unreadable_chart_contract(
    tmp_path: Path, mutate, expected: str
) -> None:
    manifest = _manifest(_one_dimensional())
    mutate(manifest)

    with pytest.raises(ValueError, match=expected):
        write_pptx(manifest, tmp_path / "invalid.pptx", tmp_path / "manifest.json")


def test_readback_detects_changed_native_value(tmp_path: Path) -> None:
    chart = _one_dimensional()
    manifest = _manifest(chart)
    out = tmp_path / "tampered.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")
    presentation = Presentation(out)
    native_chart = next(shape.chart for shape in presentation.slides[0].shapes if shape.has_chart)
    changed = ChartData()
    changed.categories = ["A", "B"]
    changed.add_series("Revenue", [12, 99])
    native_chart.replace_data(changed)
    presentation.save(out)

    violations = validate_pptx.quantitative_chart_readback_violations(out, [manifest])

    assert any(item["field"] == "charts[0].series[0].values" for item in violations)


def test_readback_detects_changed_chart_object_id(tmp_path: Path) -> None:
    manifest = _manifest(_one_dimensional())
    out = tmp_path / "tampered-id.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")
    presentation = Presentation(out)
    chart_shape = next(shape for shape in presentation.slides[0].shapes if shape.has_chart)
    chart_shape._element.xpath(".//p:cNvPr")[0].set("descr", "object_id:wrong;chart_variant:column")
    presentation.save(out)

    violations = validate_pptx.quantitative_chart_readback_violations(out, [manifest])

    assert any(item["field"] == "charts[0].object_id" for item in violations)


def test_final_deck_rebuild_preserves_quantitative_chart_readback(tmp_path: Path) -> None:
    manifests = [
        _manifest(_xy("bubble")),
        _manifest({**_one_dimensional("dot"), "target_value": 20, "actual_value": 18}),
    ]
    out = tmp_path / "final.pptx"
    write_deck(
        {"workflow_contract_version": "fixed-canvas-cm-v2", "slide": dict(SLIDE)},
        [
            {"manifest": deepcopy(manifest), "manifest_path": tmp_path / f"manifest-{index}.json"}
            for index, manifest in enumerate(manifests, start=1)
        ],
        out,
        [],
    )

    assert validate_pptx.quantitative_chart_readback_violations(out, manifests) == []


@pytest.mark.parametrize("variant", ["column", "bar", "line", "scatter", "bubble"])
def test_native_variants_refuse_shape_based_target_and_difference_marks(
    tmp_path: Path, variant: str,
) -> None:
    chart = _xy(variant) if variant in {"scatter", "bubble"} else _one_dimensional(variant)
    chart.update({"target_value": 20, "actual_value": 18})

    with pytest.raises(ValueError, match="dot"):
        write_pptx(_manifest(chart), tmp_path / f"{variant}-target.pptx", tmp_path / "manifest.json")


def test_preview_dot_matches_delivered_metadata_labels_and_horizontal_marks(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    chart = {**_one_dimensional("dot"), "target_value": 0.3, "actual_value": 0.2}
    chart["series"][0]["values"] = [0.1, 0.2]
    manifest = _manifest(chart)
    texts: list[str] = []
    lines: list[tuple[float, float, float, float]] = []
    original_text = ImageDraw.ImageDraw.text
    original_line = ImageDraw.ImageDraw.line

    def record_text(draw, xy, text, *args, **kwargs):
        texts.append(str(text))
        return original_text(draw, xy, text, *args, **kwargs)

    def record_line(draw, xy, *args, **kwargs):
        coordinates = tuple(xy)
        if len(coordinates) == 4:
            lines.append(coordinates)
        return original_line(draw, xy, *args, **kwargs)

    monkeypatch.setattr(ImageDraw.ImageDraw, "text", record_text)
    monkeypatch.setattr(ImageDraw.ImageDraw, "line", record_line)

    render_preview(manifest, tmp_path / "manifest.json", tmp_path / "preview.png")

    assert {
        "Revenue", "USD m", "FY2025", "same portfolio companies",
        "A", "B", "0.1", "0.2", "Target: 0.3", "Actual: 0.2", "Difference: -0.1",
    }.issubset(texts)
    assert sum(y1 == y2 and x1 != x2 for x1, y1, x2, y2 in lines) >= 3
    assert sum(x1 == x2 and y1 != y2 for x1, y1, x2, y2 in lines) >= 1


def test_preview_negative_bar_uses_shared_zero_baseline(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    chart = _one_dimensional("bar")
    chart["series"][0]["values"] = [-10, 20]
    rectangles: list[tuple[float, float, float, float]] = []
    original_rectangle = ImageDraw.ImageDraw.rectangle

    def record_rectangle(draw, xy, *args, **kwargs):
        if kwargs.get("fill") == "#4472C4":
            rectangles.append(tuple(xy))
        return original_rectangle(draw, xy, *args, **kwargs)

    monkeypatch.setattr(ImageDraw.ImageDraw, "rectangle", record_rectangle)

    render_preview(_manifest(chart), tmp_path / "manifest.json", tmp_path / "preview.png")

    assert len(rectangles) == 2
    assert all(right > left for left, _top, right, _bottom in rectangles)
    assert rectangles[0][2] == pytest.approx(rectangles[1][0])


def test_chart_readback_skips_opening_pptx_when_manifests_have_no_charts(tmp_path: Path) -> None:
    manifest = _manifest(_one_dimensional())
    manifest["charts"] = []
    assert validate_pptx.quantitative_chart_readback_violations(
        tmp_path / "not-a-presentation.pptx", [manifest],
    ) == []


def test_extracted_chart_canonicalizes_through_authority_manifest_build_and_readback(tmp_path: Path) -> None:
    xml = b'''<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">
      <c:chart><c:title><c:tx><c:v>Revenue</c:v></c:tx></c:title><c:plotArea><c:barChart><c:barDir val="col"/><c:ser>
      <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f><c:strCache><c:pt idx="0"><c:v>Revenue</c:v></c:pt></c:strCache></c:strRef></c:tx>
      <c:cat><c:numLit><c:pt idx="0"><c:v>2024</c:v></c:pt><c:pt idx="1"><c:v>2025</c:v></c:pt></c:numLit></c:cat>
      <c:val><c:numLit><c:pt idx="0"><c:v>12</c:v></c:pt><c:pt idx="1"><c:v>18.5</c:v></c:pt></c:numLit></c:val>
      </c:ser></c:barChart><c:valAx><c:title><c:tx><c:v>Revenue (USD m)</c:v></c:tx></c:title></c:valAx>
      </c:plotArea></c:chart></c:chartSpace>'''
    extracted = _chart_record({"page_numbers": [2], "asset_id": "word_asset_001"}, xml)
    extracted.update({"basis": "same portfolio companies", "period": "FY2025"})

    authority = select_numeric_authority([extracted])

    assert authority is not None
    assert authority["series"] == [{
        "name": "Revenue",
        "categories": ["2024", "2025"],
        "category_indices": [0, 1],
        "values": [12, 18.5],
        "value_indices": [0, 1],
    }]
    chart = {"object_id": "chart-1", "name": "Revenue chart", "box_px": [190, 90, 1142, 620], **authority}
    manifest = _manifest(chart)
    out = tmp_path / "extracted.pptx"
    write_pptx(manifest, out, tmp_path / "manifest.json")

    assert validate_pptx.quantitative_chart_readback_violations(out, [manifest]) == []
