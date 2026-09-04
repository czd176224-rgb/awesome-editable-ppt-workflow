import json
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

from editppt.runtime.build_pptx_from_manifest import normalize_manifest, write_pptx
from editppt.runtime.fixed_region_runtime import CONTENT_BOX, SLIDE
from editppt.runtime.validate_pptx import ALLOWED_SOURCE_TYPES, _connector_endpoints, _shape_arrowheads


def _manifest(width=1700, height=800):
    return {
        "workflow_contract_version": "fixed-canvas-cm-v2",
        "reconstruction_contract_version": "editable-image-v3",
        "slide": dict(SLIDE),
        "content_box": dict(CONTENT_BOX),
        "source": {"width_px": width, "height_px": height},
        "text_boxes": [{"object_id": "word-p1", "name": "body-paragraph-1", "text": "权威正文", "box_px": [80, 60, 700, 90]}],
        "tables": [{
            "object_id": "word-t1", "name": "body-table-1", "box_px": [80, 220, 1100, 300],
            "rows": [["项目", "数值"], ["投资额", "50万元"]],
            "font_size": 12, "font_color": "#000000", "cell_fill": "#FFFFFF", "cell_margin_px": 8,
        }],
        "shapes": [{"object_id": "decor-1", "name": "decorative-panel", "type": "rect", "box_px": [40, 30, 1500, 650], "fill": "#f4f4f4"}],
        "images": [],
    }


def _directed_edge_manifest(points_px, source_box=None, target_box=None):
    manifest = _manifest()
    manifest["shapes"] = [
        {"object_id": "source", "name": "source", "type": "rect", "box_px": source_box or [100, 100, 100, 100], "fill": "#FFFFFF"},
        {"object_id": "target", "name": "target", "type": "rect", "box_px": target_box or [400, 100, 100, 100], "fill": "#FFFFFF"},
        {
            "object_id": "edge:source->target", "name": "edge:source->target", "type": "line",
            "points_px": points_px, "stroke": "#6B7A90",
        },
    ]
    return manifest


def _shape_by_object_id(slide, object_id):
    return next(
        shape for shape in slide.shapes
        if shape._element.xpath(".//p:cNvPr")[0].get("descr") == f"object_id:{object_id}"
    )


def test_authentic_published_source_is_a_first_class_provenance_type():
    assert "authentic-published-source" in ALLOWED_SOURCE_TYPES


def test_v4_rejects_unexpected_16_by_9_instead_of_containing():
    with pytest.raises(ValueError, match="17:8"):
        normalize_manifest(_manifest(1600, 900))


def test_v4_builds_stable_named_text_shape_and_native_table(tmp_path: Path):
    manifest = _manifest()
    manifest_path = tmp_path / "manifest.json"
    manifest_path.write_text(json.dumps(manifest), encoding="utf-8")
    output = tmp_path / "page.pptx"
    write_pptx(manifest, output, manifest_path)

    deck = Presentation(output)
    names = {shape.name for shape in deck.slides[0].shapes}
    assert {"body-paragraph-1", "body-table-1", "decorative-panel"} <= names
    table = next(shape.table for shape in deck.slides[0].shapes if shape.has_table)
    assert table.cell(1, 1).text == "50万元"
    with zipfile.ZipFile(output) as archive:
        xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert 'descr="object_id:word-p1"' in xml
    assert 'descr="object_id:word-t1"' in xml
    assert '<a:srgbClr val="000000"' in xml
    assert '<a:srgbClr val="FFFFFF"' in xml
    assert 'marL="' in xml and 'marR="' in xml and 'marT="' in xml and 'marB="' in xml


def test_v4_native_table_rejects_unresolved_default_cell_visuals():
    manifest = _manifest()
    for field in ("font_size", "font_color", "cell_fill", "cell_margin_px"):
        manifest["tables"][0].pop(field)
    with pytest.raises(ValueError, match="table cells require explicit"):
        normalize_manifest(manifest)


@pytest.mark.parametrize("points_px,source_box,target_box", [
    ([150, 150, 399, 150], None, None),  # target endpoint misses its left edge
    ([99, 150, 450, 150], None, None),  # source endpoint misses its left edge
    ([399, 150, 201, 150], [400, 100, 100, 100], [100, 100, 100, 100]),  # flipH
])
def test_v4_sealed_directed_edge_snaps_one_pixel_misses_and_serializes_target_arrow(
    tmp_path: Path, points_px, source_box, target_box,
):
    manifest = _directed_edge_manifest(points_px, source_box, target_box)
    output = tmp_path / "page.pptx"
    write_pptx(manifest, output, tmp_path / "manifest.json")

    slide = Presentation(output).slides[0]
    source = _shape_by_object_id(slide, "source")
    target = _shape_by_object_id(slide, "target")
    edge = _shape_by_object_id(slide, "edge:source->target")
    start_x, start_y, end_x, end_y = _connector_endpoints(edge)
    assert source.left <= start_x <= source.left + source.width
    assert source.top <= start_y <= source.top + source.height
    assert target.left <= end_x <= target.left + target.width
    assert target.top <= end_y <= target.top + target.height
    assert _shape_arrowheads(edge) == {"tailEnd": "triangle"}


def test_v4_sealed_directed_edge_rejects_two_pixel_target_miss(tmp_path: Path):
    manifest = _directed_edge_manifest([200, 150, 398, 150])

    with pytest.raises(ValueError, match="outside target node"):
        write_pptx(manifest, tmp_path / "page.pptx", tmp_path / "manifest.json")


def test_v4_sealed_directed_edge_rejects_ambiguous_node_ids(tmp_path: Path):
    manifest = _directed_edge_manifest([150, 150, 450, 150])
    manifest["shapes"][0]["object_id"] = "a"
    manifest["shapes"][1]["object_id"] = "b->target"
    manifest["shapes"].extend([
        {"object_id": "a->b", "name": "alternate-source", "type": "rect", "box_px": [100, 300, 100, 100]},
        {"object_id": "target", "name": "alternate-target", "type": "rect", "box_px": [400, 300, 100, 100]},
    ])
    manifest["shapes"][2]["object_id"] = "edge:a->b->target"

    with pytest.raises(ValueError, match="missing or ambiguous"):
        write_pptx(manifest, tmp_path / "page.pptx", tmp_path / "manifest.json")


def test_v4_plain_line_keeps_geometry_and_has_no_arrowhead(tmp_path: Path):
    manifest = _manifest()
    manifest["shapes"] = [{
        "object_id": "plain-line", "name": "plain-line", "type": "line",
        "points_px": [200, 150, 399, 150], "stroke": "#6B7A90",
    }]
    output = tmp_path / "page.pptx"
    write_pptx(manifest, output, tmp_path / "manifest.json")

    edge = _shape_by_object_id(Presentation(output).slides[0], "plain-line")
    assert _connector_endpoints(edge) == tuple(
        round(value * 914400)
        for value in normalize_manifest(manifest)["shapes"][0]["points"]
    )
    assert _shape_arrowheads(edge) == {}
