#!/usr/bin/env python3
import argparse
import hashlib
import json
import posixpath
import re
import sys
import zipfile
from decimal import Decimal
from pathlib import Path
from xml.etree import ElementTree as ET

try:
    from .build_pptx_from_manifest import (
        _chart_basis_labels,
        _chart_mark_geometry,
        _number_text,
        _chart_shape_description,
        _chart_shared_text,
        _special_chart_records,
        SPECIAL_CHART_PRIMITIVES,
        normalize_manifest,
    )
except ImportError:  # direct runtime script execution through the editppt launcher
    from build_pptx_from_manifest import (
        _chart_basis_labels,
        _chart_mark_geometry,
        _number_text,
        _chart_shape_description,
        _chart_shared_text,
        _special_chart_records,
        SPECIAL_CHART_PRIMITIVES,
        normalize_manifest,
    )


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}

ALLOWED_SOURCE_TYPES = {
    "asset-sheet-separated",
    "authentic-published-source",
    "imagegen",
    "latex-rendered-formula",
    "user-provided",
    "user-approved-rasterization",
}
REQUIRED_QUALITY_CHECKS = {
    "font_size_calibrated",
    "visual_inventory_matched",
    "background_strategy_checked",
    "shape_corner_geometry_checked",
}
FOREGROUND_TERMS = {
    "badge",
    "decorative",
    "foreground",
    "hand-drawn",
    "icon",
    "illustration",
    "image block",
    "logo",
    "mark",
    "photo",
    "pictogram",
    "screenshot",
    "semantic",
    "sticker",
    "symbol",
    "trend",
    "visual object",
    "前景",
    "图标",
    "照片",
    "截图",
    "徽章",
    "贴纸",
    "语义",
    "视觉对象",
}
NON_FOREGROUND_TERMS = {
    "background",
    "clean base",
    "formula",
    "latex",
    "native structural",
    "structural shape",
    "背景",
    "公式",
    "结构",
}
ASSET_SHEET_TERMS = {
    "asset-sheet",
    "asset sheet",
    "asset_sheet",
    "image edit",
    "imagegen",
    "separated",
    "source-faithful",
    "source faithful",
    "split",
    "分离",
}
FORBIDDEN_FOREGROUND_FALLBACK_TERMS = {
    "approximate",
    "approximation",
    "approximated",
    "crop",
    "cropped",
    "direct crop",
    "direct source",
    "emoji",
    "fallback",
    "native approximation",
    "source crop",
    "source snippet",
    "text symbol",
    "warning only",
    "warning_only",
    "近似",
    "裁切",
    "裁剪",
    "降级",
}


def read_manifest(path):
    if not path:
        return {}
    return json.loads(Path(path).read_text(encoding="utf-8"))


def compact_text(value):
    if value is None:
        return ""
    if isinstance(value, str):
        return value.lower()
    if isinstance(value, (int, float, bool)):
        return str(value).lower()
    if isinstance(value, dict):
        return " ".join(compact_text(item) for item in value.values())
    if isinstance(value, (list, tuple, set)):
        return " ".join(compact_text(item) for item in value)
    return str(value).lower()


def contains_any(text, terms):
    return any(term in text for term in terms)


def visual_item_path(item):
    for key in ("path", "asset", "asset_path", "image", "image_path", "corresponding_asset"):
        value = item.get(key) if isinstance(item, dict) else None
        if isinstance(value, str) and value.strip():
            return Path(value).as_posix()
    return None


def is_foreground_visual_item(item):
    text = compact_text(item)
    if contains_any(text, NON_FOREGROUND_TERMS):
        return False
    return contains_any(text, FOREGROUND_TERMS)


def foreground_asset_contract_violations(manifest):
    violations = []
    provenance_by_path = {
        Path(entry.get("path", "")).as_posix(): entry
        for entry in manifest.get("asset_provenance", [])
        if entry.get("path")
    }

    for index, item in enumerate(manifest.get("visual_inventory", [])):
        if not isinstance(item, dict):
            continue
        text = compact_text(item)
        field = f"visual_inventory[{index}]"
        if contains_any(text, FORBIDDEN_FOREGROUND_FALLBACK_TERMS):
            violations.append(
                {
                    "field": field,
                    "reason": "foreground visual decisions must not use direct crops, native approximations, emoji/text symbols, warning-only fallbacks, or similar shortcuts",
                }
            )
        if not is_foreground_visual_item(item):
            continue
        if not contains_any(text, ASSET_SHEET_TERMS):
            violations.append(
                {
                    "field": field,
                    "reason": "foreground visual objects must explicitly use source-faithful asset-sheet separation",
                }
            )
        path = visual_item_path(item)
        if path:
            provenance = provenance_by_path.get(path, {})
            source_type = provenance.get("source_type")
            if source_type in {"user-provided", "user-approved-rasterization"}:
                violations.append(
                    {
                        "field": field,
                        "path": path,
                        "reason": "foreground visual objects cannot use user-provided/direct raster provenance; use asset-sheet separation",
                    }
                )

    for index, entry in enumerate(manifest.get("asset_provenance", [])):
        if not isinstance(entry, dict):
            continue
        text = compact_text(entry)
        source_type = entry.get("source_type")
        path = Path(entry.get("path", "")).as_posix()
        field = f"asset_provenance[{index}]"
        if source_type in {"user-provided", "user-approved-rasterization"} and contains_any(
            text, FOREGROUND_TERMS | FORBIDDEN_FOREGROUND_FALLBACK_TERMS
        ):
            violations.append(
                {
                    "field": field,
                    "path": path,
                    "reason": "foreground-like raster provenance cannot be direct user-provided/cropped source material",
                }
            )
        if contains_any(text, FORBIDDEN_FOREGROUND_FALLBACK_TERMS):
            violations.append(
                {
                    "field": field,
                    "path": path,
                    "reason": "asset provenance records a forbidden foreground fallback such as crop, approximation, or warning-only delivery",
                }
            )

    return violations


def is_full_slide_image(item, slide):
    width = float(slide.get("width", 13.333))
    height = float(slide.get("height", 7.5))
    left = float(item.get("left", 0))
    top = float(item.get("top", 0))
    image_width = float(item.get("width", 0))
    image_height = float(item.get("height", 0))
    return (
        abs(left) <= 0.02
        and abs(top) <= 0.02
        and image_width >= width * 0.98
        and image_height >= height * 0.98
    )


def page_contract_violations(manifest):
    violations = []
    slide = manifest.get("slide", {})
    images = manifest.get("images", [])
    text_boxes = manifest.get("text_boxes", [])
    provenance_by_path = {
        Path(entry.get("path", "")).as_posix(): entry
        for entry in manifest.get("asset_provenance", [])
        if entry.get("path")
    }
    for image in images:
        path = Path(image.get("path", "")).as_posix()
        provenance = provenance_by_path.get(path, {})
        source_type = provenance.get("source_type")
        if is_full_slide_image(image, slide) and Path(path).name == "source.png" and text_boxes:
            violations.append(
                {
                    "field": "images",
                    "path": path,
                    "reason": "full-slide source.png background with editable text overlays causes baked-text overlap",
                }
            )
        if (
            is_full_slide_image(image, slide)
            and source_type in {"user-provided", "user-approved-rasterization"}
            and text_boxes
        ):
            violations.append(
                {
                    "field": "asset_provenance",
                    "path": path,
                    "reason": "full-slide raster background cannot be assembled with editable text",
                }
            )

    return violations


def quality_contract_violations(manifest):
    violations = []

    if "visual_inventory" not in manifest:
        violations.append(
            {
                "field": "visual_inventory",
                "reason": "page manifest must record the non-text visual inventory, even when it is empty",
            }
        )
    elif not isinstance(manifest.get("visual_inventory"), list):
        violations.append({"field": "visual_inventory", "reason": "visual_inventory must be a list"})

    background_strategy = manifest.get("background_strategy")
    if not background_strategy:
        violations.append(
            {
                "field": "background_strategy",
                "reason": "page manifest must record how the background was rebuilt or preserved",
            }
        )

    quality_checks = manifest.get("quality_checks")
    if not isinstance(quality_checks, dict):
        violations.append({"field": "quality_checks", "reason": "quality_checks must be an object"})
    else:
        for key in sorted(REQUIRED_QUALITY_CHECKS):
            if quality_checks.get(key) is not True:
                violations.append(
                    {
                        "field": f"quality_checks.{key}",
                        "reason": "required page QA check must be explicitly true",
                    }
                )

    for index, shape in enumerate(manifest.get("shapes", [])):
        is_round_rect = shape.get("type") == "roundRect" or shape.get("preset") == "roundRect"
        if is_round_rect and not shape.get("source_corner_radius_px"):
            violations.append(
                {
                    "field": f"shapes[{index}]",
                    "reason": "roundRect requires source_corner_radius_px; use rect for source straight-corner containers",
                }
            )
        if is_round_rect and shape.get("source_corner_radius_px") and shape.get("box_px"):
            box = shape.get("box_px")
            radius = float(shape.get("source_corner_radius_px") or 0)
            min_dim = max(1.0, min(float(box[2]), float(box[3])))
            if radius > min_dim / 2:
                violations.append(
                    {
                        "field": f"shapes[{index}].source_corner_radius_px",
                        "reason": "roundRect source radius cannot exceed half of the smaller shape dimension",
                    }
                )

    violations.extend(foreground_asset_contract_violations(manifest))
    return violations


def pixel_authoring_violations(manifest):
    violations = []
    source = manifest.get("source", {})
    if not source.get("width_px") or not source.get("height_px"):
        violations.append(
            {
                "field": "source.width_px/source.height_px",
                "reason": "page manifest must record the source image pixel size",
            }
        )

    for section in ("text_boxes", "images"):
        for index, item in enumerate(manifest.get(section, [])):
            if "box_px" not in item:
                violations.append(
                    {
                        "field": f"{section}[{index}].box_px",
                        "reason": "positioned text and image objects must use source-image pixel coordinates",
                    }
                )

    for index, item in enumerate(manifest.get("shapes", [])):
        if item.get("type") == "line":
            if "points_px" not in item:
                violations.append(
                    {
                        "field": f"shapes[{index}].points_px",
                        "reason": "line shapes must use source-image pixel endpoints",
                    }
                )
        elif "box_px" not in item:
            violations.append(
                {
                    "field": f"shapes[{index}].box_px",
                    "reason": "positioned shapes must use source-image pixel coordinates",
                }
            )

    return violations


def normalize_for_validation(manifest):
    violations = pixel_authoring_violations(manifest)
    try:
        return normalize_manifest(manifest), violations
    except Exception as exc:
        violations.append({"field": "manifest", "reason": str(exc)})
        return manifest, violations


def sha256_text(value):
    return hashlib.sha256(str(value).encode("utf-8")).hexdigest()


def flatten_required_text(value):
    """Return exact text strings that should be verified in the PPTX."""
    items = []
    if value is None:
        return items
    if isinstance(value, str):
        text = value.strip()
        return [text] if text else []
    if isinstance(value, (int, float)):
        return [str(value)]
    if isinstance(value, dict):
        if "required_text" in value:
            return flatten_required_text(value.get("required_text"))
        if "items" in value:
            return flatten_required_text(value.get("items"))
        if "texts" in value:
            return flatten_required_text(value.get("texts"))
        if "text" in value:
            return flatten_required_text(value.get("text"))
        return items
    if isinstance(value, (list, tuple, set)):
        for item in value:
            items.extend(flatten_required_text(item))
    return items


def required_texts_from_manifest(manifest):
    required = []
    required.extend(flatten_required_text(manifest.get("required_text", [])))
    required.extend(flatten_required_text(manifest.get("text_inventory", [])))
    return required


def collect_text(xml_bytes):
    root = ET.fromstring(xml_bytes)
    return "".join(node.text or "" for node in root.findall(".//a:t", NS))


def collect_paragraph_text(xml_bytes):
    root = ET.fromstring(xml_bytes)
    paragraphs = []
    for paragraph in root.findall(".//a:p", NS):
        text = "".join(node.text or "" for node in paragraph.findall(".//a:t", NS))
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def collect_notes_texts(z, names):
    notes = {}
    for name in sorted(n for n in names if re.match(r"ppt/notesSlides/notesSlide\d+\.xml$", n)):
        match = re.search(r"notesSlide(\d+)\.xml$", name)
        if not match:
            continue
        notes[int(match.group(1))] = collect_paragraph_text(z.read(name))
    return notes


def _chart_xml_values(chart, tag):
    values = []
    for series in chart._element.xpath(".//c:ser"):
        dimension = series.xpath(f"./c:{tag}")
        points = dimension[0].xpath(".//c:pt") if dimension else []
        indexed = []
        for point in points:
            value = point.xpath("./c:v")
            indexed.append((int(point.get("idx")), float(value[0].text)))
        values.append([value for _index, value in sorted(indexed)])
    return values


def _shape_geometry(shape):
    xfrm = shape._element.xpath(".//a:xfrm")[0]
    off, ext = xfrm.xpath("./a:off")[0], xfrm.xpath("./a:ext")[0]
    return tuple(int(round(float(value))) for value in (off.get("x"), off.get("y"), ext.get("cx"), ext.get("cy")))


def _connector_endpoints(shape):
    left, top, width, height = _shape_geometry(shape)
    xfrm = shape._element.xpath(".//a:xfrm")[0]
    flip_h = xfrm.get("flipH") in {"1", "true"}
    flip_v = xfrm.get("flipV") in {"1", "true"}
    return (
        left + width if flip_h else left,
        top + height if flip_v else top,
        left if flip_h else left + width,
        top if flip_v else top + height,
    )


def _shape_kind(shape):
    local = shape._element.tag.rsplit("}", 1)[-1]
    if local == "cxnSp":
        return "connector"
    geometry = shape._element.xpath(".//a:prstGeom")
    return geometry[0].get("prst") if geometry else local


def _shape_arrowheads(shape):
    line = shape._element.xpath(".//a:ln")
    if not line:
        return {}
    return {
        node.tag.rsplit("}", 1)[-1]: node.get("type")
        for node in line[0]
        if node.tag.rsplit("}", 1)[-1] in {"headEnd", "tailEnd"}
    }


def quantitative_chart_readback_violations(pptx_path, manifests):
    """Read exact chart objects and direct marks back from the generated PPTX."""
    if not any(manifest.get("charts") for manifest in manifests):
        return []
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    chart_types = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "bubble": XL_CHART_TYPE.BUBBLE,
    }
    violations = []
    try:
        normalized = [normalize_manifest(manifest) for manifest in manifests]
        presentation = Presentation(pptx_path)
    except Exception as exc:
        return [{"field": "charts", "reason": str(exc)}]
    if len(presentation.slides) != len(normalized):
        return [{"field": "charts", "reason": "slide count does not match chart manifests"}]

    def add(field, expected, actual):
        if actual != expected:
            violations.append({"field": field, "expected": expected, "actual": actual})

    def validate_mark(named, chart, role, field, kind, geometry, arrowheads=None):
        shape = named.get(f"{chart['name']} {role}")
        if shape is None:
            violations.append({"field": field, "reason": "missing mark"})
            return
        description = shape._element.xpath(".//p:cNvPr")[0].get("descr")
        add(f"{field}.object_id", _chart_shape_description(chart, role), description)
        add(f"{field}.type", kind, _shape_kind(shape))
        actual_geometry = _connector_endpoints(shape) if kind == "connector" else _shape_geometry(shape)
        add(f"{field}.geometry", geometry, actual_geometry)
        if arrowheads is not None:
            add(f"{field}.arrowheads", arrowheads, _shape_arrowheads(shape))

    def validate_special_records(slide, chart, prefix):
        records = _special_chart_records(chart)
        for section in ("shapes", "text_boxes"):
            for record in records[section]:
                field = prefix + (f".{record['_field']}" if record["_field"] else "")
                matches = [shape for shape in slide.shapes if shape.name == record["name"]]
                if len(matches) != 1:
                    violations.append({"field": field, "reason": f"expected one editable object, found {len(matches)}"})
                    continue
                shape = matches[0]
                description = shape._element.xpath(".//p:cNvPr")[0].get("descr")
                add(f"{field}.object_id", f"object_id:{record['object_id']}", description)
                expected_kind = "text_box" if section == "text_boxes" else record["type"]
                actual_kind = (
                    "text_box"
                    if shape._element.xpath(".//p:cNvSpPr[@txBox='1']")
                    else _shape_kind(shape)
                )
                add(f"{field}.type", expected_kind, actual_kind)
                add(
                    f"{field}.geometry",
                    tuple(int(round(float(record[key]) * 914400)) for key in ("left", "top", "width", "height")),
                    _shape_geometry(shape),
                )
                if section == "text_boxes":
                    add(field, str(record["text"]), shape.text)

    for slide_index, (slide, manifest) in enumerate(zip(presentation.slides, normalized)):
        named = {shape.name: shape for shape in slide.shapes}
        for chart_index, expected in enumerate(manifest.get("charts", [])):
            prefix = f"slides[{slide_index}].charts[{chart_index}]" if len(normalized) > 1 else f"charts[{chart_index}]"
            if expected["rendering_primitive"] in SPECIAL_CHART_PRIMITIVES:
                validate_special_records(slide, expected, prefix)
                continue
            root = named.get(expected["name"])
            if root is None:
                violations.append({"field": prefix, "reason": "missing chart root object"})
                continue
            description = root._element.xpath(".//p:cNvPr")[0].get("descr")
            add(
                f"{prefix}.object_id",
                _chart_shape_description(expected),
                description,
            )
            add(f"{prefix}.box.left", int(expected["left"] * 914400), root.left)
            add(f"{prefix}.box.top", int(expected["top"] * 914400), root.top)
            add(f"{prefix}.box.width", int(expected["width"] * 914400), root.width)
            add(f"{prefix}.box.height", int(expected["height"] * 914400), root.height)
            variant = expected["chart_variant"]
            unit = _chart_shared_text(expected, "unit") if expected["rendering_primitive"] != "xy" else None
            if expected["rendering_primitive"] == "xy":
                unit = f"x: {expected['x_unit']} | y: {expected['y_unit']}"
                if variant == "bubble":
                    unit += f" | size: {expected['size_unit']}"
            metadata = [("Unit", unit), ("Period", expected["period"]), *_chart_basis_labels(expected)]
            for role, text in metadata:
                shape = named.get(f"{expected['name']} {role}")
                field = role.lower().replace(" ", "_")
                add(f"{prefix}.{field}", text, shape.text if shape is not None else None)

            if variant == "dot":
                title_shape = named.get(f"{expected['name']} Title")
                add(f"{prefix}.title", expected["title"], title_shape.text if title_shape is not None else None)
                point_count = sum(len(item["values"]) for item in expected["series"])
                add(
                    f"{prefix}.points",
                    point_count,
                    sum(name.startswith(f"{expected['name']} Point ") for name in named),
                )
                add(
                    f"{prefix}.connectors",
                    point_count,
                    sum(name.startswith(f"{expected['name']} Connector ") for name in named),
                )
                for series_index, item in enumerate(expected["series"], start=1):
                    series_shape = named.get(f"{expected['name']} Series {series_index}")
                    add(f"{prefix}.series[{series_index - 1}].name", item["name"], series_shape.text if series_shape is not None else None)
                point_index = 1
                mark_geometry = _chart_mark_geometry(expected)
                for series_index, item in enumerate(expected["series"]):
                    for category, value in zip(item["categories"], item["values"]):
                        validate_mark(
                            named,
                            expected,
                            f"Point {point_index}",
                            f"{prefix}.points[{point_index - 1}]",
                            "ellipse",
                            mark_geometry[f"Point {point_index}"],
                        )
                        validate_mark(
                            named,
                            expected,
                            f"Connector {point_index}",
                            f"{prefix}.connectors[{point_index - 1}]",
                            "connector",
                            mark_geometry[f"Connector {point_index}"],
                        )
                        category_shape = named.get(f"{expected['name']} Category {point_index}")
                        value_shape = named.get(f"{expected['name']} Value {point_index}")
                        add(f"{prefix}.series[{series_index}].categories[{point_index - 1}]", category, category_shape.text if category_shape is not None else None)
                        add(f"{prefix}.series[{series_index}].values[{point_index - 1}]", str(value), value_shape.text if value_shape is not None else None)
                        point_index += 1
            else:
                if not root.has_chart:
                    violations.append({"field": prefix, "reason": "expected native chart object"})
                    continue
                native = root.chart
                add(f"{prefix}.chart_variant", chart_types[variant], native.chart_type)
                add(f"{prefix}.title", expected["title"], native.chart_title.text_frame.text if native.has_title else None)
                if expected["rendering_primitive"] == "xy":
                    add(f"{prefix}.x_label", expected["x_label"], native.category_axis.axis_title.text_frame.text if native.category_axis.has_title else None)
                    add(f"{prefix}.y_label", expected["y_label"], native.value_axis.axis_title.text_frame.text if native.value_axis.has_title else None)
                add(f"{prefix}.series.names", [item["name"] for item in expected["series"]], [item.name for item in native.series])
                if variant in {"column", "bar", "line"}:
                    add(f"{prefix}.categories", expected["series"][0]["categories"], [item.label for item in native.plots[0].categories])
                    for series_index, item in enumerate(expected["series"]):
                        add(f"{prefix}.series[{series_index}].values", item["values"], list(native.series[series_index].values))
                else:
                    x_values = _chart_xml_values(native, "xVal")
                    y_values = _chart_xml_values(native, "yVal")
                    size_values = _chart_xml_values(native, "bubbleSize") if variant == "bubble" else []
                    for series_index, item in enumerate(expected["series"]):
                        add(f"{prefix}.series[{series_index}].x_values", item["x_values"], x_values[series_index])
                        add(f"{prefix}.series[{series_index}].y_values", item["y_values"], y_values[series_index])
                        if variant == "bubble":
                            add(f"{prefix}.series[{series_index}].size_values", item["size_values"], size_values[series_index])

            if expected.get("target_value") is not None:
                for role, value in (
                    ("Target", f"Target: {_number_text(expected['target_value'])}"),
                    ("Actual", f"Actual: {_number_text(expected['actual_value'])}"),
                    ("Difference", f"Difference: {_number_text(Decimal(str(expected['actual_value'])) - Decimal(str(expected['target_value'])))}"),
                ):
                    shape = named.get(f"{expected['name']} {role}")
                    add(f"{prefix}.{role.lower()}", value, shape.text if shape is not None else None)
                mark_geometry = _chart_mark_geometry(expected)
                validate_mark(
                    named, expected, "Target Line", f"{prefix}.target_line", "connector",
                    mark_geometry["Target Line"],
                )
                validate_mark(
                    named, expected, "Difference Arrow", f"{prefix}.difference_arrow", "connector",
                    mark_geometry["Difference Arrow"], {"headEnd": "triangle", "tailEnd": "triangle"},
                )
    return violations


def validate_deck(args):
    deck_path = Path(args.deck_manifest).resolve()
    deck = read_manifest(deck_path)
    root = Path(deck.get("job_dir", deck_path.parent)).resolve()
    expected_pages = int(deck.get("page_count", len(deck.get("pages", []))))
    notes_manifest = {}
    notes_path = deck.get("notes_manifest")
    if notes_path:
        notes_file = Path(notes_path)
        if not notes_file.is_absolute():
            notes_file = root / notes_file
        if notes_file.exists():
            notes_manifest = read_manifest(notes_file)

    report = {
        "pptx": str(Path(args.pptx).resolve()),
        "deck_manifest": str(deck_path),
        "expected_pages": expected_pages,
        "slides": 0,
        "page_manifests_missing": [],
        "page_validation_missing": [],
        "failed_page_validations": [],
        "page_contract_violations": [],
        "notes_expected": len(notes_manifest.get("notes", [])),
        "notes_found": 0,
        "notes_hash_mismatches": [],
        "chart_readback_violations": [],
        "missing_parts": [],
        "warnings": [],
        "passed": False,
    }

    page_manifests = []

    for page in deck.get("pages", []):
        manifest_path = Path(page.get("manifest", ""))
        validation_path = Path(page.get("validation", ""))
        if not manifest_path.is_absolute():
            manifest_path = root / manifest_path
        if not validation_path.is_absolute():
            validation_path = root / validation_path
        if not manifest_path.exists():
            report["page_manifests_missing"].append(str(manifest_path))
        else:
            try:
                raw_manifest = read_manifest(manifest_path)
                page_manifests.append(raw_manifest)
                normalized_manifest, authoring_violations = normalize_for_validation(raw_manifest)
                violations = (
                    authoring_violations
                    + page_contract_violations(normalized_manifest)
                    + quality_contract_violations(raw_manifest)
                )
                if violations:
                    report["page_contract_violations"].append(
                        {
                            "page_id": page.get("page_id"),
                            "manifest": str(manifest_path),
                            "violations": violations,
                        }
                    )
            except Exception as exc:
                report["page_contract_violations"].append(
                    {
                        "page_id": page.get("page_id"),
                        "manifest": str(manifest_path),
                        "violations": [{"field": "manifest", "reason": str(exc)}],
                    }
                )
        if not validation_path.exists():
            report["page_validation_missing"].append(str(validation_path))
        else:
            try:
                page_report = read_manifest(validation_path)
                if page_report.get("passed") is False:
                    report["failed_page_validations"].append(str(validation_path))
            except Exception as exc:
                report["failed_page_validations"].append(f"{validation_path}: {exc}")

    try:
        with zipfile.ZipFile(args.pptx) as z:
            names = z.namelist()
            report["slides"] = len([n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n)])
            for part in ("[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml", "ppt/_rels/presentation.xml.rels"):
                if part not in names:
                    report["missing_parts"].append(part)
            notes_texts = collect_notes_texts(z, names)
            report["notes_found"] = len(notes_texts)
            for entry in notes_manifest.get("notes", []):
                page_index = int(entry.get("page_index", 0))
                expected_hash = entry.get("text_sha256", sha256_text(entry.get("text", "")))
                actual = notes_texts.get(page_index)
                if actual is None:
                    report["notes_hash_mismatches"].append({"page_index": page_index, "reason": "missing notes slide"})
                elif sha256_text(actual) != expected_hash:
                    report["notes_hash_mismatches"].append({"page_index": page_index, "reason": "text hash mismatch"})
    except Exception as exc:
        report["warnings"].append(f"Unable to read pptx: {exc}")

    if page_manifests:
        report["chart_readback_violations"] = quantitative_chart_readback_violations(args.pptx, page_manifests)

    report["passed"] = (
        report["slides"] == expected_pages
        and not report["page_manifests_missing"]
        and not report["page_validation_missing"]
        and not report["failed_page_validations"]
        and not report["page_contract_violations"]
        and not report["missing_parts"]
        and not report["notes_hash_mismatches"]
        and not report["chart_readback_violations"]
    )
    output = json.dumps(report, ensure_ascii=False, indent=2)
    if args.report:
        Path(args.report).parent.mkdir(parents=True, exist_ok=True)
        Path(args.report).write_text(output + "\n", encoding="utf-8")
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except (AttributeError, OSError):
        pass
    print(output)
    raise SystemExit(0 if report["passed"] else 1)


def rel_source_part(rels_name):
    if not rels_name.endswith(".rels"):
        return posixpath.dirname(rels_name)
    directory = posixpath.dirname(rels_name)
    if directory.endswith("/_rels"):
        directory = posixpath.dirname(directory)
    source = posixpath.basename(rels_name)[:-5]
    return posixpath.normpath(posixpath.join(directory, source))


def resolve_target(rels_name, target):
    if not target or re.match(r"^[a-zA-Z][a-zA-Z0-9+.-]*:", target):
        return None
    source = rel_source_part(rels_name)
    return posixpath.normpath(posixpath.join(posixpath.dirname(source), target)).lstrip("/")


def relationship_targets(z, rels_name, names):
    if rels_name not in names:
        return []
    root = ET.fromstring(z.read(rels_name))
    targets = []
    for rel in root.findall("rel:Relationship", NS):
        mode = rel.attrib.get("TargetMode")
        target = rel.attrib.get("Target")
        resolved = resolve_target(rels_name, target)
        targets.append(
            {
                "id": rel.attrib.get("Id"),
                "type": rel.attrib.get("Type", ""),
                "target": target,
                "resolved": resolved,
                "external": mode == "External",
            }
        )
    return targets


def file_sha256(path):
    digest = hashlib.sha256()
    with Path(path).open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx")
    parser.add_argument("--manifest")
    parser.add_argument("--deck-manifest")
    parser.add_argument("--required-text", action="append", default=[])
    parser.add_argument("--report")
    args = parser.parse_args()

    if args.deck_manifest:
        validate_deck(args)

    raw_manifest = read_manifest(args.manifest)
    manifest, authoring_violations = normalize_for_validation(raw_manifest)
    manifest_base = Path(args.manifest).resolve().parent if args.manifest else Path.cwd()
    required = list(args.required_text)
    required.extend(required_texts_from_manifest(manifest))

    report = {
        "pptx": str(Path(args.pptx).resolve()),
        "zip_ok": False,
        "slides": 0,
        "images": 0,
        "editable_text_shapes": 0,
        "shape_count": 0,
        "all_text": "",
        "required_text": required,
        "missing_required_text": [],
        "missing_parts": [],
        "missing_relationship_targets": [],
        "missing_asset_provenance": [],
        "missing_manifest_images": [],
        "missing_provenance_sources": [],
        "invalid_asset_provenance": [],
        "media_hash_mismatches": [],
        "asset_provenance_checked": 0,
        "manifest_image_count": len(manifest.get("images", [])),
        "media_manifest_mismatch": False,
        "relationship_targets_checked": 0,
        "warnings": [],
        "page_contract_violations": [],
        "chart_readback_violations": [],
    }

    try:
        with zipfile.ZipFile(args.pptx) as z:
            bad = z.testzip()
            report["zip_ok"] = bad is None
            if bad:
                report["warnings"].append(f"Bad zip member: {bad}")
            names = z.namelist()
            required_parts = [
                "[Content_Types].xml",
                "_rels/.rels",
                "ppt/presentation.xml",
                "ppt/_rels/presentation.xml.rels",
            ]
            for part in required_parts:
                if part not in names:
                    report["missing_parts"].append(part)
            slide_names = sorted(n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n))
            report["slides"] = len(slide_names)
            report["images"] = len([n for n in names if n.startswith("ppt/media/")])
            report["media_manifest_mismatch"] = report["images"] != report["manifest_image_count"]
            for index, image in enumerate(manifest.get("images", []), start=1):
                image_path = image.get("path")
                if not image_path:
                    continue
                ext = Path(image_path).suffix.lower()
                if ext == ".jpeg":
                    ext = ".jpg"
                media_name = f"ppt/media/image{index}{ext}"
                source_path = Path(image_path)
                if not source_path.is_absolute():
                    source_path = manifest_base / source_path
                if media_name not in names:
                    report["media_hash_mismatches"].append(
                        {"path": image_path, "media": media_name, "reason": "missing media part"}
                    )
                    continue
                if source_path.exists():
                    manifest_hash = file_sha256(source_path)
                    media_hash = hashlib.sha256(z.read(media_name)).hexdigest()
                    if manifest_hash != media_hash:
                        report["media_hash_mismatches"].append(
                            {"path": image_path, "media": media_name, "reason": "hash mismatch"}
                        )
                else:
                    report["missing_manifest_images"].append(str(image_path))
            for slide_name in slide_names:
                rels_name = f"{posixpath.dirname(slide_name)}/_rels/{posixpath.basename(slide_name)}.rels"
                if rels_name not in names:
                    report["missing_parts"].append(rels_name)
            rel_files = [name for name in names if name.endswith(".rels")]
            for rels_name in rel_files:
                for target in relationship_targets(z, rels_name, names):
                    if target["external"] or not target["resolved"]:
                        continue
                    report["relationship_targets_checked"] += 1
                    if target["resolved"] not in names:
                        report["missing_relationship_targets"].append(
                            {
                                "rels": rels_name,
                                "id": target["id"],
                                "target": target["target"],
                                "resolved": target["resolved"],
                            }
                        )
            texts = []
            for slide_name in slide_names:
                xml = z.read(slide_name)
                root = ET.fromstring(xml)
                shapes = root.findall(".//p:sp", NS)
                report["shape_count"] += len(shapes)
                report["editable_text_shapes"] += sum(1 for shape in shapes if shape.findall(".//a:t", NS))
                texts.append(collect_text(xml))
            report["all_text"] = "\n".join(texts)
    except Exception as exc:
        report["warnings"].append(f"Unable to read pptx: {exc}")

    for text in required:
        if text and text not in report["all_text"]:
            report["missing_required_text"].append(text)

    provenance = {}
    for entry in manifest.get("asset_provenance", []):
        path = entry.get("path")
        if path:
            provenance[Path(path).as_posix()] = entry

    for image in manifest.get("images", []):
        image_path = image.get("path")
        if not image_path:
            continue
        key = Path(image_path).as_posix()
        entry = provenance.get(key)
        if not entry:
            report["missing_asset_provenance"].append(key)
            continue
        report["asset_provenance_checked"] += 1
        source_type = entry.get("source_type")
        provenance_note = entry.get("provenance_note")
        if source_type not in ALLOWED_SOURCE_TYPES:
            report["invalid_asset_provenance"].append(
                {"path": key, "field": "source_type", "value": source_type}
            )
        if not provenance_note:
            report["invalid_asset_provenance"].append(
                {"path": key, "field": "provenance_note", "value": provenance_note}
            )
        if source_type == "user-approved-rasterization" and not entry.get("approval_note"):
            report["invalid_asset_provenance"].append(
                {"path": key, "field": "approval_note", "value": entry.get("approval_note")}
            )
        source = entry.get("source")
        if not source:
            report["missing_provenance_sources"].append({"path": key, "source": source})
            continue
        source_path = Path(source)
        if not source_path.is_absolute():
            source_path = manifest_base / source_path
        if not source_path.exists():
            report["missing_provenance_sources"].append({"path": key, "source": str(source)})
    report["page_contract_violations"] = (
        authoring_violations + page_contract_violations(manifest) + quality_contract_violations(raw_manifest)
    )
    if raw_manifest.get("charts"):
        report["chart_readback_violations"] = quantitative_chart_readback_violations(args.pptx, [raw_manifest])

    report["passed"] = (
        report["zip_ok"]
        and report["slides"] >= 1
        and not report["media_manifest_mismatch"]
        and not report["missing_parts"]
        and not report["missing_relationship_targets"]
        and not report["media_hash_mismatches"]
        and not report["missing_required_text"]
        and not report["missing_asset_provenance"]
        and not report["missing_manifest_images"]
        and not report["missing_provenance_sources"]
        and not report["invalid_asset_provenance"]
        and not report["page_contract_violations"]
        and not report["chart_readback_violations"]
        and (report["editable_text_shapes"] > 0 or not required)
    )

    output = json.dumps(report, ensure_ascii=False, indent=2)
    if args.report:
        Path(args.report).parent.mkdir(parents=True, exist_ok=True)
        Path(args.report).write_text(output + "\n", encoding="utf-8")
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except (AttributeError, OSError):
        pass
    print(output)
    raise SystemExit(0 if report["passed"] else 1)


if __name__ == "__main__":
    main()
