#!/usr/bin/env python3
import argparse
import html
import json
import math
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from copy import deepcopy
from pathlib import Path

try:
    from .body_image_profile import mapping_for_source
except ImportError:  # direct runtime script execution through the editppt launcher
    from body_image_profile import mapping_for_source
try:
    from .fixed_region_runtime import require_content_box, require_contract, require_slide, validate_manifest_geometry
except ImportError:  # direct runtime script execution through the editppt launcher
    from fixed_region_runtime import require_content_box, require_contract, require_slide, validate_manifest_geometry


EMU_PER_INCH = 914400
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
ASPECT_16_9 = 16 / 9
ASPECT_TOLERANCE = 0.03
DEFAULT_TEXT_FIT_SAFETY = 0.9
DEFAULT_TEXT_LINE_HEIGHT = 1.22
DEFAULT_MIN_FONT_SIZE = 4.0


def emu(value):
    return int(round(float(value) * EMU_PER_INCH))


def hex_color(value, default="000000"):
    if not value:
        return default
    return str(value).strip().lstrip("#").upper()


def content_type_for(path):
    suffix = Path(path).suffix.lower()
    if suffix == ".png":
        return "image/png"
    if suffix in (".jpg", ".jpeg"):
        return "image/jpeg"
    if suffix == ".gif":
        return "image/gif"
    if suffix == ".svg":
        return "image/svg+xml"
    raise ValueError(f"Unsupported image type: {path}")


def image_ext(path):
    suffix = Path(path).suffix.lower()
    return ".jpg" if suffix == ".jpeg" else suffix


def xml_text(value):
    return html.escape(str(value), quote=True)


def source_size_px(manifest):
    source = manifest.get("source", {})
    width = source.get("width_px")
    height = source.get("height_px")
    if width and height:
        return float(width), float(height)
    return None


def slide_size(manifest):
    slide = manifest.get("slide", {})
    require_slide(slide)
    return float(slide["width"]), float(slide["height"])


def content_box_for_manifest(manifest):
    content_box = manifest.get("content_box")
    require_content_box(content_box)
    return {key: float(content_box[key]) for key in ("left", "top", "width", "height")}


def effective_content_box_for_manifest(manifest):
    """Map the source into the fixed body without ever changing its aspect ratio."""
    source_size = source_size_px(manifest)
    if not source_size:
        raise ValueError("Manifest uses pixel coordinates but lacks source.width_px/source.height_px")
    mapping = mapping_for_source(int(source_size[0]), int(source_size[1]))
    if (
        manifest.get("reconstruction_contract_version") == "editable-image-v3"
        and mapping["image_repair_required"]
    ):
        raise ValueError("editable-image-v3 source must satisfy the 17:8 aspect ratio within 1%")
    box_cm = mapping["effective_box_cm"]
    # The standalone legacy image-conversion interface retains proportional
    # contain. V4 is version-isolated above and never accepts this path.
    if manifest.get("reconstruction_contract_version") != "editable-image-v3" and mapping["image_repair_required"]:
        fixed = content_box_for_manifest(manifest)
        source_ratio = source_size[0] / source_size[1]
        box_ratio = fixed["width"] / fixed["height"]
        if source_ratio >= box_ratio:
            width, height = fixed["width"], fixed["width"] / source_ratio
        else:
            width, height = fixed["height"] * source_ratio, fixed["height"]
        return {
            "left": fixed["left"] + (fixed["width"] - width) / 2,
            "top": fixed["top"] + (fixed["height"] - height) / 2,
            "width": width,
            "height": height,
        }
    return {
        "left": float(box_cm["x"]) / 2.54,
        "top": float(box_cm["y"]) / 2.54,
        "width": float(box_cm["w"]) / 2.54,
        "height": float(box_cm["h"]) / 2.54,
    }


def px_to_inches(manifest, x, y, width, height):
    source_size = source_size_px(manifest)
    if not source_size:
        raise ValueError("Manifest uses pixel coordinates but lacks source.width_px/source.height_px")
    source_width, source_height = source_size
    content_box_for_manifest(manifest)
    content_box = effective_content_box_for_manifest(manifest)
    return {
        "left": content_box["left"] + float(x) / source_width * content_box["width"],
        "top": content_box["top"] + float(y) / source_height * content_box["height"],
        "width": float(width) / source_width * content_box["width"],
        "height": float(height) / source_height * content_box["height"],
    }


def normalize_position_item(manifest, item):
    item = dict(item)
    if "polygon_px" in item:
        points = [(float(point[0]), float(point[1])) for point in item["polygon_px"]]
        if points and "box_px" not in item:
            xs = [point[0] for point in points]
            ys = [point[1] for point in points]
            item["box_px"] = [min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys)]
        item["polygon"] = [
            [px_to_inches(manifest, point[0], point[1], 0, 0)["left"], px_to_inches(manifest, point[0], point[1], 0, 0)["top"]]
            for point in points
        ]
    if "box_px" in item:
        x, y, width, height = item["box_px"]
        item.update(px_to_inches(manifest, x, y, width, height))
    if "points_px" in item:
        x1, y1, x2, y2 = item["points_px"]
        left = min(float(x1), float(x2))
        top = min(float(y1), float(y2))
        width = abs(float(x2) - float(x1))
        height = abs(float(y2) - float(y1))
        item.update(px_to_inches(manifest, left, top, width, height))
        start = px_to_inches(manifest, x1, y1, 0, 0)
        end = px_to_inches(manifest, x2, y2, 0, 0)
        item["points"] = [start["left"], start["top"], end["left"], end["top"]]
        if float(x2) < float(x1):
            item["flip_h"] = True
        if float(y2) < float(y1):
            item["flip_v"] = True
    if item.get("source_corner_radius_px") is not None and "radius" not in item:
        radius = float(item.get("source_corner_radius_px") or 0)
        item["radius"] = px_to_inches(manifest, 0, 0, radius, radius)["width"]
    if item.get("cell_margin_px") is not None:
        margin = float(item["cell_margin_px"])
        mapped = px_to_inches(manifest, 0, 0, margin, margin)
        item["cell_margin_x"] = mapped["width"]
        item["cell_margin_y"] = mapped["height"]
    return item


def iter_text_lines(item):
    if item.get("paragraphs"):
        lines = []
        for paragraph in item["paragraphs"]:
            if isinstance(paragraph, str):
                lines.append(paragraph)
            else:
                runs = paragraph.get("runs")
                if runs:
                    lines.append("".join(str(run.get("text", "")) for run in runs))
                else:
                    lines.append(str(paragraph.get("text", "")))
        return lines or [""]
    if item.get("runs"):
        return ["".join(str(run.get("text", "")) for run in item["runs"])]
    return str(item.get("text", "")).splitlines() or [""]


def text_width_units(text):
    units = 0.0
    for char in str(text):
        codepoint = ord(char)
        if char.isspace():
            units += 0.32
        elif codepoint <= 0x7F:
            units += 0.55
        elif 0xFF00 <= codepoint <= 0xFFEF:
            units += 1.0
        elif 0x4E00 <= codepoint <= 0x9FFF:
            units += 1.0
        else:
            units += 0.85
    return max(units, 1.0)


def longest_unbreakable_units(text):
    tokens = [token for token in re.split(r"\s+", str(text)) if token]
    if not tokens:
        return text_width_units(text)
    return max(text_width_units(token) for token in tokens)


def is_measured_text(item):
    """True when the author marked this box as sized from `page hints` measurement."""
    return str(item.get("font_size_source", "")).strip().lower() in {"measured", "hints"}


def fitted_font_size(item, manifest):
    if item.get("fit_text") is False or manifest.get("fit_text") is False:
        return None
    if "width" not in item or "height" not in item:
        return None
    lines = iter_text_lines(item)
    requested = float(item.get("font_size", 18))
    width_pt = max(1.0, float(item.get("width", 1)) * 72)
    height_pt = max(1.0, float(item.get("height", 0.4)) * 72)
    if is_measured_text(item):
        # Box and font size were both measured from source ink; the safety
        # discount exists to absorb estimation error in hand-written boxes
        # and would systematically shrink correct text here. Clamp only at
        # the geometric limit.
        safety = 1.0
    else:
        safety = float(item.get("text_fit_safety", manifest.get("text_fit_safety", DEFAULT_TEXT_FIT_SAFETY)))
    line_height = float(item.get("line_height", manifest.get("text_line_height", DEFAULT_TEXT_LINE_HEIGHT)))
    wrap_enabled = item.get("wrap") not in (None, "", "none")
    if wrap_enabled:
        line_count = sum(max(1, math.ceil(text_width_units(line) * requested / width_pt)) for line in lines)
        width_limit = width_pt / max(longest_unbreakable_units(line) for line in lines)
    else:
        line_count = max(1, len(lines))
        width_limit = width_pt / max(text_width_units(line) for line in lines)
    height_limit = height_pt / (line_count * max(line_height, 1.0))
    max_font_size = min(width_limit, height_limit) * safety
    explicit_max = item.get("max_font_size")
    if explicit_max not in (None, ""):
        max_font_size = min(max_font_size, float(explicit_max))
    min_font_size = float(item.get("min_font_size", manifest.get("min_font_size", DEFAULT_MIN_FONT_SIZE)))
    return max(min_font_size, max_font_size)


def scale_run_font_sizes(item, ratio):
    def scale_run(run):
        if run.get("font_size") not in (None, ""):
            run["font_size"] = round(float(run["font_size"]) * ratio, 1)

    for run in item.get("runs", []):
        scale_run(run)
    for paragraph in item.get("paragraphs", []):
        if isinstance(paragraph, dict):
            for run in paragraph.get("runs", []):
                scale_run(run)


def fit_text_item(item, manifest):
    fitted = fitted_font_size(item, manifest)
    if fitted is None:
        return item
    requested = float(item.get("font_size", fitted))
    effective = min(requested, fitted)
    if effective < requested:
        item["_requested_font_size"] = requested
        item["font_size"] = round(effective, 1)
        scale_run_font_sizes(item, effective / requested)
    elif "font_size" not in item:
        item["font_size"] = round(effective, 1)
    return item


def _is_number(value):
    return type(value) in (int, float) and math.isfinite(float(value))


def _chart_categories(chart):
    series = chart.get("series", [])
    if not series:
        raise ValueError("charts[].series must be a non-empty list")
    categories = series[0].get("categories")
    if not isinstance(categories, list) or not categories or any(not isinstance(value, str) or not value for value in categories):
        raise ValueError("charts[].series[].categories must contain exact non-empty labels")
    if any(item.get("categories") != categories for item in series[1:]):
        raise ValueError("charts[].series categories must match for one native chart")
    return categories


def _chart_shared_text(chart, field):
    values = [item.get(field, chart.get(field)) for item in chart.get("series", [])]
    if not values or any(not isinstance(value, str) or not value.strip() for value in values) or len(set(values)) != 1:
        raise ValueError(f"charts[].{field} must be one explicit shared value")
    return values[0]


def _validate_chart(manifest, chart):
    if "anchor" in chart or "chart_type" in chart:
        raise ValueError("charts[] anchor/chart_type are unsupported; use fixed-canvas box_px and explicit chart_variant")
    box = chart.get("box_px")
    source = manifest.get("source", {})
    if (
        not isinstance(box, list)
        or len(box) != 4
        or any(not _is_number(value) for value in box)
        or box[0] < 0
        or box[1] < 0
        or box[2] <= 0
        or box[3] <= 0
        or box[0] + box[2] > source.get("width_px", 0)
        or box[1] + box[3] > source.get("height_px", 0)
    ):
        raise ValueError("charts[].box_px must be positive and within source bounds")
    primitive = chart.get("rendering_primitive")
    variant = chart.get("chart_variant")
    allowed = {"column_bar": {"column", "bar"}, "line_point": {"line", "dot"}, "xy": {"scatter", "bubble"}}
    if variant not in allowed.get(primitive, set()):
        raise ValueError("charts[].chart_variant must explicitly match rendering_primitive")
    for field in ("title", "period"):
        if not isinstance(chart.get(field), str) or not chart[field].strip():
            raise ValueError(f"charts[].{field} is required")
    series = chart.get("series")
    if not isinstance(series, list) or not series:
        raise ValueError("charts[].series must be a non-empty list")
    if primitive in {"column_bar", "line_point"}:
        categories = _chart_categories(chart)
        for field in ("unit", "basis"):
            _chart_shared_text(chart, field)
        for item in series:
            values = item.get("values")
            if (
                not isinstance(item.get("name"), str)
                or not item["name"].strip()
                or not isinstance(values, list)
                or len(values) != len(categories)
                or any(not _is_number(value) for value in values)
            ):
                raise ValueError("charts[].series labels and numeric values must match categories exactly")
    else:
        for prefix in ("x", "y") + (("size",) if variant == "bubble" else ()):
            for suffix in ("label", "unit", "basis"):
                field = f"{prefix}_{suffix}"
                if not isinstance(chart.get(field), str) or not chart[field].strip():
                    raise ValueError(f"charts[].{field} is required")
        for item in series:
            x_values, y_values = item.get("x_values"), item.get("y_values")
            if (
                not isinstance(item.get("name"), str)
                or not item["name"].strip()
                or not isinstance(x_values, list)
                or not isinstance(y_values, list)
                or not x_values
                or len(x_values) != len(y_values)
                or any(not _is_number(value) for value in x_values + y_values)
            ):
                raise ValueError("charts[].series x/y values must be aligned numeric lists")
            sizes = item.get("size_values")
            if variant == "bubble" and (
                not isinstance(sizes, list)
                or len(sizes) != len(x_values)
                or any(not _is_number(value) or value < 0 for value in sizes)
            ):
                raise ValueError("charts[].series bubble size_values must align and be non-negative")
    target, actual = chart.get("target_value"), chart.get("actual_value")
    if (target is None) != (actual is None) or target is not None and (
        primitive not in {"column_bar", "line_point"} or not _is_number(target) or not _is_number(actual)
    ):
        raise ValueError("charts[].target_value and actual_value require one compatible explicit numeric pair")


def normalize_manifest(manifest):
    """Return a manifest copy with pixel authoring fields resolved to inches."""
    validate_manifest_geometry(manifest)
    normalized = deepcopy(manifest)
    if normalized.get("reconstruction_contract_version") == "editable-image-v3":
        seen = set()
        for section in ("text_boxes", "tables", "images", "shapes", "charts"):
            for item in normalized.get(section, []):
                object_id = item.get("object_id")
                name = item.get("name")
                if not isinstance(object_id, str) or not object_id.strip():
                    raise ValueError(f"editable-image-v3 {section} object_id is required")
                if object_id in seen:
                    raise ValueError("editable-image-v3 object_id values must be unique")
                seen.add(object_id)
                if not isinstance(name, str) or not name.strip():
                    raise ValueError(f"editable-image-v3 {section} name is required")
                if section == "charts":
                    _validate_chart(normalized, item)
                if section == "tables":
                    required = ("font_size", "font_color", "cell_fill", "cell_margin_px")
                    if any(field not in item for field in required):
                        raise ValueError("editable-image-v3 table cells require explicit font size, foreground, fill, and margins")
                    if float(item["font_size"]) <= 0 or float(item["cell_margin_px"]) < 0:
                        raise ValueError("editable-image-v3 table cell size and margins are invalid")
                    if any(not re.fullmatch(r"#[0-9A-Fa-f]{6}", str(item[field])) for field in ("font_color", "cell_fill")):
                        raise ValueError("editable-image-v3 table cell colors must be explicit RGB")
        # Force the V4 aspect gate even for manifests without positioned objects.
        effective_content_box_for_manifest(normalized)
    normalized["text_boxes"] = [
        fit_text_item(normalize_position_item(normalized, item), normalized) for item in normalized.get("text_boxes", [])
    ]
    for key in ("images", "shapes"):
        normalized[key] = [normalize_position_item(normalized, item) for item in normalized.get(key, [])]
    normalized["tables"] = [normalize_position_item(normalized, item) for item in normalized.get("tables", [])]
    normalized["charts"] = [normalize_position_item(normalized, item) for item in normalized.get("charts", [])]
    return normalized


def preview_color(value):
    if not value or value == "none":
        return value
    value = str(value).strip()
    if value.startswith("#"):
        return value
    if len(value) == 6 and all(ch in "0123456789abcdefABCDEF" for ch in value):
        return f"#{value}"
    return value


def shape_fill(fill):
    if not fill or fill == "none":
        return '<a:noFill/>'
    return f'<a:solidFill><a:srgbClr val="{hex_color(fill)}"/></a:solidFill>'


def shape_line_xml(stroke, width, dash=None):
    if not stroke or stroke == "none":
        return '<a:ln><a:noFill/></a:ln>'
    dash_xml = f'<a:prstDash val="{xml_text(dash)}"/>' if dash else ""
    return (
        f'<a:ln w="{int(float(width or 1) * 12700)}">'
        f'<a:solidFill><a:srgbClr val="{hex_color(stroke)}"/></a:solidFill>'
        f"{dash_xml}"
        "</a:ln>"
    )


def slide_background_xml(slide):
    background = slide.get("background")
    if not background:
        return ""
    return (
        "<p:bg><p:bgPr>"
        f'<a:solidFill><a:srgbClr val="{hex_color(background, "FFFFFF")}"/></a:solidFill>'
        "<a:effectLst/></p:bgPr></p:bg>"
    )


def text_box_xml(idx, item):
    left = emu(item.get("left", 0))
    top = emu(item.get("top", 0))
    width = emu(item.get("width", 1))
    height = emu(item.get("height", 0.4))
    rotation = item.get("rotation")
    rotation_attr = f' rot="{int(float(rotation) * 60000)}"' if rotation not in (None, "") else ""
    font_size = int(float(item.get("font_size", 18)) * 100)
    font = xml_text(item.get("font", "PingFang SC"))
    align = {
        "left": "l", "center": "ctr", "right": "r", "justify": "just",
        "distributed": "dist", "l": "l", "ctr": "ctr", "r": "r",
        "just": "just", "dist": "dist",
    }.get(str(item.get("align", "left")).lower())
    if align is None:
        raise ValueError("unsupported text alignment")
    anchor = {
        "top": "t", "middle": "ctr", "mid": "ctr", "center": "ctr",
        "bottom": "b", "t": "t", "ctr": "ctr", "b": "b",
    }.get(str(item.get("valign", "top")).lower())
    if anchor is None:
        raise ValueError("unsupported vertical text alignment")
    wrap = item.get("wrap", "none")
    autofit = item.get("autofit", "none")
    autofit_xml = "<a:spAutoFit/>" if autofit == "shape" else "<a:noAutofit/>"
    paragraphs = item.get("paragraphs")
    runs = item.get("runs")

    def run_xml(run):
        run_font_size = int(float(run.get("font_size", item.get("font_size", 18))) * 100)
        run_font = xml_text(run.get("font", item.get("font", "PingFang SC")))
        run_color = hex_color(run.get("color", item.get("color", "#111111")))
        run_bold = ' b="1"' if run.get("bold", item.get("bold")) else ""
        run_italic = ' i="1"' if run.get("italic", item.get("italic")) else ""
        run_baseline = run.get("baseline")
        run_baseline_attr = f' baseline="{int(float(run_baseline))}"' if run_baseline not in (None, "") else ""
        run_text = xml_text(run.get("text", ""))
        return (
            f'<a:r><a:rPr lang="zh-CN" sz="{run_font_size}"{run_bold}{run_italic}{run_baseline_attr}>'
            f'<a:solidFill><a:srgbClr val="{run_color}"/></a:solidFill>'
            f'<a:latin typeface="{run_font}"/><a:ea typeface="{run_font}"/><a:cs typeface="{run_font}"/>'
            f'</a:rPr><a:t>{run_text}</a:t></a:r>'
        )

    def paragraph_xml(paragraph):
        if isinstance(paragraph, str):
            paragraph_runs = [{"text": paragraph}]
        else:
            paragraph_runs = paragraph.get("runs", [{"text": paragraph.get("text", "")}])
        return (
            f'<a:p><a:pPr algn="{align}"/>'
            + "".join(run_xml(run) for run in paragraph_runs)
            + f'<a:endParaRPr lang="zh-CN" sz="{font_size}"/></a:p>'
        )

    if paragraphs:
        text_body = "".join(paragraph_xml(paragraph) for paragraph in paragraphs)
    elif runs:
        text_body = paragraph_xml({"runs": runs})
    else:
        text_body = "".join(paragraph_xml(part) for part in str(item.get("text", "")).split("\n"))
    return f"""
      <p:sp>
        <p:nvSpPr><p:cNvPr id="{idx}" name="{xml_text(item.get('name') or f'TextBox {idx}')}" descr="object_id:{xml_text(item.get('object_id', ''))}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
        <p:spPr><a:xfrm{rotation_attr}><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln><a:noFill/></a:ln></p:spPr>
        <p:txBody>
          <a:bodyPr wrap="{xml_text(wrap)}" anchor="{anchor}" lIns="0" tIns="0" rIns="0" bIns="0">{autofit_xml}</a:bodyPr><a:lstStyle/>
          {text_body}
        </p:txBody>
      </p:sp>"""


def image_xml(idx, rel_id, item):
    left = emu(item.get("left", 0))
    top = emu(item.get("top", 0))
    width = emu(item.get("width", 1))
    height = emu(item.get("height", 1))
    name = xml_text(item.get("name") or item.get("alt") or Path(item.get("path", "")).stem or f"Image {idx}")
    return f"""
      <p:pic>
        <p:nvPicPr><p:cNvPr id="{idx}" name="{name}" descr="object_id:{xml_text(item.get('object_id', ''))}"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>
        <p:blipFill><a:blip r:embed="{rel_id}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>
        <p:spPr><a:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
      </p:pic>"""


def shape_xml(idx, item):
    kind = item.get("type", "rect")
    left = emu(item.get("left", 0))
    top = emu(item.get("top", 0))
    width = emu(item.get("width", 1))
    height = emu(item.get("height", 1))
    stroke_width = item.get("stroke_width", 1)
    flip_h = ' flipH="1"' if item.get("flip_h") else ""
    flip_v = ' flipV="1"' if item.get("flip_v") else ""
    fill = shape_fill(item.get("fill"))
    line = shape_line_xml(item.get("stroke", "#000000"), stroke_width, item.get("dash"))
    preset = item.get("preset")
    if item.get("polygon_px"):
        geometry = custom_polygon_geometry_xml(item)
    else:
        if not preset:
            preset = "line" if kind == "line" else "ellipse" if kind == "ellipse" else "roundRect" if kind == "roundRect" else "rect"
        geometry = preset_geometry_xml(preset, item)
    return f"""
      <p:sp>
        <p:nvSpPr><p:cNvPr id="{idx}" name="{xml_text(item.get('name') or f'{kind.title()} {idx}')}" descr="object_id:{xml_text(item.get('object_id', ''))}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
        <p:spPr><a:xfrm{flip_h}{flip_v}><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm>{geometry}{fill}{line}</p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
      </p:sp>"""


def table_xml(idx, item):
    """Build a simple native DrawingML table with stable object identity."""
    rows = item.get("rows")
    if not isinstance(rows, list) or not rows or any(not isinstance(row, list) for row in rows):
        raise ValueError("native table rows must be a non-empty rectangular array")
    columns = len(rows[0])
    if columns < 1 or any(len(row) != columns for row in rows):
        raise ValueError("native table rows must be rectangular")
    left, top = emu(item.get("left", 0)), emu(item.get("top", 0))
    width, height = emu(item.get("width", 1)), emu(item.get("height", 1))
    column_width = max(1, width // columns)
    row_height = max(1, height // len(rows))
    font_size = float(item.get("font_size", 0))
    font_color = hex_color(item.get("font_color"), "")
    cell_fill = hex_color(item.get("cell_fill"), "")
    if font_size <= 0 or not re.fullmatch(r"[0-9A-F]{6}", font_color) or not re.fullmatch(r"[0-9A-F]{6}", cell_fill):
        raise ValueError("native table cells require explicit font_size, font_color, and cell_fill")
    margin_x = emu(item.get("cell_margin_x", 0))
    margin_y = emu(item.get("cell_margin_y", 0))
    font_size_xml = max(1, int(round(font_size * 100)))
    grid = "".join(f'<a:gridCol w="{column_width}"/>' for _ in range(columns))
    row_xml = []
    for row in rows:
        cells = []
        for value in row:
            cells.append(
                '<a:tc><a:txBody><a:bodyPr/><a:lstStyle/><a:p>'
                f'<a:r><a:rPr lang="zh-CN" sz="{font_size_xml}"><a:solidFill><a:srgbClr val="{font_color}"/></a:solidFill></a:rPr><a:t>{xml_text(value)}</a:t></a:r>'
                f'<a:endParaRPr lang="zh-CN" sz="{font_size_xml}"><a:solidFill><a:srgbClr val="{font_color}"/></a:solidFill></a:endParaRPr></a:p></a:txBody>'
                f'<a:tcPr marL="{margin_x}" marR="{margin_x}" marT="{margin_y}" marB="{margin_y}"><a:solidFill><a:srgbClr val="{cell_fill}"/></a:solidFill></a:tcPr></a:tc>'
            )
        row_xml.append(f'<a:tr h="{row_height}">{"".join(cells)}</a:tr>')
    return f"""
      <p:graphicFrame>
        <p:nvGraphicFramePr><p:cNvPr id="{idx}" name="{xml_text(item.get('name') or f'Table {idx}')}" descr="object_id:{xml_text(item.get('object_id', ''))}"/><p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>
        <p:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></p:xfrm>
        <a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table"><a:tbl>
          <a:tblPr firstRow="1" bandRow="1"/><a:tblGrid>{grid}</a:tblGrid>{''.join(row_xml)}
        </a:tbl></a:graphicData></a:graphic>
      </p:graphicFrame>"""


def custom_polygon_geometry_xml(item):
    points = [(float(point[0]), float(point[1])) for point in item.get("polygon_px", [])]
    if len(points) < 3:
        return '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
    box = item.get("box_px")
    if box and len(box) == 4:
        left, top, width, height = [float(value) for value in box]
    else:
        xs = [point[0] for point in points]
        ys = [point[1] for point in points]
        left, top = min(xs), min(ys)
        width, height = max(xs) - left, max(ys) - top
    width = max(width, 1.0)
    height = max(height, 1.0)

    def rel_coord(point):
        x, y = point
        return int(round((x - left) / width * 21600)), int(round((y - top) / height * 21600))

    first_x, first_y = rel_coord(points[0])
    segments = [f'<a:moveTo><a:pt x="{first_x}" y="{first_y}"/></a:moveTo>']
    for point in points[1:]:
        x, y = rel_coord(point)
        segments.append(f'<a:lnTo><a:pt x="{x}" y="{y}"/></a:lnTo>')
    segments.append("<a:close/>")
    return (
        '<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        '<a:rect l="l" t="t" r="r" b="b"/>'
        '<a:pathLst><a:path w="21600" h="21600">'
        + "".join(segments)
        + "</a:path></a:pathLst></a:custGeom>"
    )


def round_rect_adjustment(item):
    box_px = item.get("box_px")
    if item.get("source_corner_radius_px") is not None and box_px and len(box_px) == 4:
        radius = float(item.get("source_corner_radius_px") or 0)
        min_dim = max(1.0, min(float(box_px[2]), float(box_px[3])))
    elif item.get("radius") is not None:
        radius = float(item.get("radius") or 0)
        min_dim = max(0.01, min(float(item.get("width", 1)), float(item.get("height", 1))))
    else:
        return None
    return max(0, min(50000, int(round(radius / min_dim * 100000))))


def preset_geometry_xml(preset, item):
    if preset != "roundRect":
        return f'<a:prstGeom prst="{preset}"><a:avLst/></a:prstGeom>'
    adjustment = round_rect_adjustment(item)
    if adjustment is None:
        return '<a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>'
    return (
        '<a:prstGeom prst="roundRect"><a:avLst>'
        f'<a:gd name="adj" fmla="val {adjustment}"/>'
        '</a:avLst></a:prstGeom>'
    )


def slide_xml(manifest):
    slide = manifest.get("slide", {})
    next_id = 2
    parts = []
    layered = []
    for index, item in enumerate(manifest.get("shapes", [])):
        layered.append((float(item.get("z_index", 100)), index, "shape", item, None))
    for rel_index, item in enumerate(manifest.get("images", []), start=1):
        layered.append((float(item.get("z_index", 200)), rel_index, "image", item, f"rId{rel_index + 1}"))
    for index, item in enumerate(manifest.get("text_boxes", [])):
        layered.append((float(item.get("z_index", 300)), index, "text", item, None))
    for index, item in enumerate(manifest.get("tables", [])):
        layered.append((float(item.get("z_index", 250)), index, "table", item, None))

    for _z_index, _order, kind, item, rel_id in sorted(layered, key=lambda entry: (entry[0], entry[1])):
        if kind == "shape":
            parts.append(shape_xml(next_id, item))
        elif kind == "image":
            parts.append(image_xml(next_id, rel_id, item))
        elif kind == "table":
            parts.append(table_xml(next_id, item))
        else:
            parts.append(text_box_xml(next_id, item))
        next_id += 1
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    {slide_background_xml(slide)}
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      {''.join(parts)}
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>"""


def rels_xml(manifest, media_start=1, notes_index=None):
    rels = ['<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>']
    for i, item in enumerate(manifest.get("images", []), start=1):
        target = f"../media/image{media_start + i - 1}{image_ext(item['path'])}"
        rels.append(f'<Relationship Id="rId{i + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="{target}"/>')
    if notes_index is not None:
        rels.append(
            f'<Relationship Id="rId{len(rels) + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" Target="../notesSlides/notesSlide{notes_index}.xml"/>'
        )
    return '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">' + "".join(rels) + "</Relationships>"


def notes_slide_xml(text):
    paras = "".join(
        f'<a:p><a:r><a:rPr lang="zh-CN" sz="1200"/><a:t>{xml_text(line)}</a:t></a:r><a:endParaRPr lang="zh-CN" sz="1200"/></a:p>'
        for line in str(text).splitlines()
    ) or '<a:p><a:endParaRPr lang="zh-CN" sz="1200"/></a:p>'
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Notes Placeholder"/><p:cNvSpPr txBox="1"/><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="914400"/><a:ext cx="5486400" cy="6858000"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln><a:noFill/></a:ln></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/>{paras}</p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>"""


def notes_rels_xml(slide_index):
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{REL_NS}">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{slide_index}.xml"/>
</Relationships>"""


def notes_master_xml():
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notesMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld>
  <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
</p:notesMaster>"""


def notes_master_rels_xml():
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{REL_NS}">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>"""


def content_types_xml(manifests, notes_indices=None):
    notes_indices = notes_indices or []
    defaults = {
        "rels": "application/vnd.openxmlformats-package.relationships+xml",
        "xml": "application/xml",
    }
    for manifest in manifests:
        for item in manifest.get("images", []):
            ext = image_ext(item["path"]).lstrip(".")
            defaults[ext] = content_type_for(item["path"])
    default_xml = "".join(f'<Default Extension="{ext}" ContentType="{ctype}"/>' for ext, ctype in defaults.items())
    overrides = [
        '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>',
        '<Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>',
        '<Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>',
        '<Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>',
        '<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>',
        '<Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>',
    ]
    for i in range(1, len(manifests) + 1):
        overrides.append(f'<Override PartName="/ppt/slides/slide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>')
    if notes_indices:
        overrides.append('<Override PartName="/ppt/notesMasters/notesMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"/>')
        for i in notes_indices:
            overrides.append(f'<Override PartName="/ppt/notesSlides/notesSlide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>')
    return f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">{default_xml}{"".join(overrides)}</Types>'


def is_wide_slide(width, height):
    return abs((float(width) / float(height)) / ASPECT_16_9 - 1) <= ASPECT_TOLERANCE


def slide_size_type(width, height):
    return "screen16x9" if is_wide_slide(width, height) else None


def presentation_xml(slide_count, width, height):
    slide_ids = "".join(f'<p:sldId id="{255 + i}" r:id="rId{i + 1}"/>' for i in range(1, slide_count + 1))
    size_type = slide_size_type(width, height)
    size_type_attr = f' type="{size_type}"' if size_type else ""
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
  <p:sldIdLst>{slide_ids}</p:sldIdLst>
  <p:sldSz cx="{width}" cy="{height}"{size_type_attr}/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>"""


def presentation_rels_xml(slide_count):
    rels = ['<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>']
    for i in range(1, slide_count + 1):
        rels.append(f'<Relationship Id="rId{i + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{i}.xml"/>')
    return f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="{REL_NS}">{"".join(rels)}</Relationships>'


def theme_xml():
    solid = '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
    line = '<a:ln w="9525"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
    effect = '<a:effectStyle><a:effectLst/></a:effectStyle>'
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="ImageToEditablePPT">'
        '<a:themeElements><a:clrScheme name="Office">'
        '<a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>'
        '<a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1>'
        '<a:dk2><a:srgbClr val="1F1F1F"/></a:dk2><a:lt2><a:srgbClr val="F8F8F8"/></a:lt2>'
        '<a:accent1><a:srgbClr val="0F766E"/></a:accent1><a:accent2><a:srgbClr val="E66B00"/></a:accent2>'
        '<a:accent3><a:srgbClr val="F6D365"/></a:accent3><a:accent4><a:srgbClr val="57C4B8"/></a:accent4>'
        '<a:accent5><a:srgbClr val="666666"/></a:accent5><a:accent6><a:srgbClr val="111111"/></a:accent6>'
        '<a:hlink><a:srgbClr val="0563C1"/></a:hlink><a:folHlink><a:srgbClr val="954F72"/></a:folHlink>'
        '</a:clrScheme><a:fontScheme name="PingFang">'
        '<a:majorFont><a:latin typeface="PingFang SC"/><a:ea typeface="PingFang SC"/><a:cs typeface="PingFang SC"/></a:majorFont>'
        '<a:minorFont><a:latin typeface="PingFang SC"/><a:ea typeface="PingFang SC"/><a:cs typeface="PingFang SC"/></a:minorFont>'
        '</a:fontScheme><a:fmtScheme name="Office">'
        f'<a:fillStyleLst>{solid}{solid}{solid}</a:fillStyleLst>'
        f'<a:lnStyleLst>{line}{line}{line}</a:lnStyleLst>'
        f'<a:effectStyleLst>{effect}{effect}{effect}</a:effectStyleLst>'
        f'<a:bgFillStyleLst>{solid}{solid}{solid}</a:bgFillStyleLst>'
        '</a:fmtScheme></a:themeElements></a:theme>'
    )


def write_common_parts(z, slide_count, width, height, notes_count):
    presentation_format = "Widescreen" if is_wide_slide(width, height) else "Custom"
    z.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/><Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/></Relationships>""")
    z.writestr("docProps/core.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>Image to editable PPT</dc:title></cp:coreProperties>""")
    z.writestr("docProps/app.xml", f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"><Application>Codex</Application><PresentationFormat>{presentation_format}</PresentationFormat><Slides>{slide_count}</Slides></Properties>""")
    z.writestr("ppt/presentation.xml", presentation_xml(slide_count, width, height))
    z.writestr("ppt/_rels/presentation.xml.rels", presentation_rels_xml(slide_count))
    z.writestr("ppt/slideMasters/slideMaster1.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld><p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/><p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst></p:sldMaster>""")
    z.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/></Relationships>""")
    z.writestr("ppt/slideLayouts/slideLayout1.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank" preserve="1"><p:cSld name="Blank"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sldLayout>""")
    z.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/></Relationships>""")
    z.writestr("ppt/theme/theme1.xml", theme_xml())
    if notes_count:
        z.writestr("ppt/notesMasters/notesMaster1.xml", notes_master_xml())
        z.writestr("ppt/notesMasters/_rels/notesMaster1.xml.rels", notes_master_rels_xml())


def write_pptx(manifest, out_path, manifest_path):
    validate_manifest_geometry(manifest)
    width = emu(manifest["slide"]["width"])
    height = emu(manifest["slide"]["height"])
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    normalized = normalize_manifest(manifest)
    media_index = 1
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types_xml([normalized], []))
        write_common_parts(z, 1, width, height, 0)
        z.writestr("ppt/slides/slide1.xml", slide_xml(normalized))
        z.writestr("ppt/slides/_rels/slide1.xml.rels", rels_xml(normalized, media_index, None))
        base = Path(manifest_path).resolve().parent
        for item in normalized.get("images", []):
            src = Path(item["path"])
            if not src.is_absolute():
                src = base / src
            z.write(src, f"ppt/media/image{media_index}{image_ext(src)}")
            media_index += 1
    apply_native_charts(out, [normalized])


def deck_slide_size(deck, page_entries):
    require_contract(deck.get("workflow_contract_version"))
    slide = deck.get("slide")
    require_slide(slide)
    return emu(slide["width"]), emu(slide["height"])


def write_deck(deck, page_entries, out_path, notes_entries):
    """Low-level deck writer. V18 calls it exactly once against a private temp path."""
    if not page_entries:
        raise ValueError("Deck has no pages")
    width, height = deck_slide_size(deck, page_entries)
    for entry in page_entries:
        validate_manifest_geometry(entry["manifest"])
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    notes_by_page = {int(entry.get("page_index", 0)): entry for entry in notes_entries if entry.get("text")}
    notes_indices = sorted(notes_by_page)
    normalized_entries = [{**entry, "manifest": normalize_manifest(entry["manifest"])} for entry in page_entries]
    manifests = [entry["manifest"] for entry in normalized_entries]
    media_index = 1
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types_xml(manifests, notes_indices))
        write_common_parts(z, len(page_entries), width, height, len(notes_by_page))
        for slide_index, entry in enumerate(normalized_entries, start=1):
            manifest = entry["manifest"]
            notes_index = slide_index if slide_index in notes_by_page else None
            z.writestr(f"ppt/slides/slide{slide_index}.xml", slide_xml(manifest))
            z.writestr(f"ppt/slides/_rels/slide{slide_index}.xml.rels", rels_xml(manifest, media_index, notes_index))
            base = Path(entry["manifest_path"]).resolve().parent
            media_bytes = entry.get("media_bytes") or {}
            media_sources = entry.get("media_sources") or {}
            for item in manifest.get("images", []):
                src = Path(item["path"])
                if not src.is_absolute():
                    src = base / src
                relative_name = Path(item["path"]).as_posix()
                if relative_name in media_sources:
                    with z.open(f"ppt/media/image{media_index}{image_ext(src)}", "w", force_zip64=True) as target:
                        media_sources[relative_name].copy_to(target)
                elif relative_name in media_bytes:
                    z.writestr(f"ppt/media/image{media_index}{image_ext(src)}", media_bytes[relative_name])
                else:
                    z.write(src, f"ppt/media/image{media_index}{image_ext(src)}")
                media_index += 1
            if notes_index is not None:
                note = notes_by_page[slide_index]
                notes_xml = note.get("notes_xml")
                if notes_xml and Path(notes_xml).exists():
                    z.writestr(f"ppt/notesSlides/notesSlide{notes_index}.xml", Path(notes_xml).read_bytes())
                else:
                    z.writestr(f"ppt/notesSlides/notesSlide{notes_index}.xml", notes_slide_xml(note.get("text", "")))
                z.writestr(f"ppt/notesSlides/_rels/notesSlide{notes_index}.xml.rels", notes_rels_xml(slide_index))
    apply_native_charts(out, manifests)


def officecli_executable():
    found = shutil.which("officecli")
    if found:
        return found
    bundled = Path.home() / ".codex" / "bin" / "officecli.CMD"
    if bundled.is_file():
        return str(bundled)
    raise RuntimeError("officecli executable is unavailable")


def apply_native_charts(pptx_path, manifests):
    if not any(manifest.get("charts") for manifest in manifests):
        return
    from pptx import Presentation
    from pptx.chart.data import BubbleChartData, ChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt

    chart_types = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "bubble": XL_CHART_TYPE.BUBBLE,
    }

    def identify(shape, chart, role=None):
        shape.name = chart["name"] if role is None else f"{chart['name']} {role}"
        c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
        suffix = "" if role is None else f":{role.lower().replace(' ', '-')}"
        c_nv_pr.set("descr", f"object_id:{chart['object_id']}{suffix};chart_variant:{chart['chart_variant']}")
        return shape

    def add_label(slide, chart, role, text, left, top, width, height):
        shape = identify(slide.shapes.add_textbox(left, top, width, height), chart, role)
        shape.text_frame.clear()
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.text = str(text)
        paragraph.font.size = Pt(9)
        return shape

    def unit_text(chart):
        if chart["rendering_primitive"] != "xy":
            return _chart_shared_text(chart, "unit")
        value = f"x: {chart['x_unit']} | y: {chart['y_unit']}"
        return value + (f" | size: {chart['size_unit']}" if chart["chart_variant"] == "bubble" else "")

    def add_metadata(slide, chart):
        left, top, width = (Inches(chart[key]) for key in ("left", "top", "width"))
        label_height = Inches(min(0.25, chart["height"] / 10))
        add_label(slide, chart, "Unit", unit_text(chart), left + width * 0.55, top, width * 0.45, label_height)
        add_label(slide, chart, "Period", chart["period"], left + width * 0.55, top + label_height, width * 0.45, label_height)

    def number_text(value):
        return str(value)

    def value_bounds(chart):
        values = [float(value) for item in chart["series"] for value in item["values"]]
        values.extend(float(chart[key]) for key in ("target_value", "actual_value") if chart.get(key) is not None)
        low, high = min(values + [0.0]), max(values + [0.0])
        return low, high if high != low else low + 1.0

    def value_position(chart, value, horizontal):
        low, high = value_bounds(chart)
        ratio = (float(value) - low) / (high - low)
        if horizontal:
            return Inches(chart["left"] + chart["width"] * (0.15 + 0.75 * ratio))
        return Inches(chart["top"] + chart["height"] * (0.85 - 0.7 * ratio))

    def add_arrowheads(connector):
        line = connector._element.spPr.get_or_add_ln()
        for tag in ("a:headEnd", "a:tailEnd"):
            arrow = OxmlElement(tag)
            arrow.set("type", "triangle")
            line.append(arrow)

    def add_target_marks(slide, chart):
        if chart.get("target_value") is None:
            return
        left, top, width, height = (Inches(chart[key]) for key in ("left", "top", "width", "height"))
        horizontal = chart["chart_variant"] in {"bar", "dot"}
        target = value_position(chart, chart["target_value"], horizontal)
        actual = value_position(chart, chart["actual_value"], horizontal)
        if horizontal:
            target_line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, target, top + height * 0.15, target, top + height * 0.85)
            difference = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, actual, top + height * 0.1, target, top + height * 0.1)
        else:
            target_line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, left + width * 0.15, target, left + width * 0.9, target)
            difference = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, left + width * 0.93, actual, left + width * 0.93, target)
        identify(target_line, chart, "Target Line")
        identify(difference, chart, "Difference Arrow")
        add_arrowheads(difference)
        label_width, label_height = width * 0.28, Inches(min(0.25, chart["height"] / 10))
        add_label(slide, chart, "Target", f"Target: {number_text(chart['target_value'])}", left, top + height - label_height, label_width, label_height)
        add_label(slide, chart, "Actual", f"Actual: {number_text(chart['actual_value'])}", left + label_width, top + height - label_height, label_width, label_height)
        difference_value = chart["actual_value"] - chart["target_value"]
        add_label(slide, chart, "Difference", f"Difference: {number_text(difference_value)}", left + label_width * 2, top + height - label_height, label_width, label_height)

    def add_dot(slide, chart):
        left, top, width, height = (Inches(chart[key]) for key in ("left", "top", "width", "height"))
        root = identify(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height), chart)
        root.fill.background()
        root.line.fill.background()
        points = [(series, category, value) for series in chart["series"] for category, value in zip(series["categories"], series["values"])]
        series_width = width * 0.4 / len(chart["series"])
        for series_index, series in enumerate(chart["series"], start=1):
            add_label(slide, chart, f"Series {series_index}", series["name"], left + width * 0.15 + series_width * (series_index - 1), top + height * 0.08, series_width, height * 0.08)
        row_height = height * 0.65 / max(1, len(points))
        for index, (_series, category, value) in enumerate(points, start=1):
            y = top + height * 0.18 + row_height * (index - 0.5)
            x = value_position(chart, value, True)
            connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, left + width * 0.15, y, x, y)
            identify(connector, chart, f"Connector {index}")
            diameter = Inches(min(0.14, chart["height"] / 20))
            identify(slide.shapes.add_shape(MSO_SHAPE.OVAL, x - diameter / 2, y - diameter / 2, diameter, diameter), chart, f"Point {index}")
            add_label(slide, chart, f"Category {index}", category, left, y - diameter, width * 0.12, diameter * 2)
            add_label(slide, chart, f"Value {index}", number_text(value), x + diameter, y - diameter, width * 0.12, diameter * 2)
        add_metadata(slide, chart)
        add_target_marks(slide, chart)

    presentation = Presentation(pptx_path)
    for slide_index, manifest in enumerate(manifests, start=1):
        for chart in manifest.get("charts", []):
            slide = presentation.slides[slide_index - 1]
            variant = chart["chart_variant"]
            if variant == "dot":
                add_dot(slide, chart)
                continue
            if variant in {"column", "bar", "line"}:
                data = ChartData()
                data.categories = _chart_categories(chart)
                for item in chart["series"]:
                    data.add_series(item["name"], item["values"])
            else:
                data = BubbleChartData() if variant == "bubble" else XyChartData()
                for item in chart["series"]:
                    native_series = data.add_series(item["name"])
                    for index, (x_value, y_value) in enumerate(zip(item["x_values"], item["y_values"])):
                        if variant == "bubble":
                            native_series.add_data_point(x_value, y_value, item["size_values"][index])
                        else:
                            native_series.add_data_point(x_value, y_value)
            left, top, width, height = (Inches(chart[key]) for key in ("left", "top", "width", "height"))
            chart_shape = identify(slide.shapes.add_chart(
                chart_types[variant],
                left,
                top,
                width,
                height,
                data,
            ), chart)
            native_chart = chart_shape.chart
            title = chart.get("title")
            native_chart.has_title = bool(title and title != "none")
            if native_chart.has_title:
                native_chart.chart_title.text_frame.text = str(title)
            native_chart.has_legend = chart.get("legend", "none") != "none"
            if chart["rendering_primitive"] == "xy":
                native_chart.category_axis.has_title = True
                native_chart.category_axis.axis_title.text_frame.text = chart["x_label"]
                native_chart.value_axis.has_title = True
                native_chart.value_axis.axis_title.text_frame.text = chart["y_label"]
            add_metadata(slide, chart)
            add_target_marks(slide, chart)
    presentation.save(pptx_path)


def page_entries_from_deck_manifest(deck_manifest_path):
    deck_path = Path(deck_manifest_path).resolve()
    deck = json.loads(deck_path.read_text(encoding="utf-8"))
    root = Path(deck.get("job_dir", deck_path.parent)).resolve()
    entries = []
    for page in deck.get("pages", []):
        manifest_path = Path(page.get("manifest", ""))
        if not manifest_path.is_absolute():
            manifest_path = root / manifest_path
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        entries.append({"manifest": manifest, "manifest_path": manifest_path})
    notes_path = deck.get("notes_manifest")
    notes_entries = []
    if notes_path:
        notes_file = Path(notes_path)
        if not notes_file.is_absolute():
            notes_file = root / notes_file
        if notes_file.exists():
            notes_entries = json.loads(notes_file.read_text(encoding="utf-8")).get("notes", [])
            for note in notes_entries:
                notes_xml = note.get("notes_xml")
                if notes_xml:
                    notes_xml_path = Path(notes_xml)
                    if not notes_xml_path.is_absolute():
                        notes_xml_path = root / notes_xml_path
                    note["notes_xml"] = str(notes_xml_path)
    return deck, entries, notes_entries


def output_path_from_deck_manifest(deck_manifest_path):
    deck_path = Path(deck_manifest_path).resolve()
    deck = json.loads(deck_path.read_text(encoding="utf-8"))
    root = Path(deck.get("job_dir", deck_path.parent)).resolve()
    output = Path(deck.get("output", "final/deck_edited.pptx"))
    if not output.is_absolute():
        output = root / output
    return output


def render_preview(manifest, manifest_path, out_path, *, pptx_path=None):
    from PIL import Image, ImageColor, ImageDraw, ImageFont

    manifest = normalize_manifest(manifest)
    slide = manifest.get("slide", {})
    width_in = float(slide.get("width", 13.333))
    height_in = float(slide.get("height", 7.5))
    scale = int(manifest.get("preview_scale", 120))
    canvas = Image.new("RGB", (int(width_in * scale), int(height_in * scale)), ImageColor.getrgb(slide.get("background", "#ffffff")))
    base = Path(manifest_path).resolve().parent
    draw = ImageDraw.Draw(canvas)

    def open_preview_image(src):
        if src.suffix.lower() != ".svg":
            return Image.open(src).convert("RGBA")
        convert = "/opt/homebrew/bin/magick"
        if not Path(convert).exists():
            convert = "/opt/homebrew/bin/convert"
        if not Path(convert).exists():
            print(f"Warning: cannot preview SVG without ImageMagick: {src}", file=sys.stderr)
            return None
        with tempfile.NamedTemporaryFile(suffix=".png") as handle:
            subprocess.run([convert, str(src), handle.name], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            return Image.open(handle.name).convert("RGBA")

    def render_shape(item):
        box = [item.get("left", 0) * scale, item.get("top", 0) * scale, (item.get("left", 0) + item.get("width", 1)) * scale, (item.get("top", 0) + item.get("height", 1)) * scale]
        fill = preview_color(item.get("fill"))
        outline = preview_color(item.get("stroke", "#000000"))
        width = max(1, int(float(item.get("stroke_width", 1))))
        if item.get("polygon"):
            points = [(point[0] * scale, point[1] * scale) for point in item["polygon"]]
            draw.polygon(points, fill=None if fill in (None, "none") else fill, outline=None if outline == "none" else outline)
        elif item.get("type") == "line":
            if "points" in item:
                points = [value * scale for value in item["points"]]
                draw.line(points, fill=outline, width=width)
                return
            if item.get("dash"):
                draw_dashed_line(draw, box, outline, width)
            else:
                draw.line(box, fill=outline, width=width)
        elif item.get("type") == "ellipse":
            draw.ellipse(box, fill=None if fill in (None, "none") else fill, outline=None if outline == "none" else outline, width=width)
        elif item.get("type") == "roundRect" or item.get("preset") == "roundRect":
            radius = int(float(item.get("radius", 0.12)) * scale)
            draw.rounded_rectangle(box, radius=radius, fill=None if fill in (None, "none") else fill, outline=None if outline == "none" else outline, width=width)
        elif item.get("preset") == "diamond":
            left, top, right, bottom = box
            center_x = (left + right) / 2
            center_y = (top + bottom) / 2
            points = [(center_x, top), (right, center_y), (center_x, bottom), (left, center_y)]
            draw.polygon(points, fill=None if fill in (None, "none") else fill, outline=None if outline == "none" else outline)
        else:
            draw.rectangle(box, fill=None if fill in (None, "none") else fill, outline=None if outline == "none" else outline, width=width)

    def render_image(item):
        src = Path(item["path"])
        if not src.is_absolute():
            src = base / src
        img = open_preview_image(src)
        if img is None:
            return
        img = img.resize((max(1, int(item.get("width", 1) * scale)), max(1, int(item.get("height", 1) * scale))))
        canvas.paste(img, (int(item.get("left", 0) * scale), int(item.get("top", 0) * scale)), img)

    def render_text(item):
        preview_font_scale = float(item.get("preview_font_scale", manifest.get("preview_font_scale", 1.0)))
        size = max(1, int(float(item.get("font_size", 18)) * scale / 72 * preview_font_scale))
        font_path = choose_preview_font(item.get("preview_font"))
        try:
            font = ImageFont.truetype(font_path, size=size) if font_path else ImageFont.load_default()
        except Exception:
            font = ImageFont.load_default()
        if item.get("paragraphs"):
            lines = []
            for paragraph in item["paragraphs"]:
                if isinstance(paragraph, str):
                    lines.append(paragraph)
                else:
                    lines.append("".join(str(run.get("text", "")) for run in paragraph.get("runs", [])))
            preview_text = "\n".join(lines)
        elif item.get("runs"):
            preview_text = "".join(str(run.get("text", "")) for run in item["runs"])
        else:
            preview_text = item.get("text", "")
        fill = preview_color(item.get("color", "#111111"))
        align = item.get("align", "left") if item.get("align", "left") in ("left", "center", "right") else "left"
        x = int(item.get("left", 0) * scale)
        y = int(item.get("top", 0) * scale)
        rotation = float(item.get("rotation", 0) or 0)
        if item.get("runs") and not rotation:
            cursor_x = x
            base_size = size
            for run in item["runs"]:
                run_size = max(1, int(float(run.get("font_size", item.get("font_size", 18))) * scale / 72 * preview_font_scale))
                run_font_path = choose_preview_font(run.get("preview_font") or item.get("preview_font"))
                try:
                    run_font = ImageFont.truetype(run_font_path, size=run_size) if run_font_path else ImageFont.load_default()
                except Exception:
                    run_font = font
                run_fill = preview_color(run.get("color", item.get("color", "#111111")))
                baseline = float(run.get("baseline", 0) or 0)
                run_y = y + int(-baseline / 100000 * base_size)
                run_text = str(run.get("text", ""))
                draw.text((cursor_x, run_y), run_text, fill=run_fill, font=run_font)
                cursor_x += int(draw.textlength(run_text, font=run_font))
            return
        if rotation:
            layer_w = max(1, int(item.get("width", 1) * scale))
            layer_h = max(1, int(item.get("height", 0.4) * scale))
            layer = Image.new("RGBA", (layer_w, layer_h), (0, 0, 0, 0))
            layer_draw = ImageDraw.Draw(layer)
            layer_draw.multiline_text((0, 0), preview_text, fill=fill, font=font, spacing=4, align=align)
            rotated = layer.rotate(-rotation, expand=True)
            canvas.paste(rotated, (x, y), rotated)
            return
        draw.multiline_text((x, y), preview_text, fill=fill, font=font, spacing=4, align=align)

    def render_chart(chart):
        series = chart.get("series", [])
        variant = chart["chart_variant"]
        left = int(chart["left"] * scale)
        top = int(chart["top"] * scale)
        width = int(chart["width"] * scale)
        height = int(chart["height"] * scale)
        title_height = int(0.35 * scale) if chart.get("title") not in (None, "none") else 0
        plot = (left + int(0.45 * scale), top + title_height, left + width, top + height - int(0.35 * scale))
        draw.line((plot[0], plot[1], plot[0], plot[3]), fill="#666666", width=1)
        draw.line((plot[0], plot[3], plot[2], plot[3]), fill="#666666", width=1)
        if title_height:
            draw.text((left, top), str(chart["title"]), fill="#111111", font=ImageFont.load_default())
        colors = ("#4472C4", "#ED7D31", "#A5A5A5", "#FFC000")
        if variant in {"scatter", "bubble"}:
            x_values = [float(value) for item in series for value in item["x_values"]]
            y_values = [float(value) for item in series for value in item["y_values"]]
            x_min, x_max = min(x_values), max(x_values)
            y_min, y_max = min(y_values), max(y_values)
            x_span, y_span = x_max - x_min or 1.0, y_max - y_min or 1.0
            for series_index, item in enumerate(series):
                sizes = item.get("size_values", [1] * len(item["x_values"]))
                max_size = max([float(value) for value in sizes] + [1.0])
                for x_value, y_value, size_value in zip(item["x_values"], item["y_values"], sizes):
                    x = plot[0] + (float(x_value) - x_min) / x_span * (plot[2] - plot[0])
                    y = plot[3] - (float(y_value) - y_min) / y_span * (plot[3] - plot[1])
                    radius = 4 if variant == "scatter" else max(4, int(14 * (float(size_value) / max_size) ** 0.5))
                    draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=colors[series_index % len(colors)])
            return

        categories = _chart_categories(chart)
        values = [float(value) for item in series for value in item["values"]]
        maximum = max(values + [0.0])
        minimum = min(values + [0.0])
        span = maximum - minimum or 1.0
        group_width = (plot[2] - plot[0]) / len(categories)
        bar_width = group_width * 0.7 / len(series)
        baseline = plot[1] + maximum / span * (plot[3] - plot[1])
        for series_index, item in enumerate(series):
            line_points = []
            for category_index, value in enumerate(item["values"]):
                x0 = plot[0] + category_index * group_width + group_width * 0.15 + series_index * bar_width
                y = plot[1] + (maximum - float(value)) / span * (plot[3] - plot[1])
                if variant == "bar":
                    y0 = plot[1] + (category_index + 0.15 + series_index * 0.7 / len(series)) * (plot[3] - plot[1]) / len(categories)
                    draw.rectangle((plot[0], y0, plot[0] + (float(value) - minimum) / span * (plot[2] - plot[0]), y0 + (plot[3] - plot[1]) * 0.7 / len(categories) / len(series)), fill=colors[series_index % len(colors)])
                elif variant in {"line", "dot"}:
                    point = (plot[0] + (category_index + 0.5) * group_width, y)
                    line_points.append(point)
                    draw.ellipse((point[0] - 4, point[1] - 4, point[0] + 4, point[1] + 4), fill=colors[series_index % len(colors)])
                else:
                    draw.rectangle((x0, min(y, baseline), x0 + bar_width, max(y, baseline)), fill=colors[series_index % len(colors)])
            if variant == "line" and len(line_points) > 1:
                draw.line(line_points, fill=colors[series_index % len(colors)], width=2)
        for index, category in enumerate(categories):
            draw.text((plot[0] + (index + 0.5) * group_width, plot[3] + 2), str(category), fill="#444444", anchor="ma", font=ImageFont.load_default())

    layered = []
    for index, item in enumerate(manifest.get("shapes", [])):
        layered.append((float(item.get("z_index", 100)), index, render_shape, item))
    for index, item in enumerate(manifest.get("images", [])):
        layered.append((float(item.get("z_index", 200)), index, render_image, item))
    for index, item in enumerate(manifest.get("text_boxes", [])):
        layered.append((float(item.get("z_index", 300)), index, render_text, item))
    for index, item in enumerate(manifest.get("charts", [])):
        layered.append((float(item.get("z_index", 250)), index, render_chart, item))
    for _z_index, _order, renderer, item in sorted(layered, key=lambda entry: (entry[0], entry[1])):
        renderer(item)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    canvas.save(out_path)


def choose_preview_font(preferred):
    candidates = [
        preferred,
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return candidate
    return None


def draw_dashed_line(draw, box, fill, width):
    x1, y1, x2, y2 = box
    dash = 8
    gap = 6
    if abs(y2 - y1) <= abs(x2 - x1):
        step = dash + gap
        x = min(x1, x2)
        end = max(x1, x2)
        y = y1
        while x < end:
            draw.line((x, y, min(x + dash, end), y), fill=fill, width=width)
            x += step
    else:
        step = dash + gap
        y = min(y1, y2)
        end = max(y1, y2)
        x = x1
        while y < end:
            draw.line((x, y, x, min(y + dash, end)), fill=fill, width=width)
            y += step


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("manifest", nargs="?")
    parser.add_argument("--deck-manifest")
    parser.add_argument("--out")
    parser.add_argument("--preview")
    args = parser.parse_args()
    if args.deck_manifest:
        deck, entries, notes_entries = page_entries_from_deck_manifest(args.deck_manifest)
        out = Path(args.out) if args.out else output_path_from_deck_manifest(args.deck_manifest)
        write_deck(deck, entries, out, notes_entries)
        print(f"Wrote {out}")
        return
    if not args.manifest:
        parser.error("manifest is required unless --deck-manifest is used")
    if not args.out:
        parser.error("--out is required unless --deck-manifest provides an output")
    manifest = json.loads(Path(args.manifest).read_text(encoding="utf-8"))
    write_pptx(manifest, args.out, args.manifest)
    if args.preview:
        render_preview(manifest, args.manifest, args.preview, pptx_path=args.out)
    print(f"Wrote {args.out}")
    if args.preview:
        print(f"Wrote {args.preview}")


if __name__ == "__main__":
    main()
